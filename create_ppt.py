from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.chart import XL_CHART_TYPE
import os

# Color Palette - Medical/Professional Theme
PRIMARY_BLUE = RGBColor(0x00, 0x5B, 0x96)      # Deep Blue
SECONDARY_BLUE = RGBColor(0x00, 0x8C, 0xBA)    # Teal Blue
ACCENT_GREEN = RGBColor(0x00, 0x96, 0x4B)       # Medical Green
ACCENT_ORANGE = RGBColor(0xFF, 0x8C, 0x00)      # Warning Orange
ACCENT_RED = RGBColor(0xCC, 0x00, 0x00)          # Alert Red
DARK_GRAY = RGBColor(0x33, 0x33, 0x33)
MEDIUM_GRAY = RGBColor(0x66, 0x66, 0x66)
LIGHT_GRAY = RGBColor(0xF0, 0xF0, 0xF0)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_BLUE = RGBColor(0xE8, 0xF4, 0xF8)
VERY_LIGHT_BLUE = RGBColor(0xF5, 0xFA, 0xFC)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

def add_background(slide, color):
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_shape(slide, left, top, width, height, fill_color, border_color=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.fill.solid()
        shape.line.fill.fore_color.rgb = border_color
    else:
        shape.line.fill.background()
    return shape

def add_rounded_rect(slide, left, top, width, height, fill_color, border_color=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.fill.solid()
        shape.line.fill.fore_color.rgb = border_color
    else:
        shape.line.fill.background()
    return shape

def add_text_box(slide, left, top, width, height, text, font_size=18, bold=False, color=DARK_GRAY, alignment=PP_ALIGN.LEFT, font_name='Calibri'):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.alignment = alignment
    p.font.name = font_name
    return txBox

def add_multi_text(slide, left, top, width, height, texts, font_sizes, bolds, colors, alignments, font_name='Calibri'):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, text in enumerate(texts):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = text
        p.font.size = Pt(font_sizes[i])
        p.font.bold = bolds[i]
        p.font.color.rgb = colors[i]
        p.alignment = alignments[i]
        p.font.name = font_name
    return txBox

# ==============================
# SLIDE 1: TITLE SLIDE
# ==============================
slide1 = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
add_background(slide1, PRIMARY_BLUE)

# Accent bar at top
add_shape(slide1, Inches(0), Inches(0), Inches(13.333), Inches(0.15), ACCENT_GREEN)

# Main title area
add_text_box(slide1, Inches(1), Inches(1.5), Inches(11.333), Inches(1.2),
    "PROSEDUR PEMBERSIHAN &", 40, True, WHITE, PP_ALIGN.CENTER)

add_text_box(slide1, Inches(1), Inches(2.3), Inches(11.333), Inches(1.2),
    "DEKONTAMINASI PERALATAN MEDIS", 40, True, WHITE, PP_ALIGN.CENTER)

# Decorative line
add_shape(slide1, Inches(4.5), Inches(3.6), Inches(4.333), Inches(0.06), ACCENT_GREEN)

# Subtitle
add_text_box(slide1, Inches(1.5), Inches(3.9), Inches(10.333), Inches(0.8),
    "Berdasarkan Klasifikasi Spaulding & Standar WHO/CDC", 22, False, RGBColor(0xBB, 0xDD, 0xEE), PP_ALIGN.CENTER)

# Bottom info bar
add_shape(slide1, Inches(0), Inches(6.5), Inches(13.333), Inches(1), RGBColor(0x00, 0x47, 0x7A))

add_text_box(slide1, Inches(1), Inches(6.65), Inches(5), Inches(0.7),
    "Memutus Rantai Penyebaran Infeksi", 16, False, RGBColor(0x99, 0xCC, 0xEE), PP_ALIGN.LEFT)

add_text_box(slide1, Inches(7.5), Inches(6.65), Inches(5), Inches(0.7),
    "Standar Internasional • Evidence-Based Practice", 14, False, RGBColor(0x99, 0xCC, 0xEE), PP_ALIGN.RIGHT)

# ==============================
# SLIDE 2: DAFTAR ISI
# ==============================
slide2 = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide2, WHITE)

# Header bar
add_shape(slide2, Inches(0), Inches(0), Inches(13.333), Inches(1.1), PRIMARY_BLUE)
add_text_box(slide2, Inches(0.8), Inches(0.25), Inches(11.733), Inches(0.7),
    "DAFTAR ISI", 32, True, WHITE, PP_ALIGN.LEFT)

# Content items
items = [
    ("01", "Pentingnya Dekontaminasi", "Mengapa prosedur ini vital bagi keselamatan pasien"),
    ("02", "Klasifikasi Spaulding", "Tiga kategori alat medis berdasarkan risiko penularan"),
    ("03", "Tahapan Dekontaminasi", "5 langkah wajib yang harus diikuti secara ketat"),
    ("04", "Pre-Cleaning", "Pembersihan awal segera setelah alat digunakan"),
    ("05", "Pembersihan Utama", "Scrubbing dengan detergen enzimatik di CSSD"),
    ("06", "Pembilasan & Pengeringan", "Menghilangkan residu kimia dan mempersiapkan sterilisasi"),
    ("07", "Disinfeksi & Sterilisasi", "Metode berdasarkan Klasifikasi Spaulding"),
    ("08", "Penyimpanan Bersih", "Menjaga sterilitas dengan prinsip FIFO"),
    ("09", "Keselamatan Kerja", "APD lengkap untuk melindungi petugas"),
]

y_start = 1.4
for i, (num, title, desc) in enumerate(items):
    y = y_start + (i * 0.65)
    # Number circle
    circle = slide2.shapes.add_shape(MSO_SHAPE.OVAL, Inches(1), Inches(y), Inches(0.45), Inches(0.45))
    circle.fill.solid()
    circle.fill.fore_color.rgb = PRIMARY_BLUE if i % 2 == 0 else SECONDARY_BLUE
    circle.line.fill.background()
    tf = circle.text_frame
    tf.paragraphs[0].text = num
    tf.paragraphs[0].font.size = Pt(14)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = WHITE
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    tf.word_wrap = False

    add_text_box(slide2, Inches(1.7), Inches(y - 0.02), Inches(4), Inches(0.35),
        title, 16, True, DARK_GRAY, PP_ALIGN.LEFT)
    add_text_box(slide2, Inches(1.7), Inches(y + 0.28), Inches(10), Inches(0.3),
        desc, 11, False, MEDIUM_GRAY, PP_ALIGN.LEFT)

# ==============================
# SLIDE 3: PENTINGNYA DEKONTAMINASI
# ==============================
slide3 = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide3, WHITE)

# Header
add_shape(slide3, Inches(0), Inches(0), Inches(13.333), Inches(1.1), PRIMARY_BLUE)
add_text_box(slide3, Inches(0.8), Inches(0.25), Inches(11.733), Inches(0.7),
    "01  PENTINGNYA DEKONTAMINASI", 32, True, WHITE, PP_ALIGN.LEFT)

# Main message
add_text_box(slide3, Inches(0.8), Inches(1.5), Inches(11.733), Inches(1),
    "Membersihkan dan mendekontaminasi peralatan medis bukan sekadar membuatnya terlihat bersih,\nmelainkan prosedur vital untuk memutus rantai penyebaran infeksi.", 18, False, DARK_GRAY, PP_ALIGN.CENTER)

# Three stat cards
card_data = [
    ("🦠", "HAI (Healthcare-Associated Infections)", "Infeksi nosokomial yang dapat dicegah melalui prosedur dekontaminasi yang tepat", PRIMARY_BLUE),
    ("🔬", "Biofilm", "Lapisan lendir pelindung bakteri yang terbentuk dari kotoran organik yang mengering", ACCENT_GREEN),
    ("📋", "Standar WHO & CDC", "Pedoman internasional yang mengatur prosedur dekontaminasi peralatan medis", SECONDARY_BLUE),
]

for i, (icon, title, desc, color) in enumerate(card_data):
    x = 0.8 + (i * 4.1)
    card = add_rounded_rect(slide3, Inches(x), Inches(3), Inches(3.7), Inches(3.8), VERY_LIGHT_BLUE, color)

    # Icon area
    add_text_box(slide3, Inches(x + 0.2), Inches(3.3), Inches(3.3), Inches(0.8),
        icon, 40, False, color, PP_ALIGN.CENTER)

    add_text_box(slide3, Inches(x + 0.3), Inches(4.2), Inches(3.1), Inches(0.7),
        title, 16, True, color, PP_ALIGN.CENTER)

    add_text_box(slide3, Inches(x + 0.3), Inches(5), Inches(3.1), Inches(1.5),
        desc, 13, False, MEDIUM_GRAY, PP_ALIGN.CENTER)

# ==============================
# SLIDE 4: KLASIFIKASI SPAULDING - OVERVIEW
# ==============================
slide4 = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide4, WHITE)

add_shape(slide4, Inches(0), Inches(0), Inches(13.333), Inches(1.1), PRIMARY_BLUE)
add_text_box(slide4, Inches(0.8), Inches(0.25), Inches(11.733), Inches(0.7),
    "02  KLASIFIKASI SPAULDING", 32, True, WHITE, PP_ALIGN.LEFT)

add_text_box(slide4, Inches(0.8), Inches(1.3), Inches(11.733), Inches(0.7),
    "Standar internasional yang dikeluarkan oleh WHO dan CDC untuk menentukan metode penanganan alat berdasarkan tingkat risiko penularannya.", 16, False, MEDIUM_GRAY, PP_ALIGN.CENTER)

# Three classification cards
class_data = [
    ("KRITIS", "🔴", ACCENT_RED,
     "Jaringan steril atau\npembuluh darah",
     "STERILISASI\n(Wajib membunuh semua\nmikroba & spora)",
     "Instrumen bedah\nKateter\nImplan"),
    ("SEMI-KRITIS", "🟡", ACCENT_ORANGE,
     "Selaput lendir (mukosa)\natau kulit yang terluka",
     "DISINFEKSI\nTINGKAT TINGGI (DTT)\n(Membunuh semua mikroba,\nkecuali spora dalam jumlah besar)",
     "Endoskop\nSirkuit ventilator\nSpekulum"),
    ("NON-KRITIS", "🟢", ACCENT_GREEN,
     "Kulit yang utuh/sehat",
     "DISINFEKSI\nTINGKAT RENDAH/\nMENENGAH",
     "Stetoskop\nManset tensimeter\nPermukaan bed pasien"),
]

for i, (cat, icon, color, area, treatment, examples) in enumerate(class_data):
    x = 0.6 + (i * 4.2)

    # Main card
    card = add_rounded_rect(slide4, Inches(x), Inches(2.2), Inches(3.8), Inches(4.8), WHITE, color)

    # Category header
    header = add_shape(slide4, Inches(x), Inches(2.2), Inches(3.8), Inches(0.9), color)
    add_text_box(slide4, Inches(x + 0.1), Inches(2.3), Inches(3.6), Inches(0.7),
        f"{icon}  {cat}", 22, True, WHITE, PP_ALIGN.CENTER)

    # Area
    add_text_box(slide4, Inches(x + 0.2), Inches(3.25), Inches(3.4), Inches(0.3),
        "Area Kontak:", 11, True, color, PP_ALIGN.LEFT)
    add_text_box(slide4, Inches(x + 0.2), Inches(3.55), Inches(3.4), Inches(0.8),
        area, 12, False, DARK_GRAY, PP_ALIGN.LEFT)

    # Separator
    add_shape(slide4, Inches(x + 0.3), Inches(4.35), Inches(3.2), Inches(0.02), color)

    # Treatment
    add_text_box(slide4, Inches(x + 0.2), Inches(4.45), Inches(3.4), Inches(0.3),
        "Tingkat Penanganan:", 11, True, color, PP_ALIGN.LEFT)
    add_text_box(slide4, Inches(x + 0.2), Inches(4.75), Inches(3.4), Inches(1.1),
        treatment, 11, False, DARK_GRAY, PP_ALIGN.LEFT)

    # Separator
    add_shape(slide4, Inches(x + 0.3), Inches(5.85), Inches(3.2), Inches(0.02), color)

    # Examples
    add_text_box(slide4, Inches(x + 0.2), Inches(5.95), Inches(3.4), Inches(0.3),
        "Contoh Alat:", 11, True, color, PP_ALIGN.LEFT)
    add_text_box(slide4, Inches(x + 0.2), Inches(6.25), Inches(3.4), Inches(0.8),
        examples, 11, False, MEDIUM_GRAY, PP_ALIGN.LEFT)

# ==============================
# SLIDE 5: TAHAPAN DEKONTAMINASI OVERVIEW
# ==============================
slide5 = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide5, WHITE)

add_shape(slide5, Inches(0), Inches(0), Inches(13.333), Inches(1.1), PRIMARY_BLUE)
add_text_box(slide5, Inches(0.8), Inches(0.25), Inches(11.733), Inches(0.7),
    "03  TAHAPAN DEKONTAMINASI", 32, True, WHITE, PP_ALIGN.LEFT)

add_text_box(slide5, Inches(0.8), Inches(1.3), Inches(11.733), Inches(0.7),
    "Urutan dalam proses dekontaminasi bersifat mutlak. Anda tidak boleh langsung mensterilkan alat tanpa membersihkannya terlebih dahulu.", 16, False, MEDIUM_GRAY, PP_ALIGN.CENTER)

# 5 steps with arrows
steps = [
    ("1", "Pre-Cleaning", "Pembersihan Awal", "🔵"),
    ("2", "Cleaning &\nScrubbing", "Pembersihan Utama", "🟦"),
    ("3", "Pembilasan &\nPengeringan", "Menghilangkan Residu", "🟢"),
    ("4", "Disinfeksi/\nSterilisasi", "Berdasarkan Spaulding", "🟧"),
    ("5", "Penyimpanan\nBersih", "Menjaga Sterilitas", "🟥"),
]

for i, (num, title, subtitle, color_dot) in enumerate(steps):
    x = 0.5 + (i * 2.6)

    # Step card
    card = add_rounded_rect(slide5, Inches(x), Inches(2.5), Inches(2.3), Inches(3.5), VERY_LIGHT_BLUE, PRIMARY_BLUE)

    # Number circle
    circle = slide5.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x + 0.85), Inches(2.7), Inches(0.6), Inches(0.6))
    circle.fill.solid()
    circle.fill.fore_color.rgb = PRIMARY_BLUE
    circle.line.fill.background()
    tf = circle.text_frame
    tf.paragraphs[0].text = num
    tf.paragraphs[0].font.size = Pt(22)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = WHITE
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER

    add_text_box(slide5, Inches(x + 0.15), Inches(3.5), Inches(2), Inches(0.9),
        title, 15, True, PRIMARY_BLUE, PP_ALIGN.CENTER)

    add_text_box(slide5, Inches(x + 0.15), Inches(4.5), Inches(2), Inches(0.5),
        subtitle, 11, False, MEDIUM_GRAY, PP_ALIGN.CENTER)

    # Arrow between cards
    if i < 4:
        arrow = slide5.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(x + 2.35), Inches(3.8), Inches(0.25), Inches(0.4))
        arrow.fill.solid()
        arrow.fill.fore_color.rgb = SECONDARY_BLUE
        arrow.line.fill.background()

# Warning note at bottom
warning_box = add_rounded_rect(slide5, Inches(0.8), Inches(6.3), Inches(11.733), Inches(0.8), RGBColor(0xFF, 0xF3, 0xE0), ACCENT_ORANGE)
add_text_box(slide5, Inches(1), Inches(6.4), Inches(11.333), Inches(0.6),
    "⚠️  Kotoran yang mengering dapat membentuk biofilm yang sangat sulit dibersihkan!", 15, True, ACCENT_ORANGE, PP_ALIGN.CENTER)

# ==============================
# SLIDE 6: PRE-CLEANING
# ==============================
slide6 = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide6, WHITE)

add_shape(slide6, Inches(0), Inches(0), Inches(13.333), Inches(1.1), PRIMARY_BLUE)
add_text_box(slide6, Inches(0.8), Inches(0.25), Inches(11.733), Inches(0.7),
    "04  PRE-CLEANING — PEMBERSIHAN AWAL", 32, True, WHITE, PP_ALIGN.LEFT)

# Step indicator
step_badge = add_rounded_rect(slide6, Inches(0.8), Inches(1.4), Inches(1.5), Inches(0.5), PRIMARY_BLUE)
add_text_box(slide6, Inches(0.8), Inches(1.42), Inches(1.5), Inches(0.45),
    "LANGKAH 1", 14, True, WHITE, PP_ALIGN.CENTER)

add_text_box(slide6, Inches(2.5), Inches(1.4), Inches(8), Inches(0.5),
    "Segera setelah alat digunakan", 18, False, DARK_GRAY, PP_ALIGN.LEFT)

# Main content card
card = add_rounded_rect(slide6, Inches(0.8), Inches(2.1), Inches(7.5), Inches(4.8), LIGHT_BLUE, PRIMARY_BLUE)

content_items = [
    ("Tujuan:", "Mencegah darah, cairan tubuh, atau jaringan mengering pada permukaan alat"),
    ("Waktu:", "SEGERA setelah alat digunakan — jangan ditunda!"),
    ("Metode:", "Bilas alat medis langsung di tempat penggunaan"),
    ("Cairan:", "Gunakan air mengalir atau semprotkan busa enzimatik khusus (enzymatic foam)"),
    ("Bahaya:", "Kotoran yang mengering membentuk BIOFILM — lapisan lendir pelindung bakteri yang sangat sulit dibersihkan"),
]

y_pos = 2.3
for title, desc in content_items:
    add_text_box(slide6, Inches(1.2), Inches(y_pos), Inches(6.8), Inches(0.35),
        title, 14, True, PRIMARY_BLUE, PP_ALIGN.LEFT)
    add_text_box(slide6, Inches(1.2), Inches(y_pos + 0.3), Inches(6.8), Inches(0.5),
        desc, 13, False, DARK_GRAY, PP_ALIGN.LEFT)
    y_pos += 0.9

# Right side - key point
key_box = add_rounded_rect(slide6, Inches(8.8), Inches(2.1), Inches(3.7), Inches(4.8), RGBColor(0xFF, 0xEB, 0xEE), ACCENT_RED)

add_text_box(slide6, Inches(9), Inches(2.3), Inches(3.3), Inches(0.5),
    "🔴 PENTING!", 20, True, ACCENT_RED, PP_ALIGN.CENTER)

add_text_box(slide6, Inches(9), Inches(3), Inches(3.3), Inches(3.5),
    "Pre-Cleaning adalah langkah KRITIS yang sering terlupakan.\n\n"
    "Jika dilewatkan:\n"
    "• Biofilm terbentuk\n"
    "• Sterilisasi menjadi kurang efektif\n"
    "• Risiko infeksi meningkat\n\n"
    "Selalu lakukan pre-cleaning SEBELUM memindahkan alat ke CSSD!",
    12, False, ACCENT_RED, PP_ALIGN.LEFT)

# ==============================
# SLIDE 7: PEMBERSIHAN UTAMA
# ==============================
slide7 = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide7, WHITE)

add_shape(slide7, Inches(0), Inches(0), Inches(13.333), Inches(1.1), PRIMARY_BLUE)
add_text_box(slide7, Inches(0.8), Inches(0.25), Inches(11.733), Inches(0.7),
    "05  PEMBERSIHAN UTAMA — CLEANING & SCRUBBING", 32, True, WHITE, PP_ALIGN.LEFT)

step_badge = add_rounded_rect(slide7, Inches(0.8), Inches(1.4), Inches(1.5), Inches(0.5), SECONDARY_BLUE)
add_text_box(slide7, Inches(0.8), Inches(1.42), Inches(1.5), Inches(0.45),
    "LANGKAH 2", 14, True, WHITE, PP_ALIGN.CENTER)

add_text_box(slide7, Inches(2.5), Inches(1.4), Inches(8), Inches(0.5),
    "Menggunakan APD lengkap di ruang dekontaminasi (CSSD)", 18, False, DARK_GRAY, PP_ALIGN.LEFT)

# Left column - procedure
procedure_items = [
    ("1. APD Lengkap", "Gunakan APD sebelum memulai proses pembersihan"),
    ("2. Perendaman", "Rendam alat dalam larutan detergen enzimatik sesuai dosis penunjuk kemasan"),
    ("3. Penyikatan Manual", "Lakukan penyikatan menggunakan sikat lembut di bawah permukaan air"),
    ("4. Pencegahan Aerosol", "Pastikan penyikatan dilakukan di bawah air untuk menghindari cipratan ke wajah"),
    ("5. Alat Berongga (Lumen)", "Gunakan sikat tabung yang sesuai untuk membersihkan bagian dalam"),
]

y_pos = 2.1
for title, desc in procedure_items:
    # Numbered item
    num_box = add_rounded_rect(slide7, Inches(0.8), Inches(y_pos), Inches(0.5), Inches(0.4), SECONDARY_BLUE)
    tf = num_box.text_frame
    tf.paragraphs[0].text = str(procedure_items.index((title, desc)) + 1)
    tf.paragraphs[0].font.size = Pt(14)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = WHITE
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER

    add_text_box(slide7, Inches(1.5), Inches(y_pos - 0.05), Inches(5.5), Inches(0.35),
        title, 14, True, SECONDARY_BLUE, PP_ALIGN.LEFT)
    add_text_box(slide7, Inches(1.5), Inches(y_pos + 0.3), Inches(5.5), Inches(0.45),
        desc, 12, False, DARK_GRAY, PP_ALIGN.LEFT)
    y_pos += 0.85

# Right side - Warning box
warning_card = add_rounded_rect(slide7, Inches(8), Inches(2.1), Inches(4.5), Inches(4.8), RGBColor(0xFF, 0xF8, 0xE1), ACCENT_ORANGE)

add_text_box(slide7, Inches(8.2), Inches(2.3), Inches(4.1), Inches(0.5),
    "⚠️ PERINGATAN KESELAMATAN", 16, True, ACCENT_ORANGE, PP_ALIGN.CENTER)

add_text_box(slide7, Inches(8.2), Inches(3), Inches(4.1), Inches(3.5),
    "• Selalu kenakan APD lengkap:\n"
    "  - Faceshield\n"
    "  - Gaun kedap air\n"
    "  - Apron\n"
    "  - Sarung tangan tebal\n\n"
    "• Penyikatan HARUS di bawah permukaan air\n\n"
    "• Gunakan sikat yang sesuai dengan jenis alat\n\n"
    "• Periksa dosis detergen sesuai kemasan",
    12, False, DARK_GRAY, PP_ALIGN.LEFT)

# ==============================
# SLIDE 8: PEMBILASAN & PENGERINGAN
# ==============================
slide8 = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide8, WHITE)

add_shape(slide8, Inches(0), Inches(0), Inches(13.333), Inches(1.1), PRIMARY_BLUE)
add_text_box(slide8, Inches(0.8), Inches(0.25), Inches(11.733), Inches(0.7),
    "06  PEMBILASAN & PENGERINGAN", 32, True, WHITE, PP_ALIGN.LEFT)

step_badge = add_rounded_rect(slide8, Inches(0.8), Inches(1.4), Inches(1.5), Inches(0.5), ACCENT_GREEN)
add_text_box(slide8, Inches(0.8), Inches(1.42), Inches(1.5), Inches(0.45),
    "LANGKAH 3", 14, True, WHITE, PP_ALIGN.CENTER)

add_text_box(slide8, Inches(2.5), Inches(1.4), Inches(8), Inches(0.5),
    "Menghilangkan residu kimia dan mempersiapkan alat", 18, False, DARK_GRAY, PP_ALIGN.LEFT)

# Two main sections
# Left: Pembilasan
rinse_card = add_rounded_rect(slide8, Inches(0.8), Inches(2.2), Inches(5.8), Inches(4.6), LIGHT_BLUE, ACCENT_GREEN)

add_text_box(slide8, Inches(1), Inches(2.4), Inches(5.4), Inches(0.5),
    "💧 PEMBILASAN", 20, True, ACCENT_GREEN, PP_ALIGN.LEFT)

rinse_items = [
    "Bilas seluruh alat dengan air bersih yang mengalir",
    "Diutamakan menggunakan air murni (demineralized) agar tidak meninggalkan kerak mineral",
    "Pastikan seluruh permukaan alat terbilas dengan sempurna",
    "Periksa bagian-bagian tersembunyi dan lipatan alat",
]

y_pos = 3.1
for item in rinse_items:
    add_text_box(slide8, Inches(1.3), Inches(y_pos), Inches(5.1), Inches(0.5),
        f"• {item}", 13, False, DARK_GRAY, PP_ALIGN.LEFT)
    y_pos += 0.6

# Right: Pengeringan
dry_card = add_rounded_rect(slide8, Inches(7), Inches(2.2), Inches(5.5), Inches(4.6), RGBColor(0xE8, 0xF5, 0xE9), ACCENT_GREEN)

add_text_box(slide8, Inches(7.2), Inches(2.4), Inches(5.1), Inches(0.5),
    "🌬️ PENGERINGAN", 20, True, ACCENT_GREEN, PP_ALIGN.LEFT)

dry_items = [
    "Keringkan alat sepenuhnya menggunakan kain bersih bebas serat (lint-free cloth)",
    "Alternatif: gunakan udara bertekanan medis",
    "Pastikan tidak ada area yang masih lembap",
    "Alat yang lembap dapat menurunkan efektivitas sterilisasi gas atau kimia",
]

y_pos = 3.1
for item in dry_items:
    add_text_box(slide8, Inches(7.5), Inches(y_pos), Inches(4.8), Inches(0.5),
        f"• {item}", 13, False, DARK_GRAY, PP_ALIGN.LEFT)
    y_pos += 0.6

# Bottom note
add_shape(slide8, Inches(0.8), Inches(6.2), Inches(11.733), Inches(0.06), ACCENT_GREEN)
add_text_box(slide8, Inches(0.8), Inches(6.4), Inches(11.733), Inches(0.5),
    "💡 Air demineralized sangat direkomendasikan karena tidak meninggalkan mineral deposit yang dapat mengganggu proses sterilisasi.", 13, False, MEDIUM_GRAY, PP_ALIGN.CENTER)

# ==============================
# SLIDE 9: DISINFEKSI & STERILISASI
# ==============================
slide9 = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide9, WHITE)

add_shape(slide9, Inches(0), Inches(0), Inches(13.333), Inches(1.1), PRIMARY_BLUE)
add_text_box(slide9, Inches(0.8), Inches(0.25), Inches(11.733), Inches(0.7),
    "07  DISINFEKSI & STERILISASI UTAMA", 32, True, WHITE, PP_ALIGN.LEFT)

step_badge = add_rounded_rect(slide9, Inches(0.8), Inches(1.4), Inches(1.5), Inches(0.5), ACCENT_ORANGE)
add_text_box(slide9, Inches(0.8), Inches(1.42), Inches(1.5), Inches(0.45),
    "LANGKAH 4", 14, True, WHITE, PP_ALIGN.CENTER)

add_text_box(slide9, Inches(2.5), Inches(1.4), Inches(8), Inches(0.5),
    "Metode berdasarkan Klasifikasi Spaulding", 18, False, DARK_GRAY, PP_ALIGN.LEFT)

# Three method cards
methods = [
    ("ALAT KRITIS", "🔴", ACCENT_RED,
     "STERILISASI",
     [
         "Autoklaf (sterilisasi uap panas bertekanan)",
         "Suhu: 121°C atau 134°C",
         "Sterilisasi plasma/suhu rendah untuk alat sensitif panas",
         "Membunuh SEMUA mikroba & spora",
     ]),
    ("ALAT SEMI-KRITIS", "🟡", ACCENT_ORANGE,
     "DISINFEKSI TINGKAT TINGGI (DTT)",
     [
         "Glutaraldehyde cair",
         "Hydrogen peroxide cair",
         "Waktu perendaman sesuai standar baku",
         "Membunuh semua mikroba kecuali spora dalam jumlah besar",
     ]),
    ("ALAT NON-KRITIS", "🟢", ACCENT_GREEN,
     "DISINFEKSI TINGKAT RENDAH",
     [
         "Alkohol 70%",
         "Cairan berbasis klorin",
         "Cukup seka pada permukaan alat",
         "Untuk alat yang kontak dengan kulit utuh",
     ]),
]

for i, (cat, icon, color, method, details) in enumerate(methods):
    x = 0.5 + (i * 4.2)

    # Card
    card = add_rounded_rect(slide9, Inches(x), Inches(2.1), Inches(3.9), Inches(4.8), WHITE, color)

    # Header
    header = add_shape(slide9, Inches(x), Inches(2.1), Inches(3.9), Inches(0.8), color)
    add_text_box(slide9, Inches(x + 0.1), Inches(2.2), Inches(3.7), Inches(0.6),
        f"{icon}  {cat}", 18, True, WHITE, PP_ALIGN.CENTER)

    # Method name
    add_text_box(slide9, Inches(x + 0.2), Inches(3.1), Inches(3.5), Inches(0.5),
        method, 16, True, color, PP_ALIGN.CENTER)

    # Separator
    add_shape(slide9, Inches(x + 0.3), Inches(3.65), Inches(3.3), Inches(0.03), color)

    # Details
    y_pos = 3.8
    for detail in details:
        add_text_box(slide9, Inches(x + 0.3), Inches(y_pos), Inches(3.3), Inches(0.5),
            f"• {detail}", 12, False, DARK_GRAY, PP_ALIGN.LEFT)
        y_pos += 0.5

# ==============================
# SLIDE 10: PENYIMPANAN BERSIH
# ==============================
slide10 = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide10, WHITE)

add_shape(slide10, Inches(0), Inches(0), Inches(13.333), Inches(1.1), PRIMARY_BLUE)
add_text_box(slide10, Inches(0.8), Inches(0.25), Inches(11.733), Inches(0.7),
    "08  PENYIMPANAN BERSIH", 32, True, WHITE, PP_ALIGN.LEFT)

step_badge = add_rounded_rect(slide10, Inches(0.8), Inches(1.4), Inches(1.5), Inches(0.5), RGBColor(0x7B, 0x1F, 0xA2))
add_text_box(slide10, Inches(0.8), Inches(1.42), Inches(1.5), Inches(0.45),
    "LANGKAH 5", 14, True, WHITE, PP_ALIGN.CENTER)

add_text_box(slide10, Inches(2.5), Inches(1.4), Inches(8), Inches(0.5),
    "Menjaga sterilitas setelah proses sterilisasi", 18, False, DARK_GRAY, PP_ALIGN.LEFT)

# Main content
storage_items = [
    ("📦", "Kemasan Khusus", "Alat yang sudah disterilkan harus dikemas dalam pouch khusus sterilisasi yang sesuai standar"),
    ("🗄️", "Lemari Tertutup", "Simpan di dalam lemari tertutup dengan suhu dan kelembapan terkontrol"),
    ("🔄", "Prinsip FIFO", "First In, First Out — gunakan alat yang steril paling lama terlebih dahulu"),
    ("✅", "Indikator Kimia", "Periksa indikator kimia pada kemasan untuk memastikan sterilisasi berhasil"),
    ("📋", "Dokumentasi", "Catat tanggal, nomor batch, dan hasil indikator untuk tracing jika diperlukan"),
]

for i, (icon, title, desc) in enumerate(storage_items):
    y = 2.2 + (i * 1.0)

    # Icon circle
    circle = slide10.shapes.add_shape(MSO_SHAPE.OVAL, Inches(1), Inches(y), Inches(0.6), Inches(0.6))
    circle.fill.solid()
    circle.fill.fore_color.rgb = RGBColor(0x7B, 0x1F, 0xA2)
    circle.line.fill.background()
    tf = circle.text_frame
    tf.paragraphs[0].text = icon
    tf.paragraphs[0].font.size = Pt(18)
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER

    add_text_box(slide10, Inches(2), Inches(y - 0.05), Inches(4), Inches(0.35),
        title, 16, True, RGBColor(0x7B, 0x1F, 0xA2), PP_ALIGN.LEFT)
    add_text_box(slide10, Inches(2), Inches(y + 0.3), Inches(10), Inches(0.45),
        desc, 13, False, DARK_GRAY, PP_ALIGN.LEFT)

# Right side - FIFO illustration
fifo_box = add_rounded_rect(slide10, Inches(8.5), Inches(2.2), Inches(4), Inches(4.5), VERY_LIGHT_BLUE, RGBColor(0x7B, 0x1F, 0xA2))

add_text_box(slide10, Inches(8.7), Inches(2.4), Inches(3.6), Inches(0.5),
    "🔄 PRINSIP FIFO", 18, True, RGBColor(0x7B, 0x1F, 0xA2), PP_ALIGN.CENTER)

add_text_box(slide10, Inches(8.7), Inches(3), Inches(3.6), Inches(3.5),
    "First In, First Out\n\n"
    "Alat yang steril paling lama disimpan harus digunakan terlebih dahulu.\n\n"
    "Ini memastikan:\n"
    "• Tidak ada alat yang melewati masa kadaluarsa steril\n"
    "• Rotasi stok berjalan optimal\n"
    "• Kualitas sterilisasi tetap terjaga",
    12, False, DARK_GRAY, PP_ALIGN.LEFT)

# ==============================
# SLIDE 11: KESELAMATAN KERJA (APD)
# ==============================
slide11 = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide11, WHITE)

add_shape(slide11, Inches(0), Inches(0), Inches(13.333), Inches(1.1), ACCENT_RED)
add_text_box(slide11, Inches(0.8), Inches(0.25), Inches(11.733), Inches(0.7),
    "09  KESELAMATAN KERJA — APD LENGKAP", 32, True, WHITE, PP_ALIGN.LEFT)

add_text_box(slide11, Inches(0.8), Inches(1.3), Inches(11.733), Inches(0.7),
    "Petugas wajib mengenakan Alat Pelindung Diri (APD) lengkap selama SELURUH proses dekontaminasi", 16, False, ACCENT_RED, PP_ALIGN.CENTER)

# APD items
apd_items = [
    ("🥽", "Faceshield", "Pelindung wajah dari cipratan dan aerosol berbahaya", ACCENT_RED),
    ("👕", "Gaun Kedap Air", "Mencegah kontaminasi cairan tubuh ke pakaian", ACCENT_ORANGE),
    (" apron", "Apron", "Lapisan pelindung tambahan untuk bagian depan tubuh", PRIMARY_BLUE),
    ("🧤", "Sarung Tangan Tebal", "Utility gloves khusus dekontaminasi — melindungi dari tusukan jarum/instrumen tajam", SECONDARY_BLUE),
]

for i, (icon, title, desc, color) in enumerate(apd_items):
    x = 0.6 + (i * 3.2)

    card = add_rounded_rect(slide11, Inches(x), Inches(2.3), Inches(2.9), Inches(3.8), WHITE, color)

    # Icon area
    icon_box = add_shape(slide11, Inches(x), Inches(2.3), Inches(2.9), Inches(1.2), color)
    add_text_box(slide11, Inches(x + 0.2), Inches(2.4), Inches(2.5), Inches(1),
        icon, 40, False, WHITE, PP_ALIGN.CENTER)

    add_text_box(slide11, Inches(x + 0.2), Inches(3.6), Inches(2.5), Inches(0.5),
        title, 16, True, color, PP_ALIGN.CENTER)

    add_text_box(slide11, Inches(x + 0.2), Inches(4.2), Inches(2.5), Inches(1.5),
        desc, 12, False, DARK_GRAY, PP_ALIGN.CENTER)

# Warning banner
warning_banner = add_rounded_rect(slide11, Inches(0.8), Inches(6.4), Inches(11.733), Inches(0.7), RGBColor(0xFF, 0xEB, 0xEE), ACCENT_RED)
add_text_box(slide11, Inches(1), Inches(6.45), Inches(11.333), Inches(0.6),
    "⚠️  RISIKO: Tusukan jarum, instrumen tajam, paparan cairan tubuh yang terkontaminasi!", 15, True, ACCENT_RED, PP_ALIGN.CENTER)

# ==============================
# SLIDE 12: RINGKASAN
# ==============================
slide12 = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide12, PRIMARY_BLUE)

add_text_box(slide12, Inches(0.8), Inches(0.5), Inches(11.733), Inches(1),
    "RINGKASAN", 36, True, WHITE, PP_ALIGN.CENTER)

add_shape(slide12, Inches(4.5), Inches(1.3), Inches(4.333), Inches(0.06), ACCENT_GREEN)

summary_items = [
    ("1", "Pre-Cleaning", "Segera setelah alat digunakan, bilas untuk mencegah biofilm"),
    ("2", "Cleaning & Scrubbing", "Gunakan detergen enzimatik dengan APD lengkap di CSSD"),
    ("3", "Pembilasan & Pengeringan", "Hilangkan residu kimia, keringkan sepenuhnya"),
    ("4", "Disinfeksi/Sterilisasi", "Berdasarkan Klasifikasi Spaulding (Kritis/Semi-Kritis/Non-Kritis)"),
    ("5", "Penyimpanan Bersih", "Simpan dalam kemasan steril, prinsip FIFO, cek indikator"),
]

for i, (num, title, desc) in enumerate(summary_items):
    y = 1.8 + (i * 1.0)

    # Number
    circle = slide12.shapes.add_shape(MSO_SHAPE.OVAL, Inches(1.5), Inches(y), Inches(0.55), Inches(0.55))
    circle.fill.solid()
    circle.fill.fore_color.rgb = ACCENT_GREEN
    circle.line.fill.background()
    tf = circle.text_frame
    tf.paragraphs[0].text = num
    tf.paragraphs[0].font.size = Pt(18)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = WHITE
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER

    add_text_box(slide12, Inches(2.3), Inches(y - 0.05), Inches(3.5), Inches(0.35),
        title, 18, True, WHITE, PP_ALIGN.LEFT)
    add_text_box(slide12, Inches(2.3), Inches(y + 0.3), Inches(9), Inches(0.45),
        desc, 14, False, RGBColor(0xBB, 0xDD, 0xEE), PP_ALIGN.LEFT)

# ==============================
# SLIDE 13: PENUTUP / TERIMA KASIH
# ==============================
slide13 = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide13, PRIMARY_BLUE)

add_shape(slide13, Inches(0), Inches(0), Inches(13.333), Inches(0.15), ACCENT_GREEN)

add_text_box(slide13, Inches(1), Inches(2), Inches(11.333), Inches(1),
    "TERIMA KASIH", 48, True, WHITE, PP_ALIGN.CENTER)

add_shape(slide13, Inches(4.5), Inches(3.2), Inches(4.333), Inches(0.06), ACCENT_GREEN)

add_text_box(slide13, Inches(1.5), Inches(3.5), Inches(10.333), Inches(1),
    "Memutus Rantai Penyebaran Infeksi\nMelalui Prosedur Dekontaminasi yang Tepat & Konsisten", 20, False, RGBColor(0xBB, 0xDD, 0xEE), PP_ALIGN.CENTER)

add_text_box(slide13, Inches(1.5), Inches(5), Inches(10.333), Inches(0.8),
    "\"Standar internasional WHO & CDC — Klasifikasi Spaulding\"\nEvidence-Based Practice untuk Keselamatan Pasien & Petugas", 16, False, RGBColor(0x99, 0xBB, 0xDD), PP_ALIGN.CENTER)

# Bottom bar
add_shape(slide13, Inches(0), Inches(6.5), Inches(13.333), Inches(1), RGBColor(0x00, 0x47, 0x7A))
add_text_box(slide13, Inches(1), Inches(6.65), Inches(11.333), Inches(0.7),
    " healthcare • infection prevention • quality assurance", 14, False, RGBColor(0x99, 0xCC, 0xEE), PP_ALIGN.CENTER)

# Save the presentation
output_path = "/home/runner/coba_ppt/Dekontaminasi_Peralatan_Medis.pptx"
prs.save(output_path)
print(f"✅ Presentation saved to: {output_path}")
print(f"📊 Total slides: {len(prs.slides)}")
