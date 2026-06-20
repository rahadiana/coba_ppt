#!/usr/bin/env python3
"""
Generate premium PowerPoint presentation for:
Peraturan Wali Kota Bekasi Nomor 51 Tahun 2024
Tentang Pengelolaan Pajak Reklame

Version 2 — Premium Design
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu, Cm
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR, MSO_AUTO_SIZE
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR_TYPE
from pptx.oxml.ns import qn
from pptx.oxml import parse_xml
import copy

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# ═══════════════════════════════════════════════════════════════
# COLOR PALETTE — Premium Modern
# ═══════════════════════════════════════════════════════════════
class Color:
    DARK_NAVY   = RGBColor(0x0A, 0x1F, 0x44)   # Slide backgrounds
    MID_BLUE    = RGBColor(0x1A, 0x3C, 0x6E)
    ACCENT_BLUE = RGBColor(0x2D, 0x7D, 0xD2)   # Accents
    GOLD        = RGBColor(0xF5, 0xC5, 0x18)   # Gold accent
    WARM_GOLD   = RGBColor(0xD4, 0xA0, 0x17)
    WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
    OFF_WHITE   = RGBColor(0xF5, 0xF7, 0xFA)
    DARK_TEXT    = RGBColor(0x1A, 0x1A, 0x2E)
    GRAY_TEXT   = RGBColor(0x6B, 0x72, 0x88)
    SOFT_GRAY   = RGBColor(0xE8, 0xEC, 0xF1)
    GREEN_SOFT  = RGBColor(0xE8, 0xF5, 0xE9)
    BLUE_SOFT   = RGBColor(0xE3, 0xF2, 0xFD)
    RED_ACCENT  = RGBColor(0xE7, 0x4C, 0x3C)
    TEAL        = RGBColor(0x00, 0x96, 0x88)
    ORANGE_SOFT = RGBColor(0xFF, 0xF3, 0xE0)

def add_bg_solid(slide, color=Color.WHITE):
    """Full-slide solid background."""
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid(); bg.fill.fore_color.rgb = color; bg.line.fill.background()
    sp = bg._element
    sp.getparent().remove(sp)
    # Actually don't remove; we'll use it as the bottom layer
    return bg

def add_gradient_bg(slide, color_top=Color.DARK_NAVY, color_bottom=Color.MID_BLUE):
    """Add gradient background using XML."""
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid(); bg.fill.fore_color.rgb = color_top; bg.line.fill.background()
    # Move to background via XML
    sp = bg._element
    sp.getparent().remove(sp)
    # Actually we need it, let's just remove fill later
    return bg

def add_shape(slide, shape_type, left, top, width, height, fill_color=None, line_color=None, line_width=None):
    shp = slide.shapes.add_shape(shape_type, left, top, width, height)
    if fill_color:
        shp.fill.solid(); shp.fill.fore_color.rgb = fill_color
    else:
        shp.fill.background()
    if line_color:
        shp.line.color.rgb = line_color
        if line_width: shp.line.width = line_width
    else:
        shp.line.fill.background()
    return shp

def add_rounded_rect(slide, left, top, width, height, fill_color, shadow=False):
    """Add rounded rectangle shape."""
    shp = add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height, fill_color)
    return shp

def add_textbox(slide, left, top, width, height, text, font_size=18,
                bold=False, color=Color.DARK_TEXT, alignment=PP_ALIGN.LEFT,
                font_name='Calibri', italic=False, anchor=MSO_ANCHOR.TOP):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame; tf.word_wrap = True; tf.auto_size = None
    p = tf.paragraphs[0]; p.text = text
    p.font.size = Pt(font_size); p.font.bold = bold; p.font.color.rgb = color
    p.font.name = font_name; p.font.italic = italic; p.alignment = alignment
    try: txBox.text_frame.paragraphs[0].space_after = Pt(0)
    except: pass
    return txBox

def add_rich_textbox(slide, left, top, width, height, segments, alignment=PP_ALIGN.LEFT, line_spacing=1.3):
    """Add textbox with rich text segments.
    segments = [(text, font_size, bold, color), ...]
    """
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame; tf.word_wrap = True
    for i, seg in enumerate(segments):
        if i == 0: p = tf.paragraphs[0]
        else: p = tf.add_paragraph()
        p.alignment = alignment
        p.space_after = Pt(2)
        if isinstance(seg, str):
            run = p.add_run(); run.text = seg
            run.font.size = Pt(16); run.font.color.rgb = Color.DARK_TEXT; run.font.name = 'Calibri'
        else:
            text, size, bold, color, *rest = seg if len(seg) >= 4 else (*seg, 'Calibri')
            run = p.add_run(); run.text = text
            run.font.size = Pt(size); run.font.bold = bold; run.font.color.rgb = color
            run.font.name = rest[0] if rest else 'Calibri'
    return txBox

def add_bullets(slide, left, top, width, height, items, font_size=15, color=Color.DARK_TEXT,
                bullet_char="▸", title_style=False):
    """Add bullet list with optional title style (first line bold)."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame; tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(3); p.space_before = Pt(1)
        
        if item == "":
            p.text = ""; p.space_after = Pt(6); continue
        
        # Check if item uses **bold** markers
        if "**" in item:
            parts = item.split("**")
            run = p.add_run()
            run.text = f"{bullet_char} {parts[0]}"
            run.font.size = Pt(font_size); run.font.color.rgb = color; run.font.name = 'Calibri'
            run.font.bold = True
            if len(parts) > 1:
                r2 = p.add_run()
                r2.text = parts[1]
                r2.font.size = Pt(font_size); r2.font.color.rgb = color; r2.font.name = 'Calibri'
        else:
            # Indent for sub-items
            prefix = bullet_char if not item.startswith("   ") else ""
            display_text = item.strip()
            if prefix:
                p.text = f"{prefix} {display_text}"
            else:
                p.text = f"   {display_text}"
            p.font.size = Pt(font_size); p.font.color.rgb = color; p.font.name = 'Calibri'
    return txBox

def add_card(slide, left, top, width, height, title, items, title_color=Color.ACCENT_BLUE, 
             bg_color=Color.WHITE, icon_text=""):
    """A card-style box with title and items."""
    card = add_rounded_rect(slide, left, top, width, height, bg_color)
    # Subtle shadow effect via XML
    spPr = card._element.spPr
    # Add a top accent line
    accent = add_shape(slide, MSO_SHAPE.RECTANGLE, left, top, width, Pt(4), title_color)
    
    if icon_text:
        add_textbox(slide, left + Inches(0.15), top + Inches(0.1), Inches(0.4), Inches(0.4),
                    icon_text, font_size=22, color=title_color, bold=True)
        title_left = left + Inches(0.55)
    else:
        title_left = left + Inches(0.2)
    
    add_textbox(slide, title_left, top + Inches(0.15), width - Inches(0.4), Inches(0.4),
                title, font_size=17, bold=True, color=title_color)
    
    add_bullets(slide, left + Inches(0.2), top + Inches(0.55), width - Inches(0.4), 
                height - Inches(0.7), items, font_size=13, color=Color.DARK_TEXT)
    return card

def add_page_number(slide, num, total=54):
    add_textbox(slide, Inches(12.2), Inches(7.05), Inches(1), Inches(0.35),
                f"{num} / {total}", font_size=10, color=Color.GRAY_TEXT, alignment=PP_ALIGN.RIGHT)

def add_footer_bar(slide):
    add_shape(slide, MSO_SHAPE.RECTANGLE, 0, Inches(7.3), prs.slide_width, Inches(0.2), Color.DARK_NAVY)

def add_header_bar(slide, title_text, subtitle_text=""):
    """Premium header with gradient bar and title."""
    # Top gold thin line
    add_shape(slide, MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Pt(3), Color.GOLD)
    # Navy header band
    add_shape(slide, MSO_SHAPE.RECTANGLE, 0, Pt(3), prs.slide_width, Inches(1.15), Color.DARK_NAVY)
    # Accent vertical bar
    add_shape(slide, MSO_SHAPE.RECTANGLE, Inches(0.4), Pt(3), Pt(5), Inches(1.15), Color.GOLD)
    # Title
    add_textbox(slide, Inches(0.6), Inches(0.12), Inches(11), Inches(0.55),
                title_text, font_size=26, bold=True, color=Color.WHITE)
    if subtitle_text:
        add_textbox(slide, Inches(0.6), Inches(0.62), Inches(11), Inches(0.45),
                    subtitle_text, font_size=12, color=RGBColor(0xAA, 0xBB, 0xDD))
    # Bottom separator
    add_shape(slide, MSO_SHAPE.RECTANGLE, 0, Inches(1.18), prs.slide_width, Pt(1), Color.SOFT_GRAY)

# ═══════════════════════════════════════════════════════════════
# SLIDE 1 — COVER / TITLE (Premium Design)
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
# Dark navy background with decorative elements
add_shape(slide, MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height, Color.DARK_NAVY)
# Decorative large circle (very subtle)
add_shape(slide, MSO_SHAPE.OVAL, Inches(8), -Inches(2), Inches(7), Inches(7), 
          fill_color=RGBColor(0x0D, 0x2B, 0x5A))
add_shape(slide, MSO_SHAPE.OVAL, Inches(9.5), Inches(4), Inches(5), Inches(5), 
          fill_color=RGBColor(0x0F, 0x30, 0x60))
# Gold accent lines
add_shape(slide, MSO_SHAPE.RECTANGLE, Inches(1), Inches(2.6), Inches(4), Pt(4), Color.GOLD)
add_shape(slide, MSO_SHAPE.RECTANGLE, Inches(1), Inches(5.6), Inches(3), Pt(2), Color.WARM_GOLD)

# Top metadata
add_textbox(slide, Inches(1), Inches(0.5), Inches(8), Inches(0.5),
            "BERITA DAERAH KOTA BEKASI", font_size=14, color=Color.GOLD, bold=True)
add_textbox(slide, Inches(1), Inches(0.9), Inches(5), Inches(0.4),
            "NOMOR 51  |  TAHUN 2024", font_size=12, color=RGBColor(0x88, 0x99, 0xBB))

# Main title
add_textbox(slide, Inches(1), Inches(1.7), Inches(11), Inches(0.7),
            "PERATURAN WALI KOTA BEKASI", font_size=32, bold=True, color=Color.WHITE)
add_textbox(slide, Inches(1), Inches(2.2), Inches(11), Inches(0.6),
            "NOMOR 51 TAHUN 2024", font_size=24, color=Color.GOLD, bold=True)

# Subtitle
add_textbox(slide, Inches(1), Inches(3.2), Inches(11), Inches(1.5),
            "TENTANG\nPENGELOLAAN PAJAK REKLAME",
            font_size=44, bold=True, color=Color.WHITE)

# Info
add_textbox(slide, Inches(1), Inches(4.8), Inches(8), Inches(0.4),
            "Diselenggarakan oleh Pemerintah Kota Bekasi", 
            font_size=15, color=RGBColor(0xAA, 0xBB, 0xDD))
add_textbox(slide, Inches(1), Inches(5.2), Inches(8), Inches(0.4),
            "Ditetapkan: 20 Desember 2024  |  Berlaku sejak diundangkan",
            font_size=13, color=RGBColor(0x88, 0x99, 0xBB))

# Gold bottom bar
add_shape(slide, MSO_SHAPE.RECTANGLE, 0, Inches(7.2), prs.slide_width, Inches(0.3), Color.GOLD)

# ═══════════════════════════════════════════════════════════════
# SLIDE 2 — DAFTAR ISI
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_shape(slide, MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height, Color.WHITE)
add_header_bar(slide, "DAFTAR ISI", "Peraturan Wali Kota Bekasi Nomor 51 Tahun 2024")

toc_items = [
    ("BAB I", "Ketentuan Umum"),
    ("BAB II", "Objek Pajak, Subjek Pajak, dan Wajib Pajak"),
    ("BAB III", "Masa Pajak, Tahun Pajak, dan Bagian Tahun Pajak"),
    ("BAB IV", "Pendaftaran dan Pendataan Wajib Pajak"),
    ("BAB V", "Nilai Sewa Reklame (NSR)"),
    ("BAB VI", "Perhitungan Pajak Reklame"),
    ("BAB VII", "Penetapan Besaran Pajak Terutang"),
    ("BAB VIII", "Surat Tagihan Pajak"),
    ("BAB IX", "Pembayaran dan Penyetoran"),
    ("BAB X", "Pembetulan dan Pembatalan Ketetapan"),
    ("BAB XI", "Pengembalian Kelebihan Pembayaran"),
    ("BAB XII", "Pemeriksaan Pajak"),
    ("BAB XIII", "Penagihan Pajak"),
    ("BAB XIV", "Kedaluwarsa Penagihan Pajak"),
    ("BAB XV", "Penghapusan Piutang Pajak"),
    ("BAB XVI", "Keringanan, Pengurangan, dan Pembebasan"),
    ("BAB XVII", "Kemudahan Perpajakan Daerah"),
    ("BAB XVIII", "Keberatan dan Banding"),
    ("BAB XIX", "Penghargaan"),
    ("BAB XX", "Ketentuan Penutup"),
]

# Two columns
col1 = toc_items[:10]
col2 = toc_items[10:]

y_start = Inches(1.5)
for i, (bab, title) in enumerate(col1):
    y = y_start + Inches(i * 0.48)
    num_box = add_rounded_rect(slide, Inches(0.6), y, Inches(1.0), Inches(0.38), Color.ACCENT_BLUE)
    add_textbox(slide, Inches(0.6), y + Inches(0.03), Inches(1.0), Inches(0.35),
                bab, font_size=12, bold=True, color=Color.WHITE, alignment=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(1.75), y + Inches(0.03), Inches(5), Inches(0.35),
                title, font_size=15, color=Color.DARK_TEXT)

for i, (bab, title) in enumerate(col2):
    y = y_start + Inches(i * 0.48)
    num_box = add_rounded_rect(slide, Inches(6.5), y, Inches(1.0), Inches(0.38), Color.TEAL)
    add_textbox(slide, Inches(6.5), y + Inches(0.03), Inches(1.0), Inches(0.35),
                bab, font_size=12, bold=True, color=Color.WHITE, alignment=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(7.65), y + Inches(0.03), Inches(5), Inches(0.35),
                title, font_size=15, color=Color.DARK_TEXT)

add_footer_bar(slide)
add_page_number(slide, 2)

# ═══════════════════════════════════════════════════════════════
# HELPER: Section Divider (Premium)
# ═══════════════════════════════════════════════════════════════
slide_counter = [2]  # track slides for page numbers

def section_slide(bab_title, bab_sub="", bab_num=""):
    global slide_counter
    slide_counter[0] += 1
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_shape(slide, MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height, Color.DARK_NAVY)
    # Decorative circles
    add_shape(slide, MSO_SHAPE.OVAL, Inches(10), -Inches(1), Inches(4), Inches(4),
              fill_color=RGBColor(0x0D, 0x2B, 0x5A))
    add_shape(slide, MSO_SHAPE.OVAL, Inches(-1.5), Inches(4), Inches(5), Inches(5),
              fill_color=RGBColor(0x0D, 0x2B, 0x5A))
    # Gold accent line
    add_shape(slide, MSO_SHAPE.RECTANGLE, Inches(1), Inches(2.0), Inches(2.5), Pt(4), Color.GOLD)
    
    if bab_num:
        add_textbox(slide, Inches(1), Inches(1.2), Inches(10), Inches(0.5),
                    bab_num, font_size=18, color=Color.GOLD, bold=True)
    
    add_textbox(slide, Inches(1), Inches(2.3), Inches(11), Inches(2),
                bab_title, font_size=40, bold=True, color=Color.WHITE)
    
    if bab_sub:
        add_textbox(slide, Inches(1), Inches(4.5), Inches(10), Inches(0.5),
                    bab_sub, font_size=16, color=RGBColor(0xAA, 0xBB, 0xDD))
    
    # Gold bottom bar
    add_shape(slide, MSO_SHAPE.RECTANGLE, 0, Inches(7.2), prs.slide_width, Inches(0.3), Color.GOLD)
    add_page_number(slide, slide_counter[0])
    return slide

def content_slide(title, items, sub_text="", notes="", two_col=False):
    """Standard premium content slide."""
    global slide_counter
    slide_counter[0] += 1
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_shape(slide, MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height, Color.WHITE)
    add_header_bar(slide, title, sub_text)
    
    y_start = Inches(1.5)
    if two_col:
        mid = len(items) // 2
        add_bullets(slide, Inches(0.5), y_start, Inches(5.8), Inches(5.2),
                    items[:mid], font_size=15, color=Color.DARK_TEXT)
        add_bullets(slide, Inches(6.8), y_start, Inches(5.8), Inches(5.2),
                    items[mid:], font_size=15, color=Color.DARK_TEXT)
    else:
        add_bullets(slide, Inches(0.6), y_start, Inches(12), Inches(5.2),
                    items, font_size=15, color=Color.DARK_TEXT)
    
    if notes:
        add_textbox(slide, Inches(0.6), Inches(6.6), Inches(12), Inches(0.5),
                    notes, font_size=11, color=Color.GRAY_TEXT, italic=True)
    
    add_footer_bar(slide)
    add_page_number(slide, slide_counter[0])
    return slide

def card_slide(title, cards_data, sub_text=""):
    """Slide with card layout. cards_data = [(title, [items], color), ...]"""
    global slide_counter
    slide_counter[0] += 1
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_shape(slide, MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height, Color.OFF_WHITE)
    add_header_bar(slide, title, sub_text)
    
    card_width = Inches(3.8)
    card_height = Inches(4.8)
    gap = Inches(0.3)
    total_width = len(cards_data) * card_width + (len(cards_data) - 1) * gap
    start_x = (prs.slide_width - total_width) // 2
    
    for i, (c_title, c_items, c_color) in enumerate(cards_data):
        x = start_x + i * (card_width + gap)
        add_card(slide, x, Inches(1.5), card_width, card_height, c_title, c_items, c_color)
    
    add_footer_bar(slide)
    add_page_number(slide, slide_counter[0])
    return slide

def premium_table_slide(title, headers, rows, notes="", sub_text="", col_widths=None):
    """Slide with premium styled table."""
    global slide_counter
    slide_counter[0] += 1
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_shape(slide, MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height, Color.WHITE)
    add_header_bar(slide, title, sub_text)
    
    num_rows = len(rows) + 1
    num_cols = len(headers)
    
    tbl_left = Inches(0.8)
    tbl_top = Inches(1.6)
    tbl_width = Inches(11.7)
    tbl_height = Inches(5.2)
    
    tbl_shape = slide.shapes.add_table(num_rows, num_cols, tbl_left, tbl_top, tbl_width, tbl_height)
    tbl = tbl_shape.table
    
    if col_widths:
        for i, w in enumerate(col_widths):
            tbl.columns[i].width = w
    
    # Header row
    for j, h in enumerate(headers):
        cell = tbl.cell(0, j)
        cell.text = h
        cell.fill.solid(); cell.fill.fore_color.rgb = Color.DARK_NAVY
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(12); p.font.bold = True; p.font.color.rgb = Color.WHITE
            p.font.name = 'Calibri'; p.alignment = PP_ALIGN.CENTER
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
    
    # Data rows with alternating colors
    for i, row in enumerate(rows):
        for j, val in enumerate(row):
            cell = tbl.cell(i + 1, j)
            cell.text = str(val)
            bg = Color.BLUE_SOFT if i % 2 == 0 else Color.WHITE
            cell.fill.solid(); cell.fill.fore_color.rgb = bg
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(11); p.font.color.rgb = Color.DARK_TEXT; p.font.name = 'Calibri'
                p.alignment = PP_ALIGN.CENTER if j > 0 else PP_ALIGN.LEFT
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
    
    if notes:
        add_textbox(slide, Inches(0.8), Inches(6.8), Inches(11), Inches(0.4),
                    notes, font_size=10, color=Color.GRAY_TEXT, italic=True)
    
    add_footer_bar(slide)
    add_page_number(slide, slide_counter[0])
    return slide


# ═══════════════════════════════════════════════════════════════
# BAB I: KETENTUAN UMUM
# ═══════════════════════════════════════════════════════════════
section_slide("BAB I: KETENTUAN UMUM", "Pasal 1 — 45 Definisi Istilah")

content_slide("Definisi Kunci (Pasal 1)", [
    "**Daerah** — Daerah Kota Bekasi",
    "**Pemerintah Daerah** — Wali Kota sebagai unsur penyelenggara pemerintahan daerah",
    "**Bapenda** — Badan Pendapatan Daerah Kota Bekasi",
    "**Reklame** — Benda/alat/media untuk tujuan komersial memperkenalkan, mempromosikan sesuatu",
    "**Pajak Reklame** — Pajak atas penyelenggaraan reklame",
    "**Nilai Sewa Reklame (NSR)** — Dasar pengenaan Pajak Reklame",
    "**NPWPD** — Nomor Pokok Wajib Pajak Daerah (identitas Wajib Pajak)",
    "**NOPD** — Nomor Objek Pajak Daerah (identitas objek pajak)",
    "**SKPD** — Surat Ketetapan Pajak Daerah",
    "**Wajib Pajak** — Orang pribadi/badan dengan hak dan kewajiban perpajakan",
    "**Penanggung Pajak** — Pihak yang bertanggung jawab atas pembayaran pajak",
], sub_text="Pasal 1 ayat (1) - (45)")

# ═══════════════════════════════════════════════════════════════
# BAB II: OBJEK PAJAK, SUBJEK PAJAK, DAN WAJIB PAJAK
# ═══════════════════════════════════════════════════════════════
section_slide("BAB II: OBJEK PAJAK, SUBJEK PAJAK, DAN WAJIB PAJAK", "Pasal 2–4")

card_slide("Objek Pajak Reklame — Pasal 2", [
    ("Jenis Reklame", [
        "Papan/Billboard",
        "Videotron/Megatron",
        "Kain (Spanduk, Umbul, Baliho)",
        "Melekat/Stiker",
        "Selebaran",
        "Berjalan (kendaraan)",
        "Udara (balon gas)",
        "Apung (permukaan air)",
        "Film/Slide",
        "Peragaan",
    ], Color.ACCENT_BLUE),
    ("Dikecualikan", [
        "Internet, TV, radio, media cetak",
        "Label/merek pada produk",
        "Nama usaha ≤ 1 m² di tempat usaha",
        "Reklame Pemerintah/Pemda",
        "Nama tempat ibadah & panti asuhan",
        "Reklame tanah ≤ 1 m²",
        "Kegiatan politik (masa kampanye)",
        "Sosial/keagamaan (≤ 30 hari)",
        "Olahraga KONI (≤ 30 hari)",
    ], Color.TEAL),
])

content_slide("Subjek & Wajib Pajak (Pasal 3–4)", [
    "**Subjek Pajak Reklame** (Pasal 3):",
    "   Orang pribadi atau Badan yang menggunakan Reklame",
    "",
    "**Wajib Pajak Reklame** (Pasal 4):",
    "   Orang pribadi atau Badan yang menyelenggarakan Reklame",
    "",
    "Jika reklame diselenggarakan oleh **pihak ketiga**, maka pihak ketiga tersebut menjadi Wajib Pajak.",
])

# ═══════════════════════════════════════════════════════════════
# BAB III: MASA PAJAK, TAHUN PAJAK
# ═══════════════════════════════════════════════════════════════
section_slide("BAB III: MASA PAJAK, TAHUN PAJAK, & BAGIAN TAHUN PAJAK", "Pasal 5")

card_slide("Masa & Tahun Pajak — Pasal 5", [
    ("Masa Pajak", [
        "Permanen: 12 bulan atau sesuai masa penayangan",
        "Insidentil: per hari, maks. 30 hari",
    ], Color.ACCENT_BLUE),
    ("Tahun Pajak", [
        "1 tahun kalender",
        "Atau sesuai tahun buku WP",
    ], Color.TEAL),
    ("Bagian Tahun Pajak", [
        "1 tahun pajak atas",
        "1 bulan / beberapa bulan kalender",
    ], Color.WARM_GOLD),
])

# ═══════════════════════════════════════════════════════════════
# BAB IV: PENDAFTARAN DAN PENDATAAN
# ═══════════════════════════════════════════════════════════════
section_slide("BAB IV: PENDAFTARAN & PENDATAAN WAJIB PAJAK", "Pasal 6–8")

card_slide("Pendaftaran & Pendataan — Pasal 6–8", [
    ("Pendaftaran WP (Pasal 6)", [
        "WP wajib daftarkan diri & objek pajak",
        "Form: ambil di Bapenda / online / dikirim petugas",
        "Lampirkan: KTP, NPWP, Akta, NIB, Surat Kuasa",
        "Bapenda terbitkan NPWPD",
        "Jika tidak daftar → NPWPD jabatan",
        "Juga: NOPD, nomor registrasi",
    ], Color.ACCENT_BLUE),
    ("Pendataan (Pasal 7)", [
        "Bapenda lakukan pendataan WP & objek",
        "Tujuan: memperoleh, melengkapi data",
        "Termasuk informasi geografis objek",
        "Bisa kerjasama dengan instansi/pihak ketiga",
    ], Color.TEAL),
    ("Penonaktifan (Pasal 8)", [
        "Jika WP tak lagi penuhi syarat",
        "Jabatan atau atas permohonan WP",
        "Keputusan maks. 3 bulan",
        "Syarat: tak ada tunggakan & tak ada keberatan",
        "WP nonaktif/tutup wajib lapor",
    ], Color.ORANGE_SOFT),
])

# ═══════════════════════════════════════════════════════════════
# BAB V: NILAI SEWA REKLAME
# ═══════════════════════════════════════════════════════════════
section_slide("BAB V: NILAI SEWA REKLAME (NSR)", "Pasal 9 — Dasar Pengenaan Pajak")

content_slide("Faktor Penentu NSR — Pasal 9", [
    "Dasar pengenaan Pajak Reklame adalah **Nilai Sewa Reklame (NSR)**",
    "NSR = Nilai Jual Objek Pajak Reklame + Nilai Strategis Pemasangan",
    "",
    "**7 Faktor Perhitungan NSR:**",
    "   1️⃣  Jenis reklame",
    "   2️⃣  Bahan yang digunakan",
    "   3️⃣  Lokasi penempatan (Kelas Jalan)",
    "   4️⃣  Waktu (satuan detik)",
    "   5️⃣  Jangka waktu penyelenggaraan (hari kalender)",
    "   6️⃣  Jumlah media reklame (lembar)",
    "   7️⃣  Ukuran media reklame",
], sub_text="Pasal 9 ayat (1)–(2)")

content_slide("Klasifikasi Kelas Jalan — Pasal 9 ayat (3)–(8)", [
    "🏛️  **Kelas Jalan Khusus:**",
    "     • Zona Tol | Premium I | Premium II",
    "     • Premium I: Jl. A. Yani, Cut Mutia, Juanda, Sudirman, Sultan Agung, Transyogi, KH. Noer Ali",
    "     • Premium II: Jl. Narogong, Jatiwaringin, Pekayon, Jatiasih, Jatimakmur, Joyo Martono, Chairil Anwar, Bintara",
    "",
    "🚗  **Kelas Jalan I (Kendali Ketat):**",
    "     Lebar > 3 m, pusat pelayanan/permukiman (di luar kelas khusus)",
    "",
    "🏡  **Kelas Jalan II (Kendali Sedang):**",
    "     Lebar ≤ 3 m atau jalan lingkungan perumahan",
])

# ═══════════════════════════════════════════════════════════════
# BAB VI: PERHITUNGAN PAJAK REKLAME
# ═══════════════════════════════════════════════════════════════
section_slide("BAB VI: PERHITUNGAN PAJAK REKLAME", "Pasal 10 — Tarif & Rumus")

content_slide("Rumus & Tarif Perhitungan — Pasal 10", [
    "📐  **Rumus Dasar:**",
    "     Pajak Reklame = Tarif Pajak Reklame × NSR",
    "",
    "📐  **NSR Papan/Billboard & Megatron/Videotron:**",
    "     Nilai Kelas Jalan × Ukuran (m²) × Jumlah × Jangka Waktu",
    "",
    "⚡  **Ketentuan Tambahan:**",
    "     • Reklame indoor: NSR = **50%** dari hasil perhitungan",
    "     • Ketinggian > 15 m: tambahan **20%** dari NSR",
    "     • Produk tembakau & minuman keras: tambahan **50%** dari NSR",
    "     • Perubahan reklame: pajak atas selisih perubahan",
    "     • Perubahan naskah dalam 1 badan usaha: dikecualikan",
], sub_text="Pasal 10 ayat (1)–(13)")

premium_table_slide("NSR — Reklame Papan/Billboard & Sejenisnya",
    ["Kelas Jalan", "Zona", "NSR (Rp/m²/hari)"],
    [
        ["Kelas Jalan Khusus", "Jalan Tol", "Rp 23.575"],
        ["Kelas Jalan Khusus", "Premium 1 — A. Yani, Cut Mutia, Juanda, Sudirman, dll.", "Rp 16.100"],
        ["Kelas Jalan Khusus", "Premium 2 — Narogong, Jatiwaringin, Pekayon, dll.", "Rp 14.950"],
        ["Kelas Jalan I", "Kendali Ketat (lebar > 3 m)", "Rp 13.225"],
        ["Kelas Jalan II", "Kendali Sedang (lebar ≤ 3 m)", "Rp 11.500"],
    ],
    sub_text="Pasal 10 ayat (5) huruf a — Satuan: 1 m², 1 buah, 1 hari",
    col_widths=[Inches(2.5), Inches(6.5), Inches(2.7)]
)

premium_table_slide("NSR — Reklame Megatron/Videotron & Sejenisnya",
    ["Kelas Jalan", "Zona", "NSR (/30 detik)", "NSR (/m²/tahun)"],
    [
        ["Kelas Jalan Khusus", "Jalan Tol", "Rp 17,25", "Rp 13.599.900"],
        ["Kelas Jalan Khusus", "Premium 1", "Rp 13,80", "Rp 10.879.920"],
        ["Kelas Jalan Khusus", "Premium 2", "Rp 9,20", "Rp 7.253.280"],
        ["Kelas Jalan I", "Kendali Ketat", "Rp 8,05", "Rp 6.346.620"],
        ["Kelas Jalan II", "Kendali Sedang", "Rp 5,75", "Rp 4.533.300"],
    ],
    sub_text="Pasal 10 ayat (5) huruf b — Durasi 18 jam/hari = 2.160 tayangan/hari",
    col_widths=[Inches(2.5), Inches(4.5), Inches(2.5), Inches(2.7)]
)

premium_table_slide("NSR — Reklame Kain (Spanduk, Umbul, Baliho)",
    ["Kelas Jalan", "Zona", "NSR (Rp/m²/hari)"],
    [
        ["Kelas Jalan Khusus", "Jalan Tol", "Rp 30.000"],
        ["Kelas Jalan Khusus", "Premium 1", "Rp 30.000"],
        ["Kelas Jalan Khusus", "Premium 2", "Rp 25.000"],
        ["Kelas Jalan I", "Kendali Ketat", "Rp 20.000"],
        ["Kelas Jalan II", "Kendali Sedang", "Rp 19.000"],
    ],
    sub_text="Pasal 10 ayat (5) huruf c — Satuan: 1 m², 1 buah, 1 hari",
    col_widths=[Inches(2.5), Inches(6.5), Inches(2.7)]
)

card_slide("NSR — Jenis Reklame Lainnya — Pasal 10 ayat (5) huruf d", [
    ("Reklame Stiker", [
        "Rp 7,5/cm²",
        "Min. Rp 750.000/kali",
    ], Color.ACCENT_BLUE),
    ("Reklame Melekat", [
        "Rp 750.000/m²/tahun",
        "Dinding/tembok/bangunan",
    ], Color.TEAL),
    ("Reklame Selebaran", [
        "Rp 600/lembar",
        "Min. Rp 6.000.000/kali",
    ], Color.WARM_GOLD),
    ("Reklame Berjalan", [
        "Rp 6.000/m²/hari",
        "Termasuk kendaraan",
    ], Color.RED_ACCENT),
])

card_slide("NSR — Jenis Reklame Lainnya (lanjutan)", [
    ("Reklame Udara", [
        "Rp 2.400.000/sekali",
        "Maks. 1 bulan",
    ], Color.ACCENT_BLUE),
    ("Reklame Apung", [
        "Rp 600.000/sekali",
        "Maks. 1 bulan",
    ], Color.TEAL),
    ("Reklame Film/Slide", [
        "Rp 12.000/15 detik",
        "< 15\" dibulatkan 15\"",
    ], Color.WARM_GOLD),
    ("Reklame Peragaan", [
        "Rp 480.000/penyelenggaraan",
    ], Color.RED_ACCENT),
])

# ═══════════════════════════════════════════════════════════════
# BAB VII: PENETAPAN BESARAN PAJAK TERUTANG
# ═══════════════════════════════════════════════════════════════
section_slide("BAB VII: PENETAPAN BESARAN PAJAK TERUTANG", "Pasal 11")

content_slide("Penetapan Pajak Terutang — Pasal 11", [
    "**Dasar Penetapan:** Surat pendaftaran objek pajak → SKPD",
    "",
    "**3 Skema Penetapan:**",
    "   1️⃣ WP mendaftar → SKPD diterbitkan berdasarkan data WP",
    "   2️⃣ WP tidak mendaftar → SKPD diterbitkan secara **jabatan**",
    "   3️⃣ Hasil pemeriksaan menunjukkan pajak > laporan WP",
    "       → SKPD sesuai temuan (tanpa sanksi administratif)",
    "",
    "⏳  **Batas waktu:** Paling lama 5 tahun sejak terutangnya pajak",
])

# ═══════════════════════════════════════════════════════════════
# BAB VIII: SURAT TAGIHAN PAJAK
# ═══════════════════════════════════════════════════════════════
section_slide("BAB VIII: SURAT TAGIHAN PAJAK (STPD)", "Pasal 12")

card_slide("STPD — Surat Tagihan Pajak Daerah — Pasal 12", [
    ("Dasar Penerbitan", [
        "Jangka waktu: ≤ 5 tahun sejak pajak terutang",
        "a. Pajak SKPD tidak/kurang dibayar",
        "b. Keputusan/Keberatan tidak/kurang dibayar",
        "c. Sanksi administratif bunga/denda",
    ], Color.ACCENT_BLUE),
    ("Sanksi Bunga", [
        "Huruf a: 1%/bulan (maks. 24 bulan)",
        "Huruf b: 0,6%/bulan (maks. 24 bulan)",
        "Bagian bulan dihitung 1 bulan penuh",
    ], Color.RED_ACCENT),
])

# ═══════════════════════════════════════════════════════════════
# BAB IX: PEMBAYARAN DAN PENYETORAN
# ═══════════════════════════════════════════════════════════════
section_slide("BAB IX: PEMBAYARAN DAN PENYETORAN", "Pasal 13–14")

content_slide("Tata Cara Pembayaran — Pasal 13–14", [
    "**Metode:**",
    "   • Lunas melalui Kas Daerah / tempat yang ditunjuk Wali Kota",
    "   • Prioritas: sistem pembayaran **elektronik**",
    "   • Jika elektronik tidak tersedia: **tunai**",
    "",
    "⏰  **Jatuh tempo:** Paling lama 1 bulan sejak pengiriman SKPD",
    "",
    "⚠️  **Keterlambatan:**",
    "   • Sanksi bunga **1% per bulan** (maks. 24 bulan)",
    "   • Ditagih dengan STPD, lunas ≤ 30 hari sejak pengiriman",
    "",
    "🏷️  **Tanda bukti:** Stiker reklame + bukti dari Bank Persepsi",
])

# ═══════════════════════════════════════════════════════════════
# BAB X: PEMBETULAN DAN PEMBATALAN KETETAPAN
# ═══════════════════════════════════════════════════════════════
section_slide("BAB X: PEMBETULAN & PEMBATALAN KETETAPAN", "Pasal 15–20")

content_slide("Dasar Pembetulan — Pasal 15–16", [
    "**Kesalahan yang dapat dibetulkan/dibatalkan:**",
    "   a. ✏️  **Kesalahan tulis:** nama, alamat, NPWPD, nomor surat, masa pajak, tanggal jatuh tempo, SKPD ganda",
    "   b. 🔢  **Kesalahan hitung:** penjumlahan, pengurangan, perkalian, pembagian, tarif, klasifikasi objek",
    "   c. 📜  **Kekeliruan penerapan ketentuan** perundang-undangan",
    "",
    "**Syarat permohonan (Pasal 16):**",
    "   • 1 permohonan untuk 1 ketetapan",
    "   • Diajukan tertulis (B. Indonesia) oleh WP atau kuasa",
    "   • Lampirkan: identitas WP + kuasa + asli dokumen ketetapan",
    "   • Surat kuasa bermeterai jika dikuasakan",
])

content_slide("Prosedur & Jangka Waktu — Pasal 17–20", [
    "**Pasal 17:** Permohonan tidak memenuhi syarat → tidak dipertimbangkan",
    "",
    "**Pasal 18:**",
    "   • Kepala Bapenda wajib terbitkan Surat Keputusan Pembetulan dalam **6 bulan**",
    "   • Isi keputusan: mengabulkan (menambah/mengurangi/menghapus), membatalkan, atau menolak",
    "   • Jika > 6 bulan tanpa keputusan → **dianggap dikabulkan**",
    "",
    "**Pasal 19:** Pembetulan secara **jabatan** jika Kepala Bapenda tahu ada kesalahan",
    "",
    "**Pasal 20:** Pembetulan dapat dilakukan **berulang** jika masih ada kesalahan",
])

# ═══════════════════════════════════════════════════════════════
# BAB XI: PENGEMBALIAN KELEBIHAN PEMBAYARAN
# ═══════════════════════════════════════════════════════════════
section_slide("BAB XI: PENGEMBALIAN KELEBIHAN PEMBAYARAN", "Pasal 21")

content_slide("Prosedur Pengembalian — Pasal 21", [
    "**Hak WP:** Mengajukan permohonan pengembalian kelebihan pembayaran pajak",
    "",
    "**Alur:**",
    "   1️⃣ WP mengajukan permohonan ke Wali Kota/Kepala Bapenda",
    "   2️⃣ Keputusan diterbitkan maks. **12 bulan** sejak permohonan diterima",
    "   3️⃣ Jika > 12 bulan tanpa keputusan → permohonan **dianggap dikabulkan**",
    "   4️⃣ SKPDLB diterbitkan dalam 1 bulan",
    "",
    "**Jika WP memiliki utang pajak lain** → kelebihan diperhitungkan untuk melunasi utang",
    "",
    "💰  **Pengembalian** dilakukan maks. **2 bulan** sejak SKPDLB diterbitkan",
    "💰  Jika terlambat → imbalan bunga **0,6%/bulan**",
])

# ═══════════════════════════════════════════════════════════════
# BAB XII: PEMERIKSAAN PAJAK
# ═══════════════════════════════════════════════════════════════
section_slide("BAB XII: PEMERIKSAAN PAJAK", "Pasal 22–23")

card_slide("Pemeriksaan Pajak — Pasal 22–23", [
    ("Kewenangan & Tujuan", [
        "Kepala Bapenda berwenang melakukan pemeriksaan",
        "Menguji kepatuhan WP",
        "Tujuan lain: NPWPD jabatan, penghapusan, keberatan, penagihan",
        "Ruang lingkup: tahun berjalan &/ beberapa tahun sebelumnya",
    ], Color.ACCENT_BLUE),
    ("Kewajiban WP", [
        "Perlihatkan/pinjamkan buku, catatan, dokumen",
        "Beri akses ke tempat/ruangan",
        "Beri keterangan yang diperlukan",
        "Jika tidak dipenuhi → pajak ditetapkan secara jabatan",
    ], Color.TEAL),
    ("Hak WP saat diperiksa", [
        "Minta identitas & bukti penugasan pemeriksa",
        "Minta penjelasan alasan & tujuan pemeriksaan",
        "Terima dokumen hasil pemeriksaan",
        "Beri tanggapan atas hasil pemeriksaan",
    ], Color.WARM_GOLD),
])

# ═══════════════════════════════════════════════════════════════
# BAB XIII: PENAGIHAN PAJAK
# ═══════════════════════════════════════════════════════════════
section_slide("BAB XIII: PENAGIHAN PAJAK", "Pasal 24")

content_slide("Penagihan Pajak — Pasal 24", [
    "**Dasar Penagihan Pajak:**",
    "   SKPD | SKPDKB | SKPDKBT | STPD | SK Pembetulan | SK Keberatan | Putusan Banding",
    "",
    "**Prosedur:**",
    "   • Sebelum jatuh tempo → **imbauan**",
    "   • Setelah jatuh tempo & belum dibayar → **Penagihan Pajak**",
    "      sesuai ketentuan perundang-undangan",
])

# ═══════════════════════════════════════════════════════════════
# BAB XIV: KEDALUWARSA PENAGIHAN PAJAK
# ═══════════════════════════════════════════════════════════════
section_slide("BAB XIV: KEDALUWARSA PENAGIHAN PAJAK", "Pasal 25")

content_slide("Kedaluwarsa — Pasal 25", [
    "⏳  **Jangka waktu:** Hak penagihan kedaluwarsa setelah **5 tahun** sejak terutangnya pajak",
    "⏳  Kecuali WP melakukan **tindak pidana** di bidang perpajakan daerah",
    "",
    "⏸️  **Kedaluwarsa tertangguh jika:**",
    "   a. 📬  Diterbitkan Surat Teguran dan/atau Surat Paksa",
    "       — Dihitung sejak tanggal penyampaian surat",
    "   b. ✋  Ada **pengakuan utang pajak** dari WP",
    "       — Langsung: WP menyatakan masih punya utang",
    "       — Tidak langsung: pengajuan angsuran, penundaan, atau keberatan",
    "       — Dihitung sejak tanggal pengakuan",
])

# ═══════════════════════════════════════════════════════════════
# BAB XV: PENGHAPUSAN PIUTANG PAJAK
# ═══════════════════════════════════════════════════════════════
section_slide("BAB XV: PENGHAPUSAN PIUTANG PAJAK", "Pasal 26")

content_slide("Penghapusan Piutang — Pasal 26", [
    "Piutang pajak yang **tidak mungkin ditagih lagi** karena kedaluwarsa dapat dihapuskan",
    "",
    "**Prosedur penghapusan:**",
    "   1️⃣ Bapenda melakukan **penelitian** → Berita Acara hasil penelitian",
    "   2️⃣ Tim peneliti ditetapkan dengan **Keputusan Wali Kota**",
    "   3️⃣ Kepala Bapenda menyusun **daftar usulan penghapusan**",
    "   4️⃣ Disampaikan ke Wali Kota",
    "   5️⃣ Ditetapkan dengan **Keputusan Wali Kota**",
    "",
    "**Pertimbangan:**",
    "   • Pelaksanaan penagihan sampai batas kedaluwarsa (dibuktikan dengan dokumen)",
    "   • Hasil koordinasi dengan **Inspektorat Daerah**",
])

# ═══════════════════════════════════════════════════════════════
# BAB XVI: KERINGANAN, PENGURANGAN, PEMBEBASAN
# ═══════════════════════════════════════════════════════════════
section_slide("BAB XVI: KERINGANAN, PENGURANGAN, & PEMBEBASAN", "Pasal 27")

card_slide("Keringanan & Pembebasan — Pasal 27", [
    ("Bentuk Fasilitas", [
        "Keringanan",
        "Pengurangan",
        "Pembebasan",
        "Penundaan pembayaran",
        "Atas pokok &/ sanksi pajak",
    ], Color.ACCENT_BLUE),
    ("Kondisi WP", [
        "Kemampuan membayar",
        "Tingkat likuiditas WP",
    ], Color.TEAL),
    ("Kondisi Objek Pajak", [
        "Terkena bencana alam",
        "Kebakaran",
        "Huru-hara / kerusuhan",
    ], Color.RED_ACCENT),
])

# ═══════════════════════════════════════════════════════════════
# BAB XVII: KEMUDAHAN PERPAJAKAN DAERAH
# ═══════════════════════════════════════════════════════════════
section_slide("BAB XVII: KEMUDAHAN PERPAJAKAN DAERAH", "Pasal 28")

content_slide("Kemudahan Perpajakan — Pasal 28", [
    "**Bentuk Kemudahan:**",
    "   a. 📅  **Perpanjangan batas waktu** pembayaran pajak — untuk WP keadaan kahar",
    "   b. 💳  **Fasilitas angsuran/penundaan** — untuk WP kesulitan likuiditas",
    "",
    "**Keadaan Kahar:**",
    "   🌊 Bencana alam | 🔥 Kebakaran | 🚨 Kerusuhan massal | 🦠 Wabah penyakit | Keadaan lain (keputusan Wali Kota)",
    "",
    "**Ketentuan:**",
    "   • Jangka waktu angsuran/penundaan: maks. **24 bulan**",
    "   • Bunga: **0,6%/bulan** dari jumlah pajak yang masih harus dibayar",
    "   • Wali Kota memperhatikan kepatuhan WP dalam **2 tahun terakhir**",
])

# ═══════════════════════════════════════════════════════════════
# BAB XVIII: KEBERATAN DAN BANDING
# ═══════════════════════════════════════════════════════════════
section_slide("BAB XVIII: KEBERATAN DAN BANDING", "Pasal 29–33")

card_slide("Keberatan — Pasal 29–31", [
    ("Objek Keberatan", [
        "SKPD, SKPDKB, SKPDKBT",
        "SKPDLB, SKPDN",
        "Pemotongan/pemungutan pihak ke-3",
    ], Color.ACCENT_BLUE),
    ("Syarat & Waktu", [
        "Tertulis B. Indonesia, maks. 3 bulan",
        "WP sudah bayar min. jumlah yg disetujui",
        "Keputusan maks. 12 bulan",
        "Jika tanpa keputusan → dikabulkan",
    ], Color.TEAL),
    ("Sanksi", [
        "Ditolak: denda 30%",
        "Dikabulkan: kelebihan + bunga 0,6%/bln",
        "Jangka waktu bunga maks. 24 bulan",
    ], Color.RED_ACCENT),
])

card_slide("Banding — Pasal 32–33", [
    ("Objek Banding", [
        "Surat Keputusan Keberatan",
        "Ke badan peradilan pajak",
        "Maks. 3 bulan sejak keputusan diterima",
    ], Color.ACCENT_BLUE),
    ("Efek Banding", [
        "Menangguhkan kewajiban bayar",
        "Sampai 1 bulan sejak Putusan Banding",
    ], Color.TEAL),
    ("Sanksi & Imbalan", [
        "Ditolak: denda 60%",
        "Dikabulkan: kelebihan + bunga 0,6%/bln",
        "Sanksi 30% tidak dikenakan jika banding",
    ], Color.RED_ACCENT),
])

# ═══════════════════════════════════════════════════════════════
# BAB XIX: PENGHARGAAN
# ═══════════════════════════════════════════════════════════════
section_slide("BAB XIX: PENGHARGAAN", "Pasal 34–35")

content_slide("Penghargaan Wajib Pajak Taat — Pasal 34–35", [
    "🏆  **Bentuk Penghargaan:**",
    "     Piagam penghargaan | Hadiah | dan/atau sejenisnya",
    "     Dibebankan pada APBD",
    "",
    "✅  **Kriteria Wajib Pajak Taat:**",
    "     a. Menyetorkan pajak tepat waktu minimal **1 tahun**",
    "     b. Tidak ada tunggakan **3 tahun terakhir**",
    "     c. Kontribusi signifikan dalam program Pemerintah Daerah",
    "",
    "📋  **Penetapan:**",
    "     • Keputusan Wali Kota",
    "     • Tim penilai dibentuk oleh Wali Kota",
])

# ═══════════════════════════════════════════════════════════════
# BAB XX: KETENTUAN PENUTUP
# ═══════════════════════════════════════════════════════════════
section_slide("BAB XX: KETENTUAN PENUTUP", "Pasal 36–37")

content_slide("Pencabutan & Mulai Berlaku — Pasal 36–37", [
    "**Pasal 36 — Pencabutan:**",
    "   Peraturan Wali Kota yang dicabut:",
    "   ❌  Perwal Bekasi No. 48 Tahun 2012 tentang Petunjuk Pelaksanaan Perda No. 14/2012",
    "   ❌  Perwal Bekasi No. 52 Tahun 2013 tentang Perubahan Perwal No. 48/2012",
    "",
    "**Pasal 37 — Mulai Berlaku:**",
    "   ✅  Peraturan Wali Kota ini mulai berlaku pada tanggal **diundangkan**",
    "",
    "━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━",
    "📅  Ditetapkan di Bekasi, 20 Desember 2024",
    "",
    "**Pj. WALI KOTA BEKASI,**",
    "ttd",
    "",
    "**R. GANI MUHAMAD**",
    "",
    "📅  Diundangkan di Bekasi, 20 Desember 2024",
    "**SEKRETARIS DAERAH KOTA BEKASI,**",
    "ttd",
    "**JUNAEDI**",
])

# ═══════════════════════════════════════════════════════════════
# CLOSING SLIDE
# ═══════════════════════════════════════════════════════════════
slide_counter[0] += 1
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_shape(slide, MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height, Color.DARK_NAVY)
# Decorative shapes
add_shape(slide, MSO_SHAPE.OVAL, Inches(9), -Inches(2), Inches(6), Inches(6),
          fill_color=RGBColor(0x0D, 0x2B, 0x5A))
add_shape(slide, MSO_SHAPE.OVAL, Inches(-2), Inches(4), Inches(5), Inches(5),
          fill_color=RGBColor(0x0D, 0x2B, 0x5A))
add_shape(slide, MSO_SHAPE.RECTANGLE, Inches(1), Inches(3.2), Inches(4), Pt(4), Color.GOLD)

add_textbox(slide, Inches(1), Inches(1.5), Inches(10), Inches(0.5),
            "BERITA DAERAH KOTA BEKASI", font_size=16, color=Color.GOLD, bold=True)

add_textbox(slide, Inches(1), Inches(2.2), Inches(11), Inches(1.5),
            "TERIMA KASIH", font_size=52, bold=True, color=Color.WHITE)

add_textbox(slide, Inches(1), Inches(3.8), Inches(11), Inches(1.5),
            "Peraturan Wali Kota Bekasi Nomor 51 Tahun 2024\n"
            "Tentang Pengelolaan Pajak Reklame",
            font_size=20, color=RGBColor(0xAA, 0xBB, 0xDD))

add_textbox(slide, Inches(1), Inches(5.0), Inches(11), Inches(0.8),
            "Sumber: https://jdih.bekasikota.go.id",
            font_size=13, color=RGBColor(0x88, 0x99, 0xBB))

add_shape(slide, MSO_SHAPE.RECTANGLE, 0, Inches(7.2), prs.slide_width, Inches(0.3), Color.GOLD)
add_page_number(slide, slide_counter[0])

# ═══════════════════════════════════════════════════════════════
# SAVE
# ═══════════════════════════════════════════════════════════════
output_path = "/home/runner/coba_ppt/Perwal_Bekasi_51_2024_Pajak_Reklame.pptx"
prs.save(output_path)
print(f"✅ Premium presentation saved to: {output_path}")
print(f"   Total slides: {len(prs.slides)}")
print(f"   File size: {__import__('os').path.getsize(output_path) / 1024:.1f} KB")
