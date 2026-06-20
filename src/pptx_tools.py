#!/usr/bin/env python3
"""
pptx_tools.py — Template Editing (Unpack / Edit / Pack)
=========================================================
Workflow:
    1. python pptx_tools.py unpack input.pptx unpacked/
    2. Edit XML files in unpacked/ (text, formatting, etc.)
    3. python pptx_tools.py pack unpacked/ output.pptx --original input.pptx

Atau list slide content:
    python pptx_tools.py list input.pptx

Dependencies: python-pptx (sudah included)
"""

import sys, os, zipfile, shutil, tempfile, re, glob, xml.dom.minidom
from pathlib import Path


# ═══════════════════════════════════════════════════════════
# UNPACK — extract PPTX to directory with pretty-printed XML
# ═══════════════════════════════════════════════════════════

def unpack(pptx_path, output_dir):
    """
    Extract PPTX to directory, pretty-printing all XML files.
    
    Args:
        pptx_path: path to source .pptx
        output_dir: target directory
    """
    pptx_path = os.path.abspath(pptx_path)
    output_dir = os.path.abspath(output_dir)
    
    if not os.path.exists(pptx_path):
        print(f"❌ File not found: {pptx_path}")
        return False
    
    # Remove existing output dir
    if os.path.exists(output_dir):
        shutil.rmtree(output_dir)
    
    # Extract ZIP
    with zipfile.ZipFile(pptx_path, 'r') as z:
        z.extractall(output_dir)
    
    # Pretty-print XML files
    xml_count = 0
    for root, dirs, files in os.walk(output_dir):
        for fname in files:
            if fname.endswith('.xml') or fname.endswith('.rels'):
                fpath = os.path.join(root, fname)
                try:
                    dom = xml.dom.minidom.parse(fpath)
                    pretty = dom.toprettyxml(indent="  ", encoding='utf-8')
                    # minidom adds extra newlines — clean them up
                    pretty = '\n'.join(line for line in pretty.decode('utf-8').split('\n')
                                       if line.strip() or line == '\n')
                    # Restore self-closing tags (minidom expands them)
                    pretty = re.sub(r'></\w+>', '/>', pretty)
                    with open(fpath, 'w', encoding='utf-8') as f:
                        f.write(pretty)
                    xml_count += 1
                except Exception:
                    pass  # skip binary/non-XML files
    
    # List slide files
    slides = sorted(glob.glob(os.path.join(output_dir, "ppt/slides/slide*.xml")))
    
    print(f"📦 Unpacked: {pptx_path} → {output_dir}/")
    print(f"   {xml_count} XML files pretty-printed")
    print(f"   {len(slides)} slides found")
    for s in slides:
        print(f"     {os.path.relpath(s, output_dir)}")
    
    return True


# ═══════════════════════════════════════════════════════════
# LIST — show text content per slide
# ═══════════════════════════════════════════════════════════

def list_slides(pptx_path):
    """Print text content of each slide."""
    from pptx import Presentation
    
    prs = Presentation(pptx_path)
    print(f"📋 Slides in: {pptx_path}")
    print(f"   Dimensions: {prs.slide_width.inches:.2f}\" × {prs.slide_height.inches:.2f}\"")
    print(f"   {len(prs.slides)} slides total\n")
    
    for i, slide in enumerate(prs.slides, 1):
        texts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    t = para.text.strip()
                    if t and len(t) > 1:
                        texts.append(t)
            if shape.has_table:
                for row in shape.table.rows:
                    for cell in row.cells:
                        t = cell.text.strip()
                        if t and len(t) > 1:
                            texts.append(t)
        
        print(f"── Slide {i:02d} ────────────────────────")
        for t in texts:
            print(f"  {t[:120]}{'…' if len(t) > 120 else ''}")
        if not texts:
            print(f"  (no text)")
        print()


# ═══════════════════════════════════════════════════════════
# PACK — repack directory to PPTX
# ═══════════════════════════════════════════════════════════

def pack(source_dir, output_path, original=None):
    """
    Repack directory into PPTX file.
    
    Args:
        source_dir: directory with unpacked PPTX content
        output_path: path for output .pptx
        original: optional original PPTX for reference
    """
    source_dir = os.path.abspath(source_dir)
    output_path = os.path.abspath(output_path)
    
    if not os.path.exists(source_dir):
        print(f"❌ Source dir not found: {source_dir}")
        return False
    
    # Use temp file for atomic write
    tmp_path = output_path + ".tmp"
    
    try:
        with zipfile.ZipFile(tmp_path, 'w', zipfile.ZIP_DEFLATED) as z:
            for root, dirs, files in os.walk(source_dir):
                for fname in files:
                    fpath = os.path.join(root, fname)
                    arcname = os.path.relpath(fpath, source_dir)
                    z.write(fpath, arcname)
        
        # Replace output
        if os.path.exists(output_path):
            os.remove(output_path)
        os.rename(tmp_path, output_path)
        
        size_kb = os.path.getsize(output_path) / 1024
        print(f"📦 Packed: {source_dir}/ → {output_path}")
        print(f"   {size_kb:.0f} KB")
        
        return True
        
    except Exception as e:
        print(f"❌ Pack failed: {e}")
        if os.path.exists(tmp_path):
            os.remove(tmp_path)
        return False


# ═══════════════════════════════════════════════════════════
# HELP
# ═══════════════════════════════════════════════════════════

def main():
    if len(sys.argv) < 2:
        print(__doc__)
        sys.exit(1)
    
    command = sys.argv[1]
    
    if command == "unpack" and len(sys.argv) >= 4:
        unpack(sys.argv[2], sys.argv[3])
    elif command == "list" and len(sys.argv) >= 3:
        list_slides(sys.argv[2])
    elif command == "pack" and len(sys.argv) >= 4:
        orig = sys.argv[4] if len(sys.argv) >= 5 else None
        pack(sys.argv[2], sys.argv[3], orig)
    elif command in ("-h", "--help", "help"):
        print(__doc__)
    else:
        print(f"❌ Unknown command: {command}")
        print()
        print("Usage:")
        print("  python pptx_tools.py unpack input.pptx unpacked/")
        print("  python pptx_tools.py list input.pptx")
        print("  python pptx_tools.py pack unpacked/ output.pptx [original.pptx]")
        sys.exit(1)


if __name__ == "__main__":
    main()
