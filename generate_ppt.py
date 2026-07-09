#!/usr/bin/env python3
"""
generate_ppt.py — RESUME TUPOKSI e-Government
==============================================
Pakai engine src/ppt_engine.py.
Icon pake SHAPE (bukan text/unicode) — render di mana aja.
NO circles/ovals — diamond, triangle, square, chevron.
"""

import os, sys
sys.path.insert(0, os.path.join(os.path.dirname(__file__), "src"))

from pptx.dml.color import RGBColor
from ppt_engine import Engine, MSO_SHAPE, Pt, Inches
from pptx.enum.text import PP_ALIGN as PA

# ─── Color Palette ───────────────────────────────────────────────
DARK    = RGBColor(0x0A, 0x16, 0x28)
MID     = RGBColor(0x0F, 0x2A, 0x4A)
BLUE    = RGBColor(0x1A, 0x56, 0x76)
TEAL    = RGBColor(0x0D, 0x94, 0x88)
TEAL_L  = RGBColor(0xE6, 0xF7, 0xF5)
GOLD    = RGBColor(0xF5, 0x9E, 0x0B)
GOLD_L  = RGBColor(0xFE, 0xF3, 0xC7)
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
OFF_W   = RGBColor(0xF8, 0xFA, 0xFC)
BORDER  = RGBColor(0xE2, 0xE8, 0xF0)
TDARK   = RGBColor(0x1E, 0x29, 0x3B)
TMID    = RGBColor(0x47, 0x55, 0x69)
TLIGHT  = RGBColor(0x94, 0xA3, 0xB8)


# ─── Shape Helpers ───────────────────────────────────────────────

def add_shape(slide, shape_type, l, t, w, h, fill, rotation=0):
    s = slide.shapes.add_shape(shape_type, Inches(l), Inches(t),
                                Inches(w), Inches(h))
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    s.line.fill.background()
    if rotation:
        s.rotation = rotation
    return s


def accent_bar(slide, l, t, w, h, color):
    return add_shape(slide, MSO_SHAPE.RECTANGLE, l, t, w, h, color)


def card(slide, l, t, w, h):
    s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                Inches(l), Inches(t),
                                Inches(w), Inches(h))
    s.fill.solid()
    s.fill.fore_color.rgb = WHITE
    s.line.color.rgb = BORDER
    s.line.width = Pt(0.5)
    s.adjustments[0] = 0.06
    return s


def txt(slide, l, t, w, h, text, size=12, bold=False, color=None,
        align=PA.LEFT):
    if color is None: color = TDARK
    tb = slide.shapes.add_textbox(Inches(l), Inches(t),
                                   Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text; p.font.size = Pt(size)
    p.font.bold = bold; p.font.color.rgb = color
    p.font.name = "Calibri"; p.alignment = align
    return tb


def multi_txt(slide, l, t, w, h, lines, size=12, color=None,
              spacing=2, first_bold=False, first_size=None):
    if color is None: color = TDARK
    tb = slide.shapes.add_textbox(Inches(l), Inches(t),
                                   Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.font.size = Pt(first_size if (i == 0 and first_size) else size)
        p.font.bold = (first_bold and i == 0)
        p.font.color.rgb = color; p.font.name = "Calibri"
        p.space_after = Pt(spacing)
    return tb


# ─── Shape Icons ─────────────────────────────────────────────────
# Each returns (left_offset_adjustment)

def icon_app(slide, l, t, s, color):
    """App icon: rounded square with inner square."""
    # Outer rounded square
    add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, l, t, s, s, color)
    # Inner square (lighter)
    p = 0.2
    add_shape(slide, MSO_SHAPE.RECTANGLE, l+p, t+p, s-2*p, s-2*p, WHITE)
    # Center dot
    ds = 0.15
    add_shape(slide, MSO_SHAPE.RECTANGLE, l+(s-ds)/2, t+(s-ds)/2, ds, ds, color)
    return l + s + 0.15


def icon_gear(slide, l, t, s, color):
    """Gear/SPBE icon: hexagon with center dot."""
    add_shape(slide, MSO_SHAPE.PENTAGON, l, t, s, s, color)
    add_shape(slide, MSO_SHAPE.RECTANGLE, l+s*0.3, t+s*0.3, s*0.4, s*0.4, WHITE)
    return l + s + 0.15


def icon_city(slide, l, t, s, color):
    """Building/city icon: rectangle + triangle roof."""
    add_shape(slide, MSO_SHAPE.RECTANGLE, l, t+s*0.3, s, s*0.7, color)
    add_shape(slide, MSO_SHAPE.ISOSCELES_TRIANGLE, l, t, s, s*0.5, color)
    # Window
    add_shape(slide, MSO_SHAPE.RECTANGLE, l+s*0.3, t+s*0.5, s*0.15, s*0.15, WHITE)
    return l + s + 0.15


def icon_web(slide, l, t, s, color):
    """Web/globe icon: diamond with lines."""
    add_shape(slide, MSO_SHAPE.DIAMOND, l, t, s, s, color)
    # Horizontal line
    add_shape(slide, MSO_SHAPE.RECTANGLE, l, t+s*0.42, s, s*0.16, WHITE)
    return l + s + 0.15


def icon_phone(slide, l, t, s, color):
    """Phone/mobile icon: vertical rect with screen."""
    w = s * 0.55; h = s
    x = l + (s - w)/2
    add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, x, t, w, h, color)
    # Screen area
    sc = 0.12
    add_shape(slide, MSO_SHAPE.RECTANGLE, x+sc, t+sc*2, w-2*sc, h-sc*3.5, WHITE)
    # Button at bottom
    add_shape(slide, MSO_SHAPE.RECTANGLE, x+w*0.3, t+h-sc, w*0.4, sc*0.8, WHITE)
    return l + s + 0.15


def icon_link(slide, l, t, s, color):
    """Link/chain icon: two chevrons or parallelograms."""
    add_shape(slide, MSO_SHAPE.PARALLELOGRAM, l, t, s*0.7, s, color)
    add_shape(slide, MSO_SHAPE.PARALLELOGRAM, l+s*0.3, t, s*0.7, s, WHITE)
    return l + s + 0.15


def icon_check(slide, l, t, s, color):
    """Check icon: rounded square with inner diamond."""
    add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, l, t, s, s, color)
    add_shape(slide, MSO_SHAPE.DIAMOND, l+s*0.2, t+s*0.2, s*0.6, s*0.6, WHITE)
    return l + s + 0.15


def icon_star(slide, l, t, s, color):
    """Star/diamond icon."""
    add_shape(slide, MSO_SHAPE.DIAMOND, l, t, s, s, color)
    add_shape(slide, MSO_SHAPE.DIAMOND, l+s*0.15, t+s*0.15, s*0.7, s*0.7, WHITE)
    return l + s + 0.15


def icon_arrow(slide, l, t, s, color):
    """Right arrow."""
    add_shape(slide, MSO_SHAPE.RIGHT_ARROW, l, t, s, s, color)
    return l + s + 0.15


def icon_box(slide, l, t, s, color):
    """Simple box with number inside."""
    add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, l, t, s, s, color)
    return l + s + 0.15


# ─── Layout Helpers ──────────────────────────────────────────────

def build():
    engine = Engine()
    L = engine.L

    from pptx import Presentation
    prs = Presentation()
    prs.slide_width = Inches(L.SLIDE_W)
    prs.slide_height = Inches(L.SLIDE_H)

    MX = L.MARGIN_H
    CW = L.cw
    PG = [0]

    def ns():
        return prs.slides.add_slide(prs.slide_layouts[6])

    def solid_bg(sl, c):
        bg = sl.background; bg.fill.solid(); bg.fill.fore_color.rgb = c

    def content_slide(title, sub=None):
        sl = ns()
        solid_bg(sl, WHITE)
        accent_bar(sl, 0, 0, L.SLIDE_W, 0.04, GOLD)
        accent_bar(sl, 0, 0.04, L.SLIDE_W, 1.15, DARK)
        accent_bar(sl, 0, 1.19, L.SLIDE_W, 0.03, TEAL)
        txt(sl, MX+0.15, 0.18, CW-0.3, 0.55, title, 22, bold=True, color=WHITE)
        if sub:
            txt(sl, MX+0.15, 0.72, CW-0.3, 0.35, sub, 9, color=TEAL_L)
        return sl

    def ft(sl):
        accent_bar(sl, 0, L.SLIDE_H-0.45, L.SLIDE_W, 0.45, OFF_W)
        accent_bar(sl, 0, L.SLIDE_H-0.45, L.SLIDE_W, Pt(1.5), BORDER)
        PG[0] += 1
        txt(sl, MX, L.SLIDE_H-0.38, 4, 0.22, "Sumber: RESUME TUPOKSI egov.docx", 7, color=TLIGHT)
        txt(sl, L.SLIDE_W-1.0, L.SLIDE_H-0.38, 0.8, 0.22, str(PG[0]), 8, color=TLIGHT, align=PA.RIGHT)

    # ══════════════════════════════════════════════════════════════
    #  SLIDE 1 — COVER
    # ══════════════════════════════════════════════════════════════
    sl = ns()
    solid_bg(sl, DARK)
    accent_bar(sl, 0, 0, L.SLIDE_W, 0.04, GOLD)
    accent_bar(sl, L.SLIDE_W-5, -0.5, 5.5, 5, MID)
    accent_bar(sl, L.SLIDE_W-3.8, 3.5, 4.5, 4.5, MID)
    accent_bar(sl, 0, L.SLIDE_H-0.28, L.SLIDE_W, 0.28, TEAL)

    # Large decorative app icon
    icon_app(sl, MX+0.3, 1.5, 0.8, TEAL)
    txt(sl, MX+1.5, 1.6, CW-2, 0.8, "RESUME TUPOKSI", 46, bold=True, color=WHITE)
    txt(sl, MX, 2.9, CW, 0.5, "Bidang e-Government", 22, color=TEAL_L)
    txt(sl, MX, 3.5, CW, 0.4, "Diskominfostandi Kota Bekasi", 14, color=TLIGHT)
    accent_bar(sl, MX, 4.1, 2.5, Pt(3), GOLD)
    txt(sl, MX, 4.4, CW, 0.3, "2025", 11, color=TLIGHT)
    PG[0] += 1

    # ══════════════════════════════════════════════════════════════
    #  SLIDE 2 — 3 PILAR with shape icons
    # ══════════════════════════════════════════════════════════════
    sl = content_slide("Tiga Pilar Strategis e-Government",
                       "Bidang e-Government menggerakkan 3 pilar utama transformasi digital Kota Bekasi")

    pillars_data = [
        ("01", icon_app,  "Pengembangan Aplikasi\n& Sistem Informasi",
         "Mengawal standar & kualitas\npengembangan aplikasi serta\nmenghubungkan layanan daerah\nmelalui SPLP", DARK),
        ("02", icon_gear, "Tata Kelola SPBE",
         "Merancang roadmap, arsitektur &\nkebijakan SPBE sekaligus\nmemperkuat kapasitas SDM\ndan peran Government CIO", BLUE),
        ("03", icon_city, "Pengembangan\nKota Cerdas",
         "Menyusun masterplan, membangun\nkolaborasi lintas sektor, serta\nmengevaluasi program Smart City\nsecara berkelanjutan", TEAL),
    ]

    n = len(pillars_data)
    cw = L.col_width(n, 0.35)
    gap = 0.35
    sx = (L.SLIDE_W - (n*cw + (n-1)*gap)) / 2
    ch = 4.2
    cy = 1.5

    for i, (num, icon_fn, title, desc, clr) in enumerate(pillars_data):
        cx = sx + i * (cw + gap)
        card(sl, cx, cy, cw, ch)
        accent_bar(sl, cx, cy, cw, 0.08, clr)

        # Shape icon
        icon_size = 0.55
        icon_fn(sl, cx+0.25, cy+0.3, icon_size, clr)

        # Number
        txt(sl, cx+0.95, cy+0.38, 0.5, 0.3, f"0{i+1}", 14, bold=True, color=clr)

        # Title
        txt(sl, cx+0.25, cy+1.15, cw-0.5, 0.85, title, 17, bold=True, color=DARK)

        # Separator
        accent_bar(sl, cx+0.25, cy+2.1, cw-0.5, Pt(1.5), BORDER)

        # Description
        txt(sl, cx+0.25, cy+2.3, cw-0.5, 1.6, desc, 11, color=TMID)

    ft(sl)

    # ══════════════════════════════════════════════════════════════
    #  SLIDE 3 — SISTEM & APLIKASI with shape icons
    # ══════════════════════════════════════════════════════════════
    sl = content_slide("Empat Sistem & Aplikasi Strategis",
                       "Layanan digital yang dioperasikan dan dipelihara oleh Bidang e-Government")

    systems_data = [
        (icon_web,   "Web Pemerintah Kota",   "Portal utama pemerintah kota\nsebagai pusat info & layanan publik", DARK),
        (icon_phone, "Mobile App PSW",         "Aplikasi mobile Pekan Smart City\n— layanan kota cerdas terintegrasi", BLUE),
        (icon_city,  "Web Kota Cerdas",        "Portal informasi & layanan\nprogram Smart City Kota Bekasi", TEAL),
        (icon_link,  "Aplikasi SPLP",          "Sistem Penghubung Layanan\nPemerintah — terintegrasi Kemenkomdigi", GOLD),
    ]

    cw2 = L.col_width(2, 0.4)
    rh2 = 1.8
    gh2 = 0.4; gv2 = 0.35
    gsx = (L.SLIDE_W - (2*cw2+gh2))/2
    acy = 1.55

    for i, (icon_fn, title, desc, clr) in enumerate(systems_data):
        col = i % 2; row = i // 2
        cx = gsx + col*(cw2+gh2); cy = acy + row*(rh2+gv2)
        card(sl, cx, cy, cw2, rh2)
        icon_fn(sl, cx+0.25, cy+0.3, 0.6, clr)
        txt(sl, cx+1.1, cy+0.3, cw2-1.4, 0.45, title, 18, bold=True, color=clr)
        txt(sl, cx+1.1, cy+0.85, cw2-1.4, 0.7, desc, 12, color=TMID)

    ft(sl)

    # ══════════════════════════════════════════════════════════════
    #  SLIDE 4-6 — DETAIL with icon per item
    # ══════════════════════════════════════════════════════════════

    details = [
        ("Supervisi & Pengembangan Aplikasi Daerah",
         "Lingkup kerja pengembangan sistem informasi dan aplikasi perangkat daerah",
         "✓",
         [icon_check, icon_link, icon_web, icon_gear],
         ["Supervisi, analisis, dan standarisasi\npengembangan aplikasi perangkat daerah",
          "Mengelola & mengembangkan SPLP\nsebagai tulang punggung integrasi layanan",
          "Pengelolaan domain & subdomain\npemerintah daerah secara terpusat",
          "Sosialisasi & peningkatan kapasitas\nSDM pengelola sistem informasi"]),
        ("Akselerasi Program Kota Cerdas",
         "Strategi dan kolaborasi menuju Smart City yang terintegrasi",
         "◆",
         [icon_star, icon_city, icon_app, icon_link],
         ["Menyusun strategi, rencana aksi,\ndan masterplan Kota Cerdas",
          "Membangun kolaborasi dengan pemangku\nkepentingan lintas sektor",
          "Mengoordinasikan & mengevaluasi\nprogram Kota Cerdas secara berkala",
          "Menyelaraskan rencana induk dengan\ndokumen perencanaan daerah"]),
        ("Penguatan Tata Kelola SPBE",
         "Kerangka kerja menuju SPBE yang matang",
         "⚙",
         [icon_gear, icon_arrow, icon_app, icon_star],
         ["Menyusun strategi, roadmap, arsitektur,\ndan peta rencana SPBE",
          "Melaksanakan & mengoordinasikan\nprogram SPBE lintas perangkat daerah",
          "Mengembangkan kebijakan &\ntata kelola SPBE yang adaptif",
          "Monitoring, evaluasi, & rekomendasi\nperbaikan SPBE berkelanjutan"]),
    ]

    accent_colors = [DARK, TEAL, BLUE]

    for si, (title, sub, label, icons, items) in enumerate(details):
        sl = content_slide(title, sub)
        accent = accent_colors[si]

        # Left decorative panel
        accent_bar(sl, MX, 1.5, 0.1, 5.2, accent)

        for i, (icon_fn, item) in enumerate(zip(icons, items)):
            iy = 1.8 + i * 1.15

            # Icon
            icon_fn(sl, MX+0.35, iy, 0.5, accent)

            # Vertical connector
            if i < len(icons)-1:
                accent_bar(sl, MX+0.58, iy+0.6, Pt(1.5), 0.5, BORDER)

            # Text
            lines = item.split('\n')
            multi_txt(sl, MX+1.1, iy+0.04, CW-1.4, 0.8, lines,
                      13, color=TDARK, spacing=1, first_bold=True, first_size=14)

        ft(sl)

    # ══════════════════════════════════════════════════════════════
    #  SLIDE 7 — CLOSING
    # ══════════════════════════════════════════════════════════════
    sl = ns()
    solid_bg(sl, DARK)
    accent_bar(sl, 0, 0, L.SLIDE_W, 0.04, GOLD)
    accent_bar(sl, L.SLIDE_W-4.5, -0.5, 5, 5, MID)
    accent_bar(sl, 0, L.SLIDE_H-0.28, L.SLIDE_W, 0.28, TEAL)

    # Big star icon
    icon_star(sl, L.SLIDE_W/2-0.4, 1.2, 0.8, TEAL)

    txt(sl, MX, 2.5, CW, 1.0, "Terima Kasih", 44, bold=True, color=WHITE, align=PA.CENTER)
    txt(sl, MX, 3.6, CW, 0.5, "Bidang e-Government — Diskominfostandi Kota Bekasi",
        18, color=TEAL_L, align=PA.CENTER)
    accent_bar(sl, L.SLIDE_W/2-1.0, 4.3, 2.0, Pt(2), GOLD)
    txt(sl, MX, 4.6, CW, 0.5,
        "\"Mewujudkan tata kelola pemerintahan yang cerdas,\nterpadu, dan berkelanjutan\"",
        12, color=TLIGHT, align=PA.CENTER)
    PG[0] += 1

    # ── Save ──
    out = "RESUME_TUPOKSI_egov.pptx"
    prs.save(out)
    print(f"✅ PPT selesai: {out}")
    print(f"   {len(prs.slides)} slide — all icons are SHAPES (not text)")
    return out


if __name__ == "__main__":
    build()
