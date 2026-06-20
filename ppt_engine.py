#!/usr/bin/env python3
"""
ppt_engine.py — Reusable Presentation Engine
=============================================
LayoutFrame class + draw helpers + generic archetype functions.
Terpisah dari content — bisa dipakai untuk PPT apa saja.

Arsitektur:
    ppt_engine.py       ← engine reusable (file ini)
    content_*.py        ← data konten terpisah
    generate_ppt.py     ← entry point

Usage:
    from ppt_engine import Engine
    engine = Engine()
    prs = engine.build(slides_data)
    prs.save("output.pptx")
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import math, os, colorsys


# ═══════════════════════════════════════════════════════════
# COLOR UTILITY
# ═══════════════════════════════════════════════════════════

class Colors:
    """
    Palette 60-30-10 — WCAG AA verified (no color bias).
    
    ⚡ PERUBAHAN dari palette sebelumnya:
      TEAL   #0D9488 → #0B7C72  (digelapkan: 3.7→5.1:1 on white ✅)
      WARM   #B8860B → #A0522D  (ganti hue: gold clash → sienna ~20°)
      TEXT_M #6B7288 → #8899B0  (dicerahkan: 3.8→6.2:1 on navy ✅)
      TEXT_L #9CA3AF → #64748B  (digelapkan: 2.5→4.8:1 on white ✅)
    """
    # ── 60% Dominant ──
    NAVY   = RGBColor(0x0A, 0x16, 0x28)  # bg header, cover, section
    NAVY_L = RGBColor(0x12, 0x29, 0x4A)  # decorative oval (lighter)
    NAVY_M = RGBColor(0x1B, 0x3A, 0x6B)  # medium navy
    NAVY_D = RGBColor(0x0D, 0x1F, 0x3C)  # decorative oval (darker)
    
    # ── 30% Secondary ──
    WHITE  = RGBColor(0xFF, 0xFF, 0xFF)  # card bg, text on dark
    OFF_W  = RGBColor(0xF5, 0xF7, 0xFA)  # content slide bg
    ICE    = RGBColor(0xE8, 0xED, 0xF5)  # note box, table stripe
    ICE_D  = RGBColor(0xD0, 0xD8, 0xE8)  # darker ice
    
    # ── 10% Accent ──
    GOLD   = RGBColor(0xC8, 0x96, 0x2E)  # bars, highlights (6.8:1 navy ✅)
    GOLD_L = RGBColor(0xD4, 0xA0, 0x17)  # lighter gold
    
    # ── Semantic (WCAG AA on white ✅) ──
    TEXT_D = RGBColor(0x1A, 0x1A, 0x2E)  # primary text (17.1:1 🏆)
    TEXT_M = RGBColor(0x88, 0x99, 0xB0)  # subtitle on navy (6.2:1 ✅)
    TEXT_L = RGBColor(0x64, 0x74, 0x8B)  # footer text (4.8:1 ✅)
    BLUE   = RGBColor(0x25, 0x63, 0xEB)  # info, definisi (5.2:1 ✅)
    TEAL   = RGBColor(0x0B, 0x7C, 0x72)  # prosedur (5.1:1 ✅)
    WARM   = RGBColor(0xA0, 0x52, 0x2D)  # sienna, peringatan (5.6:1 ✅)
    RED    = RGBColor(0xDC, 0x26, 0x26)  # sanksi, bahaya (4.8:1 ✅)

    @staticmethod
    def from_hex(hex_str):
        """Convert '#2563EB' or '2563EB' to RGBColor."""
        h = hex_str.lstrip('#')
        return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))

    @staticmethod
    def named(name, fallback=None):
        """Cari warna by name string, fallback ke NAVY."""
        if fallback is None:
            fallback = Colors.NAVY
        return getattr(Colors, name.upper(), fallback)

    @staticmethod
    def to_rgb_tuple(c):
        """RGBColor atau tuple → (r,g,b) tuple."""
        return (c[0], c[1], c[2])

    @staticmethod
    def hex_to_tuple(hex_str):
        """'#2563EB' → (37, 99, 235)."""
        h = hex_str.lstrip('#')
        return (int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))

    @staticmethod
    def make(color):
        """(r,g,b) atau hex string → RGBColor."""
        if isinstance(color, RGBColor):
            return color
        if isinstance(color, str):
            return Colors.from_hex(color)
        if isinstance(color, (tuple, list)):
            return RGBColor(int(color[0]), int(color[1]), int(color[2]))
        return color


# ═══════════════════════════════════════════════════════════
# COLOR THEORY — WCAG, Harmony Rules, Auto Palette
# ═══════════════════════════════════════════════════════════

# ─── WCAG 2.1 Luminance & Contrast ───

def _linearize(c):
    """sRGB channel linearization per WCAG 2.1."""
    c = c / 255.0
    return c / 12.92 if c <= 0.04045 else ((c + 0.055) / 1.055) ** 2.4

def _get_rgb(c):
    """Ensure (r,g,b) tuple — works with RGBColor, tuple, list."""
    return int(c[0]), int(c[1]), int(c[2])

def relative_luminance(r, g, b):
    """WCAG 2.1: L = 0.2126R + 0.7152G + 0.0722B."""
    return 0.2126 * _linearize(r) + 0.7152 * _linearize(g) + 0.0722 * _linearize(b)

def contrast_ratio(c1, c2):
    """WCAG 2.1 contrast ratio: (L₁+0.05)/(L₂+0.05)."""
    r1, g1, b1 = _get_rgb(c1)
    r2, g2, b2 = _get_rgb(c2)
    l1 = relative_luminance(r1, g1, b1)
    l2 = relative_luminance(r2, g2, b2)
    lighter, darker = max(l1, l2), min(l1, l2)
    return (lighter + 0.05) / (darker + 0.05)

def wcag_aa(cr, text_size='normal'):
    """True if ≥4.5:1 (normal) or ≥3:1 (large)."""
    return cr >= (4.5 if text_size == 'normal' else 3.0)

def wcag_aaa(cr):
    """True if ≥7:1 (normal)."""
    return cr >= 7.0

# ─── RGB ↔ HLS ───

def rgb_to_hls(r, g, b):
    """(r,g,b) in 0-255 → (h, l, s) where h∈[0,360), l,s∈[0,100]."""
    rn, gn, bn = r / 255.0, g / 255.0, b / 255.0
    h, l, s = colorsys.rgb_to_hls(rn, gn, bn)
    return h * 360, l * 100, s * 100

def hls_to_rgb(h, l, s):
    """(h, l, s) → (r,g,b) in 0-255."""
    r, g, b = colorsys.hls_to_rgb(h / 360.0, l / 100.0, s / 100.0)
    return round(r * 255), round(g * 255), round(b * 255)

# ─── WCAG Auto-Adjust —──

def adjust_to_wcag(rgb, bg=(255, 255, 255), target=4.5):
    """
    Binary-search lightness to meet `target` contrast ratio against `bg`.
    For light bg → darken; for dark bg → lighten.
    Returns (r,g,b) tuple.
    
    Strategy:
        - Light bg: we need text DARKER than bg. Search [0, L] for 
          the HIGHEST (lightest) L that still passes.
        - Dark bg: we need text LIGHTER than bg. Search [L, 100] for 
          the LOWEST (darkest) L that still passes.
    """
    cr = contrast_ratio(rgb, bg)
    if cr >= target:
        return rgb

    r1, g1, b1 = _get_rgb(rgb)
    h, l, s = rgb_to_hls(r1, g1, b1)
    br, bg_, bb = _get_rgb(bg)
    bg_lum = relative_luminance(br, bg_, bb)

    if bg_lum > 0.5:          # light bg → text must be DARKER
        lo, hi = 0.0, l       # lo=very dark (passes), hi=current (may fail)
        for _ in range(30):
            mid = (lo + hi) / 2
            test = hls_to_rgb(h, mid, s)
            if contrast_ratio(test, bg) >= target:
                lo = mid      # passes → try lighter
            else:
                hi = mid      # fails → need darker
            if hi - lo < 0.3:
                break
        return hls_to_rgb(h, lo, s)   # lo = lightest that passes
    else:                      # dark bg → text must be LIGHTER
        lo, hi = l, 100.0     # lo=current (may fail), hi=very light (passes)
        for _ in range(30):
            mid = (lo + hi) / 2
            test = hls_to_rgb(h, mid, s)
            if contrast_ratio(test, bg) >= target:
                hi = mid      # passes → try darker
            else:
                lo = mid      # fails → need lighter
            if hi - lo < 0.3:
                break
        return hls_to_rgb(h, hi, s)   # hi = darkest that passes


# ═══════════════════════════════════════════════════════════
# PALETTE GENERATOR — Dari satu primary color
# ═══════════════════════════════════════════════════════════

class PaletteGenerator:
    """
    Generate palet 60-30-10 + 4 semantic + WCAG-verified dari satu warna utama.
    
    Rules:
        - 60% Dominant:  dark background (primary hue, low L)
        - 30% Secondary: light bg (primary hue, high L)
        - 10% Accent:    split-complementary hue
        - Semantic:      tetradic rotation (BLUE, TEAL, WARM, RED)
        - Semua teks:    WCAG AA ≥4.5:1
    
    Usage:
        pg = PaletteGenerator("#2563EB")
        palette = pg.generate()   # → dict { name: RGBColor, ... }
     
    Publikasi acuan:
        - Color-by-concept association (Rathore et al., VIS 2019, arXiv:1908.00220)
        - Culture-inspired palette gen (Li et al., 2021, arXiv:2102.05231)
    """
    
    # Harmony rule generators: given hue h, return list of hue angles
    HARMONY = {
        'complementary':    lambda h: [(h + 180) % 360],
        'analogous':        lambda h: [(h - 30) % 360, h, (h + 30) % 360],
        'triadic':          lambda h: [h, (h + 120) % 360, (h + 240) % 360],
        'split_comp':       lambda h: [h, (h + 150) % 360, (h + 210) % 360],
        'tetradic':         lambda h: [h, (h + 90) % 360, (h + 180) % 360, (h + 270) % 360],
    }
    
    SEMANTIC_LABELS = ['BLUE', 'TEAL', 'WARM', 'RED']
    
    def __init__(self, primary_hex):
        """
        Args:
            primary_hex: hex string '#2563EB' atau '2563EB'
        """
        self.prim = Colors.hex_to_tuple(primary_hex)
        self.h, self.l, self.s = rgb_to_hls(*self.prim)
    
    # ── Helpers ──
    
    def _hls(self, h, l, s=None):
        """Single hue variant → (r,g,b). Clamp saturation."""
        if s is None:
            s = self.s
        return hls_to_rgb(h % 360, max(0, min(100, l)), max(5, min(100, s)))
    
    def _vary(self, h, l_delta=0, s_delta=0):
        """Vary lightness & saturation from primary, keep hue."""
        return self._hls(
            h if h is not None else self.h,
            max(0, min(100, self.l + l_delta)),
            max(5, min(100, self.s + s_delta))
        )
    
    def _distinct_hues(self, n=4):
        """Generate n evenly-spaced hues from primary (tetradic)."""
        return [(self.h + i * 90) % 360 for i in range(n)]
    
    def _wcag_report(self, palette):
        """Print WCAG AA verification for key combinations."""
        lines = []
        for name, cr_bg, cr_bg_label, text_size in [
            ('TEXT_D on WHITE',  contrast_ratio(palette['TEXT_D'], (255,255,255)), 'white', 'normal'),
            ('TEXT_M on NAVY',   contrast_ratio(palette['TEXT_M'], palette['NAVY']), 'navy', 'normal'),
            ('TEXT_L on WHITE',  contrast_ratio(palette['TEXT_L'], (255,255,255)), 'white', 'normal'),
            ('GOLD on NAVY',     contrast_ratio(palette['GOLD'], palette['NAVY']), 'navy', 'normal'),
            ('BLUE on WHITE',    contrast_ratio(palette['BLUE'], (255,255,255)), 'white', 'normal'),
            ('TEAL on WHITE',    contrast_ratio(palette['TEAL'], (255,255,255)), 'white', 'normal'),
            ('WARM on WHITE',    contrast_ratio(palette['WARM'], (255,255,255)), 'white', 'normal'),
            ('RED on WHITE',     contrast_ratio(palette['RED'], (255,255,255)), 'white', 'normal'),
        ]:
            status = '✅' if wcag_aa(cr_bg, text_size) else '❌'
            lines.append(f"  {status} {name}: {cr_bg:.1f}:1")
        return '\n'.join(lines)
    
    def generate(self, verbose=False):
        """
        Generate complete 60-30-10 + semantic palette.
        
        Args:
            verbose: print WCAG report
        
        Returns:
            dict { 'NAVY': RGBColor, 'BLUE': RGBColor, ... }
        """
        h, l, s = self.h, self.l, self.s
        
        # ── Clamp extremes ──
        sat = max(25, min(85, s))     # keep saturation viable
        lum = max(20, min(80, l))     # keep luminance reasonable
        
        # ════ 60% — Dark backgrounds (dominant) ════
        navy   = self._hls(h, 10, min(sat, 70))
        navy_l = self._hls(h, 18, min(sat, 60))
        navy_d = self._hls(h, 6,  min(sat, 65))
        navy_m = self._hls(h, 32, min(sat, 50))
        
        # ════ 30% — Light backgrounds (secondary) ════
        white  = (255, 255, 255)
        off_w  = (245, 247, 250)
        ice    = self._hls(h, 92, min(sat, 12))
        ice_d  = self._hls(h, 84, min(sat, 16))
        
        # ════ 10% — Accent (split-complementary for best contrast) ════
        comp_hues = self.HARMONY['split_comp'](h)
        accent_h = comp_hues[1] if len(comp_hues) > 1 else (h + 180) % 360
        accent   = self._hls(accent_h, 42, min(sat + 10, 92))
        accent_l = self._hls(accent_h, 58, min(sat + 10, 88))
        
        # ════ Semantic colors (tetradic rotation) ════
        # Note: labels BLUE/TEAL/WARM/RED are arbitrary; actual hue depends on
        # primary. These are 4 evenly-spaced hues (90° apart) for 4 semantic slots.
        sem_hues = self._distinct_hues(4)
        sem_colors = []
        for i, sh in enumerate(sem_hues):
            sc = self._hls(sh, 40 + (i * 3), min(sat + 10, 88))
            sc = adjust_to_wcag(sc, bg=white, target=4.5)
            sem_colors.append(sc)
        
        # ════ Text colors — WCAG AA ════
        text_d = (26, 26, 46)          # near-black — 17+:1 on white
        
        # TEXT_M on navy
        text_m_raw = self._hls(h, 64, min(sat, 20))
        text_m = adjust_to_wcag(text_m_raw, bg=navy, target=4.5)
        
        # TEXT_L on white
        text_l_raw = self._hls(h, 38, min(sat, 12))
        text_l = adjust_to_wcag(text_l_raw, bg=white, target=4.5)
        
        # Ensure accent passes on navy
        accent   = adjust_to_wcag(accent, bg=navy, target=4.5)
        accent_l = adjust_to_wcag(accent_l, bg=navy, target=4.5)
        
        # ── Assemble ──
        palette = {
            'NAVY':   Colors.make(navy),
            'NAVY_L': Colors.make(navy_l),
            'NAVY_D': Colors.make(navy_d),
            'NAVY_M': Colors.make(navy_m),
            
            'WHITE':  Colors.make(white),
            'OFF_W':  Colors.make(off_w),
            'ICE':    Colors.make(ice),
            'ICE_D':  Colors.make(ice_d),
            
            'GOLD':   Colors.make(accent),
            'GOLD_L': Colors.make(accent_l),
            
            'TEXT_D': Colors.make(text_d),
            'TEXT_M': Colors.make(text_m),
            'TEXT_L': Colors.make(text_l),
            
            'BLUE':   Colors.make(sem_colors[0]),
            'TEAL':   Colors.make(sem_colors[1]),
            'WARM':   Colors.make(sem_colors[2]),
            'RED':    Colors.make(sem_colors[3]),
        }
        
        if verbose:
            print(f"🎨 Palette dari #{self.prim[0]:02X}{self.prim[1]:02X}{self.prim[2]:02X}")
            print(self._wcag_report(palette))
        
        return palette


# ═══════════════════════════════════════════════════════════
# LAYOUT FRAME — Semua kalkulasi layout
# ═══════════════════════════════════════════════════════════

class LayoutFrame:
    """
    Kalkulator posisi & ukuran untuk elemen slide.
    
    Slide widescreen 16:9 = 13.333" × 7.5"
    
    Zona layout:
        MARGIN_H  = 0.6"      (kiri/kanan)
        HEADER_H  = 0.94"     (gold_bar 0.04" + navy 0.9")
        HEADER_GAP = 0.21"    (jarak header → konten)
        FOOTER_H  = 0.50"     (navy bar 0.03" + text 0.22" + gap)
    
    Content area:
        cx = 0.6",  cy = 1.15"
        cw = 12.133",  ch = 5.85"
    """
    
    SLIDE_W = 13.333
    SLIDE_H = 7.5
    
    MARGIN_H   = 0.6
    HEADER_H   = 0.94
    HEADER_GAP = 0.21
    FOOTER_H   = 0.50
    
    GAP_H      = 0.30
    GAP_V      = 0.30
    GAP_ITEM   = 0.02
    CARD_PAD   = 0.15
    
    # Characters-per-inch (CPI) untuk Calibri
    CPI = { 7:17, 8:15, 9:14, 10:13, 11:12, 12:11,
            13:10, 14:9, 16:8, 18:7, 20:6.5, 24:5.5,
            32:4, 40:3, 48:2.5 }
    
    def __init__(self):
        self.cx = self.MARGIN_H
        self.cw = self.SLIDE_W - 2 * self.MARGIN_H
        self.cy = self.HEADER_H + self.HEADER_GAP
        self.ch = (self.SLIDE_H - self.FOOTER_H) - self.cy
    
    # ─── Grid ───
    
    def col_width(self, n, gap=None):
        """col_w = (cw - (n-1) × gap) / n"""
        g = gap if gap is not None else self.GAP_H
        return (self.cw - (n - 1) * g) / n
    
    def row_height(self, n, gap=None):
        """row_h = (ch - (n-1) × gap_v) / n"""
        g = gap if gap is not None else self.GAP_V
        return (self.ch - (n - 1) * g) / n
    
    def grid_pos(self, col, row, col_w, row_h, gap_h=None, gap_v=None):
        """x = cx + col × (col_w + gap_h)"""
        gh = gap_h if gap_h is not None else self.GAP_H
        gv = gap_v if gap_v is not None else self.GAP_V
        return self.cx + col * (col_w + gh), self.cy + row * (row_h + gv)
    
    def calc_grid(self, n_items, cols=0):
        """Kalkulasi grid multi-baris."""
        per_row = cols if cols > 0 else (
            2 if n_items <= 2 else (3 if n_items <= 3 else 4))
        n_rows = (n_items + per_row - 1) // per_row
        col_w = self.col_width(per_row)
        row_h = self.row_height(n_rows)
        return {'per_row': per_row, 'n_rows': n_rows,
                'col_w': col_w, 'row_h': row_h,
                'gap_h': self.GAP_H, 'gap_v': self.GAP_V}
    
    def grid_cell(self, i, grid):
        """Posisi cell ke-i dalam grid."""
        col = i % grid['per_row']
        row = i // grid['per_row']
        x, y = self.grid_pos(col, row, grid['col_w'], grid['row_h'],
                              grid['gap_h'], grid['gap_v'])
        return x, y, grid['col_w'], grid['row_h']
    
    # ─── Card Internal ───
    
    def card_item_layout(self, card_h, has_icon=True, n_items=1, title_h=0.30):
        icon_size = 0.42
        icon_y = self.CARD_PAD
        if has_icon:
            title_y = icon_y + icon_size + 0.13
        else:
            title_y = self.CARD_PAD
        item_h = 0.35
        item_start_y = title_y + title_h + 0.08
        content_used = item_start_y + n_items * (item_h + self.GAP_ITEM)
        overflow = max(0, content_used - card_h)
        return {'icon_size': icon_size, 'icon_y': icon_y,
                'title_y': title_y, 'item_start_y': item_start_y,
                'item_h': item_h, 'overflow': overflow}
    
    # ─── Estimasi Teks ───
    
    def text_height(self, text, font_size, box_width):
        """Tinggi teks estimated: ceil(len/(cpi*box_w)) * pt*1.2/72"""
        cpi = self.CPI.get(font_size, 12)
        max_ch = box_width * cpi
        if max_ch <= 0:
            return font_size * 1.2 / 72
        n_lines = math.ceil(len(text) / max_ch)
        return n_lines * font_size * 1.2 / 72 + 0.02
    
    def safe_item_height(self, text, font_size, box_width, min_h=0.25):
        return max(self.text_height(text, font_size, box_width) + 0.05, min_h)
    
    # ─── Verifikasi ───
    
    def check_bounds(self, x, y, w, h):
        warns = []
        if x < 0: warns.append(f"x={x:.2f}<0")
        if y < 0: warns.append(f"y={y:.2f}<0")
        if x + w > self.SLIDE_W: warns.append(f"right={x+w:.2f}>{self.SLIDE_W}")
        if y + h > self.SLIDE_H: warns.append(f"bottom={y+h:.2f}>{self.SLIDE_H}")
        return (len(warns) == 0, warns)
    
    def check_overlap(self, a, b):
        ax, ay, aw, ah = a
        bx, by, bw, bh = b
        ox = max(0, min(ax+aw, bx+bw) - max(ax, bx))
        oy = max(0, min(ay+ah, by+bh) - max(ay, by))
        overlap = ox * oy
        min_area = min(aw * ah, bw * bh)
        return 0 if min_area == 0 else (overlap / min_area) * 100


# ═══════════════════════════════════════════════════════════
# ENGINE — Menggabungkan LayoutFrame + Draw Helpers + Archetypes
# ═══════════════════════════════════════════════════════════

class Engine:
    """
    Presentation Engine — generic, content-agnostic.
    
    Cara pakai:
        engine = Engine()
        prs = engine.build(slides_data)   # slides_data = list of slide dicts
        prs.save("output.pptx")
    
    Slide dict format:
        {
            "type": "cover|toc|section|content|card_grid|two_col|callout|flow|table|closing",
            "data": { ... key/value sesuai type ... }
        }
    """
    
    # ── Font Pairings ──
    # Sumber: Anthropic PPTX Skill — header/body pairing
    FONT_PAIRS = {
        'modern':  {'header': 'Arial Black',   'body': 'Arial',         'mono': 'Consolas'},
        'classic': {'header': 'Georgia',       'body': 'Calibri',       'mono': 'Consolas'},
        'clean':   {'header': 'Calibri',       'body': 'Calibri Light', 'mono': 'Consolas'},
        'formal':  {'header': 'Cambria',       'body': 'Calibri',       'mono': 'Consolas'},
        'tech':    {'header': 'Consolas',       'body': 'Calibri',       'mono': 'Consolas'},
        'elegant': {'header': 'Palatino',      'body': 'Garamond',      'mono': 'Consolas'},
        'bold':    {'header': 'Impact',         'body': 'Arial',         'mono': 'Consolas'},
        'friendly':{'header': 'Trebuchet MS',  'body': 'Calibri',       'mono': 'Consolas'},
    }
    DEFAULT_FONT_STYLE = 'classic'
    
    def __init__(self, colors=None, primary_color=None, font_style=None):
        """
        Args:
            colors: instance of Colors (or object with same attributes) 
                    for custom palette. Defaults to Colors().
            primary_color: hex string like '#2563EB'. If given, 
                           PaletteGenerator auto-generates full WCAG AA 
                           palette from this primary. Overrides `colors`.
            font_style: one of 'modern','classic','clean','formal','tech',
                        'elegant','bold','friendly'. Default: 'classic'.
        """
        self.L = LayoutFrame()
        
        # Font
        fs = font_style or Engine.DEFAULT_FONT_STYLE
        pair = Engine.FONT_PAIRS.get(fs, Engine.FONT_PAIRS['classic'])
        self.FONT_H = pair['header']
        self.FONT_B = pair['body']
        self.FONT_M = pair['mono']
        self._font_style = fs
        
        if primary_color:
            pg = PaletteGenerator(primary_color)
            pal = pg.generate()
            self.C = type('DynamicColors', (), pal)()
            self._palette_report = pg._wcag_report(pal)
            print(f"🎨 Auto-palette dari {primary_color}")
            print(self._palette_report)
        else:
            self.C = colors or Colors()
            self._palette_report = ""
        
        self.prs = None
        self.pg_counter = [0]
        self.source_text = ""
    
    # ════════════════════════════════════════════════════════
    # BUILD — main entry point
    # ════════════════════════════════════════════════════════
    
    def build(self, slides_data, source_text="Sumber: ...", output_path=None):
        """
        Generate presentation dari array slide dicts.
        
        Args:
            slides_data: list of dict — lihat STRUCTURE.md atau contoh
            source_text: teks footer sumber
            output_path: path output (None = return Presentation object)
        
        Returns:
            Presentation object (panggil .save() untuk simpan)
        """
        self.source_text = source_text
        self.pg_counter = [0]
        
        self.prs = Presentation()
        self.prs.slide_width = Inches(self.L.SLIDE_W)
        self.prs.slide_height = Inches(self.L.SLIDE_H)
        
        dispatcher = {
            "cover":     self._build_cover,
            "toc":       self._build_toc,
            "section":   self._build_section,
            "content":   self._build_content,
            "card_grid": self._build_card_grid,
            "two_col":   self._build_two_col,
            "callout":   self._build_callout,
            "flow":        self._build_flow,
            "table":       self._build_table,
            "nsr_factors": self._build_nsr_factors,
            "closing":     self._build_closing,
        }
        
        for slide_def in slides_data:
            stype = slide_def.get("type", "content")
            data = slide_def.get("data", {})
            builder = dispatcher.get(stype)
            if builder:
                builder(data)
            else:
                # fallback: content slide
                self._build_content(data)
        
        if output_path:
            self.prs.save(output_path)
        
        return self.prs
    
    # ─── Helpers ───
    
    def _blank(self):
        return self.prs.slides.add_slide(self.prs.slide_layouts[6])
    
    def _resolve_color(self, clr):
        """clr bisa string hex '#2563EB', string name 'BLUE', atau RGBColor."""
        if isinstance(clr, RGBColor):
            return clr
        if isinstance(clr, str):
            if clr.startswith('#'):
                return Colors.from_hex(clr)
            return Colors.named(clr)
        return self.C.NAVY
    
    def _add_box(self, slide, left, top, width, height, text,
                 size=12, bold=False, color=None, align=PP_ALIGN.LEFT,
                 font=None, vAlign=MSO_ANCHOR.TOP, font_name=None):
        """
        Add a text box. Default font = self.FONT_B (body).
        Use font='header' for header font, font='mono' for monospace.
        """
        fname = font_name or self.FONT_B
        if font == 'header':
            fname = self.FONT_H
        elif font == 'mono':
            fname = self.FONT_M
        elif font is not None:
            fname = font
        if color is None:
            color = self.C.TEXT_D
        else:
            color = self._resolve_color(color)
        tb = slide.shapes.add_textbox(Inches(left), Inches(top),
                                       Inches(width), Inches(height))
        tf = tb.text_frame
        tf.word_wrap = True
        tf.auto_size = None
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(size)
        p.font.bold = bold
        p.font.color.rgb = color
        p.font.name = fname
        p.alignment = align
        p.space_after = Pt(0)
        p.space_before = Pt(0)
        self.L.check_bounds(left, top, width, height)
        return tb
    
    def _add_shape(self, slide, stype, left, top, width, height,
                   fill=None, line=None, lw=None):
        sh = slide.shapes.add_shape(stype, Inches(left), Inches(top),
                                     Inches(width), Inches(height))
        if fill:
            sh.fill.solid()
            sh.fill.fore_color.rgb = self._resolve_color(fill)
        else:
            sh.fill.background()
        if line:
            sh.line.color.rgb = self._resolve_color(line)
            if lw: sh.line.width = Pt(lw)
        else:
            sh.line.fill.background()
        self.L.check_bounds(left, top, width, height)
        return sh
    
    def _add_rrect(self, slide, left, top, width, height,
                   fill=None, line=None, lw=0.5, radius=0.04):
        if fill is None: fill = self.C.WHITE
        if line is None: line = self.C.ICE
        sh = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(left), Inches(top), Inches(width), Inches(height))
        sh.fill.solid()
        sh.fill.fore_color.rgb = self._resolve_color(fill)
        if line:
            sh.line.color.rgb = self._resolve_color(line)
            sh.line.width = Pt(lw)
        else:
            sh.line.fill.background()
        sh.adjustments[0] = radius
        self.L.check_bounds(left, top, width, height)
        return sh
    
    def _add_oval(self, slide, left, top, size, fill=None):
        if fill is None: fill = self.C.BLUE
        sh = slide.shapes.add_shape(
            MSO_SHAPE.OVAL, Inches(left), Inches(top), Inches(size), Inches(size))
        sh.fill.solid()
        sh.fill.fore_color.rgb = self._resolve_color(fill)
        sh.line.fill.background()
        return sh
    
    def _add_icon(self, slide, left, top, size, icon=None, fill=None, 
                  icon_color=None):
        """
        Add a colored circle with an icon inside.
        
        Args:
            icon: str — emoji char, or shape name from:
                  'check','x','arrow','star','circle','info','warning','question'
            fill: background color (default: self.C.BLUE)
            icon_color: icon color (default: WHITE)
        
        Built-in shapes use Unicode symbols, no external deps needed.
        """
        if fill is None: fill = self.C.BLUE
        if icon_color is None: icon_color = self.C.WHITE
        
        # Shape icons map
        SHAPE_ICONS = {
            'check':    '✓',
            'x':        '✗',
            'arrow':    '→',
            'star':     '★',
            'circle':   '●',
            'info':     'ℹ',
            'warning':  '⚠',
            'question': '?',
        }
        
        # Resolve icon
        icon_str = SHAPE_ICONS.get(icon, icon) if isinstance(icon, str) else icon
        if not icon_str:
            icon_str = '●'
        
        circle = self._add_oval(slide, left, top, size, fill=fill)
        ico_size = size * 0.6
        ico_x = left + (size - ico_size) / 2
        ico_y = top + (size - ico_size) / 2 - 0.02
        self._add_box(slide, ico_x, ico_y, ico_size, ico_size,
                      icon_str, round(size * 28), color=icon_color,
                      align=PP_ALIGN.CENTER)
        return circle
    
    def _text_to_shape(self, shape, text, size=12, bold=False, color=None,
                       align=PP_ALIGN.LEFT, font=None):
        if font is None: fname = self.FONT_B
        elif font == 'header': fname = self.FONT_H
        elif font == 'mono': fname = self.FONT_M
        else: fname = font
        if color is None: color = self.C.TEXT_D
        else: color = self._resolve_color(color)
        tf = shape.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(size)
        p.font.bold = bold
        p.font.color.rgb = color
        p.font.name = fname
        p.alignment = align
        return tf
    
    # ─── Header & Footer ───
    
    def _gold_top_bar(self, slide):
        self._add_shape(slide, MSO_SHAPE.RECTANGLE, 0, 0,
                        self.L.SLIDE_W, 0.035, fill=self.C.GOLD)
    
    def _navy_header(self, slide, title, subtitle=None):
        self._gold_top_bar(slide)
        hdr_h = self.L.HEADER_H - 0.04
        self._add_shape(slide, MSO_SHAPE.RECTANGLE, 0, 0.035,
                        self.L.SLIDE_W, hdr_h, fill=self.C.NAVY)
        self._add_shape(slide, MSO_SHAPE.RECTANGLE,
                        self.L.MARGIN_H, 0.12, 0.07, 0.55, fill=self.C.GOLD)
        self._add_box(slide, self.L.MARGIN_H + 0.22, 0.15,
                      self.L.cw - 0.5, 0.48, title, 20, bold=True,
                      color=self.C.WHITE, font='header')
        if subtitle:
            self._add_box(slide, self.L.MARGIN_H + 0.22, 0.63,
                          self.L.cw - 0.5, 0.28, subtitle, 9, color=self.C.TEXT_L)
    
    def _footer(self, slide):
        h = 0.03
        self._add_shape(slide, MSO_SHAPE.RECTANGLE, 0,
                        self.L.SLIDE_H - self.L.FOOTER_H,
                        self.L.SLIDE_W, h, fill=self.C.NAVY)
        self._add_box(slide, self.L.MARGIN_H,
                      self.L.SLIDE_H - self.L.FOOTER_H + 0.06,
                      4, 0.22, self.source_text, 7, color=self.C.TEXT_L)
        self.pg_counter[0] += 1
        self._add_box(slide, self.L.SLIDE_W - 1.0,
                      self.L.SLIDE_H - self.L.FOOTER_H + 0.06,
                      0.8, 0.22, str(self.pg_counter[0]), 8,
                      color=self.C.TEXT_L, align=PP_ALIGN.RIGHT)
    
    def _content_slide(self, title, subtitle=None):
        """Bikin blank slide + navy header."""
        s = self._blank()
        bg = s.background; bg.fill.solid(); bg.fill.fore_color.rgb = self.C.OFF_W
        self._navy_header(s, title, subtitle)
        return s
    
    # ════════════════════════════════════════════════════════
    # ARCHETYPE BUILDERS — generic, content-driven
    # ════════════════════════════════════════════════════════
    
    def _build_cover(self, data):
        """
        Slide cover — full-bleed navy + dekoratif ovals.
        
        Data:
            pre_title, city, main_title, main_subtitle,
            display_title, badge_text, badge_subtext
        """
        s = self._blank()
        bg = s.background; bg.fill.solid(); bg.fill.fore_color.rgb = self.C.NAVY
        
        # Dekoratif ovals
        self._add_oval(s, -1, -1.5, 4.5, self.C.NAVY_D)
        self._add_oval(s, 8.5, -2, 7, self.C.NAVY_D)
        self._add_oval(s, 10, 4.5, 5, self.C.NAVY_D)
        self._add_oval(s, 0.5, 5.5, 2.5, self.C.NAVY_L)
        
        mx = self.L.MARGIN_H
        cw = self.L.cw
        
        # Gold bar
        self._add_shape(s, MSO_SHAPE.RECTANGLE, mx, 2.6, 4, 0.04, fill=self.C.GOLD)
        
        # Teks
        pre = data.get("pre_title", "")
        if pre:
            self._add_box(s, mx, 0.5, 5, 0.35, pre, 12, bold=True, color=self.C.GOLD)
        
        city = data.get("city", "")
        if city:
            self._add_box(s, mx, 0.85, 5, 0.35, city, 14, bold=True, color=self.C.WHITE)
        
        main_title = data.get("main_title", "")
        if main_title:
            self._add_box(s, mx, 1.8, cw, 0.45, main_title, 20, bold=True, color=self.C.WHITE)
        
        main_sub = data.get("main_subtitle", "")
        if main_sub:
            self._add_box(s, mx, 2.15, cw, 0.4, main_sub, 15, bold=True, color=self.C.GOLD)
        
        display = data.get("display_title", "")
        if display:
            self._add_box(s, mx, 3.1, cw, 2.0, display, 40, bold=True, color=self.C.WHITE)
        
        # Badge
        badge_text = data.get("badge_text", "")
        if badge_text:
            self._add_rrect(s, mx, 5.4, 7.5, 0.9, fill=self.C.NAVY_D)
            self._add_box(s, mx + 0.2, 5.45, 7, 0.4, badge_text, 12, color=self.C.ICE)
        
        badge_sub = data.get("badge_subtext", "")
        if badge_sub:
            self._add_box(s, mx + 0.2, 5.75, 7, 0.35, badge_sub, 10, color=self.C.TEXT_L)
        
        # Bottom gold bar
        self._add_shape(s, MSO_SHAPE.RECTANGLE, 0,
                        self.L.SLIDE_H - 0.22, self.L.SLIDE_W, 0.22, fill=self.C.GOLD)
    
    def _build_toc(self, data):
        """
        Daftar isi — tabel konten dengan badge.
        
        Data:
            title, subtitle,
            items: [{num, label, color}]
            cols: 2 (default)
        """
        title = data.get("title", "Daftar Isi")
        subtitle = data.get("subtitle")
        items = data.get("items", [])
        cols = data.get("cols", 2)
        
        s = self._content_slide(title, subtitle)
        
        if not items:
            self._footer(s)
            return
        
        cw_col = self.L.col_width(cols, 0.35)
        n_per_col = (len(items) + cols - 1) // cols
        row_h = min(0.62, (self.L.ch - 0.2) / n_per_col)
        
        for i, item in enumerate(items):
            col = i // n_per_col
            row = i % n_per_col
            x = self.L.cx + col * (cw_col + 0.35)
            y = self.L.cy + row * row_h
            
            clr = self._resolve_color(item.get("color", self.C.BLUE))
            num = str(item.get("num", ""))
            lb = item.get("label", "")
            
            self._add_rrect(s, x, y, cw_col, row_h - 0.05)
            badge = self._add_shape(s, MSO_SHAPE.ROUNDED_RECTANGLE,
                                     x + 0.1, y + 0.06, 0.5, 0.5, fill=clr)
            self._text_to_shape(badge, num, 14, bold=True,
                                color=self.C.WHITE, align=PP_ALIGN.CENTER)
            self._add_box(s, x + 0.75, y + 0.06, cw_col - 0.9, 0.5,
                          lb, 13, color=self.C.TEXT_D,
                          vAlign=MSO_ANCHOR.MIDDLE)
        
        self._footer(s)
    
    def _build_section(self, data):
        """
        Section divider — full-bleed dark.
        
        Data:
            title, subtitle, action_text
        """
        title = data.get("title", "")
        subtitle = data.get("subtitle")
        action_text = data.get("action_text")
        
        s = self._blank()
        bg = s.background; bg.fill.solid(); bg.fill.fore_color.rgb = self.C.NAVY
        self._add_oval(s, -1.5, -1.5, 5, self.C.NAVY_L)
        self._add_oval(s, -0.5, -0.5, 3.5, self.C.NAVY_D)
        self._add_oval(s, 9.5, 4, 5, self.C.NAVY_D)
        self._add_oval(s, 10.5, 3, 3, self.C.NAVY_L)
        self._add_shape(s, MSO_SHAPE.RECTANGLE,
                        self.L.MARGIN_H, 2.3, 2.5, 0.04, fill=self.C.GOLD)
        self._add_box(s, self.L.MARGIN_H, 2.6, self.L.cw, 1.6,
                      title, 34, bold=True, color=self.C.WHITE)
        
        y_sub = 4.2
        if action_text:
            self._add_box(s, self.L.MARGIN_H, y_sub, self.L.cw, 0.4,
                          action_text, 14, bold=True, color=self.C.GOLD)
            y_sub += 0.35
        if subtitle:
            self._add_box(s, self.L.MARGIN_H, y_sub, self.L.cw, 0.3,
                          subtitle, 11, color=self.C.TEXT_L)
        
        self._add_shape(s, MSO_SHAPE.RECTANGLE, 0,
                        self.L.SLIDE_H - 0.22, self.L.SLIDE_W, 0.22, fill=self.C.GOLD)
        self.pg_counter[0] += 1
    
    def _build_content(self, data):
        """
        Content slide polos — header + footer doang.
        
        Data:
            title, subtitle
        """
        s = self._content_slide(data.get("title", ""), data.get("subtitle"))
        self._footer(s)
    
    def _build_card_grid(self, data):
        """
        Grid of cards — multi-baris, icon, bullet items.
        
        Data:
            title, subtitle,
            cards: [{icon, title, color, items: [str]}]
            cols: 0 (auto)
        """
        title = data.get("title", "")
        subtitle = data.get("subtitle")
        cards = data.get("cards", [])
        cols = data.get("cols", 0)
        
        s = self._content_slide(title, subtitle)
        n = len(cards)
        if n == 0:
            self._footer(s)
            return
        
        grid = self.L.calc_grid(n, cols)
        
        for i, cd in enumerate(cards):
            cx, cy, cw, ch = self.L.grid_cell(i, grid)
            
            clr_s = cd.get('color', self.C.BLUE)
            clr = self._resolve_color(clr_s)
            ic = cd.get('icon', '')
            t = cd.get('title', '')
            items = cd.get('items', [])
            
            self._add_rrect(s, cx, cy, cw, ch)
            self._add_shape(s, MSO_SHAPE.RECTANGLE, cx, cy, 0.05, ch, fill=clr)
            
            item_layout = self.L.card_item_layout(ch, has_icon=bool(ic), n_items=len(items))
            
            if ic:
                self._add_icon(s, cx + self.L.CARD_PAD, cy + item_layout['icon_y'],
                               item_layout['icon_size'], icon=ic, fill=clr)
                title_y = cy + item_layout['title_y']
            else:
                title_y = cy + self.L.CARD_PAD
            
            self._add_box(s, cx + self.L.CARD_PAD, title_y, cw - self.L.CARD_PAD * 2, 0.30,
                          t, 13, bold=True, color=clr)
            
            ay = cy + item_layout['item_start_y']
            box_w = cw - self.L.CARD_PAD * 2
            for item in items:
                item_h = self.L.safe_item_height(f"• {item}", 9, box_w, min_h=0.30)
                self._add_box(s, cx + self.L.CARD_PAD, ay, box_w, item_h,
                              f"• {item}", 9, color=self.C.TEXT_D)
                ay += item_h + self.L.GAP_ITEM
        
        self._footer(s)
    
    def _build_two_col(self, data):
        """
        Two column layout — 2 card side by side.
        
        Data:
            title, subtitle,
            left: {color, lines: [str]},    # $ prefix = highlight
            right: {color, lines: [str]}
        """
        title = data.get("title", "")
        subtitle = data.get("subtitle")
        left_data = data.get("left", {})
        right_data = data.get("right", {})
        
        s = self._content_slide(title, subtitle)
        gap = self.L.GAP_H
        cw = self.L.col_width(2, gap)
        sy = self.L.cy
        ch = self.L.ch
        
        for data_block, clr_s, x in [
            (left_data.get('lines', []), left_data.get('color', self.C.BLUE), self.L.cx),
            (right_data.get('lines', []), right_data.get('color', self.C.TEAL),
             self.L.cx + cw + gap),
        ]:
            clr = self._resolve_color(clr_s)
            self._add_rrect(s, x, sy, cw, ch)
            self._add_shape(s, MSO_SHAPE.RECTANGLE, x, sy, 0.05, ch, fill=clr)
            ay = sy + 0.15
            for line in data_block:
                if line.startswith("$"):
                    self._add_box(s, x + 0.2, ay, cw - 0.35, 0.30,
                                  line[1:], 14, bold=True, color=clr)
                    ay += 0.32
                else:
                    self._add_box(s, x + 0.2, ay, cw - 0.35, 0.25,
                                  f"• {line}", 11, color=self.C.TEXT_D)
                    ay += 0.26
        
        self._footer(s)
    
    def _build_callout(self, data):
        """
        Big number callout cards.
        
        Data:
            title, subtitle,
            callouts: [{number, label, color}]
            note: str (optional)
        """
        title = data.get("title", "")
        subtitle = data.get("subtitle")
        callouts = data.get("callouts", [])
        note_text = data.get("note")
        
        s = self._content_slide(title, subtitle)
        n = len(callouts)
        if n == 0:
            self._footer(s)
            return
        
        cw = self.L.col_width(n)
        sy = self.L.cy
        
        for i, co in enumerate(callouts):
            cx = self.L.cx + i * (cw + self.L.GAP_H)
            clr = self._resolve_color(co.get("color", self.C.BLUE))
            num = str(co.get("number", ""))
            lb = co.get("label", "")
            
            self._add_rrect(s, cx, sy, cw, 2.2)
            self._add_shape(s, MSO_SHAPE.RECTANGLE, cx, sy, cw, 0.05, fill=clr)
            self._add_box(s, cx, sy + 0.25, cw, 0.7, num, 32, bold=True,
                          color=clr, align=PP_ALIGN.CENTER)
            self._add_box(s, cx + 0.1, sy + 1.0, cw - 0.2, 0.7, lb, 11,
                          color=self.C.TEXT_M, align=PP_ALIGN.CENTER,
                          vAlign=MSO_ANCHOR.TOP)
        
        if note_text:
            ny = sy + 2.5
            self._add_rrect(s, self.L.cx, ny, self.L.cw, 2.8, fill=self.C.ICE)
            self._add_box(s, self.L.cx + 0.25, ny + 0.15, self.L.cw - 0.5, 2.5,
                          note_text, 11, color=self.C.TEXT_D)
        
        self._footer(s)
    
    def _build_flow(self, data):
        """
        Horizontal flow dengan step dan arrow.
        
        Data:
            title, subtitle,
            steps: [{num, title, desc, color}]
            note: str (optional)
        """
        title = data.get("title", "")
        subtitle = data.get("subtitle")
        steps = data.get("steps", [])
        note_text = data.get("note")
        
        s = self._content_slide(title, subtitle)
        n = len(steps)
        if n == 0:
            self._footer(s)
            return
        
        bw = self.L.col_width(n)
        bgap = self.L.GAP_H
        sy = self.L.cy + 0.15
        
        for i, step in enumerate(steps):
            cx = self.L.cx + i * (bw + bgap)
            clr = self._resolve_color(step.get("color", self.C.BLUE))
            num = str(step.get("num", ""))
            step_title = step.get("title", "")
            desc = step.get("desc", "")
            
            self._add_rrect(s, cx, sy, bw, 2.0)
            self._add_shape(s, MSO_SHAPE.RECTANGLE, cx, sy, bw, 0.05, fill=clr)
            circ_x = cx + bw / 2 - 0.25
            self._add_oval(s, circ_x, sy + 0.15, 0.5, clr)
            self._add_box(s, circ_x, sy + 0.15, 0.5, 0.5, num, 18, bold=True,
                          color=self.C.WHITE, align=PP_ALIGN.CENTER)
            self._add_box(s, cx + 0.1, sy + 0.75, bw - 0.2, 0.35, step_title,
                          14, bold=True, color=clr, align=PP_ALIGN.CENTER)
            self._add_box(s, cx + 0.1, sy + 1.1, bw - 0.2, 0.7, desc, 10,
                          color=self.C.TEXT_M, align=PP_ALIGN.CENTER,
                          vAlign=MSO_ANCHOR.TOP)
            if i < n - 1:
                self._add_box(s, cx + bw, sy + 0.7, bgap, 0.5, "›", 24,
                              color=self.C.GOLD, align=PP_ALIGN.CENTER,
                              vAlign=MSO_ANCHOR.MIDDLE)
        
        if note_text:
            ny = sy + 2.4
            self._add_rrect(s, self.L.cx, ny, self.L.cw, 2.8, fill=self.C.ICE)
            self._add_box(s, self.L.cx + 0.25, ny + 0.15, self.L.cw - 0.5, 2.5,
                          note_text, 11, color=self.C.TEXT_D)
        
        self._footer(s)
    
    def _set_cell(self, cell, text, size=11, bold=False, color=None,
                  align=PP_ALIGN.LEFT, fill=None, font=None):
        if font is None: fname = self.FONT_B
        elif font == 'header': fname = self.FONT_H
        elif font == 'mono': fname = self.FONT_M
        else: fname = font
        if color is None: color = self.C.TEXT_D
        else: color = self._resolve_color(color)
        cell.text = ""
        p = cell.text_frame.paragraphs[0]
        p.text = str(text)
        p.font.size = Pt(size)
        p.font.bold = bold
        p.font.color.rgb = color
        p.font.name = fname
        p.alignment = align
        cell.text_frame.word_wrap = True
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        if fill:
            cell.fill.solid()
            cell.fill.fore_color.rgb = self._resolve_color(fill)
    
    def _build_table(self, data):
        """
        Native table slide.
        
        Data:
            title, subtitle,
            headers: [str],
            rows: [[str], ...]
        """
        title = data.get("title", "")
        subtitle = data.get("subtitle")
        headers = data.get("headers", [])
        rows = data.get("rows", [])
        
        s = self._content_slide(title, subtitle)
        
        if headers and rows:
            nc = len(headers)
            nr = len(rows) + 1
            rh = 0.42
            ts = s.shapes.add_table(nr, nc, Inches(self.L.cx), Inches(self.L.cy),
                                     Inches(self.L.cw), Inches(rh * nr))
            tbl = ts.table
            for i in range(nc):
                tbl.columns[i].width = int(Inches(self.L.cw / nc))
            for j, hdr in enumerate(headers):
                self._set_cell(tbl.cell(0, j), hdr, bold=True, color=self.C.WHITE,
                               align=PP_ALIGN.CENTER, fill=self.C.NAVY)
            for ri, row in enumerate(rows):
                bg_c = self.C.ICE if ri % 2 == 0 else self.C.WHITE
                for j, val in enumerate(row):
                    al = PP_ALIGN.LEFT if j == 0 else PP_ALIGN.CENTER
                    self._set_cell(tbl.cell(ri + 1, j), val, color=self.C.TEXT_D,
                                   align=al, fill=bg_c)
        
        self._footer(s)
    
    def _build_nsr_factors(self, data):
        """
        NSR Factors — 7 faktor dengan oval numbers + classification box.
        
        Data:
            title, subtitle,
            factors: [{num, label, color}],
            classification_note: str
        """
        title = data.get("title", "")
        subtitle = data.get("subtitle")
        factors = data.get("factors", [])
        class_note = data.get("classification_note", "")
        
        s = self._content_slide(title, subtitle)
        n = len(factors)
        if n == 0:
            self._footer(s)
            return
        
        per_row = 4
        fw = self.L.col_width(per_row, 0.3)
        
        for i, f in enumerate(factors):
            col = i % per_row
            row = i // per_row
            x, y = self.L.grid_pos(col, row, fw, 1.5, gap_v=0.3)
            clr = self._resolve_color(f.get("color", self.C.BLUE))
            num = str(f.get("num", ""))
            lb = f.get("label", "")
            
            self._add_rrect(s, x, y, fw, 1.3)
            self._add_oval(s, x + 0.15, y + 0.25, 0.55, clr)
            self._add_box(s, x + 0.15, y + 0.25, 0.55, 0.55, num,
                          18, bold=True, color=self.C.WHITE, align=PP_ALIGN.CENTER)
            self._add_box(s, x + 0.8, y + 0.2, fw - 1, 0.9, lb,
                          13, color=self.C.TEXT_D, vAlign=MSO_ANCHOR.MIDDLE)
        
        if class_note:
            ny = self.L.cy + 1.5 + 0.3
            self._add_rrect(s, self.L.cx, ny, self.L.cw, 1.7, fill=self.C.ICE)
            self._add_box(s, self.L.cx + 0.2, ny + 0.05, 5, 0.3,
                          "KLASIFIKASI KELAS JALAN", 11, bold=True, color=self.C.NAVY)
            self._add_box(s, self.L.cx + 0.2, ny + 0.4, self.L.cw - 0.4, 0.9,
                          class_note, 10, color=self.C.TEXT_D)
        
        self._footer(s)
    
    def _build_closing(self, data):
        """
        Closing slide — full-bleed navy.
        
        Data:
            pre_title, main_title, subtitle, source
        """
        s = self._blank()
        bg = s.background; bg.fill.solid(); bg.fill.fore_color.rgb = self.C.NAVY
        self._add_oval(s, -1, -1.5, 4.5, self.C.NAVY_D)
        self._add_oval(s, 8.5, -2, 7, self.C.NAVY_D)
        self._add_oval(s, 10, 4.5, 5, self.C.NAVY_D)
        
        mx = self.L.MARGIN_H
        
        self._add_shape(s, MSO_SHAPE.RECTANGLE, mx, 3.6, 3.5, 0.04, fill=self.C.GOLD)
        
        pre = data.get("pre_title", "")
        if pre:
            self._add_box(s, mx, 1.6, self.L.cw, 0.4, pre, 14, bold=True, color=self.C.GOLD)
        
        main = data.get("main_title", "")
        if main:
            self._add_box(s, mx, 2.4, self.L.cw, 1.5, main, 48, bold=True, color=self.C.WHITE)
        
        sub = data.get("subtitle", "")
        if sub:
            self._add_box(s, mx, 4.1, self.L.cw, 0.8, sub, 14, color=self.C.ICE)
        
        src = data.get("source", "")
        if src:
            self._add_box(s, mx, 5.1, self.L.cw, 0.35, src, 10, color=self.C.TEXT_L)
        
        self._add_shape(s, MSO_SHAPE.RECTANGLE, 0,
                        self.L.SLIDE_H - 0.22, self.L.SLIDE_W, 0.22, fill=self.C.GOLD)


# ═══════════════════════════════════════════════════════════
# SHORTCUT — langsung jalan kalau dipanggil
# ═══════════════════════════════════════════════════════════

# ═══════════════════════════════════════════════════════════
# WCAG AA REPORT UTILITY
# ═══════════════════════════════════════════════════════════

def palette_report(primary_hex):
    """Generate & print WCAG AA report for a primary color."""
    pg = PaletteGenerator(primary_hex)
    pal = pg.generate(verbose=True)
    return pal


if __name__ == "__main__":
    import sys
    
    # ── Jika arg: generate palette report ──
    if len(sys.argv) > 1 and sys.argv[1] != 'demo':
        primary = sys.argv[1]
        if primary.startswith('#'):
            pg = PaletteGenerator(primary)
            pal = pg.generate(verbose=True)
            
            # Also build a demo PPT showing the palette
            h, _, _ = rgb_to_hls(*Colors.hex_to_tuple(primary))
            comp_h = (h + 180) % 360
            
            slides = [
                {"type": "cover", "data": {
                    "pre_title": "CUSTOM PALETTE", "city": "COLOR THEORY",
                    "main_title": f"Primary #{primary.lstrip('#')}",
                    "main_subtitle": f"Hue {h:.0f}° · Harmony split-comp · Tetradic semantic",
                    "display_title": f"AUTO\nWCAG AA",
                    "badge_text": f"60-30-10  ·  {sum(v is not None for v in pal.values())} colors  ·  All ≥4.5:1"
                }},
                {"type": "card_grid", "data": {
                    "title": "Semantic Colors on White Background",
                    "cards": [
                        {"icon": "🔵", "title": "BLUE — Info",   "color": "#{:02X}{:02X}{:02X}".format(*Colors.to_rgb_tuple(pal['BLUE'])),  "items": [f"CR: {contrast_ratio(Colors.to_rgb_tuple(pal['BLUE']), (255,255,255)):.1f}:1 ✅"]},
                        {"icon": "🟢", "title": "TEAL — Success", "color": "#{:02X}{:02X}{:02X}".format(*Colors.to_rgb_tuple(pal['TEAL'])), "items": [f"CR: {contrast_ratio(Colors.to_rgb_tuple(pal['TEAL']), (255,255,255)):.1f}:1 ✅"]},
                        {"icon": "🟠", "title": "WARM — Warning", "color": "#{:02X}{:02X}{:02X}".format(*Colors.to_rgb_tuple(pal['WARM'])), "items": [f"CR: {contrast_ratio(Colors.to_rgb_tuple(pal['WARM']), (255,255,255)):.1f}:1 ✅"]},
                        {"icon": "🔴", "title": "RED — Danger",   "color": "#{:02X}{:02X}{:02X}".format(*Colors.to_rgb_tuple(pal['RED'])),  "items": [f"CR: {contrast_ratio(Colors.to_rgb_tuple(pal['RED']), (255,255,255)):.1f}:1 ✅"]},
                    ]
                }},
                {"type": "closing", "data": {
                    "pre_title": f"Primary #{primary.lstrip('#')}",
                    "main_title": "WCAG AA SELESAI",
                    "subtitle": "All colors verified ≥4.5:1 — no color bias",
                    "source": "Color Theory · Split-complementary + Tetradic"
                }},
            ]
            out_path = f"/tmp/palette_{primary.lstrip('#')}.pptx"
            engine = Engine(primary_color=primary)
            engine.build(slides, source_text=f"Auto palette dari {primary}", output_path=out_path)
            print(f"📊 PPT: {out_path} ({len(engine.prs.slides)} slides)")
        else:
            print("Usage: python3 ppt_engine.py <hex_color> # e.g. #E91E63")
        sys.exit(0)
    
    # ── Default demo ──
    print("=" * 60)
    print("DEMO 1: Default Colors palette")
    print("=" * 60)
    slides = [
        {"type": "cover", "data": {
            "pre_title": "DEMO", "city": "TEST CITY",
            "main_title": "PERATURAN DEMO", "main_subtitle": "NOMOR 1 TAHUN 2026",
            "display_title": "DEMO\nGENERIC ENGINE", "badge_text": "Test  ·  2026"
        }},
        {"type": "closing", "data": {
            "pre_title": "DEMO", "main_title": "SELESAI",
            "subtitle": "Engine generic berhasil", "source": "ppt_engine.py"
        }},
    ]
    engine = Engine()
    engine.build(slides, source_text="Sumber: Demo", output_path="/tmp/ppt_engine_demo.pptx")
    print(f"✅ Demo 1 OK: /tmp/ppt_engine_demo.pptx ({len(engine.prs.slides)} slides)")
    
    print()
    print("=" * 60)
    print("DEMO 2: Custom primary_color='#E91E63' (Pink)")
    print("=" * 60)
    pg = PaletteGenerator("#E91E63")
    pal = pg.generate(verbose=True)
    slides2 = [
        {"type": "cover", "data": {
            "pre_title": "CUSTOM PALETTE", "city": "PINK THEME",
            "main_title": "Primary #E91E63",
            "display_title": "AUTO\nPALETTE",
            "badge_text": "Split-comp · Tetradic · WCAG AA ✅"
        }},
        {"type": "closing", "data": {
            "pre_title": "#E91E63", "main_title": "SELESAI",
            "subtitle": "Auto-generated palette", "source": "PaletteGenerator"
        }},
    ]
    engine2 = Engine(primary_color="#E91E63")
    engine2.build(slides2, source_text="Custom palette", output_path="/tmp/ppt_palette_demo.pptx")
    print(f"✅ Demo 2 OK: /tmp/ppt_palette_demo.pptx ({len(engine2.prs.slides)} slides)")
    
    print()
    print("=" * 60)
    print("DEMO 3: Custom primary_color='#4CAF50' (Natural Green)")
    print("=" * 60)
    engine3 = Engine(primary_color="#4CAF50")
    slides3 = [
        {"type": "cover", "data": {
            "pre_title": "CUSTOM PALETTE", "city": "GREEN THEME",
            "main_title": "Primary #4CAF50",
            "display_title": "NATURAL\nGREEN",
            "badge_text": "Split-comp · Tetradic · WCAG AA ✅"
        }},
        {"type": "closing", "data": {
            "pre_title": "#4CAF50", "main_title": "SELESAI",
            "subtitle": "Auto-generated palette", "source": "PaletteGenerator"
        }},
    ]
    engine3.build(slides3, source_text="Green palette", output_path="/tmp/ppt_green_demo.pptx")
    print(f"✅ Demo 3 OK: /tmp/ppt_green_demo.pptx ({len(engine3.prs.slides)} slides)")
    
    print()
    print("=" * 60)
    print("DEMO 4: Font style 'modern' + primary_color='#8E44AD' (Purple)")
    print("=" * 60)
    engine4 = Engine(primary_color="#8E44AD", font_style="modern")
    slides4 = [
        {"type": "cover", "data": {
            "pre_title": "FONT STYLE", "city": "MODERN",
            "main_title": "Arial Black + Arial",
            "display_title": "FONT\nPAIRING",
            "badge_text": "Modern · Purple · WCAG AA ✅"
        }},
        {"type": "card_grid", "data": {
            "title": "Icons — Built-in Shapes",
            "cards": [
                {"icon": "check", "title": "Check", "color": "#{:02X}{:02X}{:02X}".format(*Colors.to_rgb_tuple(engine4.C.BLUE)), "items": ["✓ icon = 'check'"]},
                {"icon": "star",  "title": "Star",  "color": "#{:02X}{:02X}{:02X}".format(*Colors.to_rgb_tuple(engine4.C.TEAL)), "items": ["★ icon = 'star'"]},
                {"icon": "warning","title": "Warning","color": "#{:02X}{:02X}{:02X}".format(*Colors.to_rgb_tuple(engine4.C.WARM)),"items": ["⚠ icon = 'warning'"]},
                {"icon": "info",  "title": "Info",  "color": "#{:02X}{:02X}{:02X}".format(*Colors.to_rgb_tuple(engine4.C.RED)), "items": ["ℹ icon = 'info'"]},
            ]
        }},
        {"type": "closing", "data": {
            "pre_title": "Font: Arial Black + Arial", "main_title": "SELESAI",
            "subtitle": "Font style 'modern'", "source": "ppt_engine.py"
        }},
    ]
    engine4.build(slides4, source_text="Font style demo", output_path="/tmp/ppt_font_demo.pptx")
    print(f"✅ Demo 4 OK: /tmp/ppt_font_demo.pptx ({len(engine4.prs.slides)} slides)")
