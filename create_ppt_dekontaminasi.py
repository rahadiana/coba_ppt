from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# ── Color Palette ──
BG_DARK    = RGBColor(0x0A, 0x16, 0x28)   # deep navy
BG_CARD    = RGBColor(0x12, 0x29, 0x4A)   # card background
ACCENT     = RGBColor(0x0B, 0x7C, 0x72)   # teal (medical/procedures)
ACCENT2    = RGBColor(0xC8, 0x96, 0x2E)   # gold (highlights)
GREEN      = RGBColor(0x16, 0xA3, 0x4A)   # green (safe/clean)
RED        = RGBColor(0xDC, 0x26, 0x26)   # red (danger/warning)
ORANGE     = RGBColor(0xEA, 0x58, 0x0C)   # orange (caution)
TEAL       = RGBColor(0x0B, 0x7C, 0x72)   # teal accent
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0x94, 0xA3, 0xB8)
DARK_TEXT  = RGBColor(0x1A, 0x1A, 0x2E)
BLUE       = RGBColor(0x25, 0x63, 0xEB)   # blue (info)

def set_slide_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_shape_bg(slide, left, top, width, height, color, alpha=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    shape.shadow.inherit = False
    return shape

def add_circle(slide, left, top, size, color):
    shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top, size, size)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape

def add_text_box(slide, left, top, width, height, text, font_size=18, color=WHITE, bold=False, align=PP_ALIGN.LEFT, font_name='Calibri'):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = align
    return txBox

def add_rich_text_box(slide, left, top, width, height):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    return tf

# ══════════════════════════════════════════════════════════════
# SLIDE 1 — TITLE SLIDE
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
set_slide_bg(slide, BG_DARK)

# Decorative biohazard symbol (simplified as cross)
cross_v = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(6.15), Inches(0.8), Inches(0.5), Inches(1.5))
cross_v.fill.solid(); cross_v.fill.fore_color.rgb = ACCENT; cross_v.line.fill.background()
cross_h = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(5.55), Inches(1.25), Inches(1.7), Inches(0.5))
cross_h.fill.solid(); cross_h.fill.fore_color.rgb = ACCENT; cross_h.line.fill.background()

# Teal accent bar top
bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), prs.slide_width, Inches(0.08))
bar.fill.solid(); bar.fill.fore_color.rgb = ACCENT; bar.line.fill.background()

add_text_box(slide, Inches(1), Inches(2.8), Inches(11.3), Inches(1.2),
             "PROSEDUR DEKONTAMINASI", font_size=44, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
add_text_box(slide, Inches(1), Inches(3.7), Inches(11.3), Inches(1),
             "PERALATAN MEDIS BERDASARKAN STANDAR INTERNASIONAL", font_size=28, color=ACCENT, bold=True, align=PP_ALIGN.CENTER)

# Divider line
divider = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(4.5), Inches(4.9), Inches(4.3), Inches(0.04))
divider.fill.solid(); divider.fill.fore_color.rgb = LIGHT_GRAY; divider.line.fill.background()

add_text_box(slide, Inches(1), Inches(5.1), Inches(11.3), Inches(0.6),
             "Mengacu pada Standar WHO & CDC — Klasifikasi Spaulding",
             font_size=16, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)

add_text_box(slide, Inches(1), Inches(6.2), Inches(11.3), Inches(0.5),
             "Kementerian Kesehatan RI  |  CSSD (Central Sterile Supply Department)",
             font_size=14, color=ACCENT2, align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════════════
# SLIDE 2 — PENDAHULUAN
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)

bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), prs.slide_width, Inches(0.06))
bar.fill.solid(); bar.fill.fore_color.rgb = ACCENT; bar.line.fill.background()

# Section number
circle = add_circle(slide, Inches(0.6), Inches(0.4), Inches(0.7), ACCENT)
tf = circle.text_frame; tf.word_wrap = False
p = tf.paragraphs[0]; p.text = "00"; p.font.size = Pt(20); p.font.color.rgb = WHITE; p.font.bold = True; p.alignment = PP_ALIGN.CENTER
tf.paragraphs[0].space_before = Pt(6)

add_text_box(slide, Inches(1.5), Inches(0.45), Inches(5), Inches(0.6),
             "PENDAHULUAN", font_size=28, color=WHITE, bold=True)

# Content card
card = add_shape_bg(slide, Inches(0.6), Inches(1.5), Inches(12.1), Inches(5.5), BG_CARD)

tf = add_rich_text_box(slide, Inches(1.2), Inches(1.8), Inches(10.9), Inches(5))

p = tf.paragraphs[0]
p.text = "Mapa Dekontaminasi Peralatan Medis?"
p.font.size = Pt(24); p.font.color.rgb = ACCENT; p.font.bold = True

p = tf.add_paragraph(); p.space_before = Pt(20)
p.text = "Membersihkan dan mendekontaminasi peralatan medis bukan sekadar membuatnya terlihat bersih, melainkan prosedur vital untuk memutus rantai penyebaran infeksi."
p.font.size = Pt(18); p.font.color.rgb = WHITE

p = tf.add_paragraph(); p.space_before = Pt(14)
p.text = "Standar internasional yang dikeluarkan oleh:"
p.font.size = Pt(18); p.font.color.rgb = WHITE

bullets = [
    "WHO (World Health Organization)",
    "CDC (Centers for Disease Control and Prevention)",
    "Kementerian Kesehatan Republik Indonesia"
]
for b in bullets:
    p = tf.add_paragraph(); p.space_before = Pt(10); p.level = 1
    run = p.add_run(); run.text = "▸  "; run.font.size = Pt(18); run.font.color.rgb = ACCENT
    run2 = p.add_run(); run2.text = b; run2.font.size = Pt(18); run2.font.color.rgb = WHITE

p = tf.add_paragraph(); p.space_before = Pt(24)
run = p.add_run(); run.text = "Prinsip Utama: "; run.font.size = Pt(18); run.font.color.rgb = ACCENT2; run.font.bold = True
run2 = p.add_run(); run2.text = "Klasifikasi Spaulding — menentukan metode penanganan alat berdasarkan tingkat risiko penularan infeksi"
run2.font.size = Pt(17); run2.font.color.rgb = LIGHT_GRAY

# ══════════════════════════════════════════════════════════════
# SLIDE 3 — KLASIFIKASI SPAULDING (TABEL)
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)

bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), prs.slide_width, Inches(0.06))
bar.fill.solid(); bar.fill.fore_color.rgb = ACCENT; bar.line.fill.background()

circle = add_circle(slide, Inches(0.6), Inches(0.4), Inches(0.7), ACCENT)
tf = circle.text_frame
p = tf.paragraphs[0]; p.text = "01"; p.font.size = Pt(20); p.font.color.rgb = WHITE; p.font.bold = True; p.alignment = PP_ALIGN.CENTER
p.space_before = Pt(6)

add_text_box(slide, Inches(1.5), Inches(0.45), Inches(8), Inches(0.6),
             "KLASIFIKASI SPAULDING", font_size=28, color=WHITE, bold=True)

# Three category cards
cat_data = [
    ("🔴", "KRITIS", "Jaringan steril atau pembuluh darah", "STERILISASI\n(Wajib membunuh semua\nmikroba & spora)", "Instrumen bedah\nKateter\nImplan", RED),
    ("🟠", "SEMI-KRITIS", "Selaput lendir (mukosa)\natau kulit yang terluka", "DISINFEKSI TINGKAT\nTINGGI (DTT)\n(Membunuh semua mikroba,\nkecuali spora dalam jumlah besar)", "Endoskop\nSirkuit ventilator\nSpekulum", ORANGE),
    ("🟢", "NON-KRITIS", "Kulit yang utuh/sehat", "DISINFEKSI TINGKAT\nRENDAH / MENENGAH", "Stetoskop\nManset tensimeter\nPermukaan bed pasien", GREEN),
]

for i, (icon, title, area, treatment, examples, color) in enumerate(cat_data):
    left = Inches(0.6 + i * 4.1)
    card = add_shape_bg(slide, left, Inches(1.3), Inches(3.8), Inches(5.7), BG_CARD)

    # Color bar at top of card
    top_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, Inches(1.3), Inches(3.8), Inches(0.08))
    top_bar.fill.solid(); top_bar.fill.fore_color.rgb = color; top_bar.line.fill.background()

    add_text_box(slide, left + Inches(0.3), Inches(1.6), Inches(3.2), Inches(0.6),
                 icon + "  " + title, font_size=22, color=color, bold=True, align=PP_ALIGN.CENTER)

    # Area contact
    area_card = add_shape_bg(slide, left + Inches(0.2), Inches(2.3), Inches(3.4), Inches(1.0), RGBColor(0x0D, 0x1F, 0x3C))
    add_text_box(slide, left + Inches(0.3), Inches(2.35), Inches(3.2), Inches(0.9),
                 "Area Kontak:\n" + area, font_size=13, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)

    # Treatment
    treat_card = add_shape_bg(slide, left + Inches(0.2), Inches(3.5), Inches(3.4), Inches(1.5), RGBColor(0x1A, 0x2A, 0x4A))
    top_treat = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left + Inches(0.2), Inches(3.5), Inches(3.4), Inches(0.05))
    top_treat.fill.solid(); top_treat.fill.fore_color.rgb = color; top_treat.line.fill.background()
    add_text_box(slide, left + Inches(0.3), Inches(3.6), Inches(3.2), Inches(1.3),
                 treatment, font_size=13, color=WHITE, bold=True, align=PP_ALIGN.CENTER)

    # Examples
    add_text_box(slide, left + Inches(0.3), Inches(5.2), Inches(3.2), Inches(0.4),
                 "Contoh Alat:", font_size=14, color=ACCENT2, bold=True, align=PP_ALIGN.CENTER)
    add_text_box(slide, left + Inches(0.3), Inches(5.6), Inches(3.2), Inches(1.2),
                 examples, font_size=14, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════════════
# SLIDE 4 — PROSEDUR 1: PRE-CLEANING
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)

bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), prs.slide_width, Inches(0.06))
bar.fill.solid(); bar.fill.fore_color.rgb = ACCENT; bar.line.fill.background()

circle = add_circle(slide, Inches(0.6), Inches(0.4), Inches(0.7), ACCENT)
tf = circle.text_frame
p = tf.paragraphs[0]; p.text = "02"; p.font.size = Pt(20); p.font.color.rgb = WHITE; p.font.bold = True; p.alignment = PP_ALIGN.CENTER
p.space_before = Pt(6)

add_text_box(slide, Inches(1.5), Inches(0.45), Inches(10), Inches(0.6),
             "LANGKAH 1: PRE-CLEANING (Pembersihan Awal)", font_size=26, color=WHITE, bold=True)

# Main instruction card
card = add_shape_bg(slide, Inches(0.6), Inches(1.3), Inches(12.1), Inches(1.4), BG_CARD)
tf = add_rich_text_box(slide, Inches(1.2), Inches(1.5), Inches(10.9), Inches(1.0))
p = tf.paragraphs[0]
run = p.add_run(); run.text = "⏰  "; run.font.size = Pt(22); run.font.color.rgb = ORANGE
run2 = p.add_run(); run2.text = "Segera setelah alat digunakan — di tempat penggunaan!"; run2.font.size = Pt(20); run2.font.color.rgb = WHITE; run2.font.bold = True

# Steps cards
steps = [
    ("1", "BILAS LANGSUNG", "Bilas alat medis langsung di tempat penggunaan untuk mencegah darah, cairan tubuh, atau jaringan mengering.", TEAL),
    ("2", "CEGAH BIOFILM", "Kotoran yang mengering dapat membentuk biofilm (lapisan lendir pelindung bakteri) yang sangat sulit dibersihkan.", RED),
    ("3", "GUNAKAN MEDIA", "Gunakan air mengalir atau semprotkan busa enzimatik khusus (enzymatic foam).", BLUE),
]

for i, (num, title, desc, color) in enumerate(steps):
    left = Inches(0.6 + i * 4.1)
    card = add_shape_bg(slide, left, Inches(3.0), Inches(3.8), Inches(4.0), BG_CARD)
    top_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, Inches(3.0), Inches(3.8), Inches(0.06))
    top_bar.fill.solid(); top_bar.fill.fore_color.rgb = color; top_bar.line.fill.background()

    circle_num = add_circle(slide, left + Inches(1.4), Inches(3.3), Inches(1), color)
    tf_num = circle_num.text_frame
    p = tf_num.paragraphs[0]; p.text = num; p.font.size = Pt(32); p.font.color.rgb = WHITE; p.font.bold = True; p.alignment = PP_ALIGN.CENTER
    p.space_before = Pt(4)

    add_text_box(slide, left + Inches(0.3), Inches(4.5), Inches(3.2), Inches(0.5),
                 title, font_size=18, color=color, bold=True, align=PP_ALIGN.CENTER)

    add_text_box(slide, left + Inches(0.3), Inches(5.1), Inches(3.2), Inches(1.6),
                 desc, font_size=15, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════════════
# SLIDE 5 — PROSEDUR 2: PEMBERSIHAN UTAMA
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)

bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), prs.slide_width, Inches(0.06))
bar.fill.solid(); bar.fill.fore_color.rgb = ACCENT; bar.line.fill.background()

circle = add_circle(slide, Inches(0.6), Inches(0.4), Inches(0.7), ACCENT)
tf = circle.text_frame
p = tf.paragraphs[0]; p.text = "03"; p.font.size = Pt(20); p.font.color.rgb = WHITE; p.font.bold = True; p.alignment = PP_ALIGN.CENTER
p.space_before = Pt(6)

add_text_box(slide, Inches(1.5), Inches(0.45), Inches(10), Inches(0.6),
             "LANGKAH 2: PEMBERSIHAN UTAMA (Cleaning & Scrubbing)", font_size=26, color=WHITE, bold=True)

# APD warning card
warn_card = add_shape_bg(slide, Inches(0.6), Inches(1.3), Inches(12.1), Inches(1.2), RGBColor(0x3A, 0x1A, 0x1A))
tf = add_rich_text_box(slide, Inches(1.2), Inches(1.4), Inches(10.9), Inches(1.0))
p = tf.paragraphs[0]
run = p.add_run(); run.text = "⚠️  "; run.font.size = Pt(22); run.font.color.rgb = RED
run2 = p.add_run(); run2.text = "WAJIB mengenakan APD lengkap sebelum memulai proses pembersihan!"; run2.font.size = Pt(18); run2.font.color.rgb = RED; run2.font.bold = True

# Main steps
card = add_shape_bg(slide, Inches(0.6), Inches(2.8), Inches(12.1), Inches(4.3), BG_CARD)
tf = add_rich_text_box(slide, Inches(1.2), Inches(3.0), Inches(10.9), Inches(4.0))

p = tf.paragraphs[0]
p.text = "Ruang Dekontaminasi (CSSD)"
p.font.size = Pt(22); p.font.color.rgb = ACCENT; p.font.bold = True

p = tf.add_paragraph(); p.space_before = Pt(16)
p.text = "Bawa alat ke ruang dekontaminasi (CSSD — Central Sterile Supply Department)"
p.font.size = Pt(17); p.font.color.rgb = WHITE

p = tf.add_paragraph(); p.space_before = Pt(14)
p.text = "Langkah Pembersihan:"
p.font.size = Pt(18); p.font.color.rgb = ACCENT2; p.font.bold = True

steps_main = [
    ("Rendam alat dalam larutan detergen enzimatik sesuai dosis penunjuk kemasan"),
    ("Lakukan penyikatan secara manual menggunakan sikat lembut di bawah permukaan air"),
    ("Hindari cipratan cairan (aerosol) ke wajah Anda"),
    ("Untuk alat berongga (lumen), gunakan sikat tabung yang sesuai"),
]

for s in steps_main:
    p = tf.add_paragraph(); p.space_before = Pt(10); p.level = 1
    run = p.add_run(); run.text = "▸  "; run.font.size = Pt(17); run.font.color.rgb = ACCENT
    run2 = p.add_run(); run2.text = s; run2.font.size = Pt(17); run2.font.color.rgb = WHITE

# ══════════════════════════════════════════════════════════════
# SLIDE 6 — PROSEDUR 3: PEMBILASAN & PENGERINGAN
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)

bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), prs.slide_width, Inches(0.06))
bar.fill.solid(); bar.fill.fore_color.rgb = ACCENT; bar.line.fill.background()

circle = add_circle(slide, Inches(0.6), Inches(0.4), Inches(0.7), ACCENT)
tf = circle.text_frame
p = tf.paragraphs[0]; p.text = "04"; p.font.size = Pt(20); p.font.color.rgb = WHITE; p.font.bold = True; p.alignment = PP_ALIGN.CENTER
p.space_before = Pt(6)

add_text_box(slide, Inches(1.5), Inches(0.45), Inches(10), Inches(0.6),
             "LANGKAH 3: PEMBILASAN & PENGERINGAN", font_size=26, color=WHITE, bold=True)

# Two cards side by side
# Bilas card
card1 = add_shape_bg(slide, Inches(0.6), Inches(1.3), Inches(5.9), Inches(5.7), BG_CARD)
top_bar1 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.6), Inches(1.3), Inches(5.9), Inches(0.06))
top_bar1.fill.solid(); top_bar1.fill.fore_color.rgb = BLUE; top_bar1.line.fill.background()

tf1 = add_rich_text_box(slide, Inches(1.0), Inches(1.6), Inches(5.1), Inches(5.2))
p = tf1.paragraphs[0]
run = p.add_run(); run.text = "💧  "; run.font.size = Pt(28); run.font.color.rgb = BLUE
run2 = p.add_run(); run2.text = "PEMBILASAN"; run2.font.size = Pt(24); run2.font.color.rgb = BLUE; run2.font.bold = True

p = tf1.add_paragraph(); p.space_before = Pt(20)
p.text = "Menghilangkan residu kimia dari alat"
p.font.size = Pt(17); p.font.color.rgb = WHITE

p = tf1.add_paragraph(); p.space_before = Pt(16)
p.text = "▸  Bilas seluruh alat dengan air bersih yang mengalir"
p.font.size = Pt(16); p.font.color.rgb = LIGHT_GRAY

p = tf1.add_paragraph(); p.space_before = Pt(10)
p.text = "▸  Diutamakan menggunakan air murni (demineralized water)"
p.font.size = Pt(16); p.font.color.rgb = LIGHT_GRAY

p = tf1.add_paragraph(); p.space_before = Pt(10)
p.text = "▸  Air murni mencegah terbentuknya kerak mineral pada alat"
p.font.size = Pt(16); p.font.color.rgb = LIGHT_GRAY

# Keringkan card
card2 = add_shape_bg(slide, Inches(6.8), Inches(1.3), Inches(5.9), Inches(5.7), BG_CARD)
top_bar2 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(6.8), Inches(1.3), Inches(5.9), Inches(0.06))
top_bar2.fill.solid(); top_bar2.fill.fore_color.rgb = GREEN; top_bar2.line.fill.background()

tf2 = add_rich_text_box(slide, Inches(7.2), Inches(1.6), Inches(5.1), Inches(5.2))
p = tf2.paragraphs[0]
run = p.add_run(); run.text = "🌬️  "; run.font.size = Pt(28); run.font.color.rgb = GREEN
run2 = p.add_run(); run2.text = "PENGERINGAN"; run2.font.size = Pt(24); run2.font.color.rgb = GREEN; run2.font.bold = True

p = tf2.add_paragraph(); p.space_before = Pt(20)
p.text = "Memastikan alat sepenuhnya kering"
p.font.size = Pt(17); p.font.color.rgb = WHITE

p = tf2.add_paragraph(); p.space_before = Pt(16)
p.text = "▸  Gunakan kain bersih bebas serat (lint-free cloth)"
p.font.size = Pt(16); p.font.color.rgb = LIGHT_GRAY

p = tf2.add_paragraph(); p.space_before = Pt(10)
p.text = "▸  Atau gunakan udara bertekanan medis"
p.font.size = Pt(16); p.font.color.rgb = LIGHT_GRAY

p = tf2.add_paragraph(); p.space_before = Pt(10)
p.text = "▸  Alat yang lembap menurunkan efektivitas sterilisasi gas atau kimia"
p.font.size = Pt(16); p.font.color.rgb = LIGHT_GRAY

# ══════════════════════════════════════════════════════════════
# SLIDE 7 — PROSEDUR 4: DISINFEKSI / STERILISASI
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)

bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), prs.slide_width, Inches(0.06))
bar.fill.solid(); bar.fill.fore_color.rgb = ACCENT; bar.line.fill.background()

circle = add_circle(slide, Inches(0.6), Inches(0.4), Inches(0.7), ACCENT)
tf = circle.text_frame
p = tf.paragraphs[0]; p.text = "05"; p.font.size = Pt(20); p.font.color.rgb = WHITE; p.font.bold = True; p.alignment = PP_ALIGN.CENTER
p.space_before = Pt(6)

add_text_box(slide, Inches(1.5), Inches(0.45), Inches(10), Inches(0.6),
             "LANGKAH 4: DISINFEKSI / STERILISASI UTAMA", font_size=26, color=WHITE, bold=True)

# Three method cards
methods = [
    ("🔴", "ALAT KRITIS", "STERILISASI", [
        "Mesin Autoklaf (uap panas bertekanan)",
        "Suhu 121°C atau 134°C",
        "Sterilisasi plasma/suhu rendah untuk alat sensitif panas"
    ], RED),
    ("🟠", "ALAT SEMI-KRITIS", "DISINFEKSI TINGKAT TINGGI", [
        "Bahan kimia: glutaraldehyde atau hydrogen peroxide",
        "Waktu perendaman sesuai standar baku",
        "Membunuh semua mikroba kecuali spora"
    ], ORANGE),
    ("🟢", "ALAT NON-KRITIS", "DISINFEKSI TINGKAT RENDAH", [
        "Seka dengan disinfektan tingkat rendah",
        "Alkohol 70% atau cairan berbasis klorin",
        "Cukup untuk kulit utuh/sehat"
    ], GREEN),
]

for i, (icon, category, method, details, color) in enumerate(methods):
    left = Inches(0.6 + i * 4.1)
    card = add_shape_bg(slide, left, Inches(1.3), Inches(3.8), Inches(5.7), BG_CARD)
    top_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, Inches(1.3), Inches(3.8), Inches(0.08))
    top_bar.fill.solid(); top_bar.fill.fore_color.rgb = color; top_bar.line.fill.background()

    add_text_box(slide, left + Inches(0.3), Inches(1.6), Inches(3.2), Inches(0.6),
                 icon + "  " + category, font_size=20, color=color, bold=True, align=PP_ALIGN.CENTER)

    method_card = add_shape_bg(slide, left + Inches(0.2), Inches(2.4), Inches(3.4), Inches(1.0), RGBColor(0x1A, 0x2A, 0x4A))
    add_text_box(slide, left + Inches(0.3), Inches(2.5), Inches(3.2), Inches(0.8),
                 method, font_size=16, color=WHITE, bold=True, align=PP_ALIGN.CENTER)

    add_text_box(slide, left + Inches(0.3), Inches(3.6), Inches(3.2), Inches(0.4),
                 "Metode:", font_size=14, color=ACCENT2, bold=True, align=PP_ALIGN.CENTER)

    tf_detail = add_rich_text_box(slide, left + Inches(0.3), Inches(4.0), Inches(3.2), Inches(2.8))
    for d in details:
        p = tf_detail.add_paragraph(); p.space_before = Pt(8)
        p.text = "▸  " + d; p.font.size = Pt(14); p.font.color.rgb = LIGHT_GRAY

# ══════════════════════════════════════════════════════════════
# SLIDE 8 — PROSEDUR 5: PENYIMPANAN BERSIH
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)

bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), prs.slide_width, Inches(0.06))
bar.fill.solid(); bar.fill.fore_color.rgb = ACCENT; bar.line.fill.background()

circle = add_circle(slide, Inches(0.6), Inches(0.4), Inches(0.7), ACCENT)
tf = circle.text_frame
p = tf.paragraphs[0]; p.text = "06"; p.font.size = Pt(20); p.font.color.rgb = WHITE; p.font.bold = True; p.alignment = PP_ALIGN.CENTER
p.space_before = Pt(6)

add_text_box(slide, Inches(1.5), Inches(0.45), Inches(10), Inches(0.6),
             "LANGKAH 5: PENYIMPANAN BERSIH", font_size=26, color=WHITE, bold=True)

# Main instruction card
card = add_shape_bg(slide, Inches(0.6), Inches(1.3), Inches(12.1), Inches(5.7), BG_CARD)

tf = add_rich_text_box(slide, Inches(1.2), Inches(1.6), Inches(10.9), Inches(5.2))

p = tf.paragraphs[0]
p.text = "Menjaga Sterilitas Alat"
p.font.size = Pt(24); p.font.color.rgb = GREEN; p.font.bold = True

p = tf.add_paragraph(); p.space_before = Pt(20)
p.text = "Alat yang sudah disterilkan harus disimpan dengan benar untuk mempertahankan kondisi steril:"
p.font.size = Pt(18); p.font.color.rgb = WHITE

p = tf.add_paragraph(); p.space_before = Pt(16)
run = p.add_run(); run.text = "▸  "; run.font.size = Pt(18); run.font.color.rgb = GREEN
run2 = p.add_run(); run2.text = "Simpan dalam kemasan khusus (pouch) tertutup"; run2.font.size = Pt(18); run2.font.color.rgb = WHITE

p = tf.add_paragraph(); p.space_before = Pt(10)
run = p.add_run(); run.text = "▸  "; run.font.size = Pt(18); run.font.color.rgb = GREEN
run2 = p.add_run(); run2.text = "Masukkan ke dalam lemari tertutup dengan suhu dan kelembapan terkontrol"; run2.font.size = Pt(18); run2.font.color.rgb = WHITE

p = tf.add_paragraph(); p.space_before = Pt(16)
run = p.add_run(); run.text = "Prinsip FIFO: "; run.font.size = Pt(18); run.font.color.rgb = ACCENT2; run.font.bold = True
run2 = p.add_run(); run2.text = "First In, First Out — alat yang disterilkan lebih dulu digunakan terlebih dahulu"
run2.font.size = Pt(17); run2.font.color.rgb = LIGHT_GRAY

p = tf.add_paragraph(); p.space_before = Pt(20)
p.text = "Pemeriksaan Kemasan:"
p.font.size = Pt(18); p.font.color.rgb = ACCENT; p.font.bold = True

p = tf.add_paragraph(); p.space_before = Pt(10)
run = p.add_run(); run.text = "▸  "; run.font.size = Pt(18); run.font.color.rgb = ACCENT
run2 = p.add_run(); run2.text = "Periksa indikator kimia pada kemasan untuk memastikan sterilisasi berhasil"; run2.font.size = Pt(18); run2.font.color.rgb = WHITE

p = tf.add_paragraph(); p.space_before = Pt(10)
run = p.add_run(); run.text = "▸  "; run.font.size = Pt(18); run.font.color.rgb = ACCENT
run2 = p.add_run(); run2.text = "Indikator berubah warna menandakan proses sterilisasi selesai sempurna"; run2.font.size = Pt(18); run2.font.color.rgb = WHITE

# ══════════════════════════════════════════════════════════════
# SLIDE 9 — PERINGATAN KESELAMATAN APD
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)

bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), prs.slide_width, Inches(0.06))
bar.fill.solid(); bar.fill.fore_color.rgb = RED; bar.line.fill.background()

# Warning banner
warn_bg = add_shape_bg(slide, Inches(0.6), Inches(0.3), Inches(12.1), Inches(1.4), RGBColor(0x5C, 0x1A, 0x1A))
tf = add_rich_text_box(slide, Inches(1.2), Inches(0.4), Inches(10.9), Inches(1.2))
p = tf.paragraphs[0]
run = p.add_run(); run.text = "⚠️  PERINGATAN KESELAMATAN PENTING  ⚠️"
run.font.size = Pt(24); run.font.color.rgb = RED; run.font.bold = True; p.alignment = PP_ALIGN.CENTER

# APD items
apd_items = [
    ("😷", "Masker Pelindung Wajah\n(Faceshield)", "Melindungi wajah dari cipratan cairan dan aerosol berbahaya"),
    ("👔", "Gaun Kedap Air\n(Waterproof Gown)", "Mencegah kontaminasi cairan tubuh ke tubuh petugas"),
    ("🧤", "Apron Pelindung", "Lapisan tambahan untuk proteksi tubuh bagian depan"),
    ("🖐️", "Sarung Tangan Tebal\n(Utility Gloves)", "Khusus dekontaminasi — melindungi dari tusukan jarum dan instrumen tajam"),
]

for i, (icon, item, desc) in enumerate(apd_items):
    left = Inches(0.6 + i * 3.15)
    card = add_shape_bg(slide, left, Inches(2.0), Inches(2.95), Inches(5.0), BG_CARD)
    top_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, Inches(2.0), Inches(2.95), Inches(0.06))
    top_bar.fill.solid(); top_bar.fill.fore_color.rgb = RED; top_bar.line.fill.background()

    add_text_box(slide, left + Inches(0.2), Inches(2.3), Inches(2.55), Inches(0.8),
                 icon, font_size=40, align=PP_ALIGN.CENTER)

    add_text_box(slide, left + Inches(0.2), Inches(3.2), Inches(2.55), Inches(1.0),
                 item, font_size=16, color=WHITE, bold=True, align=PP_ALIGN.CENTER)

    add_text_box(slide, left + Inches(0.2), Inches(4.4), Inches(2.55), Inches(2.2),
                 desc, font_size=14, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════════════
# SLIDE 10 — RINGKASAN & PENUTUP
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)

bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), prs.slide_width, Inches(0.06))
bar.fill.solid(); bar.fill.fore_color.rgb = ACCENT; bar.line.fill.background()

add_text_box(slide, Inches(1), Inches(0.5), Inches(11.3), Inches(0.8),
             "RINGKASAN — 5 Langkah Dekontaminasi", font_size=28, color=WHITE, bold=True, align=PP_ALIGN.CENTER)

# Summary cards in a row
summary = [
    ("01", "PRE-\nCLEANING", "Bilas langsung\nsetelah pemakaian\nuntuk cegah biofilm", ORANGE),
    ("02", "PEMBERSIHAN\nUTAMA", "Rendam detergen\nenzimatik + sikat\nmanual di bawah air", ACCENT),
    ("03", "PEMBILASAN\n& PENGERINGAN", "Air bersih mengalir\n+ kain bebas serat\natau udara tekan", BLUE),
    ("04", "STERILISASI\n/ DISINFEKSI", "Sesuai klasifikasi\nSpaulding: Autoklaf\natau bahan kimia", RED),
    ("05", "PENYIMPANAN\nBERSIH", "Pouch tertutup\nFIFO + indikator\nkimia steril", GREEN),
]

for i, (num, title, desc, color) in enumerate(summary):
    left = Inches(0.4 + i * 2.55)
    card = add_shape_bg(slide, left, Inches(1.6), Inches(2.35), Inches(5.2), BG_CARD)
    top_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, Inches(1.6), Inches(2.35), Inches(0.06))
    top_bar.fill.solid(); top_bar.fill.fore_color.rgb = color; top_bar.line.fill.background()

    circle_num = add_circle(slide, left + Inches(0.65), Inches(1.9), Inches(1), color)
    tf_num = circle_num.text_frame
    p = tf_num.paragraphs[0]; p.text = num; p.font.size = Pt(28); p.font.color.rgb = WHITE; p.font.bold = True; p.alignment = PP_ALIGN.CENTER
    p.space_before = Pt(4)

    add_text_box(slide, left + Inches(0.15), Inches(3.1), Inches(2.05), Inches(1),
                 title, font_size=16, color=color, bold=True, align=PP_ALIGN.CENTER)

    add_text_box(slide, left + Inches(0.15), Inches(4.2), Inches(2.05), Inches(2.2),
                 desc, font_size=13, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)

# Bottom message
add_text_box(slide, Inches(1), Inches(7.0), Inches(11.3), Inches(0.4),
             "Standar WHO & CDC  |  Klasifikasi Spaulding  |  Patuhi Protokol Dekontaminasi!",
             font_size=16, color=ACCENT2, bold=True, align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════════════
# SAVE
# ══════════════════════════════════════════════════════════════
output_path = "/home/runner/coba_ppt/Dekontaminasi_Peralatan_Medis.pptx"
prs.save(output_path)
print(f"✅ PPT saved to: {output_path}")
