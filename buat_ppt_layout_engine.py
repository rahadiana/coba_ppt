#!/usr/bin/env python3
"""
Perwal Bekasi No 51/2024 — Pajak Reklame
VERSI LAYOUT ENGINE — Dengan LayoutFrame class untuk kalkulasi presisi

RUMUS KALKULASI LAYOUT (dokumentasi):
════════════════════════════════════════════════════════════

Slide Widescreen 16:9 = 13.333" × 7.5"

┌─ LAYOUT FRAME ─────────────────────────────────────────┐
│                                                         │
│  ← M → ════════ CONTENT AREA ════════ ← M →            │
│         ┌──────────────────────────────┐                │
│  ↑ HDR  │ [ Header Bar ]              │                │
│  ↑ GAP  │                             │                │
│         │                             │                │
│  CONT   │  ┌───┐  ┌───┐  ┌───┐       │                │
│  AREA   │  │ C │  │ C │  │ C │       │                │
│         │  │ A │  │ A │  │ A │       │                │
│         │  │ R │  │ R │  │ R │       │                │
│         │  │ D │  │ D │  │ D │       │                │
│         │  └───┘  └───┘  └───┘       │                │
│         │                             │                │
│  ↑ GAP  │ [ Source + Page # ]        │                │
│  ↑ FTR  │                             │                │
│         └──────────────────────────────┘                │
└─────────────────────────────────────────────────────────┘

FORMULA KUNCI:
════════════════════════════════════════════════════════════

1. LEBAR KOLOM untuk N item dalam grid:
   ┌────────────────────────────────────────────┐
   │  col_w = (W - (N-1) × gap) / N             │
   │                                            │
   │  Contoh: W=12.13", N=4, gap=0.3"           │
   │  col_w = (12.13 - 3×0.3) / 4              │
   │        = (12.13 - 0.9) / 4                │
   │        = 11.23 / 4 = 2.808" ✅              │
   └────────────────────────────────────────────┘

2. TINGGI BARIS untuk M baris:
   ┌────────────────────────────────────────────┐
   │  row_h = (H - (M-1) × gap_v) / M          │
   └────────────────────────────────────────────┘

3. POSISI ELEMEN di grid:
   ┌────────────────────────────────────────────┐
   │  x = content_x + col × (col_w + gap_h)    │
   │  y = content_y + row × (row_h + gap_v)    │
   └────────────────────────────────────────────┘

4. ESTIMASI TINGGI TEKS (anti-overlap):
   ┌────────────────────────────────────────────┐
   │  chars_per_line = box_w × cpi             │
   │  num_lines = ceil(len(text) / cpi / box_w) │
   │  min_height = num_lines × pt × 1.2 / 72   │
   │                                            │
   │  cpi (Calibri): 9pt=14, 11pt=12, 13pt=10  │
   └────────────────────────────────────────────┘

5. VERIFIKASI BOUND:
   ┌────────────────────────────────────────────┐
   │  x ≥ 0 AND y ≥ 0                          │
   │  AND x+w ≤ SLIDE_W                        │
   │  AND y+h ≤ SLIDE_H                        │
   └────────────────────────────────────────────┘
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os, math

# ═══════════════════════════════════════════════════════════
# COLOR PALETTE (60-30-10 Rule)
# ═══════════════════════════════════════════════════════════

class Colors:
    # 60% Dominant
    NAVY   = RGBColor(0x0A, 0x16, 0x28)
    NAVY_L = RGBColor(0x12, 0x29, 0x4A)
    NAVY_M = RGBColor(0x1B, 0x3A, 0x6B)
    NAVY_D = RGBColor(0x0D, 0x1F, 0x3C)
    # 30% Secondary
    WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
    OFF_W  = RGBColor(0xF5, 0xF7, 0xFA)
    ICE    = RGBColor(0xE8, 0xED, 0xF5)
    ICE_D  = RGBColor(0xD0, 0xD8, 0xE8)
    # 10% Accent
    GOLD   = RGBColor(0xC8, 0x96, 0x2E)
    GOLD_L = RGBColor(0xD4, 0xA0, 0x17)
    # Semantic
    TEXT_D = RGBColor(0x1A, 0x1A, 0x2E)
    TEXT_M = RGBColor(0x6B, 0x72, 0x88)
    TEXT_L = RGBColor(0x9C, 0xA3, 0xAF)
    BLUE   = RGBColor(0x25, 0x63, 0xEB)
    TEAL   = RGBColor(0x0D, 0x94, 0x88)
    WARM   = RGBColor(0xB8, 0x86, 0x0B)
    RED    = RGBColor(0xDC, 0x26, 0x26)
C = Colors  # shorthand

# ═══════════════════════════════════════════════════════════
# LAYOUT FRAME — Semua kalkulasi layout di sini
# ═══════════════════════════════════════════════════════════

class LayoutFrame:
    """
    LayoutFrame: kalkulator posisi & ukuran untuk element slide.
    
    Slide widescreen 16:9 = 13.333" × 7.5"
    
    Area layout:
    ┌──────────┬──────────────────────┬──────────┐
    │  margin   │    CONTENT AREA      │  margin  │
    │  (0.6")  │                      │  (0.6")  │
    ├──────────┴──────────────────────┴──────────┤
    │  HEADER: gold_bar(0.04") + navy(0.9")     │  ← 0.94"
    │  GAP: 0.21"                                │
    ├────────────────────────────────────────────┤
    │  CONTENT: tinggi = content_h               │
    │  ┌────┐ ┌────┐ ┌────┐ ┌────┐              │
    │  │ C1 │ │ C2 │ │ C3 │ │ C4 │              │
    │  └────┘ └────┘ └────┘ └────┘              │
    ├────────────────────────────────────────────┤
    │  FOOTER: 0.42" (source + page #)          │
    └────────────────────────────────────────────┘
    """
    
    # ─── Konstanta Slide ───
    SLIDE_W = 13.333        # inches
    SLIDE_H = 7.5           # inches
    
    # ─── Zona Layout ───
    MARGIN_H = 0.6          # margin horizontal kiri/kanan
    HEADER_H = 0.94         # tinggi header (gold_bar 0.04 + navy 0.9)
    HEADER_GAP = 0.21       # jarak header → konten
    FOOTER_H = 0.50         # tinggi footer (navy bar 0.03 + text 0.22 + gap)
    
    # ─── Spasi Default Antar Elemen ───
    GAP_H = 0.30            # gap horizontal antar card
    GAP_V = 0.30            # gap vertikal antar baris card
    GAP_ITEM = 0.02         # gap antar item dalam card
    CARD_PAD = 0.15         # padding dalam card
    
    # ─── Font Metrics (Calibri) ───
    # Characters-per-inch (CPI) untuk Calibri
    CPI = { 7: 17, 8: 15, 9: 14, 10: 13, 11: 12, 12: 11, 
            13: 10, 14: 9, 16: 8, 18: 7, 20: 6.5, 24: 5.5, 
            32: 4, 40: 3, 48: 2.5 }
    
    def __init__(self):
        """Hitung semua zona layout."""
        # Content area horizontal
        self.cx = self.MARGIN_H
        self.cw = self.SLIDE_W - 2 * self.MARGIN_H  # 12.133"
        
        # Content area vertikal
        header_bottom = self.HEADER_H
        self.cy = header_bottom + self.HEADER_GAP    # 1.15"
        footer_top = self.SLIDE_H - self.FOOTER_H    # 7.0"
        self.ch = footer_top - self.cy                # 5.85"
    
    # ─── Kalkulasi Grid ───
    
    def col_width(self, n, gap=None):
        """
        Hitung lebar per kolom untuk N item dalam grid.
        
        FORMULA:
            col_w = (cw - (n-1) × gap) / n
        
        Args:
            n: Jumlah kolom
            gap: Spasi antar kolom (default: GAP_H)
        """
        g = gap if gap is not None else self.GAP_H
        return (self.cw - (n - 1) * g) / n
    
    def row_height(self, n, gap=None):
        """
        Hitung tinggi per baris untuk N baris dalam grid.
        
        FORMULA:
            row_h = (ch - (n-1) × gap_v) / n
        """
        g = gap if gap is not None else self.GAP_V
        return (self.ch - (n - 1) * g) / n
    
    def grid_pos(self, col, row, col_w, row_h, gap_h=None, gap_v=None):
        """
        Hitung posisi x, y untuk cell grid.
        
        FORMULA:
            x = cx + col × (col_w + gap_h)
            y = cy + row × (row_h + gap_v)
        
        Returns:
            (x_inches, y_inches) — posisi dalam inches
        """
        gh = gap_h if gap_h is not None else self.GAP_H
        gv = gap_v if gap_v is not None else self.GAP_V
        x = self.cx + col * (col_w + gh)
        y = self.cy + row * (row_h + gv)
        return x, y
    
    # ─── Kalkulasi Multi-Baris Grid ───
    
    def calc_grid(self, n_items, cols=0):
        """
        Kalkulasi lengkap untuk grid multi-baris.
        
        Args:
            n_items: Jumlah total item
            cols: Jumlah kolom (0 = auto: 2/3/4 tergantung n)
        
        Returns:
            dict dengan: per_row, n_rows, col_w, row_h, gap_h, gap_v
        """
        per_row = cols if cols > 0 else (
            2 if n_items <= 2 else (3 if n_items <= 3 else 4))
        n_rows = (n_items + per_row - 1) // per_row
        col_w = self.col_width(per_row)
        row_h = self.row_height(n_rows)
        return {
            'per_row': per_row, 'n_rows': n_rows,
            'col_w': col_w, 'row_h': row_h,
            'gap_h': self.GAP_H, 'gap_v': self.GAP_V,
        }
    
    def grid_cell(self, i, grid):
        """
        Hitung posisi cell ke-i dalam grid.
        
        Args:
            i: Index item (0-based)
            grid: Hasil dari calc_grid()
        
        Returns:
            (x, y, col_w, row_h) — posisi & ukuran dalam inches
        """
        col = i % grid['per_row']
        row = i // grid['per_row']
        x, y = self.grid_pos(col, row, grid['col_w'], grid['row_h'],
                             grid['gap_h'], grid['gap_v'])
        return x, y, grid['col_w'], grid['row_h']
    
    # ─── Kalkulasi Card Internal ───
    
    def card_item_layout(self, card_h, has_icon=True, n_items=1, title_h=0.30):
        """
        Hitung layout internal dalam card.
        
        Args:
            card_h: Tinggi card dalam inches
            has_icon: Apakah ada icon circle
            n_items: Jumlah item bullet
            title_h: Tinggi title
        
        Returns:
            dict: {icon_y, icon_size, title_y, item_start_y, item_h}
        """
        icon_size = 0.42
        icon_y = self.CARD_PAD
        
        if has_icon:
            title_y = icon_y + icon_size + 0.13  # gap setelah icon
        else:
            title_y = self.CARD_PAD
        
        # Item height: minimal untuk 1 baris teks
        # CPI 9pt = 14 chars/inch
        # Asumsi teks maks 45 chars dalam 2.5" → perlu 2 baris
        item_h = 0.35  # cukup untuk 2 baris teks 9pt
        
        item_start_y = title_y + title_h + 0.08
        content_used = item_start_y + n_items * (item_h + self.GAP_ITEM)
        
        # Validasi: content_used harus ≤ card_h
        # Jika overflow, kurangi item_h atau peringatkan
        overflow = content_used - card_h
        if overflow > 0.2:  # toleransi 0.2"
            # Terlalu banyak konten untuk card — kurangi font size
            pass  # caller bisa handle
        
        return {
            'icon_size': icon_size,
            'icon_y': icon_y,
            'title_y': title_y,
            'item_start_y': item_start_y,
            'item_h': item_h,
            'overflow': max(0, content_used - card_h),
        }
    
    # ─── Estimasi Tinggi Teks ───
    
    def text_height(self, text, font_size, box_width):
        """
        Estimasi tinggi teks dalam inches.
        
        FORMULA:
            cpi = CPI[font_size]          # chars per inch
            max_chars = box_width × cpi   # chars per line
            n_lines = ceil(len(text) / max_chars)
            height = n_lines × font_size × 1.2 / 72  # inches
        
        Args:
            text: Teks yang akan di-render
            font_size: Ukuran font dalam pt
            box_width: Lebar textbox dalam inches
        
        Returns:
            Tinggi estimated dalam inches
        """
        # Cari CPI terdekat
        cpi = self.CPI.get(font_size, 12)
        max_chars_per_line = box_width * cpi
        
        if max_chars_per_line <= 0:
            return font_size * 1.2 / 72
        
        n_lines = math.ceil(len(text) / max_chars_per_line) if max_chars_per_line > 0 else 1
        line_h = font_size * 1.2 / 72  # 1.2 = line spacing
        return n_lines * line_h + 0.02  # + tiny padding
    
    def safe_item_height(self, text, font_size, box_width, min_h=0.25):
        """
        Dapatkan item height yang aman (tidak overflow).
        
        Returns:
            Tinggi dalam inches, guaranteed ≥ min_h
        """
        estimated = self.text_height(text, font_size, box_width)
        return max(estimated + 0.05, min_h)
    
    # ─── Verifikasi ───
    
    def check_bounds(self, x, y, w, h, name=""):
        """
        Verifikasi element dalam bounds slide.
        
        Returns:
            (is_safe, [list of warnings])
        """
        warnings = []
        if x < 0: warnings.append(f"x={x:.2f}<0")
        if y < 0: warnings.append(f"y={y:.2f}<0")
        if x + w > self.SLIDE_W: warnings.append(f"right={x+w:.2f}>{self.SLIDE_W}")
        if y + h > self.SLIDE_H: warnings.append(f"bottom={y+h:.2f}>{self.SLIDE_H}")
        return (len(warnings) == 0, warnings)
    
    def check_overlap(self, a, b):
        """
        Cek apakah dua rectangle overlap secara signifikan.
        
        Args:
            a, b: (x, y, w, h) tuples
        
        Returns:
            overlap_pct: persentase overlap terhadap rectangle terkecil
        """
        ax, ay, aw, ah = a
        bx, by, bw, bh = b
        
        ox = max(0, min(ax+aw, bx+bw) - max(ax, bx))
        oy = max(0, min(ay+ah, by+bh) - max(ay, by))
        overlap = ox * oy
        min_area = min(aw * ah, bw * bh)
        
        if min_area == 0:
            return 0
        return (overlap / min_area) * 100


# ═══════════════════════════════════════════════════════════
# DRAW HELPERS
# ═══════════════════════════════════════════════════════════

L = LayoutFrame()  # global layout instance

prs = Presentation()
prs.slide_width = Inches(L.SLIDE_W)
prs.slide_height = Inches(L.SLIDE_H)
pg_counter = [0]
SOURCE_TEXT = "Sumber: Perwal Bekasi No 51/2024"


def blank():
    return prs.slides.add_slide(prs.slide_layouts[6])

def add_box(slide, left, top, width, height, text, size=12, bold=False,
            color=C.TEXT_D, align=PP_ALIGN.LEFT, font="Calibri", vAlign=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font
    p.alignment = align
    p.space_after = Pt(0)
    p.space_before = Pt(0)
    L.check_bounds(left, top, width, height, f"TextBox {text[:20]}")
    return tb

def add_shape(slide, stype, left, top, width, height, fill=None, line=None, lw=None):
    sh = slide.shapes.add_shape(stype, Inches(left), Inches(top), Inches(width), Inches(height))
    if fill:
        sh.fill.solid()
        sh.fill.fore_color.rgb = fill
    else:
        sh.fill.background()
    if line:
        sh.line.color.rgb = line
        if lw: sh.line.width = Pt(lw)
    else:
        sh.line.fill.background()
    L.check_bounds(left, top, width, height, f"Shape")
    return sh

def add_rrect(slide, left, top, width, height, fill=C.WHITE, line=C.ICE, lw=0.5, radius=0.04):
    sh = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top),
                                 Inches(width), Inches(height))
    sh.fill.solid()
    sh.fill.fore_color.rgb = fill
    if line:
        sh.line.color.rgb = line
        sh.line.width = Pt(lw)
    else:
        sh.line.fill.background()
    sh.adjustments[0] = radius
    L.check_bounds(left, top, width, height, "RRect")
    return sh

def add_oval(slide, left, top, size, fill):
    sh = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(left), Inches(top), Inches(size), Inches(size))
    sh.fill.solid()
    sh.fill.fore_color.rgb = fill
    sh.line.fill.background()
    return sh

def add_text_to_shape(shape, text, size=12, bold=False, color=C.TEXT_D,
                      align=PP_ALIGN.LEFT, font="Calibri"):
    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font
    p.alignment = align
    return tf

def set_cell(cell, text, size=11, bold=False, color=C.TEXT_D,
             align=PP_ALIGN.LEFT, fill=None):
    cell.text = ""
    p = cell.text_frame.paragraphs[0]
    p.text = str(text)
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = "Calibri"
    p.alignment = align
    cell.text_frame.word_wrap = True
    cell.vertical_anchor = MSO_ANCHOR.MIDDLE
    if fill:
        cell.fill.solid()
        cell.fill.fore_color.rgb = fill

# ═══════════════════════════════════════════════════════════
# DECORATIVE / FOOTER ELEMENTS
# ═══════════════════════════════════════════════════════════

def gold_top_bar(slide):
    add_shape(slide, MSO_SHAPE.RECTANGLE, 0, 0, L.SLIDE_W, 0.035, fill=C.GOLD)

def navy_header_bg(slide, title, subtitle=None):
    """Navy header — height HEADER_H (0.94")"""
    gold_top_bar(slide)
    hdr_h = L.HEADER_H - 0.04  # 0.9"
    add_shape(slide, MSO_SHAPE.RECTANGLE, 0, 0.035, L.SLIDE_W, hdr_h, fill=C.NAVY)
    add_shape(slide, MSO_SHAPE.RECTANGLE, L.MARGIN_H, 0.12, 0.07, 0.55, fill=C.GOLD)
    add_box(slide, L.MARGIN_H + 0.22, 0.15, L.cw - 0.5, 0.48, title, 20, bold=True, color=C.WHITE)
    if subtitle:
        add_box(slide, L.MARGIN_H + 0.22, 0.63, L.cw - 0.5, 0.28, subtitle, 9, color=C.TEXT_L)

def add_footer(slide):
    h = 0.03
    add_shape(slide, MSO_SHAPE.RECTANGLE, 0, L.SLIDE_H - L.FOOTER_H, L.SLIDE_W, h, fill=C.NAVY)
    add_box(slide, L.MARGIN_H, L.SLIDE_H - L.FOOTER_H + 0.06, 4, 0.22, SOURCE_TEXT, 7, color=C.TEXT_L)
    pg_counter[0] += 1
    add_box(slide, L.SLIDE_W - 1.0, L.SLIDE_H - L.FOOTER_H + 0.06, 0.8, 0.22,
            str(pg_counter[0]), 8, color=C.TEXT_L, align=PP_ALIGN.RIGHT)

# ═══════════════════════════════════════════════════════════
# COMPLEX BUILDING BLOCKS (menggunakan LayoutFrame)
# ═══════════════════════════════════════════════════════════

def section_slide(title, subtitle=None, action_text=None):
    """Full-bleed dark section divider."""
    s = blank()
    bg = s.background; bg.fill.solid(); bg.fill.fore_color.rgb = C.NAVY
    add_oval(s, -1.5, -1.5, 5, C.NAVY_L)
    add_oval(s, -0.5, -0.5, 3.5, C.NAVY_D)
    add_oval(s, 9.5, 4, 5, C.NAVY_D)
    add_oval(s, 10.5, 3, 3, RGBColor(0x0E, 0x20, 0x3E))
    add_shape(s, MSO_SHAPE.RECTANGLE, L.MARGIN_H, 2.3, 2.5, 0.04, fill=C.GOLD)
    add_box(s, L.MARGIN_H, 2.6, L.cw, 1.6, title, 34, bold=True, color=C.WHITE)
    y_sub = 4.2
    if action_text:
        add_box(s, L.MARGIN_H, y_sub, L.cw, 0.4, action_text, 14, bold=True, color=C.GOLD)
        y_sub += 0.35
    if subtitle:
        add_box(s, L.MARGIN_H, y_sub, L.cw, 0.3, subtitle, 11, color=C.TEXT_L)
    add_shape(s, MSO_SHAPE.RECTANGLE, 0, L.SLIDE_H - 0.22, L.SLIDE_W, 0.22, fill=C.GOLD)
    pg_counter[0] += 1
    return s

def content_slide(title, subtitle=None):
    """Standard content slide with header."""
    s = blank()
    bg = s.background; bg.fill.solid(); bg.fill.fore_color.rgb = C.OFF_W
    navy_header_bg(s, title, subtitle)
    return s

def card_grid(title, cards, subtitle=None, cols=0):
    """
    Grid card dengan multi-baris. Menggunakan LayoutFrame untuk kalkulasi.
    
    Kalkulasi:
    1. calc_grid(n_items, cols) → per_row, n_rows, col_w, row_h
    2. grid_cell(i, grid) → x, y, w, h untuk tiap card
    3. card_item_layout(card_h) → posisi icon, title, items dalam card
    4. safe_item_height(text, font, box_w) → item_h anti-overflow
    """
    s = content_slide(title, subtitle)
    n = len(cards)
    grid = L.calc_grid(n, cols)
    
    for i, cd in enumerate(cards):
        cx, cy, cw, ch = L.grid_cell(i, grid)
        
        clr = cd.get('clr', C.BLUE)
        ic = cd.get('ic', '')
        t = cd.get('t', '')
        items = cd.get('items', [])
        
        # Card background
        add_rrect(s, cx, cy, cw, ch)
        # Left accent bar
        add_shape(s, MSO_SHAPE.RECTANGLE, cx, cy, 0.05, ch, fill=clr)
        
        # Layout internal card
        item_layout = L.card_item_layout(ch, has_icon=bool(ic), n_items=len(items))
        
        # Icon circle
        if ic:
            add_oval(s, cx + L.CARD_PAD, cy + item_layout['icon_y'], item_layout['icon_size'], clr)
            add_box(s, cx + L.CARD_PAD, cy + item_layout['icon_y'],
                    item_layout['icon_size'], item_layout['icon_size'],
                    ic, 13, color=C.WHITE, align=PP_ALIGN.CENTER)
            title_y = cy + item_layout['title_y']
        else:
            title_y = cy + L.CARD_PAD
        
        # Card title
        add_box(s, cx + L.CARD_PAD, title_y, cw - L.CARD_PAD * 2, 0.30,
                t, 13, bold=True, color=clr)
        
        # Items — height dihitung agar tidak overflow
        ay = cy + item_layout['item_start_y']
        box_w = cw - L.CARD_PAD * 2  # lebar textbox dalam card
        for item in items:
            item_h = L.safe_item_height(f"• {item}", 9, box_w, min_h=0.30)
            add_box(s, cx + L.CARD_PAD, ay, box_w, item_h,
                    f"• {item}", 9, color=C.TEXT_D)
            ay += item_h + L.GAP_ITEM
    
    add_footer(s)
    return s

def two_col_cards(title, left_data, right_data, subtitle=None,
                  left_color=C.BLUE, right_color=C.TEAL):
    """Two-column layout — 2 card side by side."""
    s = content_slide(title, subtitle)
    gap = L.GAP_H
    cw = L.col_width(2, gap)  # (12.133 - 0.3) / 2 = 5.9165"
    sy = L.cy
    ch = L.ch
    
    for data, clr, x in [
        (left_data, left_color, L.cx),
        (right_data, right_color, L.cx + cw + gap),
    ]:
        add_rrect(s, x, sy, cw, ch)
        add_shape(s, MSO_SHAPE.RECTANGLE, x, sy, 0.05, ch, fill=clr)
        ay = sy + 0.15
        for line in data:
            if line.startswith("$"):
                add_box(s, x + 0.2, ay, cw - 0.35, 0.30, line[1:], 14, bold=True, color=clr)
                ay += 0.32
            else:
                add_box(s, x + 0.2, ay, cw - 0.35, 0.25, f"• {line}", 11, color=C.TEXT_D)
                ay += 0.26
    add_footer(s)
    return s

def callout_slide(title, callouts, subtitle=None, note_text=None):
    """Callout cards dengan big numbers."""
    s = content_slide(title, subtitle)
    n = len(callouts)
    cw = L.col_width(n)  # proper width per kolom
    sy = L.cy
    
    for i, (num, lb, clr) in enumerate(callouts):
        cx = L.cx + i * (cw + L.GAP_H)
        add_rrect(s, cx, sy, cw, 2.2)
        add_shape(s, MSO_SHAPE.RECTANGLE, cx, sy, cw, 0.05, fill=clr)
        add_box(s, cx, sy + 0.25, cw, 0.7, num, 32, bold=True, color=clr, align=PP_ALIGN.CENTER)
        add_box(s, cx + 0.1, sy + 1.0, cw - 0.2, 0.7, lb, 11, color=C.TEXT_M,
                align=PP_ALIGN.CENTER, vAlign=MSO_ANCHOR.TOP)
    
    if note_text:
        ny = sy + 2.5
        add_rrect(s, L.cx, ny, L.cw, 2.8, fill=C.ICE)
        add_box(s, L.cx + 0.25, ny + 0.15, L.cw - 0.5, 2.5, note_text, 11, color=C.TEXT_D)
    
    add_footer(s)
    return s

def flow_slide(title, steps_data, subtitle=None, note_text=None):
    """Horizontal flow dengan arrow connectors."""
    s = content_slide(title, subtitle)
    n = len(steps_data)
    bw = L.col_width(n)  # equal width per step
    bgap = L.GAP_H
    sy = L.cy + 0.15
    
    for i, (num, step_title, desc, clr) in enumerate(steps_data):
        cx = L.cx + i * (bw + bgap)
        add_rrect(s, cx, sy, bw, 2.0)
        add_shape(s, MSO_SHAPE.RECTANGLE, cx, sy, bw, 0.05, fill=clr)
        circ_x = cx + bw / 2 - 0.25
        add_oval(s, circ_x, sy + 0.15, 0.5, clr)
        add_box(s, circ_x, sy + 0.15, 0.5, 0.5, num, 18, bold=True, color=C.WHITE,
                align=PP_ALIGN.CENTER)
        add_box(s, cx + 0.1, sy + 0.75, bw - 0.2, 0.35, step_title, 14, bold=True,
                color=clr, align=PP_ALIGN.CENTER)
        add_box(s, cx + 0.1, sy + 1.1, bw - 0.2, 0.7, desc, 10, color=C.TEXT_M,
                align=PP_ALIGN.CENTER, vAlign=MSO_ANCHOR.TOP)
        if i < n - 1:
            add_box(s, cx + bw, sy + 0.7, bgap, 0.5, "›", 24, color=C.GOLD,
                    align=PP_ALIGN.CENTER, vAlign=MSO_ANCHOR.MIDDLE)
    
    if note_text:
        ny = sy + 2.4
        add_rrect(s, L.cx, ny, L.cw, 2.8, fill=C.ICE)
        add_box(s, L.cx + 0.25, ny + 0.15, L.cw - 0.5, 2.5, note_text, 11, color=C.TEXT_D)
    
    add_footer(s)
    return s

def table_slide(title, headers, rows, subtitle=None):
    """Native table slide."""
    s = content_slide(title, subtitle)
    nc = len(headers)
    nr = len(rows) + 1
    rh = 0.42
    ts = s.shapes.add_table(nr, nc, Inches(L.cx), Inches(L.cy),
                             Inches(L.cw), Inches(rh * nr))
    tbl = ts.table
    for i in range(nc):
        tbl.columns[i].width = int(Inches(L.cw / nc))
    for j, hdr in enumerate(headers):
        set_cell(tbl.cell(0, j), hdr, bold=True, color=C.WHITE, align=PP_ALIGN.CENTER, fill=C.NAVY)
    for ri, row in enumerate(rows):
        bg_c = C.ICE if ri % 2 == 0 else C.WHITE
        for j, val in enumerate(row):
            al = PP_ALIGN.LEFT if j == 0 else PP_ALIGN.CENTER
            set_cell(tbl.cell(ri + 1, j), val, color=C.TEXT_D, align=al, fill=bg_c)
    add_footer(s)
    return s


# ═══════════════════════════════════════════════════════════
# SLIDES
# ═══════════════════════════════════════════════════════════

# ────────── COVER (1) ──────────
s = blank()
bg = s.background; bg.fill.solid(); bg.fill.fore_color.rgb = C.NAVY
add_oval(s, -1, -1.5, 4.5, C.NAVY_D)
add_oval(s, 8.5, -2, 7, C.NAVY_D)
add_oval(s, 10, 4.5, 5, C.NAVY_D)
add_oval(s, 0.5, 5.5, 2.5, C.NAVY_L)
add_shape(s, MSO_SHAPE.RECTANGLE, L.MARGIN_H, 2.6, 4, 0.04, fill=C.GOLD)
add_box(s, L.MARGIN_H, 0.5, 5, 0.35, "BERITA DAERAH", 12, bold=True, color=C.GOLD)
add_box(s, L.MARGIN_H, 0.85, 5, 0.35, "KOTA BEKASI", 14, bold=True, color=C.WHITE)
add_box(s, L.MARGIN_H, 1.8, L.cw, 0.45, "PERATURAN WALI KOTA BEKASI", 20, bold=True, color=C.WHITE)
add_box(s, L.MARGIN_H, 2.15, L.cw, 0.4, "NOMOR 51 TAHUN 2024", 15, bold=True, color=C.GOLD)
add_box(s, L.MARGIN_H, 3.1, L.cw, 2.0, "TENTANG\nPENGELOLAAN PAJAK REKLAME", 40, bold=True, color=C.WHITE)
add_rrect(s, L.MARGIN_H, 5.4, 7.5, 0.9, fill=C.NAVY_D)
add_box(s, L.MARGIN_H + 0.2, 5.45, 7, 0.4, "Pemerintah Kota Bekasi  ·  20 Desember 2024", 12, color=C.ICE)
add_box(s, L.MARGIN_H + 0.2, 5.75, 7, 0.35, "Berlaku sejak diundangkan", 10, color=C.TEXT_L)
add_shape(s, MSO_SHAPE.RECTANGLE, 0, L.SLIDE_H - 0.22, L.SLIDE_W, 0.22, fill=C.GOLD)

# ────────── DAFTAR ISI (2) ──────────
s = blank(); bg = s.background; bg.fill.solid(); bg.fill.fore_color.rgb = C.OFF_W
navy_header_bg(s, "11 Bab Kunci Mengatur Seluruh Aspek Pajak Reklame", "Perwal Bekasi No 51/2024")
toc = [
    ("1","Ketentuan Umum & Definisi",C.BLUE), ("2","Objek, Subjek & Wajib Pajak",C.TEAL),
    ("3","Masa Pajak & Tahun Pajak",C.WARM), ("4","Pendaftaran & Pendataan WP",C.NAVY_M),
    ("5","Nilai Sewa Reklame (NSR)",C.RED), ("6","Perhitungan & Tarif Pajak",C.BLUE),
    ("7","Penetapan, Tagihan & Pembayaran",C.TEAL), ("8","Pembetulan, Keberatan & Banding",C.WARM),
    ("9","Pemeriksaan, Penagihan & Penghapusan",C.NAVY_M),
    ("10","Keringanan, Kemudahan & Penghargaan",C.BLUE), ("11","Ketentuan Penutup",C.TEAL),
]
# TOC: 2 kolom × 6 baris (kebutuhan khusus — manual position)
cw_toc = L.col_width(2, 0.35)  # col width for 2 columns
sy_toc = L.cy
for i, (num, lb, clr) in enumerate(toc):
    col = 0 if i < 6 else 1
    row = i if i < 6 else i - 6
    x = L.cx + col * (cw_toc + 0.35)
    y = sy_toc + row * 0.82
    add_rrect(s, x, y, cw_toc, 0.62)
    badge = add_shape(s, MSO_SHAPE.ROUNDED_RECTANGLE, x + 0.1, y + 0.06, 0.5, 0.5, fill=clr)
    add_text_to_shape(badge, num, 14, bold=True, color=C.WHITE, align=PP_ALIGN.CENTER)
    add_box(s, x + 0.75, y + 0.06, cw_toc - 0.9, 0.5, lb, 13, color=C.TEXT_D, vAlign=MSO_ANCHOR.MIDDLE)
add_footer(s)

# ═══════════════════════════════════════════
# BAB I — KETENTUAN UMUM (3-4)
# ═══════════════════════════════════════════
section_slide("BAB I\nKETENTUAN UMUM", "Pasal 1", "7 Definisi Kunci Jadi Landasan Hukum Pajak Reklame")

# Slide 4: 7 definisi — 4+3 grid, item height auto-calculated
card_grid("7 Definisi Kunci Menjadi Landasan Pengelolaan Pajak Reklame", [
    {"ic":"🏛️","t":"Daerah","clr":C.BLUE,"items":["Kota Bekasi"]},
    {"ic":"📊","t":"Bapenda","clr":C.TEAL,"items":["Badan Pendapatan Daerah Kota Bekasi"]},
    {"ic":"📢","t":"Reklame","clr":C.WARM,"items":["Media untuk promosi & pengenalan komersial"]},
    {"ic":"💰","t":"Pajak Reklame","clr":C.NAVY_M,"items":["Pajak atas penyelenggaraan reklame"]},
    {"ic":"📐","t":"NSR","clr":C.BLUE,"items":["Nilai Sewa Reklame — dasar pengenaan pajak"]},
    {"ic":"🆔","t":"NPWPD","clr":C.TEAL,"items":["Nomor Pokok Wajib Pajak Daerah"]},
    {"ic":"👤","t":"Wajib Pajak","clr":C.WARM,"items":["Orang pribadi/badan dg hak & kewajiban pajak"]},
], subtitle="Pasal 1")

# ═══════════════════════════════════════════
# BAB II — OBJEK, SUBJEK & WAJIB PAJAK (5-7)
# ═══════════════════════════════════════════
section_slide("BAB II\nOBJEK, SUBJEK & WAJIB PAJAK", "Pasal 2–4",
              "10 Jenis Reklame Wajib Pajak + 8 Pengecualian")

card_grid("10 Jenis Reklame Wajib Pajak + 8 Jenis Dikecualikan", [
    {"ic":"📋","t":"10 Jenis Reklame","clr":C.BLUE,"items":["Papan / Billboard","Videotron / Megatron","Kain (Spanduk, Umbul, Baliho)","Melekat / Stiker","Selebaran","Berjalan (Kendaraan)","Udara (Balon Gas)","Apung","Film / Slide","Peragaan"]},
    {"ic":"🚫","t":"Dikecualikan","clr":C.TEAL,"items":["Internet, TV, radio, media cetak","Label / merek produk","Nama usaha ≤ 1 m² di tempat","Reklame Pemerintah/Pemda","Tempat ibadah & panti asuhan","Sosial & keagamaan ≤ 30 hari","Kegiatan politik (masa kampanye)","Olahraga KONI ≤ 30 hari"]},
], subtitle="Pasal 2–4")

card_grid("Subjek & Wajib Pajak: Siapa yang Terkena Kewajiban Pajak", [
    {"ic":"👤","t":"Subjek Pajak (Pasal 3)","clr":C.BLUE,"items":["Orang pribadi atau Badan","yang menggunakan Reklame"]},
    {"ic":"✋","t":"Wajib Pajak (Pasal 4)","clr":C.TEAL,"items":["Orang pribadi atau Badan","yang menyelenggarakan Reklame","Jika pihak ketiga → menjadi WP","Mendaftarkan diri & objek pajak"]},
], subtitle="Pasal 3–4")

# ═══════════════════════════════════════════
# BAB III — MASA PAJAK (8-9)
# ═══════════════════════════════════════════
section_slide("BAB III\nMASA PAJAK & TAHUN PAJAK", "Pasal 5",
              "12 Bulan atau 30 Hari — Tergantung Jenis Reklame")

callout_slide("Masa Pajak: 12 Bulan Permanen, 30 Hari Insidentil", [
    ("12", "Bulan\n(Permanen)", C.BLUE),
    ("30", "Hari\n(Insidentil)", C.TEAL),
    ("1", "Tahun Pajak\n(Kalender)", C.WARM),
    ("1", "Bulan\n(Bagian Tahun)", C.NAVY_M),
], subtitle="Pasal 5",
   note_text="• Masa Pajak Permanen: 12 bulan atau sesuai jangka waktu penayangan reklame\n"
             "• Masa Pajak Insidentil: dihitung per hari, maksimal 30 hari\n"
             "• Tahun Pajak: 1 tahun kalender atau sesuai tahun buku wajib pajak\n"
             "• Bagian Tahun Pajak: 1 bulan penuh (jika tidak mencakup satu tahun penuh)")

# ═══════════════════════════════════════════
# BAB IV — PENDAFTARAN & PENDATAAN (10-11)
# ═══════════════════════════════════════════
section_slide("BAB IV\nPENDAFTARAN & PENDATAAN WP", "Pasal 6–8",
              "Wajib Daftar atau NPWPD Jabatan — Bapenda Berwenang Mendata")

two_col_cards("Pendaftaran Wajib, Pendataan Bapenda, dan Penonaktifan WP", [
    "$PENDAFTARAN (Pasal 6)","","WP wajib mendaftarkan diri & objek pajak",
    "Formulir: ambil/online/dikirim petugas","Lampirkan: KTP, NPWP, Akta, NIB",
    "Bapenda terbitkan NPWPD","Jika tidak mendaftar → NPWPD jabatan","Juga: NOPD & nomor registrasi",
], [
    "$PENDATAAN & NONAKTIF (Pasal 7–8)","","Bapenda mendata WP & objek pajak",
    "Termasuk data geografis","Dapat kerjasama dengan instansi lain",
    "Penonaktifan: WP tak penuhi syarat","Keputusan maksimal 3 bulan","Syarat: tanpa tunggakan & keberatan",
], subtitle="Pasal 6–8")

# ═══════════════════════════════════════════
# BAB V — NILAI SEWA REKLAME (12-13)
# ═══════════════════════════════════════════
section_slide("BAB V\nNILAI SEWA REKLAME", "Pasal 9",
              "7 Faktor Penentu NSR — Jenis, Bahan, Lokasi, Waktu, dll")

s = blank(); bg = s.background; bg.fill.solid(); bg.fill.fore_color.rgb = C.OFF_W
navy_header_bg(s, "7 Faktor Penentu Nilai Sewa Reklame (NSR)", "Pasal 9")
factors = [
    ("1","Jenis Reklame",C.BLUE), ("2","Bahan",C.TEAL), ("3","Lokasi (Kelas Jalan)",C.WARM),
    ("4","Waktu Tayang (detik)",C.NAVY_M), ("5","Jangka Waktu (hari)",C.BLUE),
    ("6","Jumlah Media",C.TEAL), ("7","Ukuran (m²)",C.WARM),
]
# Manual 4+3 grid (factors adalah special case)
per_row = 4
n = len(factors)
fw = L.col_width(per_row, 0.3)
for i, (num, lb, clr) in enumerate(factors):
    col = i % per_row; row = i // per_row
    x, y = L.grid_pos(col, row, fw, 1.5, gap_v=0.3)
    add_rrect(s, x, y, fw, 1.3)
    add_oval(s, x + 0.15, y + 0.25, 0.55, clr)
    add_box(s, x + 0.15, y + 0.25, 0.55, 0.55, num, 18, bold=True, color=C.WHITE, align=PP_ALIGN.CENTER)
    add_box(s, x + 0.8, y + 0.2, fw - 1, 0.9, lb, 13, color=C.TEXT_D, vAlign=MSO_ANCHOR.MIDDLE)
# Klasifikasi jalan
box = add_rrect(s, L.cx, L.cy + 1.5 + 0.3, L.cw, 1.7, fill=C.ICE)
add_box(s, L.cx + 0.2, L.cy + 1.5 + 0.35, 5, 0.3, "KLASIFIKASI KELAS JALAN", 11, bold=True, color=C.NAVY)
add_box(s, L.cx + 0.2, L.cy + 1.5 + 0.7, L.cw - 0.4, 0.9,
        "🏛️  Kelas Jalan Khusus — Tol | Premium 1 | Premium 2\n"
        "🚗  Kelas Jalan I (Kendali Ketat) — Lebar > 3 m, pusat pelayanan\n"
        "🏡  Kelas Jalan II (Kendali Sedang) — Lebar ≤ 3 m, jalan lingkungan",
        10, color=C.TEXT_D)
add_footer(s)

# ═══════════════════════════════════════════
# BAB VI — PERHITUNGAN & TARIF (14-20)
# ═══════════════════════════════════════════
section_slide("BAB VI\nPERHITUNGAN & TARIF PAJAK", "Pasal 10",
              "Rumus: Pajak = Tarif × NSR — Dengan Berbagai Ketentuan Khusus")

callout_slide("Rumus: Pajak Reklame = Tarif (50%) × NSR", [
    ("×","Pajak =\nTarif × NSR",C.BLUE),
    ("50%","Indoor =\n50% NSR",C.TEAL),
    ("+20%","Tinggi > 15m\ntambahan 20%",C.WARM),
    ("+50%","Tembakau &\nMiras +50%",C.RED),
], subtitle="Pasal 10",
   note_text="Rumus: Pajak Reklame = Tarif Pajak × NSR\n"
             "NSR = Nilai Kelas Jalan × Ukuran (m²) × Jumlah × Jangka Waktu\n\n"
             "Ketentuan Khusus:\n"
             "• Reklame indoor: NSR 50% dari NSR normal\n"
             "• Ketinggian > 15 meter: tambahan 20%\n"
             "• Produk tembakau & minuman beralkohol: tambahan 50%\n"
             "• Perubahan naskah/revisi isi reklame: dikecualikan")

# 3 Tabel NSR
table_slide("Tabel NSR — Papan / Billboard (Rp/m²/hari)", [
    "Kelas Jalan", "Zona", "NSR (Rp/m²/hari)"
], [
    ["Kelas Jalan Khusus", "Jalan Tol", "23.575"],
    ["Kelas Jalan Khusus", "Premium 1", "16.100"],
    ["Kelas Jalan Khusus", "Premium 2", "14.950"],
    ["Kelas Jalan I", "Kendali Ketat", "13.225"],
    ["Kelas Jalan II", "Kendali Sedang", "11.500"],
], subtitle="Pasal 10")

table_slide("Tabel NSR — Megatron / Videotron", [
    "Kelas Jalan", "Zona", "NSR (/30 dtk)", "NSR (/m²/thn)"
], [
    ["Kelas Jalan Khusus", "Jalan Tol", "Rp 17,25", "13.599.900"],
    ["Kelas Jalan Khusus", "Premium 1", "Rp 13,80", "10.879.920"],
    ["Kelas Jalan Khusus", "Premium 2", "Rp 9,20", "7.253.280"],
    ["Kelas Jalan I", "Kendali Ketat", "Rp 8,05", "6.346.620"],
    ["Kelas Jalan II", "Kendali Sedang", "Rp 5,75", "4.533.300"],
], subtitle="Pasal 10")

table_slide("Tabel NSR — Kain (Spanduk/Umbul/Baliho) (Rp/m²/hari)", [
    "Kelas Jalan", "Zona", "NSR (Rp/m²/hari)"
], [
    ["Kelas Jalan Khusus", "Jalan Tol", "30.000"],
    ["Kelas Jalan Khusus", "Premium 1", "30.000"],
    ["Kelas Jalan Khusus", "Premium 2", "25.000"],
    ["Kelas Jalan I", "Kendali Ketat", "20.000"],
    ["Kelas Jalan II", "Kendali Sedang", "19.000"],
], subtitle="Pasal 10")

card_grid("NSR untuk 8 Jenis Reklame Lainnya — Bagian 1", [
    {"ic":"🏷️","t":"Stiker","clr":C.BLUE,"items":["Rp 7,5/cm²","Min. Rp 750.000/kali"]},
    {"ic":"🧱","t":"Melekat","clr":C.TEAL,"items":["Rp 750.000/m²/tahun"]},
    {"ic":"📄","t":"Selebaran","clr":C.WARM,"items":["Rp 600/lembar","Min. Rp 6.000.000/kali"]},
    {"ic":"🚌","t":"Berjalan","clr":C.RED,"items":["Rp 6.000/m²/hari","Termasuk kendaraan"]},
], subtitle="Pasal 10")

card_grid("NSR untuk 8 Jenis Reklame Lainnya — Bagian 2", [
    {"ic":"🎈","t":"Udara","clr":C.BLUE,"items":["Rp 2.400.000/sekali","Maks. 1 bulan"]},
    {"ic":"🌊","t":"Apung","clr":C.TEAL,"items":["Rp 600.000/sekali","Maks. 1 bulan"]},
    {"ic":"🎬","t":"Film / Slide","clr":C.WARM,"items":["Rp 12.000/15 detik"]},
    {"ic":"🎭","t":"Peragaan","clr":C.RED,"items":["Rp 480.000/penyelenggaraan"]},
], subtitle="Pasal 10")

# ═══════════════════════════════════════════
# BAB VII — PENETAPAN, TAGIHAN & PEMBAYARAN (21-22)
# ═══════════════════════════════════════════
section_slide("BAB VII\nPENETAPAN, TAGIHAN & PEMBAYARAN", "Pasal 11–14",
              "4 Langkah: SKPD → Bayar → Telat → STPD")

flow_slide("Alur Penetapan hingga Pembayaran: 4 Langkah Wajib Dipahami WP", [
    ("1","SKPD","Diterbitkan Bapenda\nMasa berlaku 5 tahun",C.BLUE),
    ("2","Pembayaran","Lunas 1 bulan\nsejak SKPD diterima",C.TEAL),
    ("3","Keterlambatan","Bunga 1%/bln\nDiterbitkan STPD",C.WARM),
    ("4","STPD","Harus lunas\n≤ 30 hari",C.RED),
], subtitle="Pasal 11–14",
   note_text="• Jatuh tempo: 1 bulan sejak tanggal pengiriman SKPD\n"
             "• Pembayaran: Kas Daerah / Bank Persepsi / tempat lain yang ditunjuk\n"
             "• Stiker sebagai tanda bukti pembayaran reklame\n"
             "• STPD dikenakan bunga 1%/bulan (maks. 24 bulan)")

# ═══════════════════════════════════════════
# BAB VIII — PEMBETULAN, KEBERATAN & BANDING (23-25)
# ═══════════════════════════════════════════
section_slide("BAB VIII\nPEMBETULAN, KEBERATAN & BANDING", "Pasal 15–20, 29–33",
              "Hak WP: Koreksi, Keberatan, Banding — Dengan Batas Waktu Jelas")

two_col_cards("Pembetulan: Koreksi Kesalahan Tulis, Hitung, dan Penerapan Aturan", [
    "$PEMBETULAN (Pasal 15–20)","","Kesalahan tulis: nama, alamat, NPWPD",
    "Kesalahan hitung: jumlah, tarif","Kekeliruan penerapan aturan",
    "1 permohonan = 1 ketetapan","Keputusan maksimal 6 bulan",
    "> 6 bulan tanpa putusan → dikabulkan","Dapat dilakukan berulang (Ps 20)",
    "Jenis keputusan: kabul / batal / tolak",
], [
    "$JANGKA WAKTU & PROSEDUR","","Permohonan diajukan ke Bapenda",
    "Keputusan: kabul (tambah/kurang/hapus)","Keputusan: batal | tolak",
    "Pasal 19: pembetulan jabatan","Pasal 20: berulang jika masih salah",
], subtitle="Pasal 15–20", left_color=C.BLUE, right_color=C.WARM)

two_col_cards("Keberatan & Banding: Upaya Hukum WP dalam Sengketa Pajak", [
    "$KEBERATAN (Pasal 29–31)","","Objek: SKPD, SKPDKB, SKPDKBT, dll",
    "Diajukan maks. 3 bulan sejak SKPD","Sudah bayar min. yang disetujui",
    "Keputusan maks. 12 bulan","Jika ditolak: denda 30%","Jika dikabulkan: + bunga 0,6%/bulan",
], [
    "$BANDING (Pasal 32–33)","","Objek: Surat Keputusan Keberatan",
    "Ke badan peradilan pajak","Maks. 3 bulan sejak keputusan",
    "Menangguhkan kewajiban bayar","Jika ditolak: denda 60%","Jika dikabulkan: + bunga 0,6%/bulan",
], subtitle="Pasal 29–33")

# ═══════════════════════════════════════════
# BAB IX — PEMERIKSAAN, PENAGIHAN & PENGHAPUSAN (26-28)
# ═══════════════════════════════════════════
section_slide("BAB IX\nPEMERIKSAAN, PENAGIHAN & PENGHAPUSAN", "Pasal 22–26",
              "Bapenda Berwenang Periksa — Piutang Dihapus Lewat 4 Tahapan")

card_grid("Pemeriksaan & Penagihan: 3 Pilar Penegakan Kepatuhan WP", [
    {"ic":"🔍","t":"Pemeriksaan (Ps 22–23)","clr":C.BLUE,"items":["Kepala Bapenda berwenang periksa","Menguji kepatuhan WP","WP wajib: buka buku/dokumen","Beri akses tempat & keterangan","Jika tidak → pajak ditetapkan jabatan"]},
    {"ic":"📬","t":"Penagihan (Ps 24)","clr":C.TEAL,"items":["Dasar: SKPD, SKPDKB, SKPDKBT","STPD, SK Pembetulan/Keberatan","Putusan Banding"]},
    {"ic":"⏳","t":"Kedaluwarsa (Ps 25)","clr":C.WARM,"items":["5 tahun sejak pajak terutang","Tertangguh jika ada: Surat Teguran / Paksa","Pengakuan utang dari WP"]},
], subtitle="Pasal 22–25")

flow_slide("Penghapusan Piutang Pajak: 4 Langkah dari Penelitian hingga SK", [
    ("1","Penelitian","Dilakukan Bapenda",C.BLUE),
    ("2","Penetapan","Keputusan Wali Kota",C.TEAL),
    ("3","Koordinasi","Dengan Inspektorat",C.WARM),
    ("4","SK Penghapusan","Diterbitkan",C.NAVY_M),
], subtitle="Pasal 26",
   note_text="Syarat penghapusan piutang pajak:\n"
             "• Piutang tidak mungkin ditagih lagi karena kedaluwarsa\n"
             "• Ada koordinasi dengan Inspektorat Daerah\n"
             "• Dibuktikan dengan dokumen pelaksanaan penagihan")

# ═══════════════════════════════════════════
# BAB X — KERINGANAN, KEMUDAHAN & PENGHARGAAN (29-30)
# ═══════════════════════════════════════════
section_slide("BAB X\nKERINGANAN, KEMUDAHAN & PENGHARGAAN", "Pasal 27–28, 34–35",
              "3 Fasilitas WP: Keringanan, Kemudahan Angsuran, dan Penghargaan")

card_grid("3 Fasilitas untuk WP: Keringanan, Kemudahan Angsuran, Penghargaan", [
    {"ic":"🎯","t":"Keringanan (Ps 27)","clr":C.BLUE,"items":["Keringanan / Pengurangan","Pembebasan / Penundaan","Atas pokok & sanksi pajak","WP dengan likuiditas rendah","Objek terdampak bencana/kebakaran"]},
    {"ic":"🤝","t":"Kemudahan (Ps 28)","clr":C.TEAL,"items":["Perpanjangan waktu bayar","Angsuran maks. 24 bulan","Bunga 0,6%/bulan","Keadaan kahar: bencana, wabah, kerusuhan"]},
    {"ic":"🏆","t":"Penghargaan (Ps 34–35)","clr":C.WARM,"items":["WP Taat Pajak","Bayar tepat waktu ≥ 1 tahun","Tanpa tunggakan 3 tahun","Kontribusi signifikan","Piagam / Hadiah (APBD)"]},
], subtitle="Pasal 27–28, 34–35")

# ═══════════════════════════════════════════
# BAB XI — KETENTUAN PENUTUP (31-32)
# ═══════════════════════════════════════════
section_slide("BAB XI\nKETENTUAN PENUTUP", "Pasal 36–37",
              "Perwal No 48/2012 Dicabut — Berlaku Sejak 20 Desember 2024")

two_col_cards("Perwal Lama Dicabut, Peraturan Baru Berlaku Mulai Diundangkan", [
    "$PERATURAN YANG DICABUT (Pasal 36)","","Perwal No. 48 Tahun 2012",
    "Petunjuk Pelaksanaan Perda 14/2012","Perwal No. 52 Tahun 2013 (Perubahan)",
], [
    "$MULAI BERLAKU (Pasal 37)","","Sejak diundangkan","20 Desember 2024","",
    "Pj. WALI KOTA BEKASI,","ttd.","R. GANI MUHAMAD",
], subtitle="Pasal 36–37", left_color=C.BLUE, right_color=C.NAVY_M)

# ────────── CLOSING (33) ──────────
s = blank()
bg = s.background; bg.fill.solid(); bg.fill.fore_color.rgb = C.NAVY
add_oval(s, -1, -1.5, 4.5, C.NAVY_D)
add_oval(s, 8.5, -2, 7, C.NAVY_D)
add_oval(s, 10, 4.5, 5, C.NAVY_D)
add_shape(s, MSO_SHAPE.RECTANGLE, L.MARGIN_H, 3.6, 3.5, 0.04, fill=C.GOLD)
add_box(s, L.MARGIN_H, 1.6, L.cw, 0.4, "BERITA DAERAH KOTA BEKASI", 14, bold=True, color=C.GOLD)
add_box(s, L.MARGIN_H, 2.4, L.cw, 1.5, "TERIMA KASIH", 48, bold=True, color=C.WHITE)
add_box(s, L.MARGIN_H, 4.1, L.cw, 0.8,
        "Peraturan Wali Kota Bekasi Nomor 51 Tahun 2024\nTentang Pengelolaan Pajak Reklame", 14, color=C.ICE)
add_box(s, L.MARGIN_H, 5.1, L.cw, 0.35, "Sumber: https://jdih.bekasikota.go.id", 10, color=C.TEXT_L)
add_shape(s, MSO_SHAPE.RECTANGLE, 0, L.SLIDE_H - 0.22, L.SLIDE_W, 0.22, fill=C.GOLD)

# ══════════════════════════════════════════════
# SAVE
# ══════════════════════════════════════════════
out = os.path.join(os.path.dirname(__file__), "Perwal_Bekasi_51_2024_Pajak_Reklame.pptx")
prs.save(out)
print(f"✅ OK: {out} ({len(prs.slides)} slide)")
