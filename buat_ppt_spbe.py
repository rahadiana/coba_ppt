#!/usr/bin/env python3
"""
Penerapan SPBE — Sistem Pemerintahan Berbasis Elektronik
VERSI IMPROVED — Action Titles, Source Citations, Visual Hierarchy, 7 Archetypes

Berdasarkan riset:
- McKinsey/BCG/Bain: action titles, pyramid principle, one message per slide
- Color Theory 2026: 60-30-10 rule, 4.5:1 contrast, max 5 colors
- Government design: trust, transparency, accessibility (WCAG 2.1)
- Presentation Zen: whitespace, signal vs noise, picture superiority

Arsitektur diadaptasi dari buat_ppt_improved.py (Perwal Bekasi).
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# ═══════════════════════════════════════════════════════════
# CONSTANTS
# ═══════════════════════════════════════════════════════════

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)
M = Inches(0.6)
CW = SLIDE_W - 2 * M

# ─── Color Palette (60-30-10) ───
# 60% Dominant
NAVY     = RGBColor(0x0A, 0x16, 0x28)
NAVY_L   = RGBColor(0x12, 0x29, 0x4A)
NAVY_M   = RGBColor(0x1B, 0x3A, 0x6B)

# 30% Secondary
WHITE    = RGBColor(0xFF, 0xFF, 0xFF)
OFF_W    = RGBColor(0xF5, 0xF7, 0xFA)
ICE      = RGBColor(0xE8, 0xED, 0xF5)
ICE_D    = RGBColor(0xD0, 0xD8, 0xE8)

# 10% Accent
GOLD     = RGBColor(0xC8, 0x96, 0x2E)
GOLD_L   = RGBColor(0xD4, 0xA0, 0x17)
GOLD_M   = RGBColor(0xB8, 0x86, 0x0B)

# Semantic
TEXT_D   = RGBColor(0x1A, 0x1A, 0x2E)
TEXT_M   = RGBColor(0x6B, 0x72, 0x88)
TEXT_L   = RGBColor(0x9C, 0xA3, 0xAF)
BLUE     = RGBColor(0x25, 0x63, 0xEB)
TEAL     = RGBColor(0x0D, 0x94, 0x88)
WARM     = RGBColor(0xB8, 0x86, 0x0B)
RED_C    = RGBColor(0xDC, 0x26, 0x26)
GREEN    = RGBColor(0x2E, 0x7D, 0x32)
ORANGE   = RGBColor(0xE6, 0x51, 0x00)

# ─── Presentation Setup ───
prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H
pg_counter = [0]

SOURCE_TEXT = "Sumber: Perpres No. 95/2018 • Perpres No. 132/2022 • Pedoman SPBE Kemenpan-RB"

# ═══════════════════════════════════════════════════════════
# HELPERS
# ═══════════════════════════════════════════════════════════

def blank():
    return prs.slides.add_slide(prs.slide_layouts[6])

def add_box(slide, left, top, width, height, text, size=12, bold=False,
            color=TEXT_D, align=PP_ALIGN.LEFT, font="Calibri", vAlign=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font
    p.alignment = align
    p.space_after = Pt(0)
    p.space_before = Pt(0)
    return tb

def add_shape(slide, shape_type, left, top, width, height,
              fill=None, line=None, lw=None):
    sh = slide.shapes.add_shape(shape_type, left, top, width, height)
    if fill:
        sh.fill.solid()
        sh.fill.fore_color.rgb = fill
    else:
        sh.fill.background()
    if line:
        sh.line.color.rgb = line
        if lw:
            sh.line.width = Pt(lw)
    else:
        sh.line.fill.background()
    return sh

def add_rrect(slide, left, top, width, height,
              fill=WHITE, line=ICE, lw=0.5, radius=0.04):
    sh = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    sh.fill.solid()
    sh.fill.fore_color.rgb = fill
    if line:
        sh.line.color.rgb = line
        sh.line.width = Pt(lw)
    else:
        sh.line.fill.background()
    sh.adjustments[0] = radius
    return sh

def add_oval(slide, left, top, size, fill):
    sh = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top, size, size)
    sh.fill.solid()
    sh.fill.fore_color.rgb = fill
    sh.line.fill.background()
    return sh

def add_text_to_shape(shape, text, size=12, bold=False, color=TEXT_D,
                      align=PP_ALIGN.LEFT, font="Calibri"):
    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font
    p.alignment = align
    return tf

def set_cell(cell, text, size=11, bold=False, color=TEXT_D,
             align=PP_ALIGN.LEFT, fill=None):
    cell.text = ""
    p = cell.text_frame.paragraphs[0]
    p.text = str(text)
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = "Calibri"
    p.alignment = align
    cell.text_frame.word_wrap = True
    cell.vertical_anchor = MSO_ANCHOR.MIDDLE
    if fill:
        cell.fill.solid()
        cell.fill.fore_color.rgb = fill

# ═══════════════════════════════════════════════════════════
# DECORATIVE / FOOTER ELEMENTS
# ═══════════════════════════════════════════════════════════

def gold_top_bar(slide):
    add_shape(slide, MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, Inches(0.035), fill=GOLD)

def navy_header_bg(slide, title, action_title=None, subtitle=None):
    gold_top_bar(slide)
    add_shape(slide, MSO_SHAPE.RECTANGLE, 0, Inches(0.035), SLIDE_W, Inches(0.9), fill=NAVY)
    add_shape(slide, MSO_SHAPE.RECTANGLE, M, Inches(0.12), Inches(0.07), Inches(0.6), fill=GOLD)
    y_title = Inches(0.15)
    add_box(slide, M + Inches(0.22), y_title, CW - Inches(0.5), Inches(0.48),
            title, 20, bold=True, color=WHITE)
    y_sub = y_title + Inches(0.48)
    parts = []
    if action_title:
        parts.append(action_title)
    if subtitle:
        parts.append(subtitle)
    combined = " • ".join(parts)
    add_box(slide, M + Inches(0.22), y_sub, CW - Inches(0.5), Inches(0.3),
            combined, 9, color=TEXT_L)

def navy_bottom_bar(slide):
    add_shape(slide, MSO_SHAPE.RECTANGLE, 0, SLIDE_H - Inches(0.25), SLIDE_W, Inches(0.03), fill=NAVY)

def gold_bottom_bar(slide):
    add_shape(slide, MSO_SHAPE.RECTANGLE, 0, SLIDE_H - Inches(0.25), SLIDE_W, Inches(0.25), fill=GOLD)

def add_page_num(slide):
    pg_counter[0] += 1
    add_box(slide, SLIDE_W - Inches(1.0), SLIDE_H - Inches(0.42), Inches(0.8), Inches(0.22),
            str(pg_counter[0]), 8, color=TEXT_L, align=PP_ALIGN.RIGHT)

def add_source(slide):
    add_box(slide, M, SLIDE_H - Inches(0.42), Inches(6), Inches(0.22),
            SOURCE_TEXT, 7, color=TEXT_L)

def add_footer(slide):
    navy_bottom_bar(slide)
    add_source(slide)
    add_page_num(slide)

# ═══════════════════════════════════════════════════════════
# ARCHETYPES
# ═══════════════════════════════════════════════════════════

def section_slide(title, subtitle=None, action_text=None):
    s = blank()
    bg = s.background; bg.fill.solid(); bg.fill.fore_color.rgb = NAVY
    add_oval(s, Inches(-1.5), Inches(-1.5), Inches(5), NAVY_L)
    add_oval(s, Inches(-0.5), Inches(-0.5), Inches(3.5), RGBColor(0x0D, 0x1F, 0x3C))
    add_oval(s, Inches(9.5), Inches(4), Inches(5), RGBColor(0x0D, 0x1F, 0x3C))
    add_oval(s, Inches(10.5), Inches(3), Inches(3), RGBColor(0x0E, 0x20, 0x3E))
    add_shape(s, MSO_SHAPE.RECTANGLE, M, Inches(2.3), Inches(2.5), Inches(0.04), fill=GOLD)
    add_box(s, M, Inches(2.6), CW, Inches(1.6), title, 34, bold=True, color=WHITE)
    y_sub = Inches(4.2)
    if action_text:
        add_box(s, M, y_sub, CW, Inches(0.4), action_text, 14, bold=True, color=GOLD)
        y_sub += Inches(0.35)
    if subtitle:
        add_box(s, M, y_sub, CW, Inches(0.3), subtitle, 11, color=TEXT_L)
    add_shape(s, MSO_SHAPE.RECTANGLE, 0, SLIDE_H - Inches(0.22), SLIDE_W, Inches(0.22), fill=GOLD)
    add_page_num(s)
    return s

def content_slide(action_title, subtitle=None, ref=None):
    s = blank()
    bg = s.background; bg.fill.solid(); bg.fill.fore_color.rgb = OFF_W
    parts = []
    if subtitle:
        parts.append(subtitle)
    if ref:
        parts.append(ref)
    sub_text = " • ".join(parts) if parts else None
    navy_header_bg(s, action_title, sub_text)
    return s

def card_grid(action_title, cards, subtitle=None, ref=None, cols=0):
    s = content_slide(action_title, subtitle, ref)
    n = len(cards)
    per_row = cols if cols > 0 else (2 if n <= 2 else (3 if n <= 3 else 4))
    gap = Inches(0.3)
    cw = (SLIDE_W - M * 2 - (per_row - 1) * gap) / per_row
    gap_h = Inches(0.3)
    n_rows = (n + per_row - 1) // per_row
    content_h = SLIDE_H - Inches(1.15) - Inches(0.5)
    ch = (content_h - (n_rows - 1) * gap_h) / n_rows

    max_items = max(len(cd.get('items', [])) for cd in cards) if cards else 1
    item_h = Inches(0.35)
    title_h = Inches(0.3)

    for i, cd in enumerate(cards):
        row = i // per_row
        col = i % per_row
        cx = M + col * (cw + gap)
        cy = Inches(1.15) + row * (ch + gap_h)

        clr = cd.get('clr', BLUE)
        ic = cd.get('ic', '')
        t = cd.get('t', '')
        items = cd.get('items', [])

        add_rrect(s, cx, cy, cw, ch)
        add_shape(s, MSO_SHAPE.RECTANGLE, cx, cy, Inches(0.05), ch, fill=clr)
        if ic:
            add_oval(s, cx + Inches(0.15), cy + Inches(0.15), Inches(0.42), clr)
            add_box(s, cx + Inches(0.15), cy + Inches(0.15), Inches(0.42), Inches(0.42),
                    ic, 13, color=WHITE, align=PP_ALIGN.CENTER)
            ty = cy + Inches(0.7)
        else:
            ty = cy + Inches(0.15)
        add_box(s, cx + Inches(0.15), ty, cw - Inches(0.3), title_h,
                t, 13, bold=True, color=clr)
        ay = ty + title_h + Inches(0.08)
        for item in items:
            add_box(s, cx + Inches(0.15), ay, cw - Inches(0.3), item_h,
                    f"• {item}", 9, color=TEXT_D)
            ay += item_h + Inches(0.02)

    add_footer(s)
    return s

def two_col_cards(action_title, left_data, right_data, subtitle=None, ref=None,
                  left_color=BLUE, right_color=TEAL):
    s = content_slide(action_title, subtitle, ref)
    gap = Inches(0.3)
    cw = (SLIDE_W - gap - 2 * M) / 2
    sy = Inches(1.15)
    ch = SLIDE_H - sy - Inches(0.7)

    for data, clr, x in [
        (left_data, left_color, M),
        (right_data, right_color, M + cw + gap),
    ]:
        add_rrect(s, x, sy, cw, ch)
        add_shape(s, MSO_SHAPE.RECTANGLE, x, sy, Inches(0.05), ch, fill=clr)
        ay = sy + Inches(0.15)
        for line in data:
            if line.startswith("$"):
                add_box(s, x + Inches(0.2), ay, cw - Inches(0.35), Inches(0.3),
                        line[1:], 14, bold=True, color=clr)
                ay += Inches(0.32)
            else:
                add_box(s, x + Inches(0.2), ay, cw - Inches(0.35), Inches(0.22),
                        f"• {line}", 11, color=TEXT_D)
                ay += Inches(0.24)
    add_footer(s)
    return s

def callout_slide(action_title, callouts, subtitle=None, ref=None, note_text=None):
    s = content_slide(action_title, subtitle, ref)
    n = len(callouts)
    gap = Inches(0.3)
    cw = (SLIDE_W - 2 * M - (n - 1) * gap) / n
    sx = M
    sy = Inches(1.15)

    for i, (num, lb, clr) in enumerate(callouts):
        cx = sx + i * (cw + gap)
        add_rrect(s, cx, sy, cw, Inches(2.2))
        add_shape(s, MSO_SHAPE.RECTANGLE, cx, sy, cw, Inches(0.05), fill=clr)
        add_box(s, cx, sy + Inches(0.25), cw, Inches(0.7),
                num, 32, bold=True, color=clr, align=PP_ALIGN.CENTER)
        add_box(s, cx + Inches(0.1), sy + Inches(1.0), cw - Inches(0.2), Inches(0.7),
                lb, 11, color=TEXT_M, align=PP_ALIGN.CENTER, vAlign=MSO_ANCHOR.TOP)

    if note_text:
        ny = sy + Inches(2.5)
        add_rrect(s, M, ny, CW, Inches(3.0), fill=ICE)
        add_box(s, M + Inches(0.25), ny + Inches(0.15), CW - Inches(0.5), Inches(2.6),
                note_text, 11, color=TEXT_D)

    add_footer(s)
    return s

def flow_slide(action_title, steps_data, subtitle=None, ref=None, note_text=None):
    s = content_slide(action_title, subtitle, ref)
    n = len(steps_data)
    bgap = Inches(0.35)
    bw = (SLIDE_W - 2 * M - (n - 1) * bgap) / n
    ssx = M
    sy = Inches(1.3)

    for i, (num, title, desc, clr) in enumerate(steps_data):
        cx = ssx + i * (bw + bgap)
        add_rrect(s, cx, sy, bw, Inches(2.0))
        add_shape(s, MSO_SHAPE.RECTANGLE, cx, sy, bw, Inches(0.05), fill=clr)
        circ_x = cx + bw / 2 - Inches(0.25)
        add_oval(s, circ_x, sy + Inches(0.15), Inches(0.5), clr)
        add_box(s, circ_x, sy + Inches(0.15), Inches(0.5), Inches(0.5),
                num, 18, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_box(s, cx + Inches(0.1), sy + Inches(0.75), bw - Inches(0.2), Inches(0.35),
                title, 14, bold=True, color=clr, align=PP_ALIGN.CENTER)
        add_box(s, cx + Inches(0.1), sy + Inches(1.1), bw - Inches(0.2), Inches(0.7),
                desc, 10, color=TEXT_M, align=PP_ALIGN.CENTER, vAlign=MSO_ANCHOR.TOP)
        if i < n - 1:
            add_box(s, cx + bw, sy + Inches(0.7), bgap, Inches(0.5),
                    "›", 24, color=GOLD, align=PP_ALIGN.CENTER, vAlign=MSO_ANCHOR.MIDDLE)

    if note_text:
        ny = sy + Inches(2.4)
        add_rrect(s, M, ny, CW, Inches(3.0), fill=ICE)
        add_box(s, M + Inches(0.25), ny + Inches(0.15), CW - Inches(0.5), Inches(2.6),
                note_text, 11, color=TEXT_D)

    add_footer(s)
    return s

def table_slide(action_title, headers, rows, subtitle=None, ref=None):
    s = content_slide(action_title, subtitle, ref)
    tbl_left = M
    tbl_top = Inches(1.15)
    tbl_width = CW
    rh = Inches(0.42)
    nr = len(rows) + 1
    nc = len(headers)
    ts = s.shapes.add_table(nr, nc, tbl_left, tbl_top, tbl_width, rh * nr)
    tbl = ts.table
    for i in range(nc):
        tbl.columns[i].width = int(tbl_width / nc)
    for j, hdr in enumerate(headers):
        set_cell(tbl.cell(0, j), hdr, bold=True, color=WHITE, align=PP_ALIGN.CENTER, fill=NAVY)
    for ri, row in enumerate(rows):
        bg_c = ICE if ri % 2 == 0 else WHITE
        for j, val in enumerate(row):
            al = PP_ALIGN.LEFT if j == 0 else PP_ALIGN.CENTER
            set_cell(tbl.cell(ri + 1, j), val, color=TEXT_D, align=al, fill=bg_c)
    add_footer(s)
    return s

# ═══════════════════════════════════════════════════════════
# SLIDES — PENERAPAN SPBE
# ═══════════════════════════════════════════════════════════

# ─── 1. COVER ───
s = blank()
bg = s.background; bg.fill.solid(); bg.fill.fore_color.rgb = NAVY
# Decorative ovals
add_oval(s, Inches(-1), Inches(-1.5), Inches(4.5), NAVY_L)
add_oval(s, Inches(8.5), Inches(-2), Inches(7), NAVY_L)
add_oval(s, Inches(10), Inches(4.5), Inches(5), NAVY_L)
add_oval(s, Inches(0.5), Inches(5.5), Inches(2.5), RGBColor(0x0D, 0x1F, 0x3C))
# Gold accent bar
add_shape(s, MSO_SHAPE.RECTANGLE, M, Inches(2.6), Inches(4), Inches(0.04), fill=GOLD)
# Texts
add_box(s, M, Inches(0.5), Inches(5), Inches(0.35),
        "SISTEM PEMERINTAHAN BERBASIS ELEKTRONIK", 12, bold=True, color=GOLD)
add_box(s, M, Inches(0.85), Inches(5), Inches(0.35),
        "Republik Indonesia", 14, bold=True, color=WHITE)
add_box(s, M, Inches(1.8), CW, Inches(0.45),
        "PENERAPAN SPBE", 20, bold=True, color=WHITE)
add_box(s, M, Inches(2.15), CW, Inches(0.4),
        "Transformasi Birokrasi Terintegrasi", 15, bold=True, color=GOLD)
add_box(s, M, Inches(3.1), CW, Inches(2.0),
        "Efisien • Aman • Tepercaya\nTerintegrasi", 40, bold=True, color=WHITE)
# Badge
add_rrect(s, M, Inches(5.4), Inches(7.5), Inches(0.9), fill=NAVY_L, line=None)
add_box(s, M + Inches(0.2), Inches(5.45), Inches(7), Inches(0.4),
        "Perpres No. 95 Tahun 2018  •  Perpres No. 132 Tahun 2022", 12, color=ICE)
add_box(s, M + Inches(0.2), Inches(5.75), Inches(7), Inches(0.35),
        "Kementerian PAN-RB — Pedoman SPBE Nasional", 10, color=TEXT_L)
# Gold bottom bar
add_shape(s, MSO_SHAPE.RECTANGLE, 0, SLIDE_H - Inches(0.22), SLIDE_W, Inches(0.22), fill=GOLD)

# ─── 2. DAFTAR ISI ───
section_slide("Daftar Isi", "Pokok bahasan presentasi SPBE",
              "Apa Itu SPBE • Arsitektur • 4 Domain • 6 Langkah • Prinsip Utama")

# ─── 3. APA ITU SPBE ───
section_slide("Apa Itu SPBE?",
              "Sistem Pemerintahan Berbasis Elektronik",
              "Dasar Hukum: Perpres No. 95/2018 & Perpres No. 132/2022")

# ─── 4. DEFINISI SPBE ───
card_grid("Memahami SPBE — Bukan Sekadar Digitalisasi", [
    {"ic": "⚡", "t": "Bukan Digitalisasi Biasa", "clr": GOLD,
     "items": ["Bukan memindahkan kertas ke aplikasi",
               "Transformasi menyeluruh birokrasi",
               "Menciptakan tata kelola terintegrasi"]},
    {"ic": "📜", "t": "Dasar Hukum", "clr": BLUE,
     "items": ["Perpres No. 95 Tahun 2018",
               "Perpres No. 132 Tahun 2022",
               "Pedoman dari Kemenpan-RB"]},
    {"ic": "🎯", "t": "Tujuan Akhir", "clr": TEAL,
     "items": ["Birokrasi terintegrasi & efisien",
               "Layanan prima bagi masyarakat",
               "Pemerintahan aman & tepercaya"]},
], subtitle="SPBE adalah transformasi menyeluruh — bukan proyek TIK biasa")

# ─── 5. KERANGKA ARSITEKTUR ───
callout_slide("Kerangka Arsitektur SPBE", [
    ("Layanan SPBE", "Tujuan akhir: layanan\npublik & administrasi", NAVY),
    ("Aplikasi", "Platform digital\nterpadu & interoperabel", BLUE),
    ("Data", "Satu Data Indonesia\nstandar & metadata", TEAL),
    ("Infrastruktur", "PDN, cloud, jaringan\naman & efisien", WARM),
], subtitle="Tata Kelola dan Manajemen sebagai pondasi pengikat",
note_text="Tata Kelola → perencanaan arsitektur, peta rencana, integrasi sistem\n"
          "Manajemen → pengelolaan risiko, data, keamanan, SDM, aset TIK\n\n"
          "💡 Prinsip: Seluruh komponen saling menopang untuk satu tujuan — Layanan SPBE yang prima.")

# ─── 6. 4 DOMAIN SECTION ───
section_slide("4 Domain Evaluasi SPBE",
              "Indikator Tingkat Kematangan (Maturity Index)",
              "Penilaian oleh Kemenpan-RB")

# ─── 7. 4 DOMAIN TABLE ───
table_slide("Empat Domain Evaluasi SPBE",
            ["Domain", "Fokus Utama", "Output yang Diharapkan"],
            [["Kebijakan Internal", "Regulasi, pedoman, SK formal",
              "Perbup/Perwali/Perka terkait Tim Koordinasi & Tata Kelola TIK"],
             ["Tata Kelola", "Perencanaan arsitektur, peta rencana, integrasi",
              "Dokumen Arsitektur SPBE & Peta Rencana 5 tahun"],
             ["Manajemen", "Pengelolaan risiko, data, keamanan, SDM, aset TIK",
              "Penerapan Satu Data Indonesia & Manajemen Risiko SPBE"],
             ["Layanan", "Kualitas fungsi aplikasi publik & administrasi",
              "Aplikasi terintegrasi — bukan platform terpisah-pisah"]],
            subtitle="Indikator tingkat kematangan instansi")

# ─── 8-11. DETAIL 4 DOMAIN ───
card_grid("Domain 1: Kebijakan Internal — Regulasi & SK Formal", [
    {"ic": "1", "t": "Peraturan Kepala Daerah", "clr": NAVY,
     "items": ["Terbitkan Perbup/Perwali/Perka Tim Koordinasi SPBE",
               "Landasan hukum formal implementasi"]},
    {"ic": "2", "t": "Pedoman TIK Internal", "clr": NAVY,
     "items": ["Pedoman tata kelola TIK internal instansi",
               "Selaras dengan arah strategis nasional"]},
    {"ic": "3", "t": "SK Formal Organisasi", "clr": NAVY,
     "items": ["SK pembentukan struktur organisasi SPBE",
               "Melibatkan seluruh unit terkait"]},
    {"ic": "4", "t": "Harmonisasi Regulasi", "clr": NAVY,
     "items": ["Harmonisasi antar sektor terkait",
               "Koordinasi lintas bidang"]},
], cols=4)

card_grid("Domain 2: Tata Kelola — Arsitektur & Peta Rencana", [
    {"ic": "1", "t": "Arsitektur SPBE Instansi", "clr": BLUE,
     "items": ["Proses bisnis, data, aplikasi, infrastruktur",
               "Selaras dengan Arsitektur Nasional"]},
    {"ic": "2", "t": "Peta Rencana 5 Tahun", "clr": BLUE,
     "items": ["Roadmap implementasi SPBE",
               "Tahapan dan milestone jelas"]},
    {"ic": "3", "t": "Integrasi Sistem", "clr": BLUE,
     "items": ["Integrasi lintas unit kerja",
               "Aplikasi Arsitektur SPBE Nasional"]},
    {"ic": "4", "t": "Monitoring & Evaluasi", "clr": BLUE,
     "items": ["Evaluasi berkala kemajuan",
               "Penyesuaian rencana"]},
], cols=4)

card_grid("Domain 3: Manajemen — Risiko, Data, SDM & Aset TIK", [
    {"ic": "1", "t": "Satu Data Indonesia", "clr": TEAL,
     "items": ["Standar data, metadata, interoperabilitas",
               "Referensi data nasional"]},
    {"ic": "2", "t": "Manajemen Risiko SPBE", "clr": TEAL,
     "items": ["Identifikasi dan mitigasi risiko",
               "Berbasis standar manajemen risiko"]},
    {"ic": "3", "t": "SDM TIK Profesional", "clr": TEAL,
     "items": ["Pengelolaan SDM berkelanjutan",
               "Sertifikasi kompetensi"]},
    {"ic": "4", "t": "Aset & Keamanan TIK", "clr": TEAL,
     "items": ["Inventarisasi aset TIK",
               "Kebijakan pengamanan data"]},
], cols=4)

card_grid("Domain 4: Layanan — Kualitas Aplikasi & Integrasi", [
    {"ic": "1", "t": "Platform Terpadu", "clr": GOLD,
     "items": ["Integrasi layanan publik",
               "Super Apps untuk masyarakat"]},
    {"ic": "2", "t": "Layanan Administrasi", "clr": GOLD,
     "items": ["Layanan internal efisien",
               "Kepegawaian, keuangan, dll"]},
    {"ic": "3", "t": "Interoperabilitas", "clr": GOLD,
     "items": ["Sistem saling bertukar data",
               "API terdokumentasi"]},
    {"ic": "4", "t": "Kualitas & Kepuasan", "clr": GOLD,
     "items": ["Standar kualitas layanan digital",
               "Ukur kepuasan pengguna berkala"]},
], cols=4)

# ─── 12. 6 LANGKAH SECTION ───
section_slide("6 Langkah Penerapan SPBE",
              "Urutan implementasi wajib diikuti secara runut",
              "Hindari pemborosan anggaran TIK")

# ─── 13. RINGKASAN 6 LANGKAH ───
flow_slide("Ringkasan 6 Langkah Implementasi", [
    ("1", "Tim Koordinasi", "SK Tim lintas\nsektor", NAVY),
    ("2", "Arsitektur &\nRoadmap", "Dokumen\n5 tahun", BLUE),
    ("3", "Infrastruktur\n& Data", "PDN, SPLP,\nSatu Data", TEAL),
    ("4", "Konsolidasi\nAplikasi", "Super Apps\nterpadu", WARM),
    ("5", "Keamanan\nInformasi", "ISO 27001,\nBSSN", RED_C),
    ("6", "Evaluasi &\nAudit", "Self-assessment\nAudit TIK", GOLD),
], subtitle="Tahapan implementasi SPBE yang berurutan",
note_text="📌 Urutan bersifat wajib — jangan melompati langkah. Setiap langkah adalah fondasi bagi langkah selanjutnya.")

# ─── 14. LANGKAH 1 ───
two_col_cards("Langkah 1: Pembentukan Tim Koordinasi SPBE",
    ["$Pembentukan Tim", "SK Kepala Daerah / Menteri",
     "Dipimpin Sekda / Sekjen", "Bersifat lintas sektor",
     "Bukan proyek TIK — transformasi birokrasi"],
    ["$Unsur yang Terlibat", "Perencanaan (Bappeda)",
     "Organisasi (Ortala)", "Keuangan (BPKAD)",
     "Teknis TIK (Diskominfo)"],
    subtitle="Langkah Awal / Fondasi Implementasi",
    left_color=GOLD, right_color=BLUE)

# ─── 15. LANGKAH 2 ───
two_col_cards("Langkah 2: Penyusunan Arsitektur & Peta Rencana",
    ["$Arsitektur SPBE Instansi", "Petakan Proses Bisnis dulu",
     "Data & Informasi", "Aplikasi yang dibutuhkan",
     "Infrastruktur & Keamanan"],
    ["$Peta Rencana 5 Tahun", "Gunakan Arsitektur SPBE Nasional",
     "Selaras dengan strategi nasional",
     "Milestone & target per tahun",
     "Anggaran & sumber daya"],
    subtitle="Proses Perencanaan Strategis",
    left_color=BLUE, right_color=TEAL)

# ─── 16. LANGKAH 3 ───
two_col_cards("Langkah 3: Standardisasi Infrastruktur & Integrasi Data",
    ["$Infrastruktur", "Batasi data center mandiri",
     "Migrasi ke PDN / cloud", "Hub SPLP interoperabilitas",
     "Efisiensi anggaran TIK"],
    ["$Data", "Prinsip Satu Data Indonesia",
     "Standar data & metadata", "Interoperabilitas data",
     "Referensi data nasional"],
    subtitle="Fase Eksekusi Teknis",
    left_color=TEAL, right_color=BLUE)

# ─── 17. LANGKAH 4 ───
two_col_cards("Langkah 4: Konsolidasi Aplikasi & Layanan Digital",
    ["$Yang Harus Dilakukan", "Audit seluruh aplikasi existing",
     "Moratorium aplikasi duplikatif", "Gabung ke platform terpadu",
     "Super Apps publik & internal"],
    ["$Yang Harus Dihindari", "Aplikasi baru tiap masalah",
     "Platform terpisah-pisah", "Membangun dari nol",
     "Mengabaikan interoperabilitas"],
    subtitle="Penyederhanaan Layanan Digital",
    left_color=WARM, right_color=RED_C)

# ─── 18. LANGKAH 5 ───
two_col_cards("Langkah 5: Penerapan Manajemen Keamanan Informasi",
    ["$Standar & Kolaborasi", "SMKI berbasis ISO 27001",
     "Kolaborasi dengan BSSN", "Vulnerability Assessment berkala",
     "Security by design"],
    ["$Tim & Prosedur", "CSIRT internal instansi",
     "Penetration testing rutin", "Audit keamanan informasi",
     "Insiden response plan"],
    subtitle="Aspek Keamanan & Kepatuhan",
    left_color=RED_C, right_color=GOLD)

# ─── 19. LANGKAH 6 ───
two_col_cards("Langkah 6: Pemantauan, Evaluasi, dan Audit",
    ["$Evaluasi Mandiri", "Self-assessment berkala",
     "Pantau kemajuan implementasi", "Persiapan penilaian Kemenpan-RB",
     "Gunakan hasil untuk perbaikan"],
    ["$Audit TIK", "Efisiensi infrastruktur",
     "Audit aplikasi & layanan", "Kepatuhan keamanan informasi",
     "Laporan rekomendasi perbaikan"],
    subtitle="Continuous Improvement",
    left_color=GOLD, right_color=BLUE)

# ─── 20. PRINSIP UTAMA SECTION ───
section_slide("Prinsip Utama",
              "Kunci Keberhasilan SPBE",
              "Interoperabilitas — Bukan Jumlah Aplikasi")

# ─── 21. INTEROPERABILITAS ───
callout_slide("Jebakan Terbesar SPBE", [
    ("❌", "Puluhan aplikasi terpisah\nMasyarakat bingung\nAnggaran membengkak", RED_C),
    ("✅", "Satu platform terpadu\nInteroperabilitas\nEfisiensi & sinergi", TEAL),
], subtitle="Menganggap digitalisasi = membuat aplikasi baru untuk setiap masalah",
note_text="💡 Paradigma SPBE Modern:\n"
          "• Fokus pada interoperabilitas — bagaimana sistem yang sudah ada bisa saling berbicara\n"
          "• Bertukar data tanpa memaksa masyarakat mengunduh puluhan aplikasi berbeda\n"
          "• Satu platform terpadu, banyak layanan — bukan banyak platform, satu layanan\n"
          "• SPBE adalah transformasi birokrasi, bukan proyek TIK")

# ─── 22. CLOSING ───
s = blank()
bg = s.background; bg.fill.solid(); bg.fill.fore_color.rgb = NAVY
add_oval(s, Inches(-1.5), Inches(-1.5), Inches(5), NAVY_L)
add_oval(s, Inches(9.5), Inches(4), Inches(5), NAVY_L)
add_shape(s, MSO_SHAPE.RECTANGLE, 0, Inches(3.3), SLIDE_W, Inches(0.06), fill=GOLD)

add_box(s, M, Inches(2.0), CW, Inches(1.0),
        "Terima Kasih", 40, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_box(s, M, Inches(3.5), CW, Inches(0.8),
        "Mari Wujudkan SPBE yang Terintegrasi, Efisien, Aman, dan Tepercaya",
        18, color=GOLD, align=PP_ALIGN.CENTER)
add_box(s, M, Inches(4.3), CW, Inches(0.5),
        "Bersama membangun birokrasi digital Indonesia yang melayani",
        14, color=TEXT_L, align=PP_ALIGN.CENTER)
add_shape(s, MSO_SHAPE.RECTANGLE, 0, SLIDE_H - Inches(0.22), SLIDE_W, Inches(0.22), fill=GOLD)
add_box(s, M, SLIDE_H - Inches(0.42), Inches(6), Inches(0.22),
        SOURCE_TEXT, 7, color=TEXT_L)
pg_counter[0] += 1
add_box(s, SLIDE_W - Inches(1.0), SLIDE_H - Inches(0.42), Inches(0.8), Inches(0.22),
        str(pg_counter[0]), 8, color=TEXT_L, align=PP_ALIGN.RIGHT)

# ═══════════════════════════════════════════════════════════
# SAVE
# ═══════════════════════════════════════════════════════════

output = "Penerapan_SPBE.pptx"
prs.save(output)
size = os.path.getsize(output)
print(f"✅ {output}")
print(f"   {len(prs.slides)} slides, {size:,} bytes")
print(f"   Palette: Navy #0A1628, Gold #C8962E, White, Ice #E8EDF5")
print(f"   Archetypes: cover, section, card_grid, two_col, callout, flow, table")
