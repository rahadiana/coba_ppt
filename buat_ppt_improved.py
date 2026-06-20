#!/usr/bin/env python3
"""
Perwal Bekasi No 51/2024 — Pajak Reklame
VERSI IMPROVED — Action Titles, Source Citations, Better Visual Hierarchy

Berdasarkan riset:
- McKinsey/BCG/Bain: action titles, pyramid principle, one message per slide
- Color Theory 2026: 60-30-10 rule, 4.5:1 contrast, max 5 colors
- Government design: trust, transparency, accessibility (WCAG 2.1)
- Presentation Zen: whitespace, signal vs noise, picture superiority
- Duarte: color with purpose, visual hierarchy
- Tax Policy Center: data visualization standards
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
import os

# ═══════════════════════════════════════════════════════════
# CONSTANTS
# ═══════════════════════════════════════════════════════════

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)
M = Inches(0.6)
CW = SLIDE_W - 2 * M

# ─── Color Palette (60-30-10) ───
# 60% Dominant
NAVY     = RGBColor(0x0A, 0x16, 0x28)
NAVY_L   = RGBColor(0x12, 0x29, 0x4A)
NAVY_M   = RGBColor(0x1B, 0x3A, 0x6B)

# 30% Secondary
WHITE    = RGBColor(0xFF, 0xFF, 0xFF)
OFF_W    = RGBColor(0xF5, 0xF7, 0xFA)
ICE      = RGBColor(0xE8, 0xED, 0xF5)
ICE_D    = RGBColor(0xD0, 0xD8, 0xE8)

# 10% Accent
GOLD     = RGBColor(0xC8, 0x96, 0x2E)
GOLD_L   = RGBColor(0xD4, 0xA0, 0x17)
GOLD_M   = RGBColor(0xB8, 0x86, 0x0B)

# Semantic
TEXT_D   = RGBColor(0x1A, 0x1A, 0x2E)
TEXT_M   = RGBColor(0x6B, 0x72, 0x88)
TEXT_L   = RGBColor(0x9C, 0xA3, 0xAF)
BLUE     = RGBColor(0x25, 0x63, 0xEB)
TEAL     = RGBColor(0x0D, 0x94, 0x88)
WARM     = RGBColor(0xB8, 0x86, 0x0B)
RED_C    = RGBColor(0xDC, 0x26, 0x26)
GREEN    = RGBColor(0x2E, 0x7D, 0x32)
ORANGE   = RGBColor(0xE6, 0x51, 0x00)

# ─── Presentation Setup ───
prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H
pg_counter = [0]

SOURCE_TEXT = "Sumber: Perwal Bekasi No 51/2024"

# ═══════════════════════════════════════════════════════════
# HELPERS
# ═══════════════════════════════════════════════════════════

def blank():
    """Create a blank slide."""
    return prs.slides.add_slide(prs.slide_layouts[6])

def add_box(slide, left, top, width, height, text, size=12, bold=False,
            color=TEXT_D, align=PP_ALIGN.LEFT, font="Calibri", vAlign=MSO_ANCHOR.TOP):
    """Add a text box."""
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font
    p.alignment = align
    p.space_after = Pt(0)
    p.space_before = Pt(0)
    return tb

def add_shape(slide, shape_type, left, top, width, height,
              fill=None, line=None, lw=None):
    """Add a shape with optional fill/line."""
    sh = slide.shapes.add_shape(shape_type, left, top, width, height)
    if fill:
        sh.fill.solid()
        sh.fill.fore_color.rgb = fill
    else:
        sh.fill.background()
    if line:
        sh.line.color.rgb = line
        if lw:
            sh.line.width = Pt(lw)
    else:
        sh.line.fill.background()
    return sh

def add_rrect(slide, left, top, width, height,
              fill=WHITE, line=ICE, lw=0.5, radius=0.04):
    """Add a rounded rectangle with consistent styling."""
    sh = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    sh.fill.solid()
    sh.fill.fore_color.rgb = fill
    if line:
        sh.line.color.rgb = line
        sh.line.width = Pt(lw)
    else:
        sh.line.fill.background()
    sh.adjustments[0] = radius
    return sh

def add_oval(slide, left, top, size, fill):
    """Add an oval/circle."""
    sh = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top, size, size)
    sh.fill.solid()
    sh.fill.fore_color.rgb = fill
    sh.line.fill.background()
    return sh

def add_text_to_shape(shape, text, size=12, bold=False, color=TEXT_D,
                      align=PP_ALIGN.LEFT, font="Calibri"):
    """Set text into an existing shape's text frame."""
    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font
    p.alignment = align
    return tf

def add_multiline(slide, left, top, width, height, lines, default_size=12,
                  default_color=TEXT_D):
    """lines: list of (text, size, bold, color) or strings."""
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        if isinstance(line, str):
            if line == "":
                p.text = ""
                p.font.size = Pt(6)
                p.space_after = Pt(0)
                p.space_before = Pt(0)
                continue
            text, sz, bd, clr = line, default_size, False, default_color
        else:
            text = line[0]
            sz = line[1] if len(line) > 1 else default_size
            bd = line[2] if len(line) > 2 else False
            clr = line[3] if len(line) > 3 else default_color
        p.text = text
        p.font.size = Pt(sz)
        p.font.bold = bd
        p.font.color.rgb = clr
        p.font.name = "Calibri"
        p.alignment = PP_ALIGN.LEFT
        p.space_after = Pt(1)
        p.space_before = Pt(0)
    return tb

def set_cell(cell, text, size=11, bold=False, color=TEXT_D,
             align=PP_ALIGN.LEFT, fill=None):
    """Format a table cell."""
    cell.text = ""
    p = cell.text_frame.paragraphs[0]
    p.text = str(text)
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = "Calibri"
    p.alignment = align
    cell.text_frame.word_wrap = True
    cell.vertical_anchor = MSO_ANCHOR.MIDDLE
    if fill:
        cell.fill.solid()
        cell.fill.fore_color.rgb = fill

# ═══════════════════════════════════════════════════════════
# DECORATIVE / FOOTER ELEMENTS
# ═══════════════════════════════════════════════════════════

def gold_top_bar(slide):
    """Thin gold bar across top."""
    add_shape(slide, MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, Inches(0.035), fill=GOLD)

def navy_header_bg(slide, title, action_title=None, subtitle=None):
    """Navy header band with action title + optional subtitle."""
    gold_top_bar(slide)
    add_shape(slide, MSO_SHAPE.RECTANGLE, 0, Inches(0.035), SLIDE_W, Inches(0.9), fill=NAVY)
    # Gold vertical accent
    add_shape(slide, MSO_SHAPE.RECTANGLE, M, Inches(0.12), Inches(0.07), Inches(0.6), fill=GOLD)
    # Action title (main heading — conclusion first!)
    y_title = Inches(0.15)
    add_box(slide, M + Inches(0.22), y_title, CW - Inches(0.5), Inches(0.48),
            title, 20, bold=True, color=WHITE)
    # Subtitle line
    y_sub = y_title + Inches(0.48)
    parts = []
    if action_title:
        parts.append(action_title)
    if subtitle:
        parts.append(subtitle)
    combined = " • ".join(parts)
    add_box(slide, M + Inches(0.22), y_sub, CW - Inches(0.5), Inches(0.3),
            combined, 9, color=TEXT_L)

def navy_bottom_bar(slide):
    """Navy bar at bottom of slide."""
    add_shape(slide, MSO_SHAPE.RECTANGLE, 0, SLIDE_H - Inches(0.25), SLIDE_W, Inches(0.03), fill=NAVY)

def gold_bottom_bar(slide):
    """Gold bar at very bottom."""
    add_shape(slide, MSO_SHAPE.RECTANGLE, 0, SLIDE_H - Inches(0.25), SLIDE_W, Inches(0.25), fill=GOLD)

def add_page_num(slide):
    """Add page number at bottom right."""
    pg_counter[0] += 1
    add_box(slide, SLIDE_W - Inches(1.0), SLIDE_H - Inches(0.42), Inches(0.8), Inches(0.22),
            str(pg_counter[0]), 8, color=TEXT_L, align=PP_ALIGN.RIGHT)

def add_source(slide):
    """Add source citation at bottom left."""
    add_box(slide, M, SLIDE_H - Inches(0.42), Inches(4), Inches(0.22),
            SOURCE_TEXT, 7, color=TEXT_L)

def add_footer(slide):
    """Complete footer: source + page number + bottom bar."""
    navy_bottom_bar(slide)
    add_source(slide)
    add_page_num(slide)

# ═══════════════════════════════════════════════════════════
# COMPLEX BUILDING BLOCKS
# ═══════════════════════════════════════════════════════════

def section_slide(title, subtitle=None, action_text=None):
    """Full-bleed dark section divider with decorative ovals."""
    s = blank()
    bg = s.background; bg.fill.solid(); bg.fill.fore_color.rgb = NAVY
    # Decorative ovals
    add_oval(s, Inches(-1.5), Inches(-1.5), Inches(5), NAVY_L)
    add_oval(s, Inches(-0.5), Inches(-0.5), Inches(3.5), RGBColor(0x0D, 0x1F, 0x3C))
    add_oval(s, Inches(9.5), Inches(4), Inches(5), RGBColor(0x0D, 0x1F, 0x3C))
    add_oval(s, Inches(10.5), Inches(3), Inches(3), RGBColor(0x0E, 0x20, 0x3E))
    # Gold bar
    add_shape(s, MSO_SHAPE.RECTANGLE, M, Inches(2.3), Inches(2.5), Inches(0.04), fill=GOLD)
    # Title
    add_box(s, M, Inches(2.6), CW, Inches(1.6), title, 34, bold=True, color=WHITE)
    # Action text or subtitle
    y_sub = Inches(4.2)
    if action_text:
        add_box(s, M, y_sub, CW, Inches(0.4), action_text, 14, bold=True, color=GOLD)
        y_sub += Inches(0.35)
    if subtitle:
        add_box(s, M, y_sub, CW, Inches(0.3), subtitle, 11, color=TEXT_L)
    # Gold bottom bar
    add_shape(s, MSO_SHAPE.RECTANGLE, 0, SLIDE_H - Inches(0.22), SLIDE_W, Inches(0.22), fill=GOLD)
    add_page_num(s)
    return s

def content_slide(action_title, subtitle=None, ref=None):
    """Standard content slide: gold bar → navy header → content area."""
    s = blank()
    bg = s.background; bg.fill.solid(); bg.fill.fore_color.rgb = OFF_W
    parts = []
    if subtitle:
        parts.append(subtitle)
    if ref:
        parts.append(ref)
    sub_text = " • ".join(parts) if parts else None
    navy_header_bg(s, action_title, sub_text)
    return s

def card_grid(action_title, cards, subtitle=None, ref=None, cols=0):
    """Grid of cards with multi-row support. cols=0 means auto (4 cols, 2 if n<=2)."""
    s = content_slide(action_title, subtitle, ref)
    n = len(cards)
    per_row = cols if cols > 0 else (2 if n <= 2 else (3 if n <= 3 else 4))
    gap = Inches(0.3)
    cw = (SLIDE_W - M * 2 - (per_row - 1) * gap) / per_row
    gap_h = Inches(0.3)
    n_rows = (n + per_row - 1) // per_row
    content_h = SLIDE_H - Inches(1.15) - Inches(0.5)  # ~5.85"
    ch = (content_h - (n_rows - 1) * gap_h) / n_rows

    # Calculate item height based on max items per card
    max_items = max(len(cd.get('items', [])) for cd in cards)
    item_h = Inches(0.35)  # enough for 2 lines of wrapped text
    title_h = Inches(0.3)

    for i, cd in enumerate(cards):
        row = i // per_row
        col = i % per_row
        cx = M + col * (cw + gap)
        cy = Inches(1.15) + row * (ch + gap_h)

        clr = cd.get('clr', BLUE)
        ic = cd.get('ic', '')
        t = cd.get('t', '')
        items = cd.get('items', [])

        add_rrect(s, cx, cy, cw, ch)
        # Left accent bar
        add_shape(s, MSO_SHAPE.RECTANGLE, cx, cy, Inches(0.05), ch, fill=clr)
        # Icon circle
        if ic:
            add_oval(s, cx + Inches(0.15), cy + Inches(0.15), Inches(0.42), clr)
            add_box(s, cx + Inches(0.15), cy + Inches(0.15), Inches(0.42), Inches(0.42),
                    ic, 13, color=WHITE, align=PP_ALIGN.CENTER)
            ty = cy + Inches(0.7)
        else:
            ty = cy + Inches(0.15)
        # Card title
        add_box(s, cx + Inches(0.15), ty, cw - Inches(0.3), title_h,
                t, 13, bold=True, color=clr)
        # Items with enough height for text wrapping
        ay = ty + title_h + Inches(0.08)
        for item in items:
            add_box(s, cx + Inches(0.15), ay, cw - Inches(0.3), item_h,
                    f"• {item}", 9, color=TEXT_D)
            ay += item_h + Inches(0.02)

    add_footer(s)
    return s

def two_col_cards(action_title, left_data, right_data, subtitle=None, ref=None,
                  left_color=BLUE, right_color=TEAL):
    """Two-column layout with card containers."""
    s = content_slide(action_title, subtitle, ref)
    gap = Inches(0.3)
    cw = (SLIDE_W - gap - 2 * M) / 2  # proper width: (13.333 - 0.3 - 1.2) / 2 = 5.9165"
    sy = Inches(1.15)
    ch = SLIDE_H - sy - Inches(0.7)

    for data, clr, x in [
        (left_data, left_color, M),
        (right_data, right_color, M + cw + gap),
    ]:
        add_rrect(s, x, sy, cw, ch)
        add_shape(s, MSO_SHAPE.RECTANGLE, x, sy, Inches(0.05), ch, fill=clr)
        ay = sy + Inches(0.15)
        for line in data:
            if line.startswith("$"):
                add_box(s, x + Inches(0.2), ay, cw - Inches(0.35), Inches(0.3),
                        line[1:], 14, bold=True, color=clr)
                ay += Inches(0.32)
            else:
                add_box(s, x + Inches(0.2), ay, cw - Inches(0.35), Inches(0.22),
                        f"• {line}", 11, color=TEXT_D)
                ay += Inches(0.24)
    add_footer(s)
    return s

def callout_slide(action_title, callouts, subtitle=None, ref=None, note_text=None):
    """4 callout cards with big numbers + description."""
    s = content_slide(action_title, subtitle, ref)
    n = len(callouts)
    gap = Inches(0.3)
    cw = (SLIDE_W - 2 * M - (n - 1) * gap) / n
    sx = M
    sy = Inches(1.15)

    for i, (num, lb, clr) in enumerate(callouts):
        cx = sx + i * (cw + gap)
        add_rrect(s, cx, sy, cw, Inches(2.2))
        add_shape(s, MSO_SHAPE.RECTANGLE, cx, sy, cw, Inches(0.05), fill=clr)
        add_box(s, cx, sy + Inches(0.25), cw, Inches(0.7),
                num, 32, bold=True, color=clr, align=PP_ALIGN.CENTER)
        add_box(s, cx + Inches(0.1), sy + Inches(1.0), cw - Inches(0.2), Inches(0.7),
                lb, 11, color=TEXT_M, align=PP_ALIGN.CENTER, vAlign=MSO_ANCHOR.TOP)

    if note_text:
        ny = sy + Inches(2.5)
        add_rrect(s, M, ny, CW, Inches(3.0), fill=ICE)
        add_box(s, M + Inches(0.25), ny + Inches(0.15), CW - Inches(0.5), Inches(2.6),
                note_text, 11, color=TEXT_D)

    add_footer(s)
    return s

def flow_slide(action_title, steps_data, subtitle=None, ref=None, note_text=None):
    """Horizontal flow diagram with arrow connectors."""
    s = content_slide(action_title, subtitle, ref)
    n = len(steps_data)
    bgap = Inches(0.35)
    bw = (SLIDE_W - 2 * M - (n - 1) * bgap) / n
    ssx = M
    sy = Inches(1.3)

    for i, (num, title, desc, clr) in enumerate(steps_data):
        cx = ssx + i * (bw + bgap)
        add_rrect(s, cx, sy, bw, Inches(2.0))
        add_shape(s, MSO_SHAPE.RECTANGLE, cx, sy, bw, Inches(0.05), fill=clr)
        # Number circle
        circ_x = cx + bw / 2 - Inches(0.25)
        add_oval(s, circ_x, sy + Inches(0.15), Inches(0.5), clr)
        add_box(s, circ_x, sy + Inches(0.15), Inches(0.5), Inches(0.5),
                num, 18, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        # Title
        add_box(s, cx + Inches(0.1), sy + Inches(0.75), bw - Inches(0.2), Inches(0.35),
                title, 14, bold=True, color=clr, align=PP_ALIGN.CENTER)
        # Description
        add_box(s, cx + Inches(0.1), sy + Inches(1.1), bw - Inches(0.2), Inches(0.7),
                desc, 10, color=TEXT_M, align=PP_ALIGN.CENTER, vAlign=MSO_ANCHOR.TOP)
        # Arrow connector
        if i < n - 1:
            add_box(s, cx + bw, sy + Inches(0.7), bgap, Inches(0.5),
                    "›", 24, color=GOLD, align=PP_ALIGN.CENTER, vAlign=MSO_ANCHOR.MIDDLE)

    if note_text:
        ny = sy + Inches(2.4)
        add_rrect(s, M, ny, CW, Inches(3.0), fill=ICE)
        add_box(s, M + Inches(0.25), ny + Inches(0.15), CW - Inches(0.5), Inches(2.6),
                note_text, 11, color=TEXT_D)

    add_footer(s)
    return s

def table_slide(action_title, headers, rows, subtitle=None, ref=None):
    """Native PowerPoint table slide."""
    s = content_slide(action_title, subtitle, ref)
    tbl_left = M
    tbl_top = Inches(1.15)
    tbl_width = CW
    rh = Inches(0.42)
    nr = len(rows) + 1
    nc = len(headers)
    ts = s.shapes.add_table(nr, nc, tbl_left, tbl_top, tbl_width, rh * nr)
    tbl = ts.table
    for i in range(nc):
        tbl.columns[i].width = int(tbl_width / nc)
    for j, hdr in enumerate(headers):
        set_cell(tbl.cell(0, j), hdr, bold=True, color=WHITE, align=PP_ALIGN.CENTER, fill=NAVY)
    for ri, row in enumerate(rows):
        bg_c = ICE if ri % 2 == 0 else WHITE
        for j, val in enumerate(row):
            al = PP_ALIGN.LEFT if j == 0 else PP_ALIGN.CENTER
            set_cell(tbl.cell(ri + 1, j), val, color=TEXT_D, align=al, fill=bg_c)
    add_footer(s)
    return s

# ═══════════════════════════════════════════════════════════
# SLIDES
# ═══════════════════════════════════════════════════════════

# ────────── COVER (1) ──────────
s = blank()
bg = s.background; bg.fill.solid(); bg.fill.fore_color.rgb = NAVY
# Decorative ovals
add_oval(s, Inches(-1), Inches(-1.5), Inches(4.5), RGBColor(0x0D, 0x1F, 0x3C))
add_oval(s, Inches(8.5), Inches(-2), Inches(7), RGBColor(0x0D, 0x1F, 0x3C))
add_oval(s, Inches(10), Inches(4.5), Inches(5), RGBColor(0x0D, 0x1F, 0x3C))
add_oval(s, Inches(0.5), Inches(5.5), Inches(2.5), RGBColor(0x12, 0x29, 0x4A))
# Gold accent bar
add_shape(s, MSO_SHAPE.RECTANGLE, M, Inches(2.6), Inches(4), Inches(0.04), fill=GOLD)
# Top labels
add_box(s, M, Inches(0.5), Inches(5), Inches(0.35), "BERITA DAERAH", 12, bold=True, color=GOLD)
add_box(s, M, Inches(0.85), Inches(5), Inches(0.35), "KOTA BEKASI", 14, bold=True, color=WHITE)
# Title block
add_box(s, M, Inches(1.8), CW, Inches(0.45), "PERATURAN WALI KOTA BEKASI", 20, bold=True, color=WHITE)
add_box(s, M, Inches(2.15), CW, Inches(0.4), "NOMOR 51 TAHUN 2024", 15, bold=True, color=GOLD)
add_box(s, M, Inches(3.1), CW, Inches(2.0), "TENTANG\nPENGELOLAAN PAJAK REKLAME", 40, bold=True, color=WHITE)
# Info card
crd = add_rrect(s, M, Inches(5.4), Inches(7.5), Inches(0.9), fill=RGBColor(0x0D, 0x1F, 0x3C))
add_box(s, M + Inches(0.2), Inches(5.45), Inches(7), Inches(0.4),
        "Pemerintah Kota Bekasi  ·  20 Desember 2024", 12, color=ICE)
add_box(s, M + Inches(0.2), Inches(5.75), Inches(7), Inches(0.35),
        "Berlaku sejak diundangkan", 10, color=TEXT_L)
gold_bottom_bar(s)

# ────────── DAFTAR ISI (2) ──────────
s = blank(); bg = s.background; bg.fill.solid(); bg.fill.fore_color.rgb = OFF_W
gold_top_bar(s)
add_shape(s, MSO_SHAPE.RECTANGLE, 0, Inches(0.035), SLIDE_W, Inches(0.8), fill=NAVY)
add_shape(s, MSO_SHAPE.RECTANGLE, M, Inches(0.12), Inches(0.07), Inches(0.5), fill=GOLD)
add_box(s, M + Inches(0.2), Inches(0.15), CW - Inches(1), Inches(0.35),
        "11 Bab Kunci Mengatur Seluruh Aspek Pajak Reklame", 20, bold=True, color=WHITE)
add_box(s, M + Inches(0.2), Inches(0.52), CW - Inches(1), Inches(0.25),
        "Perwal Bekasi No 51/2024", 9, color=TEXT_L)

toc = [
    ("1","Ketentuan Umum & Definisi",BLUE), ("2","Objek, Subjek & Wajib Pajak",TEAL),
    ("3","Masa Pajak & Tahun Pajak",WARM), ("4","Pendaftaran & Pendataan WP",NAVY_M),
    ("5","Nilai Sewa Reklame (NSR)",RED_C), ("6","Perhitungan & Tarif Pajak",BLUE),
    ("7","Penetapan, Tagihan & Pembayaran",TEAL), ("8","Pembetulan, Keberatan & Banding",WARM),
    ("9","Pemeriksaan, Penagihan & Penghapusan",NAVY_M),
    ("10","Keringanan, Kemudahan & Penghargaan",BLUE), ("11","Ketentuan Penutup",TEAL),
]
for i, (num, lb, clr) in enumerate(toc):
    col = 0 if i < 6 else 1
    row = i if i < 6 else i - 6
    cw_i = Inches(5.5); gap_i = Inches(0.35)
    sx_i = (SLIDE_W - 2 * cw_i - gap_i) / 2
    x = sx_i + col * (cw_i + gap_i)
    y = Inches(1.1) + row * Inches(0.82)
    add_rrect(s, x, y, cw_i, Inches(0.62))
    badge = add_shape(s, MSO_SHAPE.ROUNDED_RECTANGLE, x + Inches(0.1),
                       y + Inches(0.06), Inches(0.5), Inches(0.5), fill=clr)
    add_text_to_shape(badge, num, 14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_box(s, x + Inches(0.75), y + Inches(0.06), cw_i - Inches(0.9), Inches(0.5),
            lb, 13, color=TEXT_D, vAlign=MSO_ANCHOR.MIDDLE)
add_footer(s)

# ═══════════════════════════════════════════
# BAB I — KETENTUAN UMUM (3-4)
# ═══════════════════════════════════════════
section_slide("BAB I\nKETENTUAN UMUM", "Pasal 1", "7 Definisi Kunci Jadi Landasan Hukum Pajak Reklame")

# Slide 4: 7 Definisi
card_grid("7 Definisi Kunci Menjadi Landasan Pengelolaan Pajak Reklame", [
    {"ic":"🏛️","t":"Daerah","clr":BLUE,"items":["Kota Bekasi"]},
    {"ic":"📊","t":"Bapenda","clr":TEAL,"items":["Badan Pendapatan Daerah Kota Bekasi"]},
    {"ic":"📢","t":"Reklame","clr":WARM,"items":["Media untuk promosi & pengenalan komersial"]},
    {"ic":"💰","t":"Pajak Reklame","clr":NAVY_M,"items":["Pajak atas penyelenggaraan reklame"]},
    {"ic":"📐","t":"NSR","clr":BLUE,"items":["Nilai Sewa Reklame — dasar pengenaan pajak"]},
    {"ic":"🆔","t":"NPWPD","clr":TEAL,"items":["Nomor Pokok Wajib Pajak Daerah"]},
    {"ic":"👤","t":"Wajib Pajak","clr":WARM,"items":["Orang pribadi/badan dg hak & kewajiban pajak"]},
], subtitle="Pasal 1", ref="7 Definisi Kunci")

# ═══════════════════════════════════════════
# BAB II — OBJEK, SUBJEK & WAJIB PAJAK (5-7)
# ═══════════════════════════════════════════
section_slide("BAB II\nOBJEK, SUBJEK & WAJIB PAJAK", "Pasal 2–4",
              "10 Jenis Reklame + 8 Pengecualian + 2 Subjek Pajak")

card_grid("10 Jenis Reklame Wajib Pajak + 8 Jenis Dikecualikan", [
    {"ic":"📋","t":"10 Jenis Reklame","clr":BLUE,"items":["Papan / Billboard","Videotron / Megatron","Kain (Spanduk, Umbul, Baliho)","Melekat / Stiker","Selebaran","Berjalan (Kendaraan)","Udara (Balon Gas)","Apung","Film / Slide","Peragaan"]},
    {"ic":"🚫","t":"Dikecualikan","clr":TEAL,"items":["Internet, TV, radio, media cetak","Label / merek produk","Nama usaha ≤ 1 m² di tempat","Reklame Pemerintah/Pemda","Tempat ibadah & panti asuhan","Sosial & keagamaan ≤ 30 hari","Kegiatan politik (masa kampanye)","Olahraga KONI ≤ 30 hari"]},
], subtitle="Pasal 2–4", ref="10 Jenis + 8 Pengecualian")

card_grid("Subjek & Wajib Pajak: Siapa yang Terkena Kewajiban Pajak", [
    {"ic":"👤","t":"Subjek Pajak (Pasal 3)","clr":BLUE,"items":["Orang pribadi atau Badan","yang menggunakan Reklame"]},
    {"ic":"✋","t":"Wajib Pajak (Pasal 4)","clr":TEAL,"items":["Orang pribadi atau Badan","yang menyelenggarakan Reklame","Jika pihak ketiga → menjadi WP","Mendaftarkan diri & objek pajak"]},
], subtitle="Pasal 3–4", ref="Subjek & Wajib Pajak")

# ═══════════════════════════════════════════
# BAB III — MASA PAJAK (8-9)
# ═══════════════════════════════════════════
section_slide("BAB III\nMASA PAJAK & TAHUN PAJAK", "Pasal 5",
              "12 Bulan atau 30 Hari — Tergantung Jenis Reklame")

callout_slide("Masa Pajak: 12 Bulan Permanen, 30 Hari Insidentil", [
    ("12", "Bulan\n(Permanen)", BLUE),
    ("30", "Hari\n(Insidentil)", TEAL),
    ("1", "Tahun Pajak\n(Kalender)", WARM),
    ("1", "Bulan\n(Bagian Tahun)", NAVY_M),
], subtitle="Pasal 5", ref="Masa & Tahun Pajak",
   note_text="• Masa Pajak Permanen: 12 bulan atau sesuai jangka waktu penayangan reklame\n"
             "• Masa Pajak Insidentil: dihitung per hari, maksimal 30 hari\n"
             "• Tahun Pajak: 1 tahun kalender atau sesuai tahun buku wajib pajak\n"
             "• Bagian Tahun Pajak: 1 bulan penuh (jika tidak mencakup satu tahun penuh)")

# ═══════════════════════════════════════════
# BAB IV — PENDAFTARAN & PENDATAAN (10-11)
# ═══════════════════════════════════════════
section_slide("BAB IV\nPENDAFTARAN & PENDATAAN WP", "Pasal 6–8",
              "Wajib Daftar atau NPWPD Jabatan — Bapenda Berwenang Mendata")

two_col_cards("Pendaftaran Wajib, Pendataan Bapenda, dan Penonaktifan WP", [
    "$PENDAFTARAN (Pasal 6)",
    "",
    "WP wajib mendaftarkan diri & objek pajak",
    "Formulir: ambil/online/dikirim petugas",
    "Lampirkan: KTP, NPWP, Akta, NIB",
    "Bapenda terbitkan NPWPD",
    "Jika tidak mendaftar → NPWPD jabatan",
    "Juga: NOPD & nomor registrasi",
], [
    "$PENDATAAN & NONAKTIF (Pasal 7–8)",
    "",
    "Bapenda mendata WP & objek pajak",
    "Termasuk data geografis",
    "Dapat kerjasama dengan instansi lain",
    "Penonaktifan: WP tak penuhi syarat",
    "Keputusan maksimal 3 bulan",
    "Syarat: tanpa tunggakan & keberatan",
], subtitle="Pasal 6–8", ref="Pendaftaran & Pendataan")

# ═══════════════════════════════════════════
# BAB V — NILAI SEWA REKLAME (12-13)
# ═══════════════════════════════════════════
section_slide("BAB V\nNILAI SEWA REKLAME", "Pasal 9",
              "7 Faktor Penentu NSR — Jenis, Bahan, Lokasi, Waktu, dll")

# 7 factors grid
s = blank(); bg = s.background; bg.fill.solid(); bg.fill.fore_color.rgb = OFF_W
gold_top_bar(s)
add_shape(s, MSO_SHAPE.RECTANGLE, 0, Inches(0.035), SLIDE_W, Inches(0.9), fill=NAVY)
add_shape(s, MSO_SHAPE.RECTANGLE, M, Inches(0.12), Inches(0.07), Inches(0.6), fill=GOLD)
add_box(s, M + Inches(0.22), Inches(0.15), CW - Inches(0.5), Inches(0.48),
        "7 Faktor Penentu Nilai Sewa Reklame (NSR)", 20, bold=True, color=WHITE)
add_box(s, M + Inches(0.22), Inches(0.58), CW - Inches(0.5), Inches(0.3),
        "Pasal 9 • 7 Faktor Penentu", 9, color=TEXT_L)

factors = [
    ("1","Jenis Reklame",BLUE), ("2","Bahan",TEAL), ("3","Lokasi (Kelas Jalan)",WARM),
    ("4","Waktu Tayang\n(detik)",NAVY_M), ("5","Jangka Waktu\n(hari)",BLUE),
    ("6","Jumlah Media",TEAL), ("7","Ukuran (m²)",WARM),
]
bw = Inches(2.8); bgap = Inches(0.3); per_row = 4
tw = per_row * bw + (per_row - 1) * bgap; bsx = (SLIDE_W - tw) / 2
for i, (num, lb, clr) in enumerate(factors):
    col = i % per_row; row = i // per_row
    x = bsx + col * (bw + bgap); y = Inches(1.15) + row * Inches(1.6)
    add_rrect(s, x, y, bw, Inches(1.3))
    add_oval(s, x + Inches(0.15), y + Inches(0.25), Inches(0.55), clr)
    add_box(s, x + Inches(0.15), y + Inches(0.25), Inches(0.55), Inches(0.55),
            num, 18, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_box(s, x + Inches(0.8), y + Inches(0.2), bw - Inches(1), Inches(0.9),
            lb, 13, color=TEXT_D, vAlign=MSO_ANCHOR.MIDDLE)

# Klasifikasi jalan box
box = add_rrect(s, M, Inches(4.5), CW, Inches(1.7), fill=ICE)
add_box(s, M + Inches(0.2), Inches(4.55), Inches(5), Inches(0.3),
        "KLASIFIKASI KELAS JALAN MENENTUKAN NILAI NSR", 11, bold=True, color=NAVY)
cls_text = ("🏛️  Kelas Jalan Khusus — Tol | Premium 1 | Premium 2\n"
            "🚗  Kelas Jalan I (Kendali Ketat) — Lebar > 3 m, pusat pelayanan\n"
            "🏡  Kelas Jalan II (Kendali Sedang) — Lebar ≤ 3 m, jalan lingkungan")
add_box(s, M + Inches(0.2), Inches(4.9), CW - Inches(0.4), Inches(1.0), cls_text, 10, color=TEXT_D)
add_footer(s)

# ═══════════════════════════════════════════
# BAB VI — PERHITUNGAN & TARIF (14-20)
# ═══════════════════════════════════════════
section_slide("BAB VI\nPERHITUNGAN & TARIF PAJAK", "Pasal 10",
              "Rumus: Pajak = Tarif × NSR — Dengan Berbagai Ketentuan Khusus")

# Callout rumus
callout_slide("Rumus: Pajak Reklame = Tarif (50%) × NSR", [
    ("×","Pajak =\nTarif × NSR",BLUE),
    ("50%","Indoor =\n50% NSR",TEAL),
    ("+20%","Tinggi > 15m\ntambahan 20%",WARM),
    ("+50%","Tembakau &\nMiras +50%",RED_C),
], subtitle="Pasal 10", ref="Rumus Dasar",
   note_text="Rumus: Pajak Reklame = Tarif Pajak × NSR\n"
             "NSR = Nilai Kelas Jalan × Ukuran (m²) × Jumlah × Jangka Waktu\n\n"
             "Ketentuan Khusus:\n"
             "• Reklame indoor: NSR 50% dari NSR normal\n"
             "• Ketinggian > 15 meter: tambahan 20%\n"
             "• Produk tembakau & minuman beralkohol: tambahan 50%\n"
             "• Perubahan naskah/revisi isi reklame: dikecualikan")

# 3 Tabel NSR
table_slide("Tabel NSR — Papan / Billboard (Rp/m²/hari)", [
    "Kelas Jalan", "Zona", "NSR (Rp/m²/hari)"
], [
    ["Kelas Jalan Khusus", "Jalan Tol", "23.575"],
    ["Kelas Jalan Khusus", "Premium 1", "16.100"],
    ["Kelas Jalan Khusus", "Premium 2", "14.950"],
    ["Kelas Jalan I", "Kendali Ketat", "13.225"],
    ["Kelas Jalan II", "Kendali Sedang", "11.500"],
], subtitle="Pasal 10", ref="Papan / Billboard")

table_slide("Tabel NSR — Megatron / Videotron", [
    "Kelas Jalan", "Zona", "NSR (/30 dtk)", "NSR (/m²/thn)"
], [
    ["Kelas Jalan Khusus", "Jalan Tol", "Rp 17,25", "13.599.900"],
    ["Kelas Jalan Khusus", "Premium 1", "Rp 13,80", "10.879.920"],
    ["Kelas Jalan Khusus", "Premium 2", "Rp 9,20", "7.253.280"],
    ["Kelas Jalan I", "Kendali Ketat", "Rp 8,05", "6.346.620"],
    ["Kelas Jalan II", "Kendali Sedang", "Rp 5,75", "4.533.300"],
], subtitle="Pasal 10", ref="Megatron / Videotron")

table_slide("Tabel NSR — Kain (Spanduk/Umbul/Baliho) (Rp/m²/hari)", [
    "Kelas Jalan", "Zona", "NSR (Rp/m²/hari)"
], [
    ["Kelas Jalan Khusus", "Jalan Tol", "30.000"],
    ["Kelas Jalan Khusus", "Premium 1", "30.000"],
    ["Kelas Jalan Khusus", "Premium 2", "25.000"],
    ["Kelas Jalan I", "Kendali Ketat", "20.000"],
    ["Kelas Jalan II", "Kendali Sedang", "19.000"],
], subtitle="Pasal 10", ref="Spanduk / Umbul / Baliho")

# NSR Lainnya
card_grid("NSR untuk 8 Jenis Reklame Lainnya — Bagian 1", [
    {"ic":"🏷️","t":"Stiker","clr":BLUE,"items":["Rp 7,5/cm²","Min. Rp 750.000/kali"]},
    {"ic":"🧱","t":"Melekat","clr":TEAL,"items":["Rp 750.000/m²/tahun"]},
    {"ic":"📄","t":"Selebaran","clr":WARM,"items":["Rp 600/lembar","Min. Rp 6.000.000/kali"]},
    {"ic":"🚌","t":"Berjalan","clr":RED_C,"items":["Rp 6.000/m²/hari","Termasuk kendaraan"]},
], subtitle="Pasal 10", ref="Jenis Reklame Lainnya")

card_grid("NSR untuk 8 Jenis Reklame Lainnya — Bagian 2", [
    {"ic":"🎈","t":"Udara","clr":BLUE,"items":["Rp 2.400.000/sekali","Maks. 1 bulan"]},
    {"ic":"🌊","t":"Apung","clr":TEAL,"items":["Rp 600.000/sekali","Maks. 1 bulan"]},
    {"ic":"🎬","t":"Film / Slide","clr":WARM,"items":["Rp 12.000/15 detik"]},
    {"ic":"🎭","t":"Peragaan","clr":RED_C,"items":["Rp 480.000/penyelenggaraan"]},
], subtitle="Pasal 10", ref="Jenis Reklame Lainnya (lanjutan)")

# ═══════════════════════════════════════════
# BAB VII — PENETAPAN, TAGIHAN & PEMBAYARAN (21-22)
# ═══════════════════════════════════════════
section_slide("BAB VII\nPENETAPAN, TAGIHAN & PEMBAYARAN", "Pasal 11–14",
              "4 Langkah: SKPD → Bayar → Telat → STPD")

flow_slide("Alur Penetapan hingga Pembayaran: 4 Langkah Wajib Dipahami WP", [
    ("1","SKPD","Diterbitkan Bapenda\nMasa berlaku 5 tahun",BLUE),
    ("2","Pembayaran","Lunas 1 bulan\nsejak SKPD diterima",TEAL),
    ("3","Keterlambatan","Bunga 1%/bln\nDiterbitkan STPD",WARM),
    ("4","STPD","Harus lunas\n≤ 30 hari",RED_C),
], subtitle="Pasal 11–14", ref="Alur Penetapan → Pembayaran",
   note_text="• Jatuh tempo: 1 bulan sejak tanggal pengiriman SKPD\n"
             "• Pembayaran: Kas Daerah / Bank Persepsi / tempat lain yang ditunjuk\n"
             "• Stiker sebagai tanda bukti pembayaran reklame\n"
             "• STPD dikenakan bunga 1%/bulan (maks. 24 bulan)")

# ═══════════════════════════════════════════
# BAB VIII — PEMBETULAN, KEBERATAN & BANDING (23-25)
# ═══════════════════════════════════════════
section_slide("BAB VIII\nPEMBETULAN, KEBERATAN & BANDING", "Pasal 15–20, 29–33",
              "Hak WP: Koreksi, Keberatan, Banding — Dengan Batas Waktu Jelas")

two_col_cards("Pembetulan: Koreksi Kesalahan Tulis, Hitung, dan Penerapan Aturan", [
    "$PEMBETULAN (Pasal 15–20)",
    "",
    "Kesalahan tulis: nama, alamat, NPWPD",
    "Kesalahan hitung: jumlah, tarif",
    "Kekeliruan penerapan aturan",
    "1 permohonan = 1 ketetapan",
    "Keputusan maksimal 6 bulan",
    "> 6 bulan tanpa putusan → dikabulkan",
    "Dapat dilakukan berulang (Ps 20)",
    "Jenis keputusan: kabul / batal / tolak",
], [
    "$JANGKA WAKTU & PROSEDUR",
    "",
    "Permohonan diajukan ke Bapenda",
    "Keputusan: kabul (tambah/kurang/hapus)",
    "Keputusan: batal | tolak",
    "Pasal 19: pembetulan jabatan",
    "Pasal 20: berulang jika masih salah",
], subtitle="Pasal 15–20", ref="Pembetulan", left_color=BLUE, right_color=WARM)

two_col_cards("Keberatan & Banding: Upaya Hukum WP dalam Sengketa Pajak", [
    "$KEBERATAN (Pasal 29–31)",
    "",
    "Objek: SKPD, SKPDKB, SKPDKBT, dll",
    "Diajukan maks. 3 bulan sejak SKPD",
    "Sudah bayar min. yang disetujui",
    "Keputusan maks. 12 bulan",
    "Jika ditolak: denda 30%",
    "Jika dikabulkan: + bunga 0,6%/bulan",
], [
    "$BANDING (Pasal 32–33)",
    "",
    "Objek: Surat Keputusan Keberatan",
    "Ke badan peradilan pajak",
    "Maks. 3 bulan sejak keputusan",
    "Menangguhkan kewajiban bayar",
    "Jika ditolak: denda 60%",
    "Jika dikabulkan: + bunga 0,6%/bulan",
], subtitle="Pasal 29–33", ref="Keberatan & Banding")

# ═══════════════════════════════════════════
# BAB IX — PEMERIKSAAN, PENAGIHAN & PENGHAPUSAN (26-28)
# ═══════════════════════════════════════════
section_slide("BAB IX\nPEMERIKSAAN, PENAGIHAN & PENGHAPUSAN", "Pasal 22–26",
              "Bapenda Berwenang Periksa — Piutang Dihapus Lewat 4 Tahapan")

card_grid("Pemeriksaan & Penagihan: 3 Pilar Penegakan Kepatuhan WP", [
    {"ic":"🔍","t":"Pemeriksaan (Ps 22–23)","clr":BLUE,"items":["Kepala Bapenda berwenang periksa","Menguji kepatuhan WP","WP wajib: buka buku/dokumen","Beri akses tempat & keterangan","Jika tidak → pajak ditetapkan jabatan"]},
    {"ic":"📬","t":"Penagihan (Ps 24)","clr":TEAL,"items":["Dasar: SKPD, SKPDKB, SKPDKBT","STPD, SK Pembetulan/Keberatan","Putusan Banding"]},
    {"ic":"⏳","t":"Kedaluwarsa (Ps 25)","clr":WARM,"items":["5 tahun sejak pajak terutang","Tertangguh jika ada:","Surat Teguran / Paksa","Pengakuan utang dari WP"]},
], subtitle="Pasal 22–25", ref="Pemeriksaan & Penagihan")

flow_slide("Penghapusan Piutang Pajak: 4 Langkah dari Penelitian hingga SK", [
    ("1","Penelitian","Dilakukan Bapenda",BLUE),
    ("2","Penetapan","Keputusan Wali Kota",TEAL),
    ("3","Koordinasi","Dengan Inspektorat",WARM),
    ("4","SK Penghapusan","Diterbitkan",NAVY_M),
], subtitle="Pasal 26", ref="Penghapusan Piutang",
   note_text="Syarat penghapusan piutang pajak:\n"
             "• Piutang tidak mungkin ditagih lagi karena kedaluwarsa\n"
             "• Ada koordinasi dengan Inspektorat Daerah\n"
             "• Dibuktikan dengan dokumen pelaksanaan penagihan")

# ═══════════════════════════════════════════
# BAB X — KERINGANAN, KEMUDAHAN & PENGHARGAAN (29-30)
# ═══════════════════════════════════════════
section_slide("BAB X\nKERINGANAN, KEMUDAHAN & PENGHARGAAN", "Pasal 27–28, 34–35",
              "3 Fasilitas WP: Keringanan, Kemudahan Angsuran, dan Penghargaan")

card_grid("3 Fasilitas untuk WP: Keringanan, Kemudahan Angsuran, Penghargaan", [
    {"ic":"🎯","t":"Keringanan (Ps 27)","clr":BLUE,"items":["Keringanan / Pengurangan","Pembebasan / Penundaan","Atas pokok & sanksi pajak","WP dengan likuiditas rendah","Objek terdampak bencana/kebakaran"]},
    {"ic":"🤝","t":"Kemudahan (Ps 28)","clr":TEAL,"items":["Perpanjangan waktu bayar","Angsuran maks. 24 bulan","Bunga 0,6%/bulan","Keadaan kahar: bencana, wabah, kerusuhan"]},
    {"ic":"🏆","t":"Penghargaan (Ps 34–35)","clr":WARM,"items":["WP Taat Pajak","Bayar tepat waktu ≥ 1 tahun","Tanpa tunggakan 3 tahun","Kontribusi signifikan","Piagam / Hadiah (APBD)"]},
], subtitle="Pasal 27–28, 34–35", ref="Keringanan, Kemudahan & Penghargaan")

# ═══════════════════════════════════════════
# BAB XI — KETENTUAN PENUTUP (31-32)
# ═══════════════════════════════════════════
section_slide("BAB XI\nKETENTUAN PENUTUP", "Pasal 36–37",
              "Perwal No 48/2012 Dicabut — Berlaku Sejak 20 Desember 2024")

two_col_cards("Perwal Lama Dicabut, Peraturan Baru Berlaku Mulai Diundangkan", [
    "$PERATURAN YANG DICABUT (Pasal 36)",
    "",
    "Perwal No. 48 Tahun 2012",
    "Petunjuk Pelaksanaan Perda 14/2012",
    "Perwal No. 52 Tahun 2013 (Perubahan)",
], [
    "$MULAI BERLAKU (Pasal 37)",
    "",
    "Sejak diundangkan",
    "20 Desember 2024",
    "",
    "Pj. WALI KOTA BEKASI,",
    "ttd.",
    "R. GANI MUHAMAD",
], subtitle="Pasal 36–37", ref="Penutup", left_color=BLUE, right_color=NAVY_M)

# ────────── CLOSING (33) ──────────
s = blank()
bg = s.background; bg.fill.solid(); bg.fill.fore_color.rgb = NAVY
add_oval(s, Inches(-1), Inches(-1.5), Inches(4.5), RGBColor(0x0D, 0x1F, 0x3C))
add_oval(s, Inches(8.5), Inches(-2), Inches(7), RGBColor(0x0D, 0x1F, 0x3C))
add_oval(s, Inches(10), Inches(4.5), Inches(5), RGBColor(0x0D, 0x1F, 0x3C))
add_shape(s, MSO_SHAPE.RECTANGLE, M, Inches(3.6), Inches(3.5), Inches(0.04), fill=GOLD)
add_box(s, M, Inches(1.6), CW, Inches(0.4), "BERITA DAERAH KOTA BEKASI", 14, bold=True, color=GOLD)
add_box(s, M, Inches(2.4), CW, Inches(1.5), "TERIMA KASIH", 48, bold=True, color=WHITE)
add_box(s, M, Inches(4.1), CW, Inches(0.8),
        "Peraturan Wali Kota Bekasi Nomor 51 Tahun 2024\nTentang Pengelolaan Pajak Reklame", 14, color=ICE)
add_box(s, M, Inches(5.1), CW, Inches(0.35),
        "Sumber: https://jdih.bekasikota.go.id", 10, color=TEXT_L)
gold_bottom_bar(s)

# ══════════════════════════════════════════════
# SAVE
# ══════════════════════════════════════════════
out = os.path.join(os.path.dirname(__file__), "Perwal_Bekasi_51_2024_Pajak_Reklame.pptx")
prs.save(out)
print(f"✅ OK: {out} ({len(prs.slides)} slide)")
