#!/usr/bin/env python3
"""
Generate a professional PowerPoint presentation about SPBE (Sistem Pemerintahan Berbasis Elektronik)
Based on Perpres No. 95/2018 and Perpres No. 132/2022
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import copy

# ── Colour Palette ──────────────────────────────────────────────
DARK_BLUE   = RGBColor(0x00, 0x2B, 0x5C)   # header / footer bands
MED_BLUE    = RGBColor(0x00, 0x5B, 0x96)   # accent shapes
LIGHT_BLUE  = RGBColor(0xD6, 0xE8, 0xF7)   # light backgrounds
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
BLACK       = RGBColor(0x33, 0x33, 0x33)
GRAY        = RGBColor(0x66, 0x66, 0x66)
ACCENT_GOLD = RGBColor(0xF0, 0xA5, 0x00)   # highlights
TABLE_HEAD  = RGBColor(0x00, 0x4B, 0x87)
TABLE_ALT   = RGBColor(0xE8, 0xF0, 0xF8)

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
W = prs.slide_width
H = prs.slide_height

# ── Helper functions ────────────────────────────────────────────

def _rect(slide, left, top, width, height, fill_color, line_color=None):
    """Add a filled rectangle."""
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill_color
    if line_color:
        shp.line.color.rgb = line_color
    else:
        shp.line.fill.background()
    return shp

def _footer(slide, text="Sistem Pemerintahan Berbasis Elektronik (SPBE)"):
    """Standard footer band."""
    bar = _rect(slide, 0, H - Inches(0.5), W, Inches(0.5), DARK_BLUE)
    tf = bar.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(10)
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    tf.paragraphs[0].space_before = Pt(4)

def _section_header_bar(slide, title_text):
    """Dark blue bar at top with title."""
    bar = _rect(slide, 0, 0, W, Inches(1.0), DARK_BLUE)
    # gold accent line
    _rect(slide, 0, Inches(0.95), W, Inches(0.05), ACCENT_GOLD)
    txBox = slide.shapes.add_textbox(Inches(0.6), Inches(0.12), Inches(12), Inches(0.8))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = title_text
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = WHITE

def _subtitle_box(slide, left, top, width, height, text, color=GRAY):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(14)
    p.font.color.rgb = color
    p.font.italic = True

def _bullet_text(slide, left, top, width, height, items, font_size=14, color=BLACK, bold_first=False):
    """Add a textbox with bullet points."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.space_after = Pt(6)
        p.level = 0
        # bullet
        pPr = p._pPr
        if pPr is None:
            from lxml import etree
            nsmap = {'a': 'http://schemas.openxmlformats.org/drawingml/2006/main'}
            pPr = etree.SubElement(p._p, '{http://schemas.openxmlformats.org/drawingml/2006/main}pPr')
        from lxml import etree
        nsmap = {'a': 'http://schemas.openxmlformats.org/drawingml/2006/main'}
        buChar = etree.SubElement(pPr, '{http://schemas.openxmlformats.org/drawingml/2006/main}buChar')
        buChar.set('char', '•')
        if bold_first and i == 0:
            p.font.bold = True
    return txBox

def _icon_shape(slide, left, top, size, color, label, label_below=True):
    """Circle with number/text inside."""
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top, size, size)
    circle.fill.solid()
    circle.fill.fore_color.rgb = color
    circle.line.fill.background()
    tf = circle.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = label
    p.font.size = Pt(int(size / Inches(1) * 18)) if len(label) <= 3 else Pt(11)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    tf.paragraphs[0].space_before = Pt(0)
    return circle

def _add_table(slide, left, top, width, height, headers, rows, col_widths=None):
    """Add a styled table."""
    num_rows = len(rows) + 1
    num_cols = len(headers)
    table_shape = slide.shapes.add_table(num_rows, num_cols, left, top, width, height)
    table = table_shape.table
    
    if col_widths:
        for i, cw in enumerate(col_widths):
            table.columns[i].width = cw
    
    # Header row
    for ci, h in enumerate(headers):
        cell = table.cell(0, ci)
        cell.text = h
        cell.fill.solid()
        cell.fill.fore_color.rgb = TABLE_HEAD
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(12)
            p.font.bold = True
            p.font.color.rgb = WHITE
            p.alignment = PP_ALIGN.CENTER
    
    # Data rows
    for ri, row in enumerate(rows):
        for ci, val in enumerate(row):
            cell = table.cell(ri + 1, ci)
            cell.text = val
            if ri % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = TABLE_ALT
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = WHITE
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(11)
                p.font.color.rgb = BLACK
            cell.text_frame.paragraphs[0].space_before = Pt(3)
    
    return table_shape

def _card(slide, left, top, width, height, title, body_items, icon_color=MED_BLUE):
    """Card-like box with title and bullet body."""
    # Card background
    card = _rect(slide, left, top, width, height, WHITE, RGBColor(0xCC, 0xCC, 0xCC))
    card.shadow.inherit = False
    
    # Title bar inside card
    title_bar = _rect(slide, left, top, width, Inches(0.45), icon_color)
    txBox = slide.shapes.add_textbox(left + Inches(0.15), top + Inches(0.05), width - Inches(0.3), Inches(0.4))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = WHITE
    
    # Body
    _bullet_text(slide, left + Inches(0.15), top + Inches(0.5), width - Inches(0.3),
                 height - Inches(0.5), body_items, font_size=11, color=BLACK)


# ════════════════════════════════════════════════════════════════
# SLIDE 1 — TITLE
# ════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank

# Full blue background
_rect(slide, 0, 0, W, H, DARK_BLUE)
# Accent stripe
_rect(slide, 0, H * 0.45, W, Inches(0.06), ACCENT_GOLD)
_rect(slide, 0, H * 0.47, W, Inches(0.02), MED_BLUE)

# Title
txBox = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(11), Inches(1.5))
tf = txBox.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "PENERAPAN SISTEM PEMERINTAHAN\nBERBASIS ELEKTRONIK (SPBE)"
p.font.size = Pt(36)
p.font.bold = True
p.font.color.rgb = WHITE
p.alignment = PP_ALIGN.CENTER
p.line_spacing = Pt(44)

# Subtitle
txBox2 = slide.shapes.add_textbox(Inches(2), Inches(3.3), Inches(9), Inches(0.8))
tf2 = txBox2.text_frame
tf2.word_wrap = True
p2 = tf2.paragraphs[0]
p2.text = "Transformasi Birokrasi Menuju Pemerintahan yang Terintegrasi, Efisien, Aman, dan Tepercaya"
p2.font.size = Pt(16)
p2.font.color.rgb = LIGHT_BLUE
p2.alignment = PP_ALIGN.CENTER

# Regulation info
txBox3 = slide.shapes.add_textbox(Inches(2), Inches(4.3), Inches(9), Inches(1.0))
tf3 = txBox3.text_frame
tf3.word_wrap = True
p3 = tf3.paragraphs[0]
p3.text = "Perpres No. 95 Tahun 2018  •  Perpres No. 132 Tahun 2022"
p3.font.size = Pt(14)
p3.font.color.rgb = ACCENT_GOLD
p3.alignment = PP_ALIGN.CENTER
p3b = tf3.add_paragraph()
p3b.text = "Kementerian Pendayagunaan Aparatur Negara dan Reformasi Birokrasi"
p3b.font.size = Pt(12)
p3b.font.color.rgb = LIGHT_BLUE
p3b.alignment = PP_ALIGN.CENTER

_footer(slide)


# ════════════════════════════════════════════════════════════════
# SLIDE 2 — PENGANTAR / APA ITU SPBE?
# ════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])

_section_header_bar(slide, "Apa Itu SPBE?")

# Left column - Definition
txBox = slide.shapes.add_textbox(Inches(0.6), Inches(1.3), Inches(6), Inches(1.2))
tf = txBox.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "Sistem Pemerintahan Berbasis Elektronik (SPBE)"
p.font.size = Pt(18)
p.font.bold = True
p.font.color.rgb = DARK_BLUE
p2 = tf.add_paragraph()
p2.text = "Penyelenggaraan pemerintahan yang memanfaatkan teknologi informasi dan komunikasi untuk memberikan layanan kepada pengguna SPBE."
p2.font.size = Pt(13)
p2.font.color.rgb = GRAY
p2.space_before = Pt(8)

# Key points - 3 cards
y_card = Inches(2.7)
card_w = Inches(3.7)
card_h = Inches(2.0)
gap = Inches(0.3)
x_start = Inches(0.6)

cards_data = [
    ("Bukan Sekadar Digitalisasi", [
        "Bukan memindahkan kertas ke aplikasi",
        "Transformasi menyeluruh birokrasi",
        "Menciptakan tata kelola terintegrasi"
    ]),
    ("Dasar Hukum", [
        "Perpres No. 95 Tahun 2018",
        "Perpres No. 132 Tahun 2022 (penyempurnaan)",
        "Pedoman dari Kemenpan-RB"
    ]),
    ("Tujuan Akhir", [
        "Birokrasi terintegrasi & efisien",
        "Layanan prima masyarakat",
        "Pemerintahan aman & tepercaya"
    ])
]

for i, (title, items) in enumerate(cards_data):
    x = x_start + i * (card_w + gap)
    _card(slide, x, y_card, card_w, card_h, title, items, MED_BLUE if i != 1 else ACCENT_GOLD)

# Bottom insight
txBox4 = slide.shapes.add_textbox(Inches(0.6), Inches(5.0), Inches(12), Inches(0.8))
tf4 = txBox4.text_frame
tf4.word_wrap = True
p4 = tf4.paragraphs[0]
p4.text = "💡 "  # We'll just use text
p4.font.size = Pt(12)
p4.font.color.rgb = GRAY
p4b = tf4.add_paragraph()
p4b.text = "SPBE adalah transformasi menyeluruh — bukan proyek TIK biasa. Digitalisasi tanpa perubahan tata kelola hanya akan menciptakan ego sektoral digital."
p4b.font.size = Pt(12)
p4b.font.bold = True
p4b.font.color.rgb = DARK_BLUE

_footer(slide)


# ════════════════════════════════════════════════════════════════
# SLIDE 3 — KERANGKA ARSITEKTUR SPBE
# ════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])

_section_header_bar(slide, "Kerangka Arsitektur SPBE")

# Architecture layers - concentric / layered boxes
layers = [
    ("LAYANAN SPBE", Inches(4.0), Inches(1.5), Inches(5.3), Inches(0.7), DARK_BLUE, WHITE),
    ("Aplikasi & Layanan Digital", Inches(3.2), Inches(2.3), Inches(6.9), Inches(0.7), MED_BLUE, WHITE),
    ("Data & Informasi", Inches(2.4), Inches(3.1), Inches(8.5), Inches(0.7), RGBColor(0x00, 0x70, 0xC0), WHITE),
    ("Infrastruktur & Keamanan", Inches(1.6), Inches(3.9), Inches(10.1), Inches(0.7), RGBColor(0x00, 0x88, 0xD0), WHITE),
]

for text, x, y, w, h, bg, fc in layers:
    shp = _rect(slide, x, y, w, h, bg)
    tx = slide.shapes.add_textbox(x, y, w, h)
    tf = tx.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = fc
    p.alignment = PP_ALIGN.CENTER
    p.space_before = Pt(int(h / Inches(1) * 6))

# Side pillars: Tata Kelola (left) and Manajemen (right)
pillar_w = Inches(0.8)
pillar_h = Inches(3.8)
pillar_y = Inches(1.5)

# Left pillar - Tata Kelola
shp_l = _rect(slide, Inches(0.5), pillar_y, pillar_w, pillar_h, ACCENT_GOLD)
tx_l = slide.shapes.add_textbox(Inches(0.5), pillar_y, pillar_w, pillar_h)
tf_l = tx_l.text_frame
tf_l.word_wrap = True
p_l = tf_l.paragraphs[0]
p_l.text = "TATA\nKELOLA"
p_l.font.size = Pt(14)
p_l.font.bold = True
p_l.font.color.rgb = WHITE
p_l.alignment = PP_ALIGN.CENTER

# Right pillar - Manajemen
shp_r = _rect(slide, Inches(12.0), pillar_y, pillar_w, pillar_h, ACCENT_GOLD)
tx_r = slide.shapes.add_textbox(Inches(12.0), pillar_y, pillar_w, pillar_h)
tf_r = tx_r.text_frame
tf_r.word_wrap = True
p_r = tf_r.paragraphs[0]
p_r.text = "MANA-\nJEMEN"
p_r.font.size = Pt(14)
p_r.font.bold = True
p_r.font.color.rgb = WHITE
p_r.alignment = PP_ALIGN.CENTER

# Bottom description
bottom_items = [
    "Tata Kelola dan Manajemen menjadi pondasi luar yang mengikat seluruh komponen arsitektur SPBE.",
    "Proses bisnis, data, aplikasi, infrastruktur, dan keamanan saling menopang untuk satu tujuan: Layanan SPBE yang prima.",
    "Layanan Publik → untuk masyarakat  |  Layanan Administrasi → untuk internal pemerintah"
]
_bullet_text(slide, Inches(0.5), Inches(5.0), Inches(12.3), Inches(1.5), bottom_items, font_size=12, color=GRAY)

_footer(slide)


# ════════════════════════════════════════════════════════════════
# SLIDE 4 — 4 DOMAIN EVALUASI SPBE (TABLE OVERVIEW)
# ════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])

_section_header_bar(slide, "4 Domain Utama Evaluasi SPBE — Indikator Tingkat Kematangan")

# Table
headers = ["Domain", "Fokus Utama", "Output yang Diharapkan"]
rows = [
    ["Kebijakan Internal", "Regulasi, pedoman, SK formal di internal instansi", "Perbup/Perwali/Perka terkait Tim Koordinasi & Tata Kelola TIK"],
    ["Tata Kelola", "Perencanaan arsitektur, peta rencana, integrasi sistem", "Dokumen Arsitektur SPBE & Peta Rencana 5 tahun"],
    ["Manajemen", "Pengelolaan risiko, data, keamanan, SDM, aset TIK", "Penerapan Satu Data Indonesia & Manajemen Risiko SPBE"],
    ["Layanan", "Kualitas fungsi aplikasi pelayanan publik & administrasi", "Aplikasi terintegrasi (bukan platform terpisah-pisah)"],
]

tbl = _add_table(slide, Inches(0.6), Inches(1.4), Inches(12.1), Inches(3.5), headers, rows,
                 col_widths=[Inches(2.2), Inches(4.8), Inches(5.1)])

# Note
txBox = slide.shapes.add_textbox(Inches(0.6), Inches(5.2), Inches(12), Inches(0.8))
tf = txBox.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "⚠️  Sebelum menyusun strategi, instansi Anda akan dinilai berdasarkan 4 domain evaluasi ini sebagai indikator tingkat kematangan (Maturity Index) SPBE."
p.font.size = Pt(12)
p.font.bold = True
p.font.color.rgb = DARK_BLUE

_footer(slide)


# ════════════════════════════════════════════════════════════════
# SLIDES 5-8 — DETAIL SETIAP DOMAIN (cards)
# ════════════════════════════════════════════════════════════════

domain_details = [
    ("Kebijakan Internal", "Regulasi, Pedoman, dan SK Formal", DARK_BLUE, [
        "Menyusun dan menetapkan Peraturan Kepala Daerah / Menteri tentang Tim Koordinasi SPBE",
        "Menyusun pedoman tata kelola TIK internal instansi",
        "Memastikan setiap kebijakan TIK selaras dengan arah strategis SPBE nasional",
        "Menerbitkan SK formal untuk pembentukan struktur organisasi SPBE",
        "Melakukan harmonisasi regulasi antar sektor terkait",
    ]),
    ("Tata Kelola", "Perencanaan Arsitektur & Peta Rencana", MED_BLUE, [
        "Menyusun Dokumen Arsitektur SPBE Instansi (proses bisnis, data, aplikasi, infrastruktur)",
        "Membuat Peta Rencana (Roadmap) 5 tahun implementasi SPBE",
        "Memastikan arsitektur instansi selaras dengan Arsitektur SPBE Nasional",
        "Merencanakan integrasi sistem lintas unit kerja",
        "Menggunakan aplikasi Arsitektur SPBE Nasional sebagai acuan",
    ]),
    ("Manajemen", "Pengelolaan Risiko, Data, SDM & Aset TIK", RGBColor(0x00, 0x70, 0xC0), [
        "Menerapkan prinsip Satu Data Indonesia (SDI) dalam pengelolaan data",
        "Menyusun dan menerapkan Manajemen Risiko SPBE",
        "Mengelola SDM TIK secara profesional dan berkelanjutan",
        "Melakukan inventarisasi dan optimalisasi aset TIK",
        "Menyusun kebijakan pengamanan data dan informasi",
    ]),
    ("Layanan", "Kualitas Fungsi Aplikasi & Integrasi", ACCENT_GOLD, [
        "Mengintegrasikan aplikasi pelayanan publik dalam satu platform terpadu",
        "Menyediakan layanan administrasi internal yang efisien",
        "Memastikan aplikasi dapat saling bertukar data (interoperabilitas)",
        "Menerapkan standar kualitas layanan digital",
        "Mengukur kepuasan pengguna (masyarakat & internal) secara berkala",
    ]),
]

for idx, (title, subtitle, color, items) in enumerate(domain_details):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    _section_header_bar(slide, f"Domain {idx + 1}: {title}")
    
    # Subtitle
    _subtitle_box(slide, Inches(0.6), Inches(1.2), Inches(8), Inches(0.4), subtitle)
    
    # Cards
    for i, item in enumerate(items):
        row = i // 2
        col = i % 2
        x = Inches(0.6) + col * Inches(6.2)
        y = Inches(1.8) + row * Inches(1.3)
        
        # Numbered circle
        num_size = Inches(0.4)
        circle = _icon_shape(slide, x, y + Inches(0.05), num_size, color, str(i + 1))
        
        # Text
        txBox = slide.shapes.add_textbox(x + Inches(0.55), y, Inches(5.3), Inches(0.8))
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = item
        p.font.size = Pt(13)
        p.font.color.rgb = BLACK
        p.space_after = Pt(4)
    
    # Output hint
    txBox_out = slide.shapes.add_textbox(Inches(0.6), Inches(5.5), Inches(12), Inches(0.5))
    tf_out = txBox_out.text_frame
    tf_out.word_wrap = True
    p_out = tf_out.paragraphs[0]
    outputs = ["Output: Perbup/Perwali/Perka", "Output: Dokumen Arsitektur & Peta Rencana",
               "Output: Satu Data Indonesia & Manajemen Risiko", "Output: Aplikasi Terintegrasi"]
    p_out.text = "📋 " + outputs[idx]
    p_out.font.size = Pt(12)
    p_out.font.bold = True
    p_out.font.color.rgb = color
    
    _footer(slide)


# ════════════════════════════════════════════════════════════════
# SLIDE 9 — LANGKAH 1: TIM KOORDINASI
# ════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])

_section_header_bar(slide, "Langkah 1: Pembentukan Tim Koordinasi SPBE")

# Step number
_icon_shape(slide, Inches(0.6), Inches(1.3), Inches(0.8), ACCENT_GOLD, "1")

txBox = slide.shapes.add_textbox(Inches(1.6), Inches(1.3), Inches(10), Inches(0.5))
tf = txBox.text_frame
p = tf.paragraphs[0]
p.text = "Langkah Awal / Fondasi Implementasi SPBE"
p.font.size = Pt(18)
p.font.bold = True
p.font.color.rgb = DARK_BLUE

# Key points
points = [
    "Tim Koordinasi SPBE Instansi disahkan melalui SK Kepala Daerah / Menteri",
    "Dipimpin oleh Sekretaris Daerah atau Sekretaris Jenderal (Sekjen)",
    "Bersifat lintas sektor — melibatkan berbagai unsur strategis:",
]
sub_points = [
    "  • Perencanaan (Bappeda/Bappenas)",
    "  • Organisasi & Tata Laksana (Ortala)",
    "  • Keuangan (BPKAD / Biro Keuangan)",
    "  • Teknis TIK (Diskominfo / Pusdatin)",
]

_bullet_text(slide, Inches(0.6), Inches(2.3), Inches(12), Inches(2.5), points, font_size=14, color=BLACK)
_bullet_text(slide, Inches(1.0), Inches(4.0), Inches(11), Inches(1.5), sub_points, font_size=13, color=MED_BLUE)

# Emphasis box
em = _rect(slide, Inches(0.6), Inches(5.5), Inches(12), Inches(0.7), LIGHT_BLUE, MED_BLUE)
txEm = slide.shapes.add_textbox(Inches(0.8), Inches(5.55), Inches(11.5), Inches(0.6))
tfEm = txEm.text_frame
tfEm.word_wrap = True
pEm = tfEm.paragraphs[0]
pEm.text = "💡  Prinsip: Tim harus bersifat lintas sektor — bukan hanya domain Diskominfo / TIK semata. SPBE bukan proyek TIK, melainkan transformasi birokrasi menyeluruh."
pEm.font.size = Pt(12)
pEm.font.bold = True
pEm.font.color.rgb = DARK_BLUE

_footer(slide)


# ════════════════════════════════════════════════════════════════
# SLIDE 10 — LANGKAH 2: ARSITEKTUR & PETA RENCANA
# ════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])

_section_header_bar(slide, "Langkah 2: Penyusunan Arsitektur & Peta Rencana (Roadmap)")

_icon_shape(slide, Inches(0.6), Inches(1.3), Inches(0.8), ACCENT_GOLD, "2")

txBox = slide.shapes.add_textbox(Inches(1.6), Inches(1.3), Inches(10), Inches(0.5))
tf = txBox.text_frame
p = tf.paragraphs[0]
p.text = "Proses Perencanaan Strategis SPBE"
p.font.size = Pt(18)
p.font.bold = True
p.font.color.rgb = DARK_BLUE

points = [
    "Petakan Proses Bisnis Instansi terlebih dahulu — sebelum menyentuh teknologi",
    "Gunakan aplikasi Arsitektur SPBE Nasional sebagai acuan penyusunan",
    "Susun Dokumen Arsitektur SPBE Instansi:",
    "  • Proses Bisnis  • Data & Informasi  • Aplikasi  • Infrastruktur  • Keamanan",
    "Susun Peta Rencana (Roadmap) implementasi — biasanya berdurasi 5 tahun",
    "Pastikan keselarasan dengan arah strategis SPBE Nasional"
]
_bullet_text(slide, Inches(0.6), Inches(2.3), Inches(12), Inches(3.0), points, font_size=14, color=BLACK)

em = _rect(slide, Inches(0.6), Inches(5.5), Inches(12), Inches(0.7), LIGHT_BLUE, MED_BLUE)
txEm = slide.shapes.add_textbox(Inches(0.8), Inches(5.55), Inches(11.5), Inches(0.6))
tfEm = txEm.text_frame
tfEm.word_wrap = True
pEm = tfEm.paragraphs[0]
pEm.text = "📌  Prinsip: Jangan mulai dari teknologi — mulailah dari proses bisnis. Teknologi hanyalah alat untuk mewujudkan proses yang efisien."
pEm.font.size = Pt(12)
pEm.font.bold = True
pEm.font.color.rgb = DARK_BLUE

_footer(slide)


# ════════════════════════════════════════════════════════════════
# SLIDE 11 — LANGKAH 3: STANDARDISASI INFRASTRUKTUR
# ════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])

_section_header_bar(slide, "Langkah 3: Standardisasi Infrastruktur & Integrasi Data")

_icon_shape(slide, Inches(0.6), Inches(1.3), Inches(0.8), ACCENT_GOLD, "3")

txBox = slide.shapes.add_textbox(Inches(1.6), Inches(1.3), Inches(10), Inches(0.5))
tf = txBox.text_frame
p = tf.paragraphs[0]
p.text = "Fase Eksekusi Teknis"
p.font.size = Pt(18)
p.font.bold = True
p.font.color.rgb = DARK_BLUE

points = [
    "Batasi pembangunan pusat data (data center) mandiri — biaya besar & tidak efisien",
    "Mulai migrasi ke Pusat Data Nasional (PDN) atau cloud tersertifikasi",
    "Implementasikan sistem berbagi pakai melalui Hub SPLP:",
    "  • SPLP = Sistem Penghubung Layanan Pemerintah (interkoneksi antar sistem)",
    "Pastikan data selaras dengan prinsip Satu Data Indonesia (SDI):",
    "  • Standar data, metadata, interoperabilitas, dan referensi data"
]
_bullet_text(slide, Inches(0.6), Inches(2.3), Inches(12), Inches(3.0), points, font_size=14, color=BLACK)

em = _rect(slide, Inches(0.6), Inches(5.5), Inches(12), Inches(0.7), LIGHT_BLUE, MED_BLUE)
txEm = slide.shapes.add_textbox(Inches(0.8), Inches(5.55), Inches(11.5), Inches(0.6))
tfEm = txEm.text_frame
tfEm.word_wrap = True
pEm = tfEm.paragraphs[0]
pEm.text = "💰  Prinsip: Tidak perlu membangun data center sendiri — manfaatkan PDN dan komputasi cloud untuk efisiensi anggaran."
pEm.font.size = Pt(12)
pEm.font.bold = True
pEm.font.color.rgb = DARK_BLUE

_footer(slide)


# ════════════════════════════════════════════════════════════════
# SLIDE 12 — LANGKAH 4: KONSOLIDASI APLIKASI
# ════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])

_section_header_bar(slide, "Langkah 4: Konsolidasi Aplikasi & Layanan Digital")

_icon_shape(slide, Inches(0.6), Inches(1.3), Inches(0.8), ACCENT_GOLD, "4")

txBox = slide.shapes.add_textbox(Inches(1.6), Inches(1.3), Inches(10), Inches(0.5))
tf = txBox.text_frame
p = tf.paragraphs[0]
p.text = "Penyederhanaan Layanan Digital"
p.font.size = Pt(18)
p.font.bold = True
p.font.color.rgb = DARK_BLUE

points = [
    "Lakukan audit menyeluruh terhadap seluruh aplikasi yang ada di instansi",
    "Terapkan moratorium pembuatan aplikasi baru yang bersifat duplikatif",
    "Gabungkan aplikasi sektoral ke dalam satu platform layanan digital terpadu:",
    "  • Super Apps untuk layanan administrasi internal (kepegawaian, keuangan, dll)",
    "  • Super Apps untuk layanan publik (perizinan, pengaduan, informasi)",
    "Pastikan setiap aplikasi memiliki API yang terdokumentasi untuk interoperabilitas"
]
_bullet_text(slide, Inches(0.6), Inches(2.3), Inches(12), Inches(3.0), points, font_size=14, color=BLACK)

em = _rect(slide, Inches(0.6), Inches(5.5), Inches(12), Inches(0.7), LIGHT_BLUE, MED_BLUE)
txEm = slide.shapes.add_textbox(Inches(0.8), Inches(5.55), Inches(11.5), Inches(0.6))
tfEm = txEm.text_frame
tfEm.word_wrap = True
pEm = tfEm.paragraphs[0]
pEm.text = "⚠️  Prinsip: Stop membuat aplikasi baru untuk setiap masalah! Fokus pada integrasi dan interoperabilitas sistem yang sudah ada."
pEm.font.size = Pt(12)
pEm.font.bold = True
pEm.font.color.rgb = DARK_BLUE

_footer(slide)


# ════════════════════════════════════════════════════════════════
# SLIDE 13 — LANGKAH 5: KEAMANAN INFORMASI
# ════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])

_section_header_bar(slide, "Langkah 5: Penerapan Manajemen Keamanan Informasi")

_icon_shape(slide, Inches(0.6), Inches(1.3), Inches(0.8), ACCENT_GOLD, "5")

txBox = slide.shapes.add_textbox(Inches(1.6), Inches(1.3), Inches(10), Inches(0.5))
tf = txBox.text_frame
p = tf.paragraphs[0]
p.text = "Aspek Keamanan & Kepatuhan"
p.font.size = Pt(18)
p.font.bold = True
p.font.color.rgb = DARK_BLUE

points = [
    "Terapkan Sistem Manajemen Keamanan Informasi (SMKI) berbasis ISO 27001 secara bertahap",
    "Berkolaborasi dengan BSSN (Badan Siber dan Sandi Negara):",
    "  • Vulnerability Assessment (penilaian kerentanan) secara berkala",
    "  • Konsultasi keamanan siber dan kriptografi",
    "Bentuk CSIRT (Computer Security Incident Response Team) internal instansi",
    "Lakukan uji penetrasi (penetration testing) dan audit keamanan rutin"
]
_bullet_text(slide, Inches(0.6), Inches(2.3), Inches(12), Inches(3.0), points, font_size=14, color=BLACK)

em = _rect(slide, Inches(0.6), Inches(5.5), Inches(12), Inches(0.7), LIGHT_BLUE, MED_BLUE)
txEm = slide.shapes.add_textbox(Inches(0.8), Inches(5.55), Inches(11.5), Inches(0.6))
tfEm = txEm.text_frame
tfEm.word_wrap = True
pEm = tfEm.paragraphs[0]
pEm.text = "🔒  Prinsip: Keamanan bukan lapisan terakhir — keamanan harus menjadi pertimbangan sejak tahap perencanaan (security by design)."
pEm.font.size = Pt(12)
pEm.font.bold = True
pEm.font.color.rgb = DARK_BLUE

_footer(slide)


# ════════════════════════════════════════════════════════════════
# SLIDE 14 — LANGKAH 6: PEMANTAUAN & EVALUASI
# ════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])

_section_header_bar(slide, "Langkah 6: Pemantauan, Evaluasi, dan Audit")

_icon_shape(slide, Inches(0.6), Inches(1.3), Inches(0.8), ACCENT_GOLD, "6")

txBox = slide.shapes.add_textbox(Inches(1.6), Inches(1.3), Inches(10), Inches(0.5))
tf = txBox.text_frame
p = tf.paragraphs[0]
p.text = "Continuous Improvement"
p.font.size = Pt(18)
p.font.bold = True
p.font.color.rgb = DARK_BLUE

points = [
    "Lakukan evaluasi mandiri (self-assessment) secara berkala untuk memantau kemajuan",
    "Persiapkan diri sebelum penilaian formal dari Kemenpan-RB",
    "Audit TIK (Teknologi Informasi dan Komunikasi) wajib dijalankan:",
    "  • Menilai efisiensi infrastruktur yang berjalan",
    "  • Mengaudit aplikasi dan layanan digital",
    "  • Memeriksa kepatuhan keamanan informasi",
    "Gunakan hasil evaluasi sebagai dasar perbaikan berkelanjutan (continuous improvement)"
]
_bullet_text(slide, Inches(0.6), Inches(2.3), Inches(12), Inches(3.0), points, font_size=14, color=BLACK)

em = _rect(slide, Inches(0.6), Inches(5.5), Inches(12), Inches(0.7), LIGHT_BLUE, MED_BLUE)
txEm = slide.shapes.add_textbox(Inches(0.8), Inches(5.55), Inches(11.5), Inches(0.6))
tfEm = txEm.text_frame
tfEm.word_wrap = True
pEm = tfEm.paragraphs[0]
pEm.text = "📊  Prinsip: Evaluasi bukan formalitas — gunakan hasilnya untuk perbaikan nyata. Ukur, evaluasi, perbaiki, ulangi."
pEm.font.size = Pt(12)
pEm.font.bold = True
pEm.font.color.rgb = DARK_BLUE

_footer(slide)


# ════════════════════════════════════════════════════════════════
# SLIDE 15 — RINGKASAN 6 LANGKAH
# ════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])

_section_header_bar(slide, "Ringkasan 6 Langkah Penerapan SPBE")

steps = [
    ("1", "Tim Koordinasi", "SK Tim lintas sektor", DARK_BLUE),
    ("2", "Arsitektur & Roadmap", "Dokumen arsitektur 5 tahun", MED_BLUE),
    ("3", "Infrastruktur & Data", "PDN, SPLP, Satu Data", RGBColor(0x00, 0x70, 0xC0)),
    ("4", "Konsolidasi Aplikasi", "Super Apps terintegrasi", RGBColor(0x00, 0x88, 0xD0)),
    ("5", "Keamanan Informasi", "ISO 27001, CSIRT, BSSN", RGBColor(0x00, 0x96, 0xD8)),
    ("6", "Evaluasi & Audit", "Self-assessment & Audit TIK", ACCENT_GOLD),
]

card_w2 = Inches(3.8)
card_h2 = Inches(2.0)
gap_x = Inches(0.25)
gap_y = Inches(0.25)
start_x = Inches(0.6)
start_y = Inches(1.3)

for idx, (num, title, desc, color) in enumerate(steps):
    col = idx % 3
    row = idx // 3
    x = start_x + col * (card_w2 + gap_x)
    y = start_y + row * (card_h2 + gap_y)
    
    _card(slide, x, y, card_w2, card_h2, f"Langkah {num}: {title}", [desc], color)

_footer(slide)


# ════════════════════════════════════════════════════════════════
# SLIDE 16 — PRINSIP UTAMA & PENUTUP
# ════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])

_section_header_bar(slide, "Prinsip Utama: Interoperabilitas")

# Main message
txBox = slide.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(11.5), Inches(1.2))
tf = txBox.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "Jebakan Terbesar SPBE:"
p.font.size = Pt(20)
p.font.bold = True
p.font.color.rgb = ACCENT_GOLD
p2 = tf.add_paragraph()
p2.text = "Menganggap digitalisasi berarti 'membuat aplikasi baru untuk setiap masalah'"
p2.font.size = Pt(18)
p2.font.bold = True
p2.font.color.rgb = DARK_BLUE
p2.space_before = Pt(6)

# The right paradigm
txBox2 = slide.shapes.add_textbox(Inches(0.8), Inches(3.0), Inches(11.5), Inches(1.5))
tf2 = txBox2.text_frame
tf2.word_wrap = True
p3 = tf2.paragraphs[0]
p3.text = "Paradigma SPBE Modern"
p3.font.size = Pt(18)
p3.font.bold = True
p3.font.color.rgb = MED_BLUE

paradigm_items = [
    "Fokus pada interoperabilitas — bagaimana sistem yang sudah ada bisa saling berbicara",
    "Bertukar data tanpa memaksa masyarakat mengunduh puluhan aplikasi berbeda",
    "Satu platform terpadu, banyak layanan — bukan banyak platform, satu layanan",
    "Efisiensi anggaran melalui berbagi pakai infrastruktur dan aplikasi"
]
_bullet_text(slide, Inches(0.8), Inches(3.7), Inches(11.5), Inches(2.0), paradigm_items, font_size=14, color=BLACK)

# Closing
txBox3 = slide.shapes.add_textbox(Inches(0.8), Inches(5.6), Inches(11.5), Inches(0.8))
tf3 = txBox3.text_frame
tf3.word_wrap = True
p4 = tf3.paragraphs[0]
p4.text = "SPBE adalah transformasi birokrasi, bukan proyek TIK. Keberhasilannya diukur dari layanan yang dirasakan masyarakat, bukan jumlah aplikasi yang dibuat."
p4.font.size = Pt(14)
p4.font.bold = True
p4.font.color.rgb = GRAY
p4.font.italic = True

_footer(slide)


# ════════════════════════════════════════════════════════════════
# SLIDE 17 — SLIDE AKHIR
# ════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])

_rect(slide, 0, 0, W, H, DARK_BLUE)
_rect(slide, 0, H * 0.44, W, Inches(0.06), ACCENT_GOLD)

txBox = slide.shapes.add_textbox(Inches(1), Inches(2.0), Inches(11), Inches(1.5))
tf = txBox.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "Terima Kasih"
p.font.size = Pt(40)
p.font.bold = True
p.font.color.rgb = WHITE
p.alignment = PP_ALIGN.CENTER

txBox2 = slide.shapes.add_textbox(Inches(2), Inches(3.5), Inches(9), Inches(1.0))
tf2 = txBox2.text_frame
tf2.word_wrap = True
p2 = tf2.paragraphs[0]
p2.text = "Mari Wujudkan SPBE yang Terintegrasi, Efisien, Aman, dan Tepercaya"
p2.font.size = Pt(18)
p2.font.color.rgb = ACCENT_GOLD
p2.alignment = PP_ALIGN.CENTER

p2b = tf2.add_paragraph()
p2b.text = "Bersama membangun birokrasi digital Indonesia yang melayani"
p2b.font.size = Pt(14)
p2b.font.color.rgb = LIGHT_BLUE
p2b.alignment = PP_ALIGN.CENTER
p2b.space_before = Pt(12)

# Sources
txBox3 = slide.shapes.add_textbox(Inches(2), Inches(5.5), Inches(9), Inches(0.8))
tf3 = txBox3.text_frame
tf3.word_wrap = True
p3 = tf3.paragraphs[0]
p3.text = "Sumber: Perpres No. 95/2018 • Perpres No. 132/2022 • Pedoman SPBE Kemenpan-RB"
p3.font.size = Pt(10)
p3.font.color.rgb = GRAY
p3.alignment = PP_ALIGN.CENTER

_footer(slide)


# ════════════════════════════════════════════════════════════════
# SAVE
# ════════════════════════════════════════════════════════════════
output_path = "/home/runner/coba_ppt/Penerapan_SPBE.pptx"
prs.save(output_path)
print(f"✅ Presentation saved to: {output_path}")
print(f"   Total slides: {len(prs.slides)}")
