#!/usr/bin/env python3
"""
Perwal Bekasi No 51/2024 — Pajak Reklame
Versi: python-pptx (PowerPoint compatible)
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import copy

# ─── CONSTANTS ───
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)
M = Inches(0.5)
CW = SLIDE_W - 2 * M

NAVY = RGBColor(0x0A, 0x16, 0x28)
NAVY_L = RGBColor(0x12, 0x29, 0x4A)
NAVY_M = RGBColor(0x1B, 0x3A, 0x6B)
GOLD = RGBColor(0xC8, 0x96, 0x2E)
GOLD_L = RGBColor(0xD4, 0xA0, 0x17)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
OFF_W = RGBColor(0xF5, 0xF7, 0xFA)
ICE = RGBColor(0xE8, 0xED, 0xF5)
TX_D = RGBColor(0x1A, 0x1A, 0x2E)
TX_G = RGBColor(0x6B, 0x72, 0x88)
TX_M = RGBColor(0x9C, 0xA3, 0xAF)
BL = RGBColor(0x25, 0x63, 0xEB)
TEAL = RGBColor(0x0D, 0x94, 0x88)
WARM = RGBColor(0xB8, 0x86, 0x0B)
RED = RGBColor(0xDC, 0x26, 0x26)

prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H

pg_num = [0]

def blank_slide():
    layout = prs.slide_layouts[6]  # blank
    return prs.slides.add_slide(layout)

def add_textbox(slide, left, top, width, height, text, font_size=12, bold=False, color=TX_D, align=PP_ALIGN.LEFT, font_name="Calibri", valign=MSO_ANCHOR.TOP, line_spacing=1.0):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = align
    p.space_after = Pt(0)
    p.space_before = Pt(0)
    tf.paragraphs[0].space_after = Pt(0)
    return txBox

def add_shape(slide, shape_type, left, top, width, height, fill_color=None, line_color=None, line_width=None):
    shape = slide.shapes.add_shape(shape_type, left, top, width, height)
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if line_color:
        shape.line.color.rgb = line_color
        if line_width:
            shape.line.width = Pt(line_width)
    else:
        shape.line.fill.background()
    return shape

def add_rounded_rect(slide, left, top, width, height, fill_color=WHITE, line_color=ICE, line_width=0.5):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(line_width)
    else:
        shape.line.fill.background()
    # Adjust corner radius
    shape.adjustments[0] = 0.04
    return shape

def add_oval(slide, left, top, size, fill_color):
    shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top, size, size)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    return shape

def add_text_into_shape(shape, text, font_size=12, bold=False, color=TX_D, align=PP_ALIGN.LEFT, font_name="Calibri"):
    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = align
    return tf

def add_multiline_textbox(slide, left, top, width, height, lines, default_size=12, default_color=TX_D):
    """lines: list of (text, size, bold, color, bullet) or just strings"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    
    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        
        if isinstance(line, str):
            if line == "":
                p.text = ""
                p.font.size = Pt(6)
                p.space_after = Pt(0)
                p.space_before = Pt(0)
                continue
            text = line
            size = default_size
            bold = False
            color = default_color
            bullet = True
        else:
            text = line[0]
            size = line[1] if len(line) > 1 else default_size
            bold = line[2] if len(line) > 2 else False
            color = line[3] if len(line) > 3 else default_color
            bullet = line[4] if len(line) > 4 else True
        
        p.text = text
        p.font.size = Pt(size)
        p.font.bold = bold
        p.font.color.rgb = color
        p.font.name = "Calibri"
        p.alignment = PP_ALIGN.LEFT
        p.space_after = Pt(1)
        p.space_before = Pt(0)
    
    return txBox

def gold_bar(slide, top=Inches(0)):
    add_shape(slide, MSO_SHAPE.RECTANGLE, 0, top, SLIDE_W, Inches(0.03), fill_color=GOLD)

def navy_bottom(slide):
    add_shape(slide, MSO_SHAPE.RECTANGLE, 0, SLIDE_H - Inches(0.22), SLIDE_W, Inches(0.03), fill_color=NAVY)

def gold_bottom(slide):
    add_shape(slide, MSO_SHAPE.RECTANGLE, 0, SLIDE_H - Inches(0.22), SLIDE_W, Inches(0.22), fill_color=GOLD)

def page_num(slide):
    pg_num[0] += 1
    add_textbox(slide, SLIDE_W - Inches(0.9), SLIDE_H - Inches(0.4), Inches(0.7), Inches(0.25),
                str(pg_num[0]), 8, color=TX_M, align=PP_ALIGN.RIGHT)

def header_slide(slide, title, subtitle=None):
    gold_bar(slide)
    add_shape(slide, MSO_SHAPE.RECTANGLE, 0, Inches(0.03), SLIDE_W, Inches(0.8), fill_color=NAVY)
    add_shape(slide, MSO_SHAPE.RECTANGLE, M, Inches(0.12), Inches(0.07), Inches(0.5), fill_color=GOLD)
    add_textbox(slide, M + Inches(0.2), Inches(0.12), CW - Inches(1), Inches(0.42),
                title, 20, bold=True, color=WHITE)
    if subtitle:
        add_textbox(slide, M + Inches(0.2), Inches(0.52), CW - Inches(1), Inches(0.25),
                    subtitle, 10, color=TX_M)

def section_slide(title, subtitle=None):
    s = blank_slide()
    bg = s.background
    bg.fill.solid()
    bg.fill.fore_color.rgb = NAVY
    
    add_oval(s, Inches(-1.5), Inches(-1.5), Inches(5), NAVY_L)
    # Use a slightly different shade - transparency not directly supported, use lighter color
    add_oval(s, Inches(-1.5), Inches(-1.5), Inches(5), RGBColor(0x0D, 0x1F, 0x3C))
    add_oval(s, Inches(9.5), Inches(4), Inches(5), RGBColor(0x0D, 0x1F, 0x3C))
    
    add_shape(s, MSO_SHAPE.RECTANGLE, M, Inches(2.3), Inches(2.5), Inches(0.04), fill_color=GOLD)
    add_textbox(s, M, Inches(2.6), CW, Inches(1.8), title, 34, bold=True, color=WHITE)
    if subtitle:
        add_textbox(s, M, Inches(4.3), CW - Inches(2), Inches(0.4), subtitle, 13, color=TX_M)
    gold_bottom(s)
    page_num(s)
    return s

# ═══════════════════════════════════════════════════
# SLIDES
# ═══════════════════════════════════════════════════

# ─── COVER ───
s = blank_slide()
bg = s.background; bg.fill.solid(); bg.fill.fore_color.rgb = NAVY
add_oval(s, Inches(-1), Inches(-1.5), Inches(4.5), RGBColor(0x0D, 0x1F, 0x3C))
add_oval(s, Inches(8.5), Inches(-2), Inches(7), RGBColor(0x0D, 0x1F, 0x3C))
add_oval(s, Inches(10), Inches(4.5), Inches(5), RGBColor(0x0D, 0x1F, 0x3C))
add_shape(s, MSO_SHAPE.RECTANGLE, M, Inches(2.6), Inches(4), Inches(0.04), fill_color=GOLD)
add_textbox(s, M, Inches(0.5), Inches(5), Inches(0.35), "BERITA DAERAH", 12, bold=True, color=GOLD)
add_textbox(s, M, Inches(0.85), Inches(5), Inches(0.35), "KOTA BEKASI", 14, bold=True, color=WHITE)
add_textbox(s, M, Inches(1.8), CW, Inches(0.45), "PERATURAN WALI KOTA BEKASI", 20, bold=True, color=WHITE)
add_textbox(s, M, Inches(2.15), CW, Inches(0.4), "NOMOR 51 TAHUN 2024", 15, bold=True, color=GOLD)
add_textbox(s, M, Inches(3.1), CW, Inches(2.0), "TENTANG\nPENGELOLAAN PAJAK REKLAME", 40, bold=True, color=WHITE)
s2 = add_shape(s, MSO_SHAPE.ROUNDED_RECTANGLE, M, Inches(5.4), Inches(7.5), Inches(0.9), fill_color=RGBColor(0x0D, 0x1F, 0x3C))
add_textbox(s, M + Inches(0.2), Inches(5.45), Inches(7), Inches(0.4), "Pemerintah Kota Bekasi  ·  20 Desember 2024", 12, color=ICE)
add_textbox(s, M + Inches(0.2), Inches(5.75), Inches(7), Inches(0.35), "Berlaku sejak diundangkan", 10, color=TX_M)
gold_bottom(s)

# ─── DAFTAR ISI ───
s = blank_slide(); bg = s.background; bg.fill.solid(); bg.fill.fore_color.rgb = OFF_W
header_slide(s, "DAFTAR ISI", "Perwal Bekasi No 51/2024")
toc = [
    ("1","Ketentuan Umum & Definisi",BL),("2","Objek, Subjek & Wajib Pajak",TEAL),("3","Masa Pajak & Tahun Pajak",WARM),
    ("4","Pendaftaran & Pendataan WP",NAVY_M),("5","Nilai Sewa Reklame (NSR)",RED),("6","Perhitungan & Tarif Pajak",BL),
    ("7","Penetapan, Tagihan & Pembayaran",TEAL),("8","Pembetulan, Keberatan & Banding",WARM),("9","Pemeriksaan, Penagihan & Penghapusan",NAVY_M),
    ("10","Keringanan, Kemudahan & Penghargaan",BL),("11","Ketentuan Penutup",TEAL),
]
cw2 = Inches(5.5); gap2 = Inches(0.35); sx2 = (SLIDE_W - 2*cw2 - gap2) / 2
for i, (num, lb, clr) in enumerate(toc):
    col = 0 if i < 6 else 1; row = i if i < 6 else i - 6
    x = sx2 + col * (cw2 + gap2); y = Inches(1.1) + row * Inches(0.85)
    add_rounded_rect(s, x, y, cw2, Inches(0.65))
    badge = add_shape(s, MSO_SHAPE.ROUNDED_RECTANGLE, x + Inches(0.1), y + Inches(0.08), Inches(0.5), Inches(0.5), fill_color=clr)
    add_text_into_shape(badge, num, 14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_textbox(s, x + Inches(0.75), y + Inches(0.08), cw2 - Inches(0.9), Inches(0.5), lb, 13, color=TX_D, valign=MSO_ANCHOR.MIDDLE)
navy_bottom(s); page_num(s)

# ═══════════════════════════════════════════
# BAB I
# ═══════════════════════════════════════════
section_slide("BAB I\nKETENTUAN UMUM", "Pasal 1 — Definisi Kunci")

s = blank_slide(); bg = s.background; bg.fill.solid(); bg.fill.fore_color.rgb = OFF_W
header_slide(s, "Definisi Penting")
defs = [
    ("🏛️","Daerah","Kota Bekasi"),("📊","Bapenda","Badan Pendapatan Daerah Kota Bekasi"),
    ("📢","Reklame","Media untuk promosi & pengenalan komersial"),("💰","Pajak Reklame","Pajak atas penyelenggaraan reklame"),
    ("📐","NSR","Nilai Sewa Reklame — dasar pengenaan pajak"),("🆔","NPWPD","Nomor Pokok Wajib Pajak Daerah"),
    ("👤","Wajib Pajak","Orang pribadi/badan dengan hak & kewajiban pajak"),
]
for i, (ic, term, defn) in enumerate(defs):
    y = Inches(1.1) + i * Inches(0.78)
    add_rounded_rect(s, M, y, CW, Inches(0.65))
    add_oval(s, M + Inches(0.12), y + Inches(0.08), Inches(0.5), ICE)
    add_textbox(s, M + Inches(0.12), y + Inches(0.1), Inches(0.5), Inches(0.45), ic, 12, align=PP_ALIGN.CENTER)
    add_textbox(s, M + Inches(0.75), y + Inches(0.04), Inches(2.5), Inches(0.3), term, 13, bold=True, color=TX_D)
    add_textbox(s, M + Inches(0.75), y + Inches(0.32), CW - Inches(1), Inches(0.3), defn, 11, color=TX_G)
navy_bottom(s); page_num(s)

# ═══════════════════════════════════════════
# BAB II
# ═══════════════════════════════════════════
section_slide("BAB II\nOBJEK, SUBJEK & WAJIB PAJAK", "Pasal 2–4")

# Grid helper
def card_grid(title, cards, subtitle=None):
    s = blank_slide(); bg = s.background; bg.fill.solid(); bg.fill.fore_color.rgb = OFF_W
    header_slide(s, title, subtitle)
    n = len(cards)
    if n == 2: cw = Inches(5.85)
    elif n == 3: cw = Inches(3.85)
    else: cw = Inches(2.85)
    gap = Inches(0.3)
    total_w = n * cw + (n - 1) * gap
    sx = (SLIDE_W - total_w) / 2
    sy = Inches(1.1)
    ch = Inches(5.4)
    
    for i, card_data in enumerate(cards):
        cx = sx + i * (cw + gap)
        clr = card_data.get('clr', BL)
        ic = card_data.get('ic', None)
        t = card_data.get('t', '')
        items = card_data.get('items', [])
        
        add_rounded_rect(s, cx, sy, cw, ch)
        # Left accent
        add_shape(s, MSO_SHAPE.RECTANGLE, cx, sy, Inches(0.05), ch, fill_color=clr)
        
        if ic:
            add_oval(s, cx + Inches(0.15), sy + Inches(0.15), Inches(0.4), clr)
            add_textbox(s, cx + Inches(0.15), sy + Inches(0.15), Inches(0.4), Inches(0.4), ic, 14, color=WHITE, align=PP_ALIGN.CENTER)
            ty = sy + Inches(0.65)
        else:
            ty = sy + Inches(0.15)
        
        add_textbox(s, cx + Inches(0.15), ty, cw - Inches(0.3), Inches(0.35), t, 14, bold=True, color=clr)
        
        # Bullet items
        lines = []
        for item in items:
            lines.append((f"• {item}", 11, False, TX_D, False))
        # ... we'll handle this inline
        ay = ty + Inches(0.4)
        for item in items:
            add_textbox(s, cx + Inches(0.15), ay, cw - Inches(0.3), Inches(0.25), f"• {item}", 11, color=TX_D)
            ay += Inches(0.25)
    
    navy_bottom(s); page_num(s)
    return s

# Objek Pajak Reklame
card_grid("Objek Pajak Reklame", [
    {"ic":"📋", "t":"10 Jenis Reklame", "clr":BL, "items":["Papan / Billboard","Videotron / Megatron","Kain (Spanduk, Umbul, Baliho)","Melekat / Stiker","Selebaran","Berjalan (Kendaraan)","Udara (Balon Gas)","Apung","Film / Slide","Peragaan"]},
    {"ic":"🚫", "t":"Dikecualikan", "clr":TEAL, "items":["Internet, TV, radio, media cetak","Label / merek produk","Nama usaha ≤ 1 m² di tempat","Reklame Pemerintah/Pemda","Tempat ibadah & panti asuhan","Sosial & keagamaan ≤ 30 hari","Kegiatan politik (masa kampanye)","Olahraga KONI ≤ 30 hari"]},
])

# Subjek & Wajib Pajak
card_grid("Subjek & Wajib Pajak", [
    {"ic":"👤", "t":"Subjek Pajak (Pasal 3)", "clr":BL, "items":["Orang pribadi atau Badan","yang menggunakan Reklame"]},
    {"ic":"✋", "t":"Wajib Pajak (Pasal 4)", "clr":TEAL, "items":["Orang pribadi atau Badan","yang menyelenggarakan Reklame","Jika pihak ketiga → menjadi WP"]},
])

# ═══════════════════════════════════════════
# BAB III
# ═══════════════════════════════════════════
section_slide("BAB III\nMASA PAJAK & TAHUN PAJAK", "Pasal 5")

# Callout slide
s = blank_slide(); bg = s.background; bg.fill.solid(); bg.fill.fore_color.rgb = OFF_W
header_slide(s, "Masa & Tahun Pajak", "Pasal 5")
callouts = [
    ("12", "Bulan (Permanen)", BL), ("30", "Hari (Insidentil)", TEAL),
    ("1", "Tahun Pajak", WARM), ("1", "Bulan (Bagian Tahun)", NAVY_M),
]
n = len(callouts); cw = Inches(2.8); gap = Inches(0.35)
tw = n * cw + (n - 1) * gap; sx = (SLIDE_W - tw) / 2
for i, (num, lb, clr) in enumerate(callouts):
    cx = sx + i * (cw + gap)
    add_rounded_rect(s, cx, Inches(1.3), cw, Inches(2.1))
    add_shape(s, MSO_SHAPE.RECTANGLE, cx, Inches(1.3), cw, Inches(0.05), fill_color=clr)
    add_textbox(s, cx, Inches(1.6), cw, Inches(0.7), num, 32, bold=True, color=clr, align=PP_ALIGN.CENTER)
    add_textbox(s, cx + Inches(0.1), Inches(2.35), cw - Inches(0.2), Inches(0.6), lb, 11, color=TX_G, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)

# Extra note
box = add_rounded_rect(s, M, Inches(3.8), CW, Inches(2.5), ICE, None)
extra_text = "• Masa Pajak Permanen: 12 bulan atau sesuai jangka waktu penayangan reklame\n• Masa Pajak Insidentil: dihitung per hari, maksimal 30 hari\n• Tahun Pajak: 1 tahun kalender atau sesuai tahun buku wajib pajak\n• Bagian Tahun Pajak: 1 bulan penuh (jika tidak mencakup satu tahun penuh)"
add_textbox(s, M + Inches(0.2), Inches(3.9), CW - Inches(0.4), Inches(2.3), extra_text, 12, color=TX_D)
navy_bottom(s); page_num(s)

# ═══════════════════════════════════════════
# BAB IV
# ═══════════════════════════════════════════
section_slide("BAB IV\nPENDAFTARAN & PENDATAAN WP", "Pasal 6–8")

def two_col(title, left_lines, right_lines, subtitle=None, right_color=TEAL):
    s = blank_slide(); bg = s.background; bg.fill.solid(); bg.fill.fore_color.rgb = OFF_W
    header_slide(s, title, subtitle)
    cw = Inches(5.8); gap = Inches(0.35)
    pairs = [
        (left_lines, BL, M),
        (right_lines, right_color, M + cw + gap),
    ]
    for lines, clr, x in pairs:
        add_rounded_rect(s, x, Inches(1.1), cw, Inches(5.4))
        add_shape(s, MSO_SHAPE.RECTANGLE, x, Inches(1.1), Inches(0.05), Inches(5.4), fill_color=clr)
        first = True
        ay = Inches(1.3)
        for line in lines:
            if line.startswith("$"):
                add_textbox(s, x + Inches(0.2), ay, cw - Inches(0.35), Inches(0.3), line[1:], 13, bold=True, color=clr)
            else:
                add_textbox(s, x + Inches(0.2), ay, cw - Inches(0.35), Inches(0.25), f"• {line}", 11, color=TX_D)
            ay += Inches(0.28)
            first = False
    navy_bottom(s); page_num(s)
    return s

two_col("Pendaftaran & Pendataan", [
    "$PENDAFTARAN (Pasal 6)","","WP wajib mendaftarkan diri & objek pajak","Formulir: ambil/online/dikirim petugas","Lampirkan: KTP, NPWP, Akta, NIB","Bapenda terbitkan NPWPD","Jika tidak mendaftar → NPWPD jabatan","Juga: NOPD & nomor registrasi",
], [
    "$PENDATAAN & NONAKTIF (Pasal 7–8)","","Bapenda mendata WP & objek pajak","Termasuk data geografis","Dapat kerjasama dengan instansi lain","Penonaktifan: WP tak penuhi syarat","Keputusan maksimal 3 bulan","Syarat: tanpa tunggakan & keberatan",
])

# ═══════════════════════════════════════════
# BAB V
# ═══════════════════════════════════════════
section_slide("BAB V\nNILAI SEWA REKLAME", "Pasal 9 — Dasar Pengenaan Pajak")

s = blank_slide(); bg = s.background; bg.fill.solid(); bg.fill.fore_color.rgb = OFF_W
header_slide(s, "Faktor Penentu NSR")
factors = [("1","Jenis Reklame",BL),("2","Bahan",TEAL),("3","Lokasi (Kelas Jalan)",WARM),("4","Waktu Tayang (detik)",NAVY_M),("5","Jangka Waktu (hari)",BL),("6","Jumlah Media",TEAL),("7","Ukuran (m²)",WARM)]
bw = Inches(2.8); bgap = Inches(0.3); per_row = 4
total_bw = per_row * bw + (per_row - 1) * bgap; bsx = (SLIDE_W - total_bw) / 2
for i, (num, lb, clr) in enumerate(factors):
    col = i % per_row; row = i // per_row
    x = bsx + col * (bw + bgap); y = Inches(1.1) + row * Inches(1.6)
    add_rounded_rect(s, x, y, bw, Inches(1.3))
    add_oval(s, x + Inches(0.15), y + Inches(0.25), Inches(0.55), clr)
    add_textbox(s, x + Inches(0.15), y + Inches(0.25), Inches(0.55), Inches(0.55), num, 18, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_textbox(s, x + Inches(0.8), y + Inches(0.2), bw - Inches(1), Inches(0.9), lb, 13, color=TX_D, valign=MSO_ANCHOR.MIDDLE)
# Bottom note
box = add_rounded_rect(s, M, Inches(4.4), CW, Inches(1.7), ICE, None)
add_textbox(s, M + Inches(0.2), Inches(4.45), Inches(5), Inches(0.3), "KLASIFIKASI KELAS JALAN", 11, bold=True, color=NAVY)
cls_text = "🏛️  Kelas Jalan Khusus — Tol | Premium 1 | Premium 2\n🚗  Kelas Jalan I (Kendali Ketat) — Lebar > 3 m, pusat pelayanan\n🏡  Kelas Jalan II (Kendali Sedang) — Lebar ≤ 3 m, jalan lingkungan"
add_textbox(s, M + Inches(0.2), Inches(4.8), CW - Inches(0.4), Inches(1.0), cls_text, 10, color=TX_D)
navy_bottom(s); page_num(s)

# ═══════════════════════════════════════════
# BAB VI
# ═══════════════════════════════════════════
section_slide("BAB VI\nPERHITUNGAN & TARIF PAJAK", "Pasal 10")

# Callout rumus
s = blank_slide(); bg = s.background; bg.fill.solid(); bg.fill.fore_color.rgb = OFF_W
header_slide(s, "Rumus Dasar Perhitungan", "Pasal 10")
rumus = [("×","Pajak = Tarif × NSR",BL),("50%","Indoor = 50% NSR",TEAL),("+20%","Tinggi > 15 m",WARM),("+50%","Tembakau & Miras",RED)]
for i, (num, lb, clr) in enumerate(rumus):
    cx = (SLIDE_W - 4*Inches(2.8) - 3*Inches(0.35))/2 + i*(Inches(2.8)+Inches(0.35))
    add_rounded_rect(s, cx, Inches(1.3), Inches(2.8), Inches(2.1))
    add_shape(s, MSO_SHAPE.RECTANGLE, cx, Inches(1.3), Inches(2.8), Inches(0.05), fill_color=clr)
    add_textbox(s, cx, Inches(1.6), Inches(2.8), Inches(0.7), num, 32, bold=True, color=clr, align=PP_ALIGN.CENTER)
    add_textbox(s, cx + Inches(0.1), Inches(2.35), Inches(2.6), Inches(0.6), lb, 11, color=TX_G, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
box = add_rounded_rect(s, M, Inches(3.8), CW, Inches(2.5), ICE, None)
add_textbox(s, M + Inches(0.2), Inches(3.9), CW - Inches(0.4), Inches(2.3),
    "Rumus: Pajak Reklame = Tarif Pajak × NSR\nNSR = Nilai Kelas Jalan × Ukuran (m²) × Jumlah × Jangka Waktu\n\nKetentuan Khusus:\n• Reklame indoor: NSR 50% dari NSR normal\n• Ketinggian > 15 meter: tambahan 20%\n• Produk tembakau & minuman beralkohol: tambahan 50%\n• Perubahan naskah/revisi isi reklame: dikecualikan", 12, color=TX_D)
navy_bottom(s); page_num(s)

# ─── NATIVE TABLE ───
def set_cell(cell, text, font_size=11, bold=False, color=TX_D, align=PP_ALIGN.LEFT, fill=None):
    cell.text = ""
    p = cell.text_frame.paragraphs[0]
    p.text = str(text)
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = "Calibri"
    p.alignment = align
    cell.text_frame.word_wrap = True
    cell.vertical_anchor = MSO_ANCHOR.MIDDLE
    if fill:
        cell.fill.solid()
        cell.fill.fore_color.rgb = fill

def table_slide(title, headers, rows, subtitle=None):
    s = blank_slide(); bg = s.background; bg.fill.solid(); bg.fill.fore_color.rgb = OFF_W
    header_slide(s, title, subtitle)
    
    tbl_left = M
    tbl_top = Inches(1.15)
    tbl_width = CW
    row_height = Inches(0.42)
    n_rows = len(rows) + 1  # +1 for header
    n_cols = len(headers)
    
    table_shape = s.shapes.add_table(n_rows, n_cols, tbl_left, tbl_top, tbl_width, row_height * n_rows)
    table = table_shape.table
    
    # Distribute columns evenly (in EMU)
    for i in range(n_cols):
        table.columns[i].width = int(tbl_width / n_cols)
    
    # Header row
    for j, hdr in enumerate(headers):
        set_cell(table.cell(0, j), hdr, bold=True, color=WHITE, align=PP_ALIGN.CENTER, fill=NAVY)
    
    # Data rows
    for ri, row in enumerate(rows):
        bg_color = ICE if ri % 2 == 0 else WHITE
        for j, val in enumerate(row):
            align = PP_ALIGN.LEFT if j == 0 else PP_ALIGN.CENTER
            set_cell(table.cell(ri + 1, j), val, color=TX_D, align=align, fill=bg_color)
    
    navy_bottom(s); page_num(s)
    return s

table_slide("NSR — Papan / Billboard", ["Kelas Jalan", "Zona", "NSR (Rp/m²/hari)"], [
    ["Kelas Jalan Khusus", "Jalan Tol", "23.575"], ["Kelas Jalan Khusus", "Premium 1", "16.100"],
    ["Kelas Jalan Khusus", "Premium 2", "14.950"], ["Kelas Jalan I", "Kendali Ketat", "13.225"],
    ["Kelas Jalan II", "Kendali Sedang", "11.500"],
])

table_slide("NSR — Megatron / Videotron", ["Kelas Jalan", "Zona", "NSR (/30 dtk)", "NSR (/m²/thn)"], [
    ["Kelas Jalan Khusus", "Jalan Tol", "Rp 17,25", "13.599.900"], ["Kelas Jalan Khusus", "Premium 1", "Rp 13,80", "10.879.920"],
    ["Kelas Jalan Khusus", "Premium 2", "Rp 9,20", "7.253.280"], ["Kelas Jalan I", "Kendali Ketat", "Rp 8,05", "6.346.620"],
    ["Kelas Jalan II", "Kendali Sedang", "Rp 5,75", "4.533.300"],
])

table_slide("NSR — Kain (Spanduk/Umbul/Baliho)", ["Kelas Jalan", "Zona", "NSR (Rp/m²/hari)"], [
    ["Kelas Jalan Khusus", "Jalan Tol", "30.000"], ["Kelas Jalan Khusus", "Premium 1", "30.000"],
    ["Kelas Jalan Khusus", "Premium 2", "25.000"], ["Kelas Jalan I", "Kendali Ketat", "20.000"],
    ["Kelas Jalan II", "Kendali Sedang", "19.000"],
])

# Card grids for NSR lainnya
card_grid("NSR — Jenis Reklame Lainnya", [
    {"ic":"🏷️", "t":"Stiker", "clr":BL, "items":["Rp 7,5/cm²","Min. Rp 750.000/kali"]},
    {"ic":"🧱", "t":"Melekat", "clr":TEAL, "items":["Rp 750.000/m²/tahun"]},
    {"ic":"📄", "t":"Selebaran", "clr":WARM, "items":["Rp 600/lembar","Min. Rp 6.000.000/kali"]},
    {"ic":"🚌", "t":"Berjalan", "clr":RED, "items":["Rp 6.000/m²/hari","Termasuk kendaraan"]},
])
card_grid("NSR — Jenis Lainnya (lanjutan)", [
    {"ic":"🎈", "t":"Udara", "clr":BL, "items":["Rp 2.400.000/sekali","Maks. 1 bulan"]},
    {"ic":"🌊", "t":"Apung", "clr":TEAL, "items":["Rp 600.000/sekali","Maks. 1 bulan"]},
    {"ic":"🎬", "t":"Film / Slide", "clr":WARM, "items":["Rp 12.000/15 detik"]},
    {"ic":"🎭", "t":"Peragaan", "clr":RED, "items":["Rp 480.000/penyelenggaraan"]},
])

# ═══════════════════════════════════════════
# BAB VII
# ═══════════════════════════════════════════
section_slide("BAB VII\nPENETAPAN, TAGIHAN & PEMBAYARAN", "Pasal 11–14")

# Flow slide
s = blank_slide(); bg = s.background; bg.fill.solid(); bg.fill.fore_color.rgb = OFF_W
header_slide(s, "Alur Penetapan → Tagihan → Pembayaran", "Pasal 11–14")
steps = [("1","SKPD","Diterbitkan Bapenda\nMasa berlaku 5 tahun",BL),("2","Pembayaran","Lunas 1 bulan\nsejak SKPD diterima",TEAL),("3","Keterlambatan","Bunga 1%/bln\nDiterbitkan STPD",WARM),("4","STPD","Harus lunas\n≤ 30 hari",RED)]
bw = Inches(2.6); bgap = Inches(0.5)
total_sw = len(steps)*bw + (len(steps)-1)*bgap; ssx = (SLIDE_W - total_sw)/2
for i, (num, title, desc, clr) in enumerate(steps):
    cx = ssx + i*(bw+bgap)
    add_rounded_rect(s, cx, Inches(1.8), bw, Inches(2.0))
    add_shape(s, MSO_SHAPE.RECTANGLE, cx, Inches(1.8), bw, Inches(0.05), fill_color=clr)
    add_oval(s, cx + bw/2 - Inches(0.25), Inches(1.95), Inches(0.5), clr)
    add_textbox(s, cx + bw/2 - Inches(0.25), Inches(1.95), Inches(0.5), Inches(0.5), num, 18, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_textbox(s, cx + Inches(0.1), Inches(2.55), bw - Inches(0.2), Inches(0.35), title, 13, bold=True, color=clr, align=PP_ALIGN.CENTER)
    add_textbox(s, cx + Inches(0.1), Inches(2.9), bw - Inches(0.2), Inches(0.7), desc, 10, color=TX_G, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.TOP)
    if i < len(steps) - 1:
        add_textbox(s, cx + bw, Inches(2.4), bgap, Inches(0.5), "›", 24, color=GOLD, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
box = add_rounded_rect(s, M, Inches(4.3), CW, Inches(1.7), ICE, None)
add_textbox(s, M + Inches(0.2), Inches(4.35), CW - Inches(0.4), Inches(1.5),
    "• Jatuh tempo: 1 bulan sejak tanggal pengiriman SKPD\n• Pembayaran: Kas Daerah / Bank Persepsi / tempat lain yang ditunjuk\n• Stiker sebagai tanda bukti pembayaran reklame\n• STPD dikenakan bunga 1%/bulan (maks. 24 bulan)", 11, color=TX_D)
navy_bottom(s); page_num(s)

# ═══════════════════════════════════════════
# BAB VIII
# ═══════════════════════════════════════════
section_slide("BAB VIII\nPEMBETULAN, KEBERATAN & BANDING", "Pasal 15–20, 29–33")

two_col("Pembetulan Ketetapan", [
    "$PEMBETULAN (Pasal 15–20)","","Kesalahan tulis: nama, alamat, NPWPD","Kesalahan hitung: jumlah, tarif","Kekeliruan penerapan aturan","1 permohonan = 1 ketetapan","Keputusan maksimal 6 bulan","> 6 bulan tanpa putusan → dikabulkan","Dapat dilakukan berulang (Ps 20)","Jenis keputusan: kabul / batal / tolak",
], [
    "$JANGKA WAKTU & SANKSI","","Permohonan diajukan ke Bapenda","Keputusan: kabul (tambah/kurang/hapus)","Keputusan: batal | tolak","Pasal 19: pembetulan jabatan","Pasal 20: berulang jika masih salah",
], right_color=WARM)

two_col("Keberatan & Banding", [
    "$KEBERATAN (Pasal 29–31)","","Objek: SKPD, SKPDKB, SKPDKBT, dll","Diajukan maks. 3 bulan sejak SKPD","Sudah bayar min. yang disetujui","Keputusan maks. 12 bulan","Jika ditolak: denda 30%","Jika dikabulkan: + bunga 0,6%/bulan",
], [
    "$BANDING (Pasal 32–33)","","Objek: Surat Keputusan Keberatan","Ke badan peradilan pajak","Maks. 3 bulan sejak keputusan","Menangguhkan kewajiban bayar","Jika ditolak: denda 60%","Jika dikabulkan: + bunga 0,6%/bulan",
])

# ═══════════════════════════════════════════
# BAB IX
# ═══════════════════════════════════════════
section_slide("BAB IX\nPEMERIKSAAN, PENAGIHAN & PENGHAPUSAN", "Pasal 22–26")

card_grid("Pemeriksaan & Penagihan", [
    {"ic":"🔍", "t":"Pemeriksaan (Ps 22–23)", "clr":BL, "items":["Kepala Bapenda berwenang periksa","Menguji kepatuhan WP","WP wajib: buka buku/dokumen","Beri akses tempat & keterangan","Jika tidak → pajak ditetapkan jabatan"]},
    {"ic":"📬", "t":"Penagihan (Ps 24)", "clr":TEAL, "items":["Dasar: SKPD, SKPDKB, SKPDKBT","STPD, SK Pembetulan/Keberatan","Putusan Banding"]},
    {"ic":"⏳", "t":"Kedaluwarsa (Ps 25)", "clr":WARM, "items":["5 tahun sejak pajak terutang","Tertangguh jika ada:","Surat Teguran / Paksa","Pengakuan utang dari WP"]},
])

# Penghapusan flow
s = blank_slide(); bg = s.background; bg.fill.solid(); bg.fill.fore_color.rgb = OFF_W
header_slide(s, "Penghapusan Piutang Pajak (Pasal 26)")
steps = [("1","Penelitian","Dilakukan Bapenda",BL),("2","Penetapan","Keputusan Wali Kota",TEAL),("3","Koordinasi","Dengan Inspektorat",WARM),("4","SK Penghapusan","Diterbitkan",NAVY_M)]
bw = Inches(2.6); bgap = Inches(0.5)
total_sw = len(steps)*bw + (len(steps)-1)*bgap; ssx = (SLIDE_W - total_sw)/2
for i, (num, title, desc, clr) in enumerate(steps):
    cx = ssx + i*(bw+bgap)
    add_rounded_rect(s, cx, Inches(1.8), bw, Inches(2.0))
    add_shape(s, MSO_SHAPE.RECTANGLE, cx, Inches(1.8), bw, Inches(0.05), fill_color=clr)
    add_oval(s, cx + bw/2 - Inches(0.25), Inches(1.95), Inches(0.5), clr)
    add_textbox(s, cx + bw/2 - Inches(0.25), Inches(1.95), Inches(0.5), Inches(0.5), num, 18, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_textbox(s, cx + Inches(0.1), Inches(2.55), bw - Inches(0.2), Inches(0.35), title, 13, bold=True, color=clr, align=PP_ALIGN.CENTER)
    add_textbox(s, cx + Inches(0.1), Inches(2.9), bw - Inches(0.2), Inches(0.7), desc, 10, color=TX_G, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.TOP)
    if i < len(steps) - 1:
        add_textbox(s, cx + bw, Inches(2.4), bgap, Inches(0.5), "›", 24, color=GOLD, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
box = add_rounded_rect(s, M, Inches(4.3), CW, Inches(1.7), ICE, None)
add_textbox(s, M + Inches(0.2), Inches(4.35), CW - Inches(0.4), Inches(1.5),
    "Syarat penghapusan piutang pajak:\n• Piutang tidak mungkin ditagih lagi karena kedaluwarsa\n• Ada koordinasi dengan Inspektorat Daerah\n• Dibuktikan dengan dokumen pelaksanaan penagihan", 11, color=TX_D)
navy_bottom(s); page_num(s)

# ═══════════════════════════════════════════
# BAB X
# ═══════════════════════════════════════════
section_slide("BAB X\nKERINGANAN, KEMUDAHAN & PENGHARGAAN", "Pasal 27–28, 34–35")

card_grid("Fasilitas & Penghargaan", [
    {"ic":"🎯", "t":"Keringanan (Ps 27)", "clr":BL, "items":["Keringanan / Pengurangan","Pembebasan / Penundaan","Atas pokok & sanksi pajak","WP dengan likuiditas rendah","Objek terdampak bencana/kebakaran"]},
    {"ic":"🤝", "t":"Kemudahan (Ps 28)", "clr":TEAL, "items":["Perpanjangan waktu bayar","Angsuran maks. 24 bulan","Bunga 0,6%/bulan","Keadaan kahar: bencana, wabah, kerusuhan"]},
    {"ic":"🏆", "t":"Penghargaan (Ps 34–35)", "clr":WARM, "items":["WP Taat Pajak","Bayar tepat waktu ≥ 1 tahun","Tanpa tunggakan 3 tahun","Kontribusi signifikan","Piagam / Hadiah (APBD)"]},
])

# ═══════════════════════════════════════════
# BAB XI
# ═══════════════════════════════════════════
section_slide("BAB XI\nKETENTUAN PENUTUP", "Pasal 36–37")

two_col("Pencabutan & Mulai Berlaku", [
    "$PERATURAN YANG DICABUT (Pasal 36)","","Perwal No. 48 Tahun 2012","Petunjuk Pelaksanaan Perda 14/2012","Perwal No. 52 Tahun 2013 (Perubahan)","",
], [
    "$MULAI BERLAKU (Pasal 37)","","Sejak diundangkan","20 Desember 2024","","Pj. WALI KOTA BEKASI,","ttd.","R. GANI MUHAMAD",
], right_color=NAVY_M)

# ─── CLOSING ───
s = blank_slide()
bg = s.background; bg.fill.solid(); bg.fill.fore_color.rgb = NAVY
add_oval(s, Inches(-1), Inches(-1.5), Inches(4.5), RGBColor(0x0D, 0x1F, 0x3C))
add_oval(s, Inches(8.5), Inches(-2), Inches(7), RGBColor(0x0D, 0x1F, 0x3C))
add_oval(s, Inches(10), Inches(4.5), Inches(5), RGBColor(0x0D, 0x1F, 0x3C))
add_shape(s, MSO_SHAPE.RECTANGLE, M, Inches(3.6), Inches(3.5), Inches(0.04), fill_color=GOLD)
add_textbox(s, M, Inches(1.6), CW, Inches(0.4), "BERITA DAERAH KOTA BEKASI", 14, bold=True, color=GOLD)
add_textbox(s, M, Inches(2.4), CW, Inches(1.5), "TERIMA KASIH", 48, bold=True, color=WHITE)
add_textbox(s, M, Inches(4.1), CW, Inches(0.8), "Peraturan Wali Kota Bekasi Nomor 51 Tahun 2024\nTentang Pengelolaan Pajak Reklame", 14, color=ICE)
add_textbox(s, M, Inches(5.1), CW, Inches(0.35), "Sumber: https://jdih.bekasikota.go.id", 10, color=TX_M)
gold_bottom(s)

# ══════════════════════════════════════════════
# SAVE
# ══════════════════════════════════════════════
import os
out = os.path.join(os.path.dirname(__file__), "Perwal_Bekasi_51_2024_Pajak_Reklame.pptx")
prs.save(out)
print(f"✅ OK: {out} ({len(prs.slides)} slide)")
