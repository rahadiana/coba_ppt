#!/usr/bin/env python3
"""
Generate PowerPoint presentation for:
Peraturan Wali Kota Bekasi Nomor 51 Tahun 2024
Tentang Pengelolaan Pajak Reklame
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import datetime

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# ─── Color Palette ───
BLUE_DARK   = RGBColor(0x00, 0x2B, 0x5C)
BLUE_MID    = RGBColor(0x00, 0x4E, 0x8C)
BLUE_LIGHT  = RGBColor(0x00, 0x7B, 0xBF)
GOLD        = RGBColor(0xE8, 0xB0, 0x00)
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
BLACK       = RGBColor(0x00, 0x00, 0x00)
GRAY        = RGBColor(0x55, 0x55, 0x55)
LIGHT_GRAY  = RGBColor(0xF0, 0xF0, 0xF0)
RED_ACCENT  = RGBColor(0xC0, 0x39, 0x2B)
GREEN_TABLE = RGBColor(0xE8, 0xF5, 0xE9)
BLUE_TABLE  = RGBColor(0xE3, 0xF2, 0xFD)

def add_bg_shape(slide, color=BLUE_DARK):
    """Add full-slide background rectangle."""
    shp = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height
    )
    shp.fill.solid()
    shp.fill.fore_color.rgb = color
    shp.line.fill.background()
    # send to back — we need to manually re-order; we'll add it first
    return shp

def add_accent_bar(slide, left=0, top=0, width=Inches(0.15), height=None, color=GOLD):
    """Add a thin accent bar."""
    if height is None:
        height = prs.slide_height
    shp = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, left, top, width, height
    )
    shp.fill.solid()
    shp.fill.fore_color.rgb = color
    shp.line.fill.background()
    return shp

def add_text_box(slide, left, top, width, height, text, font_size=18,
                 bold=False, color=WHITE, alignment=PP_ALIGN.LEFT, 
                 font_name='Calibri', line_spacing=1.2):
    """Add a text box with single paragraph."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = alignment
    p.space_after = Pt(line_spacing * 2)
    return txBox

def add_bullet_text(slide, left, top, width, height, items, font_size=16,
                    color=BLACK, font_name='Calibri', bold_first=False, 
                    bullet_char="•", line_spacing=1.3):
    """Add a text box with bullet-pointed items."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        
        # Handle bold prefix with "**text**" syntax
        if "**" in item:
            parts = item.split("**")
            run = p.add_run()
            run.text = f"{bullet_char} {parts[0]}"
            run.font.size = Pt(font_size)
            run.font.color.rgb = color
            run.font.name = font_name
            run.font.bold = True
            if len(parts) > 1:
                run2 = p.add_run()
                run2.text = parts[1]
                run2.font.size = Pt(font_size)
                run2.font.color.rgb = color
                run2.font.name = font_name
        else:
            p.text = f"{bullet_char} {item}"
            p.font.size = Pt(font_size)
            p.font.color.rgb = color
            p.font.name = font_name
        
        p.space_after = Pt(4)
    return txBox

def make_title_slide():
    """Slide 1: Cover / Title"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    # Background
    add_bg_shape(slide, BLUE_DARK)
    add_accent_bar(slide, left=Inches(0.5), color=GOLD)
    
    # Subtitle / metadata
    add_text_box(slide, Inches(1.2), Inches(1.0), Inches(10), Inches(0.6),
                 "BERITA DAERAH KOTA BEKASI", font_size=18, color=GOLD,
                 bold=True, alignment=PP_ALIGN.LEFT)
    
    add_text_box(slide, Inches(1.2), Inches(1.7), Inches(10), Inches(1.5),
                 "PERATURAN WALI KOTA BEKASI\nNOMOR 51 TAHUN 2024",
                 font_size=36, color=WHITE, bold=True, alignment=PP_ALIGN.LEFT)
    
    add_text_box(slide, Inches(1.2), Inches(3.5), Inches(10), Inches(1.2),
                 "TENTANG\nPENGELOLAAN PAJAK REKLAME",
                 font_size=44, color=GOLD, bold=True, alignment=PP_ALIGN.LEFT)
    
    add_text_box(slide, Inches(1.2), Inches(5.2), Inches(10), Inches(0.5),
                 "Diselenggarakan oleh Pemerintah Kota Bekasi", 
                 font_size=16, color=WHITE, alignment=PP_ALIGN.LEFT)
    
    add_text_box(slide, Inches(1.2), Inches(5.7), Inches(10), Inches(0.5),
                 "Ditetapkan: 20 Desember 2024 | Berlaku sejak diundangkan",
                 font_size=14, color=RGBColor(0xAA, 0xAA, 0xAA), alignment=PP_ALIGN.LEFT)
    
    # Bottom decorative line
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1.2), Inches(6.5), Inches(5), Pt(3))
    shp.fill.solid()
    shp.fill.fore_color.rgb = GOLD
    shp.line.fill.background()

def make_section_slide(judul_bab, subjudul="", bab_num=""):
    """Slide type for section / chapter dividers."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg_shape(slide, BLUE_DARK)
    add_accent_bar(slide, left=Inches(0.5), color=GOLD)
    
    if bab_num:
        add_text_box(slide, Inches(1.2), Inches(1.8), Inches(10), Inches(0.6),
                     bab_num, font_size=20, color=GOLD, bold=True)
    
    add_text_box(slide, Inches(1.2), Inches(2.5), Inches(10), Inches(2),
                 judul_bab, font_size=40, color=WHITE, bold=True)
    
    if subjudul:
        add_text_box(slide, Inches(1.2), Inches(4.5), Inches(10), Inches(1),
                     subjudul, font_size=18, color=RGBColor(0xCC, 0xCC, 0xCC))

def make_content_slide(title, items, notes="", sub_text=None):
    """Standard content slide with title and bullet points."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # White background with top bar
    add_bg_shape(slide, WHITE)
    # Top header band
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(0.08))
    shp.fill.solid(); shp.fill.fore_color.rgb = GOLD; shp.line.fill.background()
    
    shp2 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(0.08), Inches(0.08), Inches(1.2))
    shp2.fill.solid(); shp2.fill.fore_color.rgb = BLUE_MID; shp2.line.fill.background()
    
    # Title
    add_text_box(slide, Inches(0.5), Inches(0.25), Inches(12), Inches(0.9),
                 title, font_size=28, bold=True, color=BLUE_DARK)
    
    # Thin separator
    shp3 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(1.05), Inches(2), Pt(3))
    shp3.fill.solid(); shp3.fill.fore_color.rgb = GOLD; shp3.line.fill.background()
    
    # Sub text if any
    y_start = Inches(1.3)
    if sub_text:
        add_text_box(slide, Inches(0.6), y_start, Inches(12), Inches(0.6),
                     sub_text, font_size=15, color=GRAY)
        y_start = Inches(1.9)
    
    # Bullet items
    add_bullet_text(slide, Inches(0.6), y_start, Inches(12), Inches(5.2),
                    items, font_size=17, color=BLACK)
    
    if notes:
        add_text_box(slide, Inches(0.6), Inches(6.5), Inches(12), Inches(0.7),
                     notes, font_size=12, color=GRAY)
    
    # Page footer
    shp4 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(7.35), prs.slide_width, Inches(0.15))
    shp4.fill.solid(); shp4.fill.fore_color.rgb = BLUE_DARK; shp4.line.fill.background()

def make_table_slide(title, headers, rows, notes="", col_widths=None):
    """Content slide with a table."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg_shape(slide, WHITE)
    
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(0.08))
    shp.fill.solid(); shp.fill.fore_color.rgb = GOLD; shp.line.fill.background()
    
    shp2 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(0.08), Inches(0.08), Inches(1.2))
    shp2.fill.solid(); shp2.fill.fore_color.rgb = BLUE_MID; shp2.line.fill.background()
    
    add_text_box(slide, Inches(0.5), Inches(0.25), Inches(12), Inches(0.9),
                 title, font_size=26, bold=True, color=BLUE_DARK)
    
    shp3 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(1.0), Inches(2), Pt(3))
    shp3.fill.solid(); shp3.fill.fore_color.rgb = GOLD; shp3.line.fill.background()
    
    num_rows = len(rows) + 1
    num_cols = len(headers)
    table_left = Inches(0.5)
    table_top = Inches(1.3)
    table_width = Inches(12.3)
    table_height = Inches(5.5)
    
    table_shape = slide.shapes.add_table(num_rows, num_cols, table_left, table_top, table_width, table_height)
    table = table_shape.table
    
    if col_widths:
        for i, w in enumerate(col_widths):
            table.columns[i].width = w
    
    # Header row
    for j, h in enumerate(headers):
        cell = table.cell(0, j)
        cell.text = h
        cell.fill.solid()
        cell.fill.fore_color.rgb = BLUE_DARK
        for paragraph in cell.text_frame.paragraphs:
            paragraph.font.size = Pt(13)
            paragraph.font.bold = True
            paragraph.font.color.rgb = WHITE
            paragraph.alignment = PP_ALIGN.CENTER
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
    
    # Data rows
    for i, row in enumerate(rows):
        for j, val in enumerate(row):
            cell = table.cell(i + 1, j)
            cell.text = str(val)
            if i % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = BLUE_TABLE
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = WHITE
            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.size = Pt(12)
                paragraph.font.color.rgb = BLACK
                paragraph.alignment = PP_ALIGN.CENTER if j > 0 else PP_ALIGN.LEFT
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
    
    if notes:
        add_text_box(slide, Inches(0.5), Inches(6.8), Inches(12), Inches(0.5),
                     notes, font_size=11, color=GRAY)
    
    shp5 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(7.35), prs.slide_width, Inches(0.15))
    shp5.fill.solid(); shp5.fill.fore_color.rgb = BLUE_DARK; shp5.line.fill.background()

# ═══════════════════════════════════════════════════
# SLIDE 1: TITLE
# ═══════════════════════════════════════════════════
make_title_slide()

# ═══════════════════════════════════════════════════
# SLIDE 2: DAFTAR ISI
# ═══════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg_shape(slide, WHITE)
shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(0.08))
shp.fill.solid(); shp.fill.fore_color.rgb = GOLD; shp.line.fill.background()
add_text_box(slide, Inches(0.5), Inches(0.25), Inches(12), Inches(0.8),
             "DAFTAR ISI", font_size=30, bold=True, color=BLUE_DARK)
shp3 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(1.0), Inches(2), Pt(3))
shp3.fill.solid(); shp3.fill.fore_color.rgb = GOLD; shp3.line.fill.background()

toc_items = [
    "BAB I   : Ketentuan Umum",
    "BAB II  : Objek Pajak, Subjek Pajak, dan Wajib Pajak",
    "BAB III : Masa Pajak, Tahun Pajak, dan Bagian Tahun Pajak",
    "BAB IV  : Pendaftaran dan Pendataan Wajib Pajak",
    "BAB V   : Nilai Sewa Reklame (NSR)",
    "BAB VI  : Perhitungan Pajak Reklame",
    "BAB VII : Penetapan Besaran Pajak Terutang",
    "BAB VIII: Surat Tagihan Pajak",
    "BAB IX  : Pembayaran dan Penyetoran",
    "BAB X   : Pembetulan dan Pembatalan Ketetapan",
    "BAB XI  : Pengembalian Kelebihan Pembayaran",
    "BAB XII : Pemeriksaan Pajak",
    "BAB XIII : Penagihan Pajak",
    "BAB XIV : Kedaluwarsa Penagihan Pajak",
    "BAB XV  : Penghapusan Piutang Pajak",
    "BAB XVI : Keringanan, Pengurangan, dan Pembebasan",
    "BAB XVII: Kemudahan Perpajakan Daerah",
    "BAB XVIII: Keberatan dan Banding",
    "BAB XIX : Penghargaan",
    "BAB XX  : Ketentuan Penutup"
]
add_bullet_text(slide, Inches(0.8), Inches(1.3), Inches(11), Inches(5.8),
                toc_items, font_size=15, color=BLACK, bullet_char="", line_spacing=1.5)

shp5 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(7.35), prs.slide_width, Inches(0.15))
shp5.fill.solid(); shp5.fill.fore_color.rgb = BLUE_DARK; shp5.line.fill.background()

# ═══════════════════════════════════════════════════
# BAB I: KETENTUAN UMUM (Pasal 1)
# ═══════════════════════════════════════════════════
make_section_slide("BAB I: KETENTUAN UMUM", "Pasal 1 — Definisi dan Istilah dalam Peraturan ini", "BAB I")

# Key definitions
make_content_slide("Definisi Penting (Pasal 1)",
    [
        "**Daerah** — Daerah Kota Bekasi",
        "**Pemerintah Daerah** — Wali Kota sebagai unsur penyelenggara pemerintahan daerah",
        "**Bapenda** — Badan Pendapatan Daerah Kota Bekasi",
        "**Reklame** — Benda/alat/media untuk tujuan komersial memperkenalkan, menganjurkan, mempromosikan sesuatu",
        "**Pajak Reklame** — Pajak atas penyelenggaraan reklame",
        "**Nilai Sewa Reklame (NSR)** — Dasar pengenaan Pajak Reklame (nilai jual objek + nilai strategis)",
        "**NPWPD** — Nomor Pokok Wajib Pajak Daerah (identitas Wajib Pajak)",
        "**NOPD** — Nomor Objek Pajak Daerah (identitas objek pajak)",
        "**SKPD** — Surat Ketetapan Pajak Daerah",
        "**Wajib Pajak** — Orang pribadi atau badan yang mempunyai hak dan kewajiban perpajakan"
    ],
    notes="Pasal 1 memuat 45 definisi. Hanya definisi kunci ditampilkan."
)

# ═══════════════════════════════════════════════════
# BAB II: OBJEK PAJAK, SUBJEK PAJAK, DAN WAJIB PAJAK
# ═══════════════════════════════════════════════════
make_section_slide("BAB II: OBJEK PAJAK, SUBJEK PAJAK, DAN WAJIB PAJAK",
                   "Pasal 2–4", "BAB II")

make_content_slide("Objek Pajak Reklame (Pasal 2)",
    [
        "**Objek Pajak:** Semua penyelenggaraan reklame, meliputi:",
        "   a. Reklame Papan/Billboard/Videotron/Megatron",
        "   b. Reklame Kain (spanduk, umbul-umbul, baliho)",
        "   c. Reklame Melekat/Stiker",
        "   d. Reklame Selebaran",
        "   e. Reklame Berjalan (termasuk pada kendaraan)",
        "   f. Reklame Udara (balon gas, dll)",
        "   g. Reklame Apung (terapung di permukaan air)",
        "   h. Reklame Film/Slide",
        "   i. Reklame Peragaan"
    ],
    sub_text="Pasal 2 ayat (1) & (2)"
)

make_content_slide("Pengecualian Objek Pajak (Pasal 2 ayat 3)",
    [
        "**Tidak termasuk objek Pajak Reklame:**",
        "   a. Reklame melalui internet, televisi, radio, wartaharian/mingguan/bulanan",
        "   b. Label/merek produk yang melekat pada barang diperdagangkan",
        "   c. Nama pengenal usaha/profesi (luas ≤ 1 m²) di bangunan tempat usaha",
        "   d. Reklame yang diselenggarakan oleh Pemerintah/Pemerintah Daerah",
        "   e. Reklame nama tempat ibadah dan panti asuhan",
        "   f. Reklame kepemilikan/peruntukan tanah (luas ≤ 1 m²)",
        "   g. Reklame kegiatan politik, sosial, keagamaan (tanpa iklan komersial)",
        "      — Politik: selama masa kampanye KPU",
        "      — Sosial/keagamaan: maks. 30 hari kalender",
        "   h. Reklame pertandingan olahraga (di bawah KONI): maks. 30 hari"
    ],
    sub_text="Pasal 2 ayat (3)"
)

make_content_slide("Subjek & Wajib Pajak Reklame (Pasal 3–4)",
    [
        "**Subjek Pajak Reklame** (Pasal 3):",
        "   Orang pribadi atau Badan yang menggunakan Reklame",
        "",
        "**Wajib Pajak Reklame** (Pasal 4):",
        "   Orang pribadi atau Badan yang menyelenggarakan Reklame",
        "",
        "Jika reklame diselenggarakan oleh **pihak ketiga**, maka pihak ketiga tersebut menjadi Wajib Pajak."
    ]
)

# ═══════════════════════════════════════════════════
# BAB III: MASA PAJAK, TAHUN PAJAK, BAGIAN TAHUN PAJAK
# ═══════════════════════════════════════════════════
make_section_slide("BAB III: MASA PAJAK, TAHUN PAJAK, DAN BAGIAN TAHUN PAJAK",
                   "Pasal 5", "BAB III")

make_content_slide("Masa & Tahun Pajak (Pasal 5)",
    [
        "**Masa Pajak Reklame:**",
        "   — Permanen: 12 (dua belas) bulan dan/atau sesuai masa penayangan",
        "   — Insidentil: dalam satuan hari, paling lama 30 (tiga puluh) hari",
        "",
        "**Tahun Pajak Reklame:**",
        "   1 (satu) tahun kalender (atau sesuai tahun buku Wajib Pajak)",
        "",
        "**Bagian Tahun Pajak Reklame:**",
        "   Jangka waktu 1 tahun pajak atas 1 bulan atau beberapa bulan kalender"
    ]
)

# ═══════════════════════════════════════════════════
# BAB IV: PENDAFTARAN DAN PENDATAAN
# ═══════════════════════════════════════════════════
make_section_slide("BAB IV: PENDAFTARAN DAN PENDATAAN WAJIB PAJAK",
                   "Pasal 6–8", "BAB IV")

make_content_slide("Pendaftaran Wajib Pajak (Pasal 6)",
    [
        "**Wajib Pajak** wajib mendaftarkan diri dan/atau objek pajak ke Wali Kota melalui Kepala Bapenda",
        "",
        "**Cara memperoleh formulir pendaftaran:**",
        "   — Ambil sendiri di Bapenda/UPTD Bapenda",
        "   — Dikirim oleh petugas Bapenda",
        "   — Akses daring (online) di laman resmi Bapenda",
        "",
        "**Persyaratan lampiran:**",
        "   — Fotokopi identitas (KTP/SIM/Paspor) | NPWP | Akta pendirian | NIB",
        "   — Surat kuasa bermeterai (jika dikuasakan)",
        "",
        "Setelah memenuhi syarat, Kepala Bapenda menerbitkan NPWPD.",
        "Jika WP tidak mendaftar, NPWPD diterbitkan secara jabatan."
    ]
)

make_content_slide("Pendataan & Penonaktifan (Pasal 7–8)",
    [
        "**Pendataan (Pasal 7):**",
        "   — Bapenda melakukan pendataan Wajib Pajak & objek Pajak Reklame",
        "   — Untuk memperoleh, melengkapi data & penatausahaan",
        "   — Termasuk informasi geografis objek Pajak Reklame",
        "   — Bapenda dapat bekerjasama dengan instansi terkait/pihak ketiga",
        "",
        "**Penonaktifan/penghapusan NPWPD (Pasal 8):**",
        "   — Jika WP tidak lagi memenuhi persyaratan subjektif & objektif",
        "   — Dapat dilakukan secara jabatan atau atas permohonan WP",
        "   — Keputusan paling lama 3 bulan sejak permohonan diterima lengkap",
        "   — Syarat: tidak ada tunggakan & tidak sedang mengajukan keberatan/banding",
        "   — WP yang tidak aktif sementara/tutup permanen wajib melapor"
    ]
)

# ═══════════════════════════════════════════════════
# BAB V: NILAI SEWA REKLAME
# ═══════════════════════════════════════════════════
make_section_slide("BAB V: NILAI SEWA REKLAME (NSR)",
                   "Pasal 9 — Dasar Pengenaan Pajak Reklame", "BAB V")

make_content_slide("Faktor Penentu Nilai Sewa Reklame (Pasal 9)",
    [
        "Dasar pengenaan Pajak Reklame adalah **Nilai Sewa Reklame (NSR)**",
        "",
        "**Faktor perhitungan NSR:**",
        "   1. Jenis reklame",
        "   2. Bahan yang digunakan",
        "   3. Lokasi penempatan (Kelas Jalan)",
        "   4. Waktu (dalam satuan detik)",
        "   5. Jangka waktu penyelenggaraan (hari kalender)",
        "   6. Jumlah media reklame (lembar)",
        "   7. Ukuran media reklame"
    ],
    sub_text="Pasal 9 ayat (1)–(2)"
)

make_content_slide("Klasifikasi Kelas Jalan (Pasal 9 ayat 3–8)",
    [
        "**Kelas Jalan Khusus:**",
        "   — Zona Khusus Jalan Tol",
        "   — Zona Khusus Premium I: Jl. A. Yani, Cut Mutia, Ir. H. Juanda, Jend. Sudirman, Sultan Agung, Transyogi, KH. Noer Ali",
        "   — Zona Khusus Premium II: Jl. Narogong Siliwangi, Jatiwaringin, Pekayon Jatiasih, Jatiasih Pondokgede, Jatimakmur, Joyo Martono, Chairil Anwar, Bintara",
        "",
        "**Kelas Jalan I (Kendali Ketat):**",
        "   Jalan lebar > 3 m, menghubungkan pusat pelayanan/permukiman (di luar kelas khusus)",
        "",
        "**Kelas Jalan II (Kendali Sedang):**",
        "   Jalan lingkungan lebar ≤ 3 m atau jalan dalam perumahan"
    ]
)

# ═══════════════════════════════════════════════════
# BAB VI: PERHITUNGAN PAJAK REKLAME
# ═══════════════════════════════════════════════════
make_section_slide("BAB VI: PERHITUNGAN PAJAK REKLAME",
                   "Pasal 10 — Tarif & Rumus Perhitungan", "BAB VI")

make_content_slide("Rumus Perhitungan Pajak Reklame (Pasal 10)",
    [
        "**Rumus Dasar:**",
        "   Pajak Reklame = Tarif Pajak Reklame × NSR",
        "",
        "**Rumus NSR Papan/Billboard:**",
        "   Nilai kelas jalan × Ukuran (m²) × Jumlah × Jangka Waktu",
        "",
        "**Rumus NSR Megatron/Videotron:**",
        "   Nilai kelas jalan × Ukuran (m²) × Jumlah × Jangka Waktu",
        "",
        "**Ketentuan tambahan:**",
        "   — Reklame indoor: NSR = 50% dari hasil perhitungan",
        "   — Ketinggian > 15 m: tambahan 20% dari NSR",
        "   — Produk tembakau & minuman keras: tambahan 50% dari NSR",
        "   — Perubahan naskah/bentuk/ukuran: dikenakan pajak atas selisihnya"
    ]
)

make_content_slide("Perubahan & Penutupan Naskah (Pasal 10 ayat 10–13)",
    [
        "**Perubahan Reklame:**",
        "   — Perubahan pemasangan terhadap naskah, bentuk, dan ukuran objek pajak → dikenakan Pajak Reklame",
        "   — Penghitungan hanya atas selisih dari perubahan",
        "   — Perubahan naskah dikecualikan jika masih dalam satu badan usaha",
        "",
        "**Penutupan Naskah:**",
        "   — Jika objek pajak sudah menayangkan naskah tetapi subjek pajak tidak diketahui",
        "   → Bapenda dapat melakukan penutupan naskah reklame"
    ]
)

# NSR TABLES
make_table_slide("Tabel NSR — Reklame Papan/Billboard (Rp/m²/hari)",
    ["Kelas Jalan", "Zona", "NSR (Rp/m²)"],
    [
        ["Kelas Jalan Khusus", "Jalan Tol", "23.575"],
        ["Kelas Jalan Khusus", "Premium 1", "16.100"],
        ["Kelas Jalan Khusus", "Premium 2", "14.950"],
        ["Kelas Jalan I", "Kendali Ketat", "13.225"],
        ["Kelas Jalan II", "Kendali Sedang", "11.500"],
    ],
    notes="Berdasarkan Pasal 10 ayat (5) huruf a. Satuan: 1 m², 1 buah, 1 hari."
)

make_table_slide("Tabel NSR — Reklame Megatron/Videotron",
    ["Kelas Jalan", "Zona", "NSR (Rp/30 detik)", "NSR (Rp/m²/tahun)"],
    [
        ["Kelas Jalan Khusus", "Jalan Tol", "17,25", "13.599.900"],
        ["Kelas Jalan Khusus", "Premium 1", "13,80", "10.879.920"],
        ["Kelas Jalan Khusus", "Premium 2", "9,20", "7.253.280"],
        ["Kelas Jalan I", "Kendali Ketat", "8,05", "6.346.620"],
        ["Kelas Jalan II", "Kendali Sedang", "5,75", "4.533.300"],
    ],
    notes="Berdasarkan Pasal 10 ayat (5) huruf b. Durasi 18 jam/hari = 2.160 tayangan/hari."
)

make_table_slide("Tabel NSR — Reklame Kain (Spanduk/Umbul/Baliho) (Rp/m²/hari)",
    ["Kelas Jalan", "Zona", "NSR (Rp/m²)"],
    [
        ["Kelas Jalan Khusus", "Jalan Tol", "30.000"],
        ["Kelas Jalan Khusus", "Premium 1", "30.000"],
        ["Kelas Jalan Khusus", "Premium 2", "25.000"],
        ["Kelas Jalan I", "Kendali Ketat", "20.000"],
        ["Kelas Jalan II", "Kendali Sedang", "19.000"],
    ],
    notes="Berdasarkan Pasal 10 ayat (5) huruf c. Satuan: 1 m², 1 buah, 1 hari."
)

make_content_slide("Tabel NSR — Jenis Reklame Lainnya (Pasal 10 ayat 5d)",
    [
        "**1. Reklame Stiker:** Rp7,5/cm² (min. Rp750.000 setiap kali penyelenggaraan)",
        "**2. Reklame Melekat (dinding/tembok):** Rp750.000/m²/tahun",
        "**3. Reklame Selebaran:** Rp600/lembar (min. Rp6.000.000 setiap kali)",
        "**4. Reklame Berjalan (kendaraan):** Rp6.000/m²/hari",
        "**5. Reklame Udara:** Rp2.400.000 sekali peragaan (maks. 1 bulan)",
        "**6. Reklame Apung:** Rp600.000 sekali peragaan (maks. 1 bulan)",
        "**7. Reklame Film/Slide:** Rp12.000/15 detik (kurang dari 15\" dibulatkan)",
        "**8. Reklame Peragaan:** Rp480.000 per penyelenggaraan"
    ]
)

# ═══════════════════════════════════════════════════
# BAB VII: PENETAPAN BESARAN PAJAK TERUTANG
# ═══════════════════════════════════════════════════
make_section_slide("BAB VII: PENETAPAN BESARAN PAJAK TERUTANG",
                   "Pasal 11", "BAB VII")

make_content_slide("Penetapan Pajak Terutang (Pasal 11)",
    [
        "Pejabat Yang Ditunjuk menetapkan Pajak Reklame terutang berdasarkan:",
        "   — Surat pendaftaran objek pajak dengan menggunakan SKPD",
        "",
        "**Jika WP tidak mendaftar:** SKPD diterbitkan secara jabatan berdasarkan data yang dimiliki Daerah",
        "",
        "**Jika hasil pemeriksaan menunjukkan pajak lebih besar:**",
        "   SKPD diterbitkan sesuai temuan (tanpa sanksi administratif)",
        "",
        "**Batas waktu penetapan:** Paling lama 5 (lima) tahun sejak terutangnya pajak"
    ]
)

# ═══════════════════════════════════════════════════
# BAB VIII: SURAT TAGIHAN PAJAK
# ═══════════════════════════════════════════════════
make_section_slide("BAB VIII: SURAT TAGIHAN PAJAK (STPD)",
                   "Pasal 12", "BAB VIII")

make_content_slide("STPD — Surat Tagihan Pajak Daerah (Pasal 12)",
    [
        "Pejabat Yang Ditunjuk dapat menerbitkan STPD dalam jangka waktu paling lama **5 tahun** sejak terutangnya pajak",
        "",
        "**STPD diterbitkan apabila:**",
        "   a. Pajak dalam SKPD tidak/kurang dibayar setelah jatuh tempo",
        "   b. Surat Keputusan Pembetulan/Keberatan/Putusan Banding tidak/kurang dibayar",
        "   c. Wajib Pajak dikenakan sanksi administratif bunga/denda",
        "",
        "**Bunga untuk huruf a:** 1% per bulan (maks. 24 bulan)",
        "**Bunga untuk huruf b:** 0,6% per bulan (maks. 24 bulan)"
    ]
)

# ═══════════════════════════════════════════════════
# BAB IX: PEMBAYARAN DAN PENYETORAN
# ═══════════════════════════════════════════════════
make_section_slide("BAB IX: PEMBAYARAN DAN PENYETORAN",
                   "Pasal 13–14", "BAB IX")

make_content_slide("Tata Cara Pembayaran (Pasal 13–14)",
    [
        "**Pembayaran/Penyetoran Pajak Reklame:**",
        "   — Dilakukan secara lunas melalui Kas Daerah atau tempat yang ditunjuk Wali Kota",
        "   — Menggunakan sistem pembayaran berbasis elektronik (jika belum tersedia → tunai)",
        "   — **Jatuh tempo:** paling lama 1 (satu) bulan sejak pengiriman SKPD",
        "",
        "**Keterlambatan:**",
        "   — Sanksi bunga 1% per bulan (maks. 24 bulan) + STPD",
        "   — STPD harus dilunasi maks. 30 hari sejak pengiriman",
        "",
        "**Tanda bukti:**",
        "   — Wajib Pajak diberikan stiker sebagai tanda pembayaran",
        "   — Bukti pembayaran sah dari Bank Persepsi/tempat lain yang ditunjuk",
        "",
        "**Tempat:** Bank Persepsi atau tempat lain yang ditunjuk Wali Kota"
    ]
)

# ═══════════════════════════════════════════════════
# BAB X: PEMBETULAN DAN PEMBATALAN KETETAPAN
# ═══════════════════════════════════════════════════
make_section_slide("BAB X: PEMBETULAN DAN PEMBATALAN KETETAPAN",
                   "Pasal 15–20", "BAB X")

make_content_slide("Pembetulan/Pembatalan Ketetapan (Pasal 15–16)",
    [
        "**Dapat dilakukan jika terdapat:**",
        "   a. Kesalahan tulis (nama, alamat, NPWPD, nomor surat, masa pajak, dll.)",
        "   b. Kesalahan hitung (penjumlahan, pengurangan, perkalian, pembagian, tarif)",
        "   c. Kekeliruan penerapan ketentuan perundang-undangan",
        "",
        "**Persyaratan permohonan (Pasal 16):**",
        "   — 1 permohonan untuk 1 ketetapan",
        "   — Diajukan oleh WP atau kuasanya secara tertulis (Bahasa Indonesia)",
        "   — Melampirkan identitas WP + kuasa + asli dokumen ketetapan"
    ]
)

make_content_slide("Prosedur Pembetulan (Pasal 17–20)",
    [
        "**Pasal 17:** Permohonan yang tidak memenuhi syarat → tidak dipertimbangkan",
        "",
        "**Pasal 18:**",
        "   — Kepala Bapenda wajib menerbitkan Surat Keputusan Pembetulan dalam 6 bulan",
        "   — Keputusan: mengabulkan (menambah/mengurangi/menghapus), membatalkan, atau menolak",
        "   — Jika > 6 bulan tanpa keputusan → permohonan dianggap dikabulkan",
        "",
        "**Pasal 19:** Pembetulan secara jabatan jika Kepala Bapenda mengetahui kesalahan",
        "",
        "**Pasal 20:** Pembetulan dapat dilakukan kembali jika masih ada kesalahan"
    ]
)

# ═══════════════════════════════════════════════════
# BAB XI: PENGEMBALIAN KELEBIHAN PEMBAYARAN
# ═══════════════════════════════════════════════════
make_section_slide("BAB XI: PENGEMBALIAN KELEBIHAN PEMBAYARAN",
                   "Pasal 21", "BAB XI")

make_content_slide("Pengembalian Kelebihan Pembayaran (Pasal 21)",
    [
        "WP dapat mengajukan permohonan pengembalian kelebihan pembayaran pajak",
        "",
        "**Proses:**",
        "   — Keputusan diterbitkan maks. 12 bulan sejak permohonan diterima",
        "   — Jika > 12 bulan tanpa keputusan → permohonan dianggap dikabulkan",
        "   → SKPDLB diterbitkan dalam 1 bulan",
        "",
        "**Jika WP memiliki utang pajak lain:** kelebihan diperhitungkan untuk melunasi utang tsb.",
        "",
        "**Pengembalian** dilakukan maks. 2 bulan sejak SKPDLB diterbitkan",
        "Jika terlambat → imbalan bunga 0,6%/bulan"
    ]
)

# ═══════════════════════════════════════════════════
# BAB XII: PEMERIKSAAN PAJAK
# ═══════════════════════════════════════════════════
make_section_slide("BAB XII: PEMERIKSAAN PAJAK",
                   "Pasal 22–23", "BAB XII")

make_content_slide("Pemeriksaan Pajak (Pasal 22–23)",
    [
        "**Kewenangan:** Kepala Bapenda berwenang melakukan pemeriksaan untuk menguji kepatuhan WP",
        "",
        "**Pemeriksaan kepatuhan dilakukan jika:**",
        "   — WP mengajukan pengembalian/kompensasi kelebihan pembayaran",
        "   — Ada data konkret yang menunjukkan pajak kurang dibayar",
        "   — WP terpilih berdasarkan analisis risiko",
        "",
        "**Pemeriksaan untuk tujuan lain:**",
        "   — Pemberian/penghapusan NPWPD, penyelesaian keberatan, pencocokan data, penagihan",
        "",
        "**Kewajiban WP saat diperiksa:**",
        "   — Memperlihatkan buku/catatan/dokumen",
        "   — Memberi akses ke tempat/ruangan",
        "   — Memberikan keterangan yang diperlukan"
    ]
)

# ═══════════════════════════════════════════════════
# BAB XIII: PENAGIHAN PAJAK
# ═══════════════════════════════════════════════════
make_section_slide("BAB XIII: PENAGIHAN PAJAK",
                   "Pasal 24", "BAB XIII")

make_content_slide("Penagihan Pajak (Pasal 24)",
    [
        "**Dasar Penagihan Pajak:**",
        "   SKPD, SKPDKB, SKPDKBT, STPD, Surat Keputusan Pembetulan, Surat Keputusan Keberatan, dan Putusan Banding",
        "",
        "**Prosedur:**",
        "   — Sebelum jatuh tempo: dapat dilakukan imbauan",
        "   — Setelah jatuh tempo dan belum dilunasi: Penagihan Pajak sesuai ketentuan perundang-undangan"
    ]
)

# ═══════════════════════════════════════════════════
# BAB XIV: KEDALUWARSA PENAGIHAN PAJAK
# ═══════════════════════════════════════════════════
make_section_slide("BAB XIV: KEDALUWARSA PENAGIHAN PAJAK",
                   "Pasal 25", "BAB XIV")

make_content_slide("Kedaluwarsa Penagihan (Pasal 25)",
    [
        "Hak melakukan penagihan pajak **kedaluwarsa setelah 5 tahun** sejak terutangnya pajak",
        "Kecuali WP melakukan tindak pidana di bidang perpajakan daerah",
        "",
        "**Kedaluwarsa tertangguh jika:**",
        "   a. Diterbitkan Surat Teguran dan/atau Surat Paksa",
        "   b. Ada pengakuan utang pajak dari WP (langsung/tidak langsung)",
        "      — Langsung: WP menyatakan masih punya utang",
        "      — Tidak langsung: pengajuan angsuran, penundaan, atau keberatan"
    ]
)

# ═══════════════════════════════════════════════════
# BAB XV: PENGHAPUSAN PIUTANG PAJAK
# ═══════════════════════════════════════════════════
make_section_slide("BAB XV: PENGHAPUSAN PIUTANG PAJAK",
                   "Pasal 26", "BAB XV")

make_content_slide("Penghapusan Piutang Pajak (Pasal 26)",
    [
        "Piutang pajak yang **tidak mungkin ditagih lagi** karena kedaluwarsa dapat dihapuskan",
        "",
        "**Prosedur:**",
        "   — Bapenda melakukan penelitian → membuat berita acara",
        "   — Tim peneliti ditetapkan dengan Keputusan Wali Kota",
        "   — Kepala Bapenda menyusun daftar usulan penghapusan",
        "   — Disampaikan ke Wali Kota",
        "   — Ditetapkan dengan Keputusan Wali Kota",
        "",
        "**Pertimbangan:**",
        "   — Pelaksanaan penagihan sampai batas kedaluwarsa",
        "   — Hasil koordinasi dengan Inspektorat Daerah"
    ]
)

# ═══════════════════════════════════════════════════
# BAB XVI: KERINGANAN, PENGURANGAN, DAN PEMBEBASAN
# ═══════════════════════════════════════════════════
make_section_slide("BAB XVI: KERINGANAN, PENGURANGAN, DAN PEMBEBASAN",
                   "Pasal 27", "BAB XVI")

make_content_slide("Keringanan, Pengurangan, & Pembebasan (Pasal 27)",
    [
        "Wali Kota/Kepala Bapenda dapat memberikan:",
        "   — Keringanan",
        "   — Pengurangan",
        "   — Pembebasan",
        "   — Penundaan pembayaran",
        "atas pokok dan/atau sanksi pajak",
        "",
        "**Kondisi WP:** kemampuan membayar / tingkat likuiditas",
        "**Kondisi Objek Pajak:** terdampak bencana alam, kebakaran, huru-hara, kerusuhan",
        "",
        "**Ditetapkan dengan Keputusan Wali Kota**"
    ]
)

# ═══════════════════════════════════════════════════
# BAB XVII: KEMUDAHAN PERPAJAKAN DAERAH
# ═══════════════════════════════════════════════════
make_section_slide("BAB XVII: KEMUDAHAN PERPAJAKAN DAERAH",
                   "Pasal 28", "BAB XVII")

make_content_slide("Kemudahan Perpajakan Daerah (Pasal 28)",
    [
        "**Bentuk kemudahan:**",
        "   a. **Perpanjangan batas waktu pembayaran** — untuk WP yang mengalami keadaan kahar",
        "   b. **Fasilitas angsuran/penundaan pembayaran** — untuk WP dengan kesulitan likuiditas",
        "",
        "**Keadaan kahar meliputi:**",
        "   Bencana alam, kebakaran, kerusuhan massal, wabah penyakit, atau keadaan lain berdasarkan pertimbangan Wali Kota",
        "",
        "**Jangka waktu angsuran/penundaan:** maksimal 24 bulan",
        "**Bunga:** 0,6% per bulan dari jumlah pajak yang masih harus dibayar",
        "",
        "Wali Kota memperhatikan kepatuhan WP dalam 2 tahun terakhir"
    ]
)

# ═══════════════════════════════════════════════════
# BAB XVIII: KEBERATAN DAN BANDING
# ═══════════════════════════════════════════════════
make_section_slide("BAB XVIII: KEBERATAN DAN BANDING",
                   "Pasal 29–33", "BAB XVIII")

make_content_slide("Keberatan (Pasal 29–31)",
    [
        "**WP dapat mengajukan keberatan** terhadap SKPD, SKPDKB, SKPDKBT, SKPDLB, SKPDN",
        "",
        "**Syarat:**",
        "   — Diajukan tertulis dalam Bahasa Indonesia, maks. **3 bulan** sejak SKPD dikirim",
        "   — WP harus sudah membayar paling sedikit sejumlah yang disetujui",
        "   — Jika tidak memenuhi syarat → tidak dianggap sebagai surat keberatan",
        "",
        "**Keputasan:** diterbitkan maks. **12 bulan** sejak surat diterima",
        "   — Jika tidak ada keputusan dalam 12 bulan → keberatan dianggap diterima",
        "",
        "**Sanksi jika keberatan ditolak:** denda 30% dari jumlah pajak"
    ]
)

make_content_slide("Banding (Pasal 32–33)",
    [
        "**WP dapat mengajukan banding** ke badan peradilan pajak atas Surat Keputusan Keberatan",
        "",
        "**Syarat:**",
        "   — Paling lama **3 bulan** sejak keputusan diterima",
        "   — Diajukan tertulis dalam Bahasa Indonesia dengan alasan jelas",
        "",
        "**Efek:** menangguhkan kewajiban bayar pajak sampai 1 bulan sejak Putusan Banding",
        "",
        "**Sanksi jika banding ditolak:** denda 60% dari jumlah pajak berdasarkan Putusan Banding",
        "**Jika dikabulkan:** kelebihan dikembalikan + imbalan bunga 0,6%/bulan (maks. 24 bulan)"
    ]
)

# ═══════════════════════════════════════════════════
# BAB XIX: PENGHARGAAN
# ═══════════════════════════════════════════════════
make_section_slide("BAB XIX: PENGHARGAAN",
                   "Pasal 34–35", "BAB XIX")

make_content_slide("Penghargaan Wajib Pajak Taat (Pasal 34–35)",
    [
        "**Bentuk penghargaan:**",
        "   Piagam penghargaan, hadiah, dan/atau sejenisnya",
        "   Dibebankan pada APBD",
        "",
        "**Kriteria Wajib Pajak Taat Pajak:**",
        "   a. Menyetorkan pajak tepat waktu minimal 1 tahun & tidak ada tunggakan 3 tahun terakhir",
        "   b. Kontribusi signifikan dalam mendukung program Pemerintah Daerah",
        "",
        "**Penetapan:**",
        "   — Ditetapkan dengan Keputusan Wali Kota",
        "   — Tim penilai dibentuk oleh Wali Kota"
    ]
)

# ═══════════════════════════════════════════════════
# BAB XX: KETENTUAN PENUTUP
# ═══════════════════════════════════════════════════
make_section_slide("BAB XX: KETENTUAN PENUTUP",
                   "Pasal 36–37", "BAB XX")

make_content_slide("Ketentuan Penutup (Pasal 36–37)",
    [
        "**Pasal 36 — Pencabutan:**",
        "Peraturan Wali Kota yang dicabut dan dinyatakan tidak berlaku:",
        "   1. Perwal Bekasi No. 48 Tahun 2012 tentang Petunjuk Pelaksanaan Perda No. 14/2012",
        "   2. Perwal Bekasi No. 52 Tahun 2013 tentang Perubahan atas Perwal No. 48/2012",
        "",
        "**Pasal 37 — Mulai Berlaku:**",
        "Peraturan Wali Kota ini mulai berlaku pada tanggal **diundangkan**",
        "",
        "Ditetapkan di Bekasi",
        "pada tanggal **20 Desember 2024**",
        "",
        "**Pj. WALI KOTA BEKASI,**",
        "ttd",
        "**R. GANI MUHAMAD**"
    ]
)

# ═══════════════════════════════════════════════════
# CLOSING SLIDE
# ═══════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg_shape(slide, BLUE_DARK)
add_accent_bar(slide, left=Inches(0.5), color=GOLD)

add_text_box(slide, Inches(1.2), Inches(1.0), Inches(10), Inches(0.6),
             "BERITA DAERAH KOTA BEKASI", font_size=18, color=GOLD, bold=True)

add_text_box(slide, Inches(1.2), Inches(1.8), Inches(10), Inches(1.5),
             "TERIMA KASIH", font_size=48, color=WHITE, bold=True)

add_text_box(slide, Inches(1.2), Inches(3.3), Inches(10), Inches(2),
             "Peraturan Wali Kota Bekasi Nomor 51 Tahun 2024\n"
             "Tentang Pengelolaan Pajak Reklame\n\n"
             "Sumber: https://jdih.bekasikota.go.id",
             font_size=18, color=RGBColor(0xCC, 0xCC, 0xCC))

shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1.2), Inches(5.8), Inches(5), Pt(3))
shp.fill.solid(); shp.fill.fore_color.rgb = GOLD; shp.line.fill.background()

# ═══════════════════════════════════════════════════
# SAVE
# ═══════════════════════════════════════════════════
output_path = "/home/runner/coba_ppt/Perwal_Bekasi_51_2024_Pajak_Reklame.pptx"
prs.save(output_path)
print(f"✅ Presentation saved to: {output_path}")
print(f"   Total slides: {len(prs.slides)}")
